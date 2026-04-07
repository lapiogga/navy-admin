"""
해병대 행정포탈 시스템 - 18개 서브시스템 분석/설계 PPT 자동 생성
각 서브시스템당 5페이지: 표지, 업무개요, 주요기능, 입출력/규칙, 시스템구조
"""
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 색상 정의
NAVY = RGBColor(0x00, 0x33, 0x66)
DARK_NAVY = RGBColor(0x00, 0x1F, 0x3F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
GRAY = RGBColor(0x66, 0x66, 0x66)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
ACCENT_BLUE = RGBColor(0x1A, 0x73, 0xE8)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
BORDER_GRAY = RGBColor(0xDD, 0xDD, 0xDD)

# 서브시스템 정보
SUBSYSTEMS = [
    {"code": "SYS01", "name": "초과근무관리체계", "processes": 99, "complexity": "상", "pages": 23,
     "keywords": "신청/결재/월말결산/당직/개인설정",
     "overview": "해병대 장병의 초과근무 신청/결재/정산/통계를 관리하는 체계.\n신청서 작성부터 월말결산, 당직업무, 개인설정까지 초과근무 전 과정을 체계적으로 처리한다.",
     "users": "일반 사용자(신청자), 결재자, 부대관리자, 체계관리자, 당직자",
     "flows": ["신청서 작성 > 결재 요청 > 승인/반려", "일괄처리 신청서 작성/승인", "월말결산 마감기한 입력 > 결산 > 마감",
               "당직개소/초과근무자 관리", "부대 설정 (최대인정시간/근무시간/공휴일/결재선)"],
     "menus": [("게시판", "공지사항, 질의응답"), ("신청서 관리", "신청서 작성/결재/일괄처리/월말결산"),
               ("현황조회", "나의 근무현황/부재관리/월말결산/자료출력"), ("부대관리", "부대인원/최대인정시간/근무시간/공휴일/결재선"),
               ("당직업무", "초과근무자/당직개소/개인별 승인"), ("개인설정", "당직개소/부서 설정"), ("관리자", "공통기능 모듈")],
     "input": "신청서 종류(사전/사후), 근무일, 초과근무 시간, 신청자 IP, 당직개소, 초과근무 내용",
     "output": "월말결산 엑셀, 초과근무 현황 그래프, 부대인원 엑셀",
     "rules": ["결재 이후 신청서 수정/삭제는 관리자만 가능", "최대인정시간 예외처리자는 소속 변경 시 기본값 자동 복구",
               "당직 PC MAC 주소 확인 별도 Client 필요"],
     "page_list": "OtRequestPage, OtApprovalPage, OtBulkPage, OtBulkApprovalPage, OtMonthlyClosingPage 등 23개"},
    {"code": "SYS02", "name": "설문종합관리체계", "processes": 31, "complexity": "중", "pages": 7,
     "keywords": "설문생성/참여/승인/결과분석",
     "overview": "해병대 내 설문조사를 체계적으로 관리하는 시스템.\n설문 생성/문항 등록/승인/배포/참여/결과 분석까지 설문 전 주기를 관리한다.",
     "users": "설문 작성자(일반), 설문 참여자(일반), 체계관리자(승인/관리)",
     "flows": ["설문조사 설정 등록 > 문항 등록 > 제출(승인요청)", "승인 대기 > 승인/반려", "활성화 설문 > 참여", "지난 설문 결과 조회/출력"],
     "menus": [("게시판", "공지사항, 질의응답"), ("설문관리", "나의 설문/설문참여/지난 설문/체계관리"), ("관리자", "공통기능 모듈")],
     "input": "설문명, 설문내용, 설문대상, 대상자 수, 설문기간, 결과공개, 대상 계급/부대/직책/성별, 첨부파일",
     "output": "설문 결과 차트(Chart), 결과 엑셀",
     "rules": ["설문조사 대상자 여부 체크 필요", "결과 공개 설문만 지난 설문 결과 조회 가능", "문항 등록은 엑셀 양식 업로드 방식"],
     "page_list": "MySurveyPage, SurveyParticipationPage, SurveyAdminPage 등 7개"},
    {"code": "SYS03", "name": "성과관리체계", "processes": 76, "complexity": "상", "pages": 16,
     "keywords": "과제계층/실적입력/평가/그래프",
     "overview": "해병대 부대(서)별 성과관리를 지원하는 체계.\n기준정보 관리, 연간과제 등록/실적 입력/승인/평가, 평가결과 조회까지 성과관리 전 주기를 관리한다.",
     "users": "과제 담당자(일반), 결재자, 평가자, 체계관리자",
     "flows": ["기준년도/평가조직 설정", "지휘방침 > 추진중점 > 중과제 계층 구조 관리", "소과제/상세과제 등록 > 실적 입력 > 상신",
               "과제실적 승인/반려", "과제/개인 실적 평가 > 결과 조회"],
     "menus": [("대시보드", "업무 달성률/추진율 그래프"), ("기준정보관리", "기준년도/평가조직/개인실적/과제관리"),
               ("연간과제관리", "과제등록/실적입력/승인/평가"), ("평가결과", "결과조회/입력현황/과제검색"), ("관리자", "공통기능 모듈")],
     "input": "기준년도, 평가조직, 과제 계층(지휘방침~상세과제), 업무실적 내용",
     "output": "업무 달성률/추진율 차트, 평가결과/과제목록 엑셀",
     "rules": ["과제 계층: 지휘방침 > 추진중점과제 > 중과제 > 소과제 > 상세과제", "중과제/소과제 엑셀 일괄 등록 지원"],
     "page_list": "PerfPolicyPage, PerfMainTaskPage, PerfEvalResultPage 등 16개"},
    {"code": "SYS04", "name": "인증서발급신청체계", "processes": 14, "complexity": "하", "pages": 3,
     "keywords": "발급신청/승인/등록대장",
     "overview": "해병대 장병의 인증서(국방전자서명인증서) 발급을 온라인으로 신청/승인/관리하는 체계.",
     "users": "신청자(일반), 관리자(승인)",
     "flows": ["발급 신청 > 현황 조회 > 회수", "승인대기 목록 > 상세 > 승인/반려", "발급 목록/엑셀/통계"],
     "menus": [("게시판", "공지사항, 질의응답"), ("인증서발급신청", "신청/승인관리/등록대장"), ("관리자", "공통기능 모듈")],
     "input": "인증서 구분, 신청구분, 용도, 사유, 신청자정보(계급/성명/군번/소속/이메일/전화), 정보활용동의",
     "output": "발급현황 엑셀, 발급 통계",
     "rules": ["승인대기 상태에서만 회수 가능", "반려 시 반송의견 필수"],
     "page_list": "CertificateApplyPage, CertificateApprovalPage, CertificateRegisterPage"},
    {"code": "SYS05", "name": "행정규칙포탈체계", "processes": 15, "complexity": "하", "pages": 4,
     "keywords": "현행규정/예규/지시문서",
     "overview": "해군 현행규정, 예규, 지시문서를 통합 관리하는 포탈.\n부/실/단별 규정을 조직도 기반으로 조회하고, 예규와 지시문서를 게시판 형태로 관리한다.",
     "users": "일반 사용자(열람), 관리자(등록/수정)",
     "flows": ["조직도 전시 > 규정 조회 > 열람/다운로드", "예규 목록 조회 > CRUD", "지시문서 CRUD"],
     "menus": [("게시판", "공지사항, 질의응답"), ("해군규정", "현행규정"), ("예규", "해군본부/예하부대"), ("지시", "지시문서"), ("관리자", "공통기능 모듈")],
     "input": "규정명, 규정 파일, 예규 정보, 지시문서 정보",
     "output": "규정 전문 열람, 파일 다운로드",
     "rules": ["현행규정 열람 시 전문 전시 + 파일 다운로드"],
     "page_list": "RegulationListPage, PrecedentHQPage, PrecedentUnitPage, DirectiveListPage"},
    {"code": "SYS06", "name": "해병대규정관리체계", "processes": 30, "complexity": "중", "pages": 4,
     "keywords": "해병대규정/예규/지시/DocViewer",
     "overview": "해병대 고유 규정, 예규, 지시문서를 관리하는 체계.\nSYS05와 유사 구조이나 해병대사령부 예규와 규정예고/자료실 게시판이 추가된다.",
     "users": "일반 사용자(열람), 관리자(등록/수정)",
     "flows": ["해병대규정 조직도 > 규정 열람(Document Viewer)", "해병대사령부 예규 CRUD", "예하부대 조직도/링크", "지시문서 CRUD"],
     "menus": [("게시판", "공지사항, 규정예고, 자료실"), ("해병대규정", "현행규정"), ("예규", "해병대사령부/예하부대"),
               ("지시", "지시문서"), ("관리자", "공통기능 모듈")],
     "input": "규정명, 규정 파일, 예규 정보, 지시문서 정보",
     "output": "규정 전문 열람(Document Viewer), 파일 다운로드",
     "rules": ["현행규정 열람 시 Document Viewer 필요", "게시판 3개(공지/규정예고/자료실)"],
     "page_list": "RegulationListPage, PrecedentHQPage, PrecedentUnitPage, DirectiveListPage"},
    {"code": "SYS07", "name": "군사자료관리체계", "processes": 40, "complexity": "중", "pages": 9,
     "keywords": "비밀문서/활용/통계/해기단",
     "overview": "해기단에서 관리하는 군사자료(비밀문서 포함)의 목록 관리, 활용(대출/열람/지출) 처리, 통계 관리 체계.",
     "users": "일반 사용자(열람/검색), 관리자(등록/활용처리/통계)",
     "flows": ["군사자료 목록/상세 > 등록/수정/삭제/일괄등록", "평가심의 대상 > 결과 입력", "대출/열람/지출 처리 > 반납", "통계자료 조회/저장/출력"],
     "menus": [("게시판", "공지사항, 질의응답"), ("군사자료 관리", "군사자료 관리/활용/통계자료"), ("해기단자료", "해기단자료 CRUD"), ("관리자", "공통기능 모듈")],
     "input": "비밀등급, 보관장소, 관리번호, 문서구분, 이관일자, 보존기간 등 17개 필드",
     "output": "군사자료 목록 엑셀, 각종 기록부/통계 양식",
     "rules": ["삭제 시 삭제 사유 필수 / 삭제 로그 저장", "데이터 무결성 검증 필수", "파기/보존기간 연장 기능"],
     "page_list": "MilDataListPage, MilDataUsagePage, MilDataStatsPage, HaegidanListPage 등 9개"},
    {"code": "SYS08", "name": "부대계보관리체계", "processes": 59, "complexity": "상", "pages": 12,
     "keywords": "계보/주요활동/직위자/부대기",
     "overview": "해병대 부대의 계보(주요활동, 주요직위자, 제원/계승부대, 부대기/마크, 부대기록부) 정보를 체계적으로 관리하는 시스템.",
     "users": "계보담당자, 중간결재자, 확인관, 부대관리자, 체계관리자",
     "flows": ["권한 신청 > 관리자 승인/반려", "주요활동 CRUD > 결재", "주요직위자 CRUD", "제원/계승부대, 부대기/마크 관리", "통계/출력"],
     "menus": [("권한신청", "권한신청/관리/조회"), ("주요활동", "관리/결재"), ("주요직위자", "관리"),
               ("제원/계승부대", "관리"), ("부대기/부대마크", "관리"), ("통계 및 출력", "입력 통계"), ("관리자", "공통기능 모듈")],
     "input": "주요활동(연혁부호/비밀여부/카테고리/대표사진), 주요직위자(계승번호/역대/군번/병과/표어), 제원(임무/소재지)",
     "output": "각 항목 엑셀/프린트, 각종 통계",
     "rules": ["사용자 소속 변경 시 권한 자동 해지", "반려 시 반려 사유 필수", "일괄등록 시 데이터 무결성 검증"],
     "page_list": "UnitKeyPersonPage, UnitActivityPage, UnitLineageTreePage, UnitFlagPage 등 12개"},
    {"code": "SYS09", "name": "영현보훈체계", "processes": 35, "complexity": "중", "pages": 18,
     "keywords": "사망자/상이자/확인서/통계",
     "overview": "해병대 전공사상자(사망자/상이자) 정보 관리 및 전공사상심사위원회 결과 관리 체계.\n각종 확인서/보고서 출력과 다양한 통계를 제공한다.",
     "users": "담당자(자료 입력), 관리자(통계/출력)",
     "flows": ["사망자 현황 CRUD", "상이자 현황 CRUD", "전공사상심사 결과 CRUD", "확인서/보고서/통계 출력"],
     "menus": [("게시판", "공지사항"), ("자료입력", "사망자/상이자/전공사상심사"), ("통계 및 자료 출력", "확인서/보고서/통계/명부"), ("관리자", "공통기능 모듈")],
     "input": "사망자(군번/성명/주민번호/사망구분 등), 상이자(보훈청명/병명/전공상여부), 심사(심사차수/일자/분류)",
     "output": "국가유공자 확인서, 현황 보고서, 각종 통계 엑셀, 명부",
     "rules": ["양식은 전공사상자 처리 훈령 참조", "검색조건: 군구분, 군번, 성명, 주민번호, 계급, 소속"],
     "page_list": "DeceasedPage, InjuredPage, ReviewPage, CertPages, StatsPages 등 18개"},
    {"code": "SYS10", "name": "주말버스예약관리체계", "processes": 44, "complexity": "중", "pages": 13,
     "keywords": "예약/좌석배정/배차/타군사용자",
     "overview": "해병대 장병의 주말버스 예약/좌석 배정/승차권 발급/위규자 관리/타군 사용자 관리를 수행하는 체계.",
     "users": "일반 사용자(예약), 관리자(배차/좌석관리), 타군 사용자",
     "flows": ["예약자 검색 > 예약 > 승차권 발급/반납", "대기자 > 좌석 배정(수동/자동)", "배차 정보 CRUD", "위규자/타군사용자 관리"],
     "menus": [("게시판", "공지사항, 질의응답"), ("주말버스", "예약/대기자/예약현황/배차/예약시간/사용현황/위규자/타군사용자"), ("관리자", "공통기능 모듈")],
     "input": "운행일자, 출발지/도착지, 시간, 예약자 정보, 배차 정보, 타군사용자(군구분/부대/군번/직책/전화/메일)",
     "output": "승차권 프린트, 예약현황/대기자 엑셀, 사용현황/위규자 양식",
     "rules": ["타군 사용자 별도 로그인 화면 필요", "타군 패스워드 초기화 메일 발송", "예약시간은 구간별/계급별/일자별"],
     "page_list": "BusReservationPage, BusSchedulePage, BusDispatchPage, ExternalUserPage 등 13개"},
    {"code": "SYS11", "name": "연구자료종합관리체계", "processes": 19, "complexity": "하", "pages": 5,
     "keywords": "연구자료/카테고리/다운로드이력",
     "overview": "해군 연구용역 자료를 카테고리별(10개)로 등록/관리하고, 대시보드에서 연구현황/소개자료를 전시하는 체계.",
     "users": "일반 사용자(열람/등록), 관리자",
     "flows": ["대시보드 연구현황/소개자료 전시", "카테고리별 게시물 목록/검색/상세/CRUD", "관리자 CRUD/다운로드 이력 조회"],
     "menus": [("대시보드", "연구용역 현황/소개자료"), ("게시판", "공지사항, 질의응답"), ("연구자료종합관리", "연구자료/관리자"), ("관리자", "공통기능 모듈")],
     "input": "제목, 연구분야(10개 카테고리), 연구년도, 예산, 연구자(군내/외부), 진행사항, 과제담당관, 연구계획서/보고서",
     "output": "연구자료 목록 엑셀, 연구용역 현황 게시표",
     "rules": ["열람은 모든 사용자, 수정/삭제는 작성자/관리자만", "첨부파일 다운로드 시 이력 저장(군번/계급/성명/IP)"],
     "page_list": "ResearchMainPage, ResearchListPage, ResearchAdminPage, ResearchFilePage"},
    {"code": "SYS12", "name": "지시건의사항관리체계", "processes": 32, "complexity": "중", "pages": 5,
     "keywords": "지시/건의/추진현황매트릭스",
     "overview": "대통령/국방부장관/지휘관의 지시사항과 건의사항의 추진현황을 매트릭스 형태로 관리하고 조치사항을 추적하는 체계.",
     "users": "담당자(등록/조치), 관리자(추진현황 조회)",
     "flows": ["지시사항 추진현황 매트릭스 조회", "지시사항/건의사항 CRUD", "조치사항 등록/수정"],
     "menus": [("게시판", "공지사항, 질의응답"), ("지시사항", "대통령/국방부장관/지휘관 지시"), ("건의사항", "지휘관 건의"), ("관리자", "공통기능 모듈")],
     "input": "지시자, 지시일자, 수명부대, 추진상태, 지시내용, 종류(문서/구두), 조치(담당자/계획/내용/첨부)",
     "output": "추진현황 매트릭스(총괄/부서별), 엑셀",
     "rules": ["추진현황은 매트릭스 형태", "지시사항 4개 섹션(대통령/국방부장관/지휘관/건의) 동일 구조"],
     "page_list": "DirectiveListPage, ProposalListPage, DirectiveProgressPage, ProposalProgressPage, DirectiveAdminPage"},
    {"code": "SYS13", "name": "지식관리체계", "processes": 23, "complexity": "하", "pages": 6,
     "keywords": "지식등록/추천/평점/댓글/통계",
     "overview": "해병대 장병이 업무 지식을 등록/공유/열람하고, 추천/평점/댓글 기능으로 활용도를 높이는 지식관리 시스템.",
     "users": "일반 사용자(등록/열람), 관리자(숨김/삭제)",
     "flows": ["나의 지식 CRUD", "지식 검색/상세 > 추천/평점/댓글", "유형별/부대별/작성자별 통계"],
     "menus": [("게시판", "공지사항, 질의응답"), ("지식관리", "나의 지식/지식 관리"), ("지식열람", "지식열람"), ("지식통계", "통계"), ("관리자", "공통기능 모듈")],
     "input": "제목, 지식유형, 출처(생산/카피), 키워드(다수), 내용, 첨부파일",
     "output": "유형별/부대별/작성자별 통계, 부대별 목록 엑셀",
     "rules": ["부대별 통계는 등록자 소속 기준", "정렬: 작성 수/추천 수/평점/조회 수"],
     "page_list": "MyKnowledgePage, KnowledgeListPage, KnowledgeDetailPage, KnowledgeAdminPage, KnowledgeStatsPage"},
    {"code": "SYS14", "name": "나의 제언", "processes": 16, "complexity": "하", "pages": 3,
     "keywords": "제언등록/상태관리/댓글/서식",
     "overview": "해병대 장병이 업무 개선 의견이나 제안을 등록하고, 관리자가 진행상태/조치유형을 관리하며 댓글로 소통하는 제언 시스템.",
     "users": "일반 사용자(제언 등록), 관리자(상태관리/서식관리)",
     "flows": ["대시보드 제언 목록 조회", "제언 등록/수정/삭제", "관리자 상태변경(등록>접수>진행>완료/반려)", "댓글 CRUD"],
     "menus": [("대시보드", "제언 목록"), ("게시판", "공지사항"), ("나의제언", "제언확인/서식관리"), ("관리자", "공통기능 모듈")],
     "input": "제목, 내용, 첨부파일 (소속/제언자/직책/전화번호는 로그인정보 연동)",
     "output": "제언 목록 엑셀",
     "rules": ["진행상태: 등록>접수>진행>완료/반려", "조치유형: 정책반영/업무추진/기추진/업무참고", "반려 시 반려사유 입력 필수"],
     "page_list": "SuggestionMainPage, SuggestionListPage, SuggestionAdminPage"},
    {"code": "SYS15", "name": "보안일일결산체계", "processes": 138, "complexity": "최상", "pages": 27,
     "keywords": "비밀매체/결산/평가/점검항목",
     "overview": "해병대 보안업무의 핵심 체계. 비밀/저장매체/보안자재 보유 현황, 개인/사무실 보안일일결산, 보안수준평가, 보안교육, 결재 등을 관리한다.",
     "users": "일반 사용자, 일일보안점검관, 부서관리자, 평가관, 결재자, 체계관리자",
     "flows": ["저장매체/비밀/보안자재 보유현황 CRUD", "인계/인수 처리/결재", "개인/사무실 보안일일결산 실시", "일일보안점검관 당직표 관리",
               "개인보안수준 수시/정기 평가", "부재처리/보안교육", "결재대기/완료", "결산종합현황/관리설정"],
     "menus": [("비밀/매체관리", "저장매체/비밀/보안자재/예고문/인계인수"), ("보안일일결산", "개인/사무실결산/점검관/수준평가/부재/교육"),
               ("결재", "결재대기/완료"), ("결산종합현황", "비밀매체/개인/사무실/부재"), ("결산체계 관리", "개인설정/점검항목/휴무일/알림시간/로그/예외처리"),
               ("관리자", "공통기능 모듈")],
     "input": "저장매체(매체구분/등급/증감수량), 비밀(등급/증감/수량), 부재(대상/기간/내용), 보안교육(구분/일자/교관/내용)",
     "output": "저장매체 관리대장, 비밀관리기록부, 결산 캘린더, 이력 엑셀",
     "rules": ["사무실결산 미실시자/부재자 사유 입력 필수", "당직표 결재자 기본 지정/변경 가능", "정기평가 수정은 2차 결재 이전까지만"],
     "page_list": "SecretMediaPage, PersonalSecDailyPage, OfficeSecDailyPage, DutyOfficerPage 등 27개"},
    {"code": "SYS16", "name": "회의실예약관리체계", "processes": 21, "complexity": "하", "pages": 5,
     "keywords": "예약신청/승인/회의실관리",
     "overview": "부대(서)별 회의실 예약/관리 체계. 사용자가 예약을 신청하고 관리자가 승인/반려하며, 회의실 정보와 고정 일정을 관리한다.",
     "users": "일반 사용자(예약 신청), 관리자(승인/회의실 관리)",
     "flows": ["회의실 선택 > 예약 신청", "내예약 목록/수정/취소", "조건별 회의현황 조회", "관리자 승인/승인취소", "회의실/고정일정 CRUD"],
     "menus": [("게시판", "공지사항, 질의응답"), ("회의실 예약관리", "예약신청/내예약/회의현황/예약관리/회의실관리"), ("관리자", "공통기능 모듈")],
     "input": "회의실 관리부대, 회의실, 회의명, 회의일시, 회의등급, 주관부서",
     "output": "예약 목록 엑셀/프린트",
     "rules": ["검색조건: 회의실 관리부대, 회의실, 회의일자"],
     "page_list": "MeetingReservePage, MyReservationPage, MeetingStatusPage, ReservationMgmtPage, MeetingRoomMgmtPage"},
    {"code": "SYS17", "name": "검열결과관리체계", "processes": 25, "complexity": "중", "pages": 9,
     "keywords": "검열계획/조치과제/결재/추진현황",
     "overview": "해병대 검열계획 수립, 검열결과(조치과제) 등록, 조치결과 입력/결재, 추진현황 종합 조회 체계.",
     "users": "검열관(계획/과제 등록), 대상부대 담당자(조치결과), 결재자, 관리자",
     "flows": ["검열부대 지정", "검열계획 작성/삭제", "조치과제 등록 > 조치결과 입력", "접수(승인)/반송(반려)", "종합/세부 추진현황"],
     "menus": [("게시판", "공지사항, 질의응답"), ("검열결과 관리", "검열부대/계획/결과/결재/추진현황"), ("관리자", "공통기능 모듈")],
     "input": "검열계획(연도/계획명/기간/대상부대), 조치과제(주관부서/검열관/분야/처분종류/과제명), 과제처리(진행상태/문제점/조치계획)",
     "output": "조치과제 엑셀, 추진현황 종합/세부(그래프)",
     "rules": ["검색조건: 검열연도, 계획명, 대상부대, 기간, 진행상태, 검열분야"],
     "page_list": "InspectionPlanPage, InspectionResultPage, InspectionApprovalPage, InspectionProgressPage 등 9개"},
    {"code": "SYS18", "name": "직무기술서관리체계", "processes": 47, "complexity": "중", "pages": 7,
     "keywords": "조직진단/직무기술서/결재/표준시간",
     "overview": "해병대 조직진단 대상 관리, 개인/부서 직무기술서 작성/결재/조회, 관리자 검토, 표준업무시간 관리 체계.",
     "users": "일반 사용자(작성), 부서관리자(대리작성/부서직무기술서), 결재자, 조직진단 관리자",
     "flows": ["조직진단 대상 CRUD", "개인/부서 직무기술서 작성/수정", "결재대기 > 결재/반려", "관리자 검토/의견/반송", "표준업무시간 관리"],
     "menus": [("게시판", "공지사항, 질의응답, 자료실"), ("직무기술서 관리", "조직진단/작성/결재/조회(관리자)/표준업무시간"), ("관리자", "공통기능 모듈")],
     "input": "조직진단(진단명/부대/기간/대상자/제외대상자), 직무기술서(엑셀 업로드/결재자), 표준업무시간(신분별/적용기간)",
     "output": "개인/부서 직무기술서 프린트, 목록 프린트",
     "rules": ["조직진단 수정/삭제는 진단기간 이전까지만", "적용여부 자동 표시(만료/적용중/예정)", "초과근무 실적 연동"],
     "page_list": "OrgDiagnosisPage, JobDescListPage, JobDescFormPage, JobDescApprovalPage, JobDescAdminPage, StandardWorkTimePage"},
]


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    """사각형 도형 추가"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=12, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT):
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def add_bullet_text(text_frame, items, font_size=10, color=DARK_GRAY, level=0):
    """불릿 리스트 텍스트 추가"""
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.level = level
        p.space_before = Pt(2)
        p.space_after = Pt(2)


def create_cover_slide(prs, sys_info):
    """슬라이드 1: 표지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃

    # 상단 군청색 배경
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(3.2), fill_color=NAVY)

    # 시스템 코드
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.6),
                 sys_info["code"], font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT)

    # 시스템명
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(8.4), Inches(1.2),
                 sys_info["name"], font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

    # 분석/설계 자료
    add_text_box(slide, Inches(0.8), Inches(2.1), Inches(8.4), Inches(0.6),
                 "분석/설계 자료", font_size=16, bold=False, color=RGBColor(0xCC, 0xDD, 0xFF), alignment=PP_ALIGN.LEFT)

    # 하단 정보 영역
    # 프로세스 수
    add_shape(slide, Inches(0.8), Inches(3.8), Inches(2.5), Inches(1.4), fill_color=LIGHT_BLUE)
    add_text_box(slide, Inches(0.8), Inches(3.9), Inches(2.5), Inches(0.4),
                 "단위 프로세스", font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(4.3), Inches(2.5), Inches(0.8),
                 f"{sys_info['processes']}개", font_size=28, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)

    # 복잡도
    add_shape(slide, Inches(3.7), Inches(3.8), Inches(2.5), Inches(1.4), fill_color=LIGHT_BLUE)
    add_text_box(slide, Inches(3.7), Inches(3.9), Inches(2.5), Inches(0.4),
                 "복잡도", font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.7), Inches(4.3), Inches(2.5), Inches(0.8),
                 sys_info["complexity"], font_size=28, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)

    # 구현 페이지
    add_shape(slide, Inches(6.6), Inches(3.8), Inches(2.5), Inches(1.4), fill_color=LIGHT_BLUE)
    add_text_box(slide, Inches(6.6), Inches(3.9), Inches(2.5), Inches(0.4),
                 "구현 페이지", font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(6.6), Inches(4.3), Inches(2.5), Inches(0.8),
                 f"{sys_info['pages']}개", font_size=28, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)

    # 키워드
    add_text_box(slide, Inches(0.8), Inches(5.5), Inches(8.4), Inches(0.4),
                 f"핵심 키워드: {sys_info['keywords']}", font_size=11, color=GRAY, alignment=PP_ALIGN.LEFT)

    # 하단 푸터
    add_shape(slide, Inches(0), Inches(6.8), Inches(10), Inches(0.7), fill_color=DARK_NAVY)
    add_text_box(slide, Inches(0.8), Inches(6.9), Inches(8.4), Inches(0.4),
                 "해병대 행정포탈 시스템 | 2026-04-07", font_size=9, color=RGBColor(0x99, 0xAA, 0xCC), alignment=PP_ALIGN.CENTER)


def create_overview_slide(prs, sys_info):
    """슬라이드 2: 업무 개요"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(0.8), fill_color=NAVY)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.5),
                 f"{sys_info['code']} {sys_info['name']} - 업무 개요", font_size=18, bold=True, color=WHITE)

    # 업무 개요
    add_text_box(slide, Inches(0.5), Inches(1.0), Inches(4.2), Inches(0.4),
                 "업무 개요", font_size=13, bold=True, color=NAVY)
    add_shape(slide, Inches(0.5), Inches(1.4), Inches(4.2), Inches(0.03), fill_color=NAVY)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.2), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for line in sys_info["overview"].split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)

    # 주요 사용자
    add_text_box(slide, Inches(0.5), Inches(3.2), Inches(4.2), Inches(0.4),
                 "주요 사용자", font_size=13, bold=True, color=NAVY)
    add_shape(slide, Inches(0.5), Inches(3.6), Inches(4.2), Inches(0.03), fill_color=NAVY)
    add_text_box(slide, Inches(0.5), Inches(3.7), Inches(4.2), Inches(0.6),
                 sys_info["users"], font_size=10, color=DARK_GRAY)

    # 핵심 업무 흐름
    add_text_box(slide, Inches(5.2), Inches(1.0), Inches(4.3), Inches(0.4),
                 "핵심 업무 흐름", font_size=13, bold=True, color=NAVY)
    add_shape(slide, Inches(5.2), Inches(1.4), Inches(4.3), Inches(0.03), fill_color=NAVY)

    txBox = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(4.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    flows = sys_info["flows"]
    for i, flow in enumerate(flows):
        p = tf.add_paragraph()
        p.text = f"{i+1}. {flow}"
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(3)
        p.space_after = Pt(3)

    # 푸터
    add_shape(slide, Inches(0), Inches(7.0), Inches(10), Inches(0.5), fill_color=LIGHT_GRAY)
    add_text_box(slide, Inches(0.5), Inches(7.05), Inches(9), Inches(0.3),
                 f"해병대 행정포탈 | {sys_info['code']} {sys_info['name']} | 2/5", font_size=8, color=GRAY, alignment=PP_ALIGN.RIGHT)


def create_features_slide(prs, sys_info):
    """슬라이드 3: 주요 기능"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(0.8), fill_color=NAVY)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.5),
                 f"{sys_info['code']} {sys_info['name']} - 주요 기능", font_size=18, bold=True, color=WHITE)

    # 테이블 생성
    menus = sys_info["menus"]
    rows = len(menus) + 1
    cols = 2
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.1), Inches(9), Inches(0))
    table = table_shape.table

    # 컬럼 너비
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(6.5)

    # 헤더 행
    for col_idx, header in enumerate(["대메뉴", "주요 기능"]):
        cell = table.cell(0, col_idx)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

    # 데이터 행
    for row_idx, (menu, desc) in enumerate(menus, 1):
        cell0 = table.cell(row_idx, 0)
        cell0.text = menu
        for p in cell0.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = DARK_NAVY
            p.alignment = PP_ALIGN.CENTER

        cell1 = table.cell(row_idx, 1)
        cell1.text = desc
        for p in cell1.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_GRAY

        # 행 배경색 (짝수행)
        if row_idx % 2 == 0:
            cell0.fill.solid()
            cell0.fill.fore_color.rgb = LIGHT_BLUE
            cell1.fill.solid()
            cell1.fill.fore_color.rgb = LIGHT_BLUE

    # 구현 페이지 목록
    y_pos = 1.1 + (rows * 0.4) + 0.3
    if y_pos < 5.5:
        add_text_box(slide, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.3),
                     "구현 페이지", font_size=11, bold=True, color=NAVY)
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.3), Inches(9), Inches(0.8),
                     sys_info["page_list"], font_size=9, color=GRAY)

    # 푸터
    add_shape(slide, Inches(0), Inches(7.0), Inches(10), Inches(0.5), fill_color=LIGHT_GRAY)
    add_text_box(slide, Inches(0.5), Inches(7.05), Inches(9), Inches(0.3),
                 f"해병대 행정포탈 | {sys_info['code']} {sys_info['name']} | 3/5", font_size=8, color=GRAY, alignment=PP_ALIGN.RIGHT)


def create_io_rules_slide(prs, sys_info):
    """슬라이드 4: 입출력 데이터 및 규칙"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(0.8), fill_color=NAVY)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.5),
                 f"{sys_info['code']} {sys_info['name']} - 입출력 데이터 및 규칙", font_size=18, bold=True, color=WHITE)

    # 입력 데이터
    add_shape(slide, Inches(0.5), Inches(1.1), Inches(4.2), Inches(2.5), fill_color=LIGHT_BLUE, line_color=BORDER_GRAY)
    add_text_box(slide, Inches(0.7), Inches(1.2), Inches(3.8), Inches(0.3),
                 "입력 데이터 (Input)", font_size=12, bold=True, color=NAVY)
    add_text_box(slide, Inches(0.7), Inches(1.55), Inches(3.8), Inches(2.0),
                 sys_info["input"], font_size=9, color=DARK_GRAY)

    # 출력 데이터
    add_shape(slide, Inches(5.3), Inches(1.1), Inches(4.2), Inches(2.5), fill_color=LIGHT_BLUE, line_color=BORDER_GRAY)
    add_text_box(slide, Inches(5.5), Inches(1.2), Inches(3.8), Inches(0.3),
                 "출력 데이터 (Output)", font_size=12, bold=True, color=NAVY)
    add_text_box(slide, Inches(5.5), Inches(1.55), Inches(3.8), Inches(2.0),
                 sys_info["output"], font_size=9, color=DARK_GRAY)

    # 규칙 / 예외사항
    add_text_box(slide, Inches(0.5), Inches(3.9), Inches(9), Inches(0.3),
                 "규칙 / 예외사항", font_size=13, bold=True, color=NAVY)
    add_shape(slide, Inches(0.5), Inches(4.25), Inches(9), Inches(0.03), fill_color=NAVY)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(9), Inches(2.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, rule in enumerate(sys_info["rules"]):
        p = tf.add_paragraph()
        p.text = f"  {rule}"
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    # 푸터
    add_shape(slide, Inches(0), Inches(7.0), Inches(10), Inches(0.5), fill_color=LIGHT_GRAY)
    add_text_box(slide, Inches(0.5), Inches(7.05), Inches(9), Inches(0.3),
                 f"해병대 행정포탈 | {sys_info['code']} {sys_info['name']} | 4/5", font_size=8, color=GRAY, alignment=PP_ALIGN.RIGHT)


def create_architecture_slide(prs, sys_info):
    """슬라이드 5: 시스템 구조"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(0.8), fill_color=NAVY)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.5),
                 f"{sys_info['code']} {sys_info['name']} - 시스템 구조", font_size=18, bold=True, color=WHITE)

    # 기술 스택
    add_text_box(slide, Inches(0.5), Inches(1.0), Inches(4.2), Inches(0.3),
                 "기술 스택", font_size=12, bold=True, color=NAVY)
    add_shape(slide, Inches(0.5), Inches(1.3), Inches(4.2), Inches(0.03), fill_color=NAVY)

    tech_items = [
        "React 18 + TypeScript 5 + Vite 5",
        "Ant Design 5 (ProTable, ProForm, ProLayout)",
        "Zustand 5 (Client State) + TanStack Query 5 (Server State)",
        "React Router 7 (SPA Routing)",
        "MSW 2 (Mock API) + Faker.js 9",
        "Tailwind CSS 3",
    ]
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.2), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_bullet_text(tf, tech_items, font_size=9, color=DARK_GRAY)

    # 공통 컴포넌트
    add_text_box(slide, Inches(5.3), Inches(1.0), Inches(4.2), Inches(0.3),
                 "공통 컴포넌트", font_size=12, bold=True, color=NAVY)
    add_shape(slide, Inches(5.3), Inches(1.3), Inches(4.2), Inches(0.03), fill_color=NAVY)

    comp_items = [
        "DataTable: 군청색 라인 테이블 (상단 2px, 하단 1px)",
        "SearchForm: 100px 높이 검색영역 컨테이너",
        "CrudForm: 등록/수정 폼 (text/select/date/file 등)",
        "militaryPersonColumn: 군번/계급/성명 3항목 표시",
        "StatusBadge: 상태 뱃지 (색상/라벨 매핑)",
        "SubsystemHomePage: 서브시스템 메인 화면",
        "AdminRoutes: 공통기능 관리자 라우트",
    ]
    txBox = slide.shapes.add_textbox(Inches(5.3), Inches(1.4), Inches(4.2), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_bullet_text(tf, comp_items, font_size=9, color=DARK_GRAY)

    # 아키텍처 패턴
    add_text_box(slide, Inches(0.5), Inches(3.6), Inches(9), Inches(0.3),
                 "아키텍처 패턴", font_size=12, bold=True, color=NAVY)
    add_shape(slide, Inches(0.5), Inches(3.9), Inches(9), Inches(0.03), fill_color=NAVY)

    # 아키텍처 다이어그램 (박스 형태)
    boxes = [
        (Inches(0.5), Inches(4.2), Inches(2.0), Inches(0.8), "사용자 (브라우저)", LIGHT_BLUE),
        (Inches(3.0), Inches(4.2), Inches(2.0), Inches(0.8), f"React SPA\n/{sys_info['code'].lower()}/*", LIGHT_BLUE),
        (Inches(5.5), Inches(4.2), Inches(2.0), Inches(0.8), "MSW\n(Mock API)", LIGHT_BLUE),
        (Inches(8.0), Inches(4.2), Inches(1.5), Inches(0.8), "향후\nBackend", RGBColor(0xFF, 0xF0, 0xE0)),
    ]
    for left, top, width, height, text, color in boxes:
        shape = add_shape(slide, left, top, width, height, fill_color=color, line_color=NAVY)
        shape.text_frame.word_wrap = True
        p = shape.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(9)
        p.font.color.rgb = DARK_NAVY
        p.alignment = PP_ALIGN.CENTER

    # 화살표 대신 텍스트로 흐름 표시
    add_text_box(slide, Inches(2.3), Inches(4.4), Inches(0.7), Inches(0.3), "-->", font_size=14, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(4.8), Inches(4.4), Inches(0.7), Inches(0.3), "-->", font_size=14, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.3), Inches(4.4), Inches(0.7), Inches(0.3), "-->", font_size=14, bold=True, color=GRAY, alignment=PP_ALIGN.CENTER)

    # 라우트 구조
    add_text_box(slide, Inches(0.5), Inches(5.3), Inches(9), Inches(0.3),
                 "라우트 구조", font_size=12, bold=True, color=NAVY)
    add_shape(slide, Inches(0.5), Inches(5.6), Inches(9), Inches(0.03), fill_color=NAVY)

    sys_code_lower = sys_info["code"].lower().replace("sys", "sys")
    route_text = f"/ (메인 포탈) --> /{sys_code_lower}/* (서브시스템 레이아웃) --> /{sys_code_lower}/N/N (기능 페이지) | /{sys_code_lower}/admin/* (관리자)"
    add_text_box(slide, Inches(0.5), Inches(5.7), Inches(9), Inches(0.5),
                 route_text, font_size=9, color=DARK_GRAY)

    # 총 규모
    add_shape(slide, Inches(0.5), Inches(6.3), Inches(9), Inches(0.5), fill_color=LIGHT_BLUE, line_color=BORDER_GRAY)
    add_text_box(slide, Inches(0.7), Inches(6.35), Inches(8.6), Inches(0.3),
                 f"전체 시스템: 18개 서브시스템 | 845개 프로세스 | 213개 페이지 | {sys_info['code']}: {sys_info['processes']}개 프로세스 | {sys_info['pages']}개 페이지",
                 font_size=9, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)

    # 푸터
    add_shape(slide, Inches(0), Inches(7.0), Inches(10), Inches(0.5), fill_color=LIGHT_GRAY)
    add_text_box(slide, Inches(0.5), Inches(7.05), Inches(9), Inches(0.3),
                 f"해병대 행정포탈 | {sys_info['code']} {sys_info['name']} | 5/5", font_size=8, color=GRAY, alignment=PP_ALIGN.RIGHT)


def generate_ppt(sys_info, output_dir):
    """서브시스템별 PPT 생성"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    create_cover_slide(prs, sys_info)
    create_overview_slide(prs, sys_info)
    create_features_slide(prs, sys_info)
    create_io_rules_slide(prs, sys_info)
    create_architecture_slide(prs, sys_info)

    filename = f"{sys_info['code']}_{sys_info['name']}_분석설계.pptx"
    filepath = os.path.join(output_dir, filename)
    prs.save(filepath)
    return filepath


def main():
    output_dir = os.path.join(os.path.dirname(__file__), "ppt")
    os.makedirs(output_dir, exist_ok=True)

    print(f"PPT 출력 폴더: {output_dir}")
    print(f"총 {len(SUBSYSTEMS)}개 서브시스템 PPT 생성 시작...\n")

    for sys_info in SUBSYSTEMS:
        filepath = generate_ppt(sys_info, output_dir)
        print(f"  [완료] {sys_info['code']} {sys_info['name']} -> {os.path.basename(filepath)}")

    print(f"\n총 {len(SUBSYSTEMS)}개 PPT 생성 완료! (각 5페이지, 총 {len(SUBSYSTEMS)*5}페이지)")
    print(f"저장 위치: {output_dir}")


if __name__ == "__main__":
    main()
