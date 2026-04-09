# -*- coding: utf-8 -*-
"""
SYS01 초과근무관리체계 - 플로우차트 / 스토리보드 PPT 생성
 - 1장: 플로우차트 (sample_flowchart.png 스타일)
 - 2장: 스토리보드 (sample_storyboard.png 스타일)
 - 폰트: KoPub돋움체 10/12/16/20pt
 - 투명 배경, 복잡한 질감/색상 없음
 - 하나의 오브젝트 내 설명 가능 시 단일 오브젝트로 구성
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import math
import os

# ── 색상 (투명 배경 스타일) ──
NAVY      = RGBColor(0x00, 0x33, 0x66)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY     = RGBColor(0x33, 0x33, 0x33)
MGRAY     = RGBColor(0x66, 0x66, 0x66)
LGRAY     = RGBColor(0xF2, 0xF2, 0xF2)
BDR       = RGBColor(0xBB, 0xBB, 0xBB)
ORANGE    = RGBColor(0xF5, 0xA6, 0x23)
LBLUE     = RGBColor(0xE8, 0xF0, 0xFE)
LRED      = RGBColor(0xFD, 0xED, 0xED)
RED_BD    = RGBColor(0xE7, 0x4C, 0x3C)
RED_T     = RGBColor(0xCC, 0x33, 0x33)
GREEN_T   = RGBColor(0x27, 0xAE, 0x60)
CONN_GRAY = RGBColor(0xDD, 0xDD, 0xDD)

FONT = 'KoPub돋움체 Medium'
I = Inches


# ════════════════════════════════════════
#  헬퍼 함수
# ════════════════════════════════════════

def sf(r, sz, b=False, c=DGRAY):
    """폰트 설정"""
    r.font.name = FONT
    r.font.size = Pt(sz)
    r.font.bold = b
    r.font.color.rgb = c


def _anchor_ctr(tf):
    try:
        tf._txBody.bodyPr.attrib['anchor'] = 'ctr'
    except Exception:
        pass


def _make(slide, typ, l, t, w, h, fill=None, bd=None, bw=1):
    """도형 생성 공통"""
    s = slide.shapes.add_shape(typ, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if bd:
        s.line.fill.solid()
        s.line.color.rgb = bd
        s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    return s


def rct(s, l, t, w, h, f=None, bd=None, bw=1):
    return _make(s, MSO_SHAPE.RECTANGLE, l, t, w, h, f, bd, bw)


def rnd(s, l, t, w, h, f=None, bd=None):
    return _make(s, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h, f, bd)


def ovl(s, l, t, w, h, f=None, bd=None):
    return _make(s, MSO_SHAPE.OVAL, l, t, w, h, f, bd)


def dia(s, l, t, w, h):
    return _make(s, MSO_SHAPE.FLOWCHART_DECISION, l, t, w, h, WHITE, BDR)


def stxt(sh, txt, sz=10, b=False, c=DGRAY, a=PP_ALIGN.CENTER):
    """도형 내부 텍스트 (단일 오브젝트 유지)"""
    tf = sh.text_frame
    tf.word_wrap = True
    _anchor_ctr(tf)
    p = tf.paragraphs[0]
    p.alignment = a
    r = p.add_run()
    r.text = txt
    sf(r, sz, b, c)


def tb(s, l, t, w, h, txt, sz=10, b=False, c=DGRAY, a=PP_ALIGN.LEFT):
    """독립 텍스트 박스"""
    box = s.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.alignment = a
    r = p.add_run()
    r.text = txt
    sf(r, sz, b, c)
    return box


# ── 연결선 / 화살표 ──

def vln(s, cx, y1, y2, c=DGRAY):
    w = Pt(1.5)
    rct(s, cx - w // 2, y1, w, y2 - y1, f=c)


def hln(s, x1, x2, cy, c=DGRAY):
    h = Pt(1.5)
    rct(s, x1, cy - h // 2, x2 - x1, h, f=c)


def adn(s, cx, y, c=DGRAY, sz=I(0.1)):
    sh = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, cx - sz // 2, y, sz, sz)
    sh.rotation = 180.0
    sh.fill.solid()
    sh.fill.fore_color.rgb = c
    sh.line.fill.background()


def art(s, x, cy, c=DGRAY, sz=I(0.1)):
    sh = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, x, cy - sz // 2, sz, sz)
    sh.rotation = 90.0
    sh.fill.solid()
    sh.fill.fore_color.rgb = c
    sh.line.fill.background()


def cd(s, cx, y1, y2, c=DGRAY):
    """수직 연결 (선 + 화살표)"""
    vln(s, cx, y1, y2 - I(0.08), c)
    adn(s, cx, y2 - I(0.1), c)


def cr(s, x1, x2, cy, c=DGRAY):
    """수평 연결 (선 + 화살표)"""
    hln(s, x1, x2 - I(0.08), cy, c)
    art(s, x2 - I(0.1), cy, c)


def rarr(s, x, cy, w=I(0.3), h=I(0.15)):
    """우측 화살표 도형"""
    sh = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, cy - h // 2, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = NAVY
    sh.line.fill.background()


def fline(slide, x1, y1, x2, y2, c=DGRAY, w=Pt(1), dash=False):
    """임의 각도 직선 (freeform 또는 connector)"""
    s = None
    try:
        ff = slide.shapes.build_freeform(x1, y1)
        ff.add_line_segments([(x2, y2)])
        s = ff.convert_to_shape()
        s.fill.background()
    except (AttributeError, TypeError):
        pass
    if s is None:
        try:
            s = slide.shapes.add_connector(1, x1, y1, x2, y2)
        except (AttributeError, TypeError):
            return None
    s.line.color.rgb = c
    s.line.width = w
    if dash:
        try:
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            s.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        except (ImportError, AttributeError):
            pass
    return s


# ── 공통 슬라이드 요소 (각각 단일 오브젝트) ──

def title_bar(slide, title):
    """타이틀 바 (단일 오브젝트: 네이비 사각형 + 내부 텍스트)"""
    s = rct(slide, I(0), I(0), I(13.33), I(0.5), f=NAVY)
    tf = s.text_frame
    tf.margin_left = Pt(30)
    _anchor_ctr(tf)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"SYS01 초과근무관리체계 \u2014 {title}"
    sf(r, 20, True, WHITE)


def footer(slide, pg, total):
    """푸터 (단일 오브젝트)"""
    s = rct(slide, I(0), I(7.15), I(13.33), I(0.35), f=LGRAY)
    tf = s.text_frame
    tf.margin_right = Pt(20)
    _anchor_ctr(tf)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"해병대 행정포탈  |  SYS01 초과근무관리체계  |  {pg}/{total}  |  2026-04-09"
    sf(r, 10, False, MGRAY)


def annot(slide, x, y, w_inch, title, items):
    """
    주석 박스 — 단일 오브젝트 (둥근 사각형 + 멀티 패러그래프 텍스트)
    sample_flowchart.png의 좌/우 설명 박스 스타일
    """
    lh = 0.155
    h = 0.15 + (1 + len(items)) * lh + 0.08

    s = rnd(slide, x, y, I(w_inch), I(h), bd=BDR)
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(4)

    # 제목 (첫 번째 패러그래프)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = f"[{title}]"
    sf(r, 10, True, NAVY)
    p.space_after = Pt(3)

    # 항목들 (후속 패러그래프)
    for item in items:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"\u2022 {item}"
        sf(r, 10, False, DGRAY)
        p.space_before = Pt(1)
        p.space_after = Pt(1)

    return h


def section(slide, x, y, w, h, header, items):
    """
    스토리보드 섹션 — 2개 오브젝트 (콘텐츠 사각형 + 오렌지 헤더)
    sample_storyboard.png의 메뉴 박스 스타일
    items: [(번호, 텍스트, [서브항목...]), ...]
    """
    # 오브젝트 1: 테두리 사각형 + 내부 아이템 텍스트 (단일 오브젝트)
    s = rct(slide, x, y, w, h, bd=BDR)
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(4)
    tf.margin_top = I(0.30)   # 오렌지 헤더 아래로 밀기
    tf.margin_bottom = Pt(4)

    first = True
    for num, text, subs in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        r = p.add_run()
        r.text = f"{num}  {text}"
        sf(r, 10, True, DGRAY)
        p.space_before = Pt(2)
        p.space_after = Pt(0)

        for sub in subs:
            p = tf.add_paragraph()
            r = p.add_run()
            r.text = f"   \u2022 {sub}"
            sf(r, 10, False, MGRAY)
            p.space_before = Pt(0)
            p.space_after = Pt(0)

    # 오브젝트 2: 오렌지 헤더 (단일 오브젝트: 사각형 + 내부 텍스트)
    hs = rct(slide, x, y, w, I(0.26), f=ORANGE)
    stxt(hs, header, 12, True, WHITE)


# ════════════════════════════════════════
#  슬라이드 1: 플로우차트
# ════════════════════════════════════════

def create_flowchart(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(slide, "플로우차트 (Flowchart)")
    tb(slide, I(0.4), I(0.55), I(12.5), I(0.25),
       "초과근무 신청 \u2192 결재 \u2192 당직확인 \u2192 월말결산 프로세스 흐름도", 16, False, MGRAY)
    footer(slide, 1, 4)

    # ── 치수 ──
    CX = I(6.3)                          # 중앙 흐름 X
    PW, PH = I(2.2), I(0.28)            # 프로세스 박스
    DW, DH = I(1.8), I(0.45)            # 다이아몬드
    TW, TH = I(1.4), I(0.28)            # 터미널

    def cl(w):
        return CX - w // 2

    # Y 좌표
    y = [I(v) for v in [
        0.90,   # 0  시작
        1.30,   # 1  접속
        1.75,   # 2  신청구분? (diamond)
        2.45,   # 3  신청서 작성
        2.90,   # 4  유효성 검증 (diamond)
        3.50,   # 5  결재 상신
        3.93,   # 6  중간결재 (diamond)
        4.55,   # 7  최종결재 (diamond)
        5.15,   # 8  당직확인
        5.55,   # 9  시간확정
        5.95,   # 10 월말결산
        6.35,   # 11 종료
    ]]

    # ── 0: 시작 (터미널) ──
    s = ovl(slide, cl(TW), y[0], TW, TH, f=DGRAY)
    stxt(s, "시작", 12, True, WHITE)
    cd(slide, CX, y[0] + TH, y[1])

    # ── 1: 접속 ──
    s = rct(slide, cl(PW), y[1], PW, PH, bd=BDR)
    stxt(s, "초과근무 신청서 관리 접속", 10)
    cd(slide, CX, y[1] + PH, y[2])

    # ── 2: 신청구분? (다이아몬드) ──
    d = dia(slide, cl(DW), y[2], DW, DH)
    stxt(d, "신청구분?", 10, True)

    dcy = y[2] + DH // 2
    dlx = cl(DW)
    drx = cl(DW) + DW

    # 사전신청 라벨 (단일 오브젝트)
    lbw, lbh = I(1.6), I(0.35)
    lbx = dlx - lbw - I(0.15)
    lby = dcy - lbh // 2
    s = rnd(slide, lbx, lby, lbw, lbh, f=LBLUE, bd=BDR)
    stxt(s, "사전신청\n(근무일 전 작성)", 10)
    hln(slide, lbx + lbw, dlx, dcy, MGRAY)

    # 사후신청 라벨 (단일 오브젝트)
    rbx = drx + I(0.15)
    s = rnd(slide, rbx, lby, lbw, lbh, f=LBLUE, bd=BDR)
    stxt(s, "사후신청\n(사유 입력 필수)", 10)
    hln(slide, drx, rbx, dcy, MGRAY)

    cd(slide, CX, y[2] + DH, y[3])

    # ── 3: 신청서 작성/제출 ──
    s = rct(slide, cl(PW), y[3], PW, PH, bd=BDR)
    stxt(s, "신청서 작성 / 제출", 10)
    cd(slide, CX, y[3] + PH, y[4])

    # ── 4: 유효성 검증 (다이아몬드) ──
    d = dia(slide, cl(DW), y[4], DW, DH)
    stxt(d, "유효성\n검증", 10, True)

    d4cy = y[4] + DH // 2
    d4rx = cl(DW) + DW

    # No → 시스템 알럿 (단일 오브젝트)
    ax = d4rx + I(0.3)
    aw, ah = I(1.6), I(0.3)
    s = rnd(slide, ax, d4cy - ah // 2, aw, ah, f=LRED, bd=RED_BD)
    stxt(s, "[시스템 알럿]\n입력 정보 확인 요청", 10, False, RED_T)
    cr(slide, d4rx, ax, d4cy, MGRAY)
    tb(slide, d4rx + I(0.02), d4cy - I(0.22), I(0.25), I(0.14), "No", 10, True, RED_T)
    tb(slide, CX + I(0.1), y[4] + DH - I(0.02), I(0.3), I(0.14), "Yes", 10, True, GREEN_T)

    cd(slide, CX, y[4] + DH, y[5])

    # ── 5: 결재 상신 ──
    s = rct(slide, cl(PW), y[5], PW, PH, bd=BDR)
    stxt(s, "결재 상신", 10)
    cd(slide, CX, y[5] + PH, y[6])

    # ── 6: 중간결재자 (다이아몬드) ──
    d = dia(slide, cl(DW), y[6], DW, DH)
    stxt(d, "중간결재자\n검토", 10, True)

    d6cy = y[6] + DH // 2
    d6rx = cl(DW) + DW
    rjx = d6rx + I(0.3)
    rjw, rjh = I(1.5), I(0.25)

    # 반려 (단일 오브젝트)
    s = rnd(slide, rjx, d6cy - rjh // 2, rjw, rjh, f=LRED, bd=RED_BD)
    stxt(s, "반려사유 \u2192 신청자 통보", 10, False, RED_T)
    cr(slide, d6rx, rjx, d6cy, MGRAY)
    tb(slide, d6rx + I(0.02), d6cy - I(0.22), I(0.3), I(0.14), "반려", 10, False, RED_T)
    tb(slide, CX + I(0.1), y[6] + DH - I(0.02), I(0.3), I(0.14), "승인", 10, False, GREEN_T)

    cd(slide, CX, y[6] + DH, y[7])

    # ── 7: 최종결재자 (다이아몬드) ──
    d = dia(slide, cl(DW), y[7], DW, DH)
    stxt(d, "최종결재자\n검토", 10, True)

    d7cy = y[7] + DH // 2
    d7rx = cl(DW) + DW

    # 반려 (단일 오브젝트)
    s = rnd(slide, rjx, d7cy - rjh // 2, rjw, rjh, f=LRED, bd=RED_BD)
    stxt(s, "반려사유 \u2192 신청자 통보", 10, False, RED_T)
    cr(slide, d7rx, rjx, d7cy, MGRAY)
    tb(slide, d7rx + I(0.02), d7cy - I(0.22), I(0.3), I(0.14), "반려", 10, False, RED_T)
    tb(slide, CX + I(0.1), y[7] + DH - I(0.02), I(0.3), I(0.14), "승인", 10, False, GREEN_T)

    cd(slide, CX, y[7] + DH, y[8])

    # ── 8: 당직확인 ──
    s = rct(slide, cl(PW), y[8], PW, PH, bd=BDR)
    stxt(s, "당직자 확인 / 출퇴근 기록", 10)
    cd(slide, CX, y[8] + PH, y[9])

    # ── 9: 시간확정 ──
    s = rct(slide, cl(PW), y[9], PW, PH, bd=BDR)
    stxt(s, "초과근무 시간 확정 (자동계산)", 10)
    cd(slide, CX, y[9] + PH, y[10])

    # ── 10: 월말결산 ──
    s = rct(slide, cl(PW), y[10], PW, PH, bd=BDR)
    stxt(s, "월말결산 집계 반영", 10)
    cd(slide, CX, y[10] + PH, y[11])

    # ── 11: 종료 (터미널) ──
    s = ovl(slide, cl(TW), y[11], TW, TH, f=DGRAY)
    stxt(s, "종료", 12, True, WHITE)

    # ════════════════════════════════
    #  좌측 주석 (각각 단일 오브젝트)
    # ════════════════════════════════

    LX = I(0.2)
    LW = 3.2  # inches

    # 주석 1: 신청서 입력 항목
    annot(slide, LX, I(1.15), LW, "신청서 입력 항목", [
        "신청구분: 사전신청 / 사후신청",
        "초과근무 종류: 평일 / 공휴일 / 훈련",
        "근무일자: YYYY-MM-DD (달력 선택)",
        "시작시간 / 종료시간 (30분 단위)",
        "초과근무시간: 자동계산 (종료-시작)",
        "초과근무 내용 (텍스트 입력)",
        "당직개소 선택 (드롭다운)",
        "신청자 IP 자동 수집",
    ])
    hln(slide, LX + I(LW), cl(PW), y[3] + PH // 2, CONN_GRAY)

    # 주석 2: 당직업무 세부
    annot(slide, LX, I(4.90), LW, "당직업무 세부", [
        "당직개소별 초과근무자 실시간 관리",
        "출근 / 퇴근 시간 자동 기록",
        "당직개소 변경 시 승인 프로세스 필요",
        "IP 기반 위치 확인 / 이상접속 차단",
    ])
    hln(slide, LX + I(LW), cl(PW), y[8] + PH // 2, CONN_GRAY)

    # ════════════════════════════════
    #  우측 주석 (각각 단일 오브젝트)
    # ════════════════════════════════

    RX = I(9.5)
    RW = 3.6  # inches

    # 주석 3: 유효성 검증 규칙
    annot(slide, RX, I(2.65), RW, "유효성 검증 규칙", [
        "필수값 미입력 시 알럿 표시",
        "최대인정시간 초과 시 경고 표시",
        "사후신청: 사유 입력 필수 검증",
        "공휴일/훈련 별도 시간계산 적용",
        "IP 중복 접속 체크",
    ])

    # 주석 4: 결재 규칙
    annot(slide, RX, I(3.92), RW, "결재 규칙", [
        "중간결재 \u2192 최종결재 순차 진행",
        "반려 시 사유 입력 필수",
        "신청자 실시간 알림 (결재 현황)",
        "일괄결재 지원 (복수 건 동시 처리)",
        "대리결재 설정 가능",
    ])

    # 주석 5: 월말결산 프로세스
    annot(slide, RX, I(5.35), RW, "월말결산 프로세스", [
        "월 단위 자동 마감 처리",
        "부대별 / 개인별 통계 집계",
        "결재 후 확정 (수정 불가)",
        "Excel / PDF 자료 출력 지원",
    ])
    hln(slide, cl(PW) + PW, RX, y[10] + PH // 2, CONN_GRAY)


# ════════════════════════════════════════
#  슬라이드 2: 스토리보드
# ════════════════════════════════════════

def create_storyboard(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(slide, "스토리보드 (Storyboard)")
    tb(slide, I(0.4), I(0.55), I(12.5), I(0.25),
       "7개 메뉴, 27개 하위기능, 99개 단위 프로세스", 16, False, MGRAY)
    footer(slide, 2, 4)

    # ── Entry Flow (각 박스 = 단일 오브젝트) ──
    ey = I(0.88)
    eh = I(0.33)
    ecy = ey + eh // 2

    # 메인 포탈
    e1x, e1w = I(3.0), I(1.6)
    s = rnd(slide, e1x, ey, e1w, eh, bd=BDR)
    stxt(s, "메인 포탈", 12, True, DGRAY)
    rarr(slide, e1x + e1w + I(0.05), ecy, I(0.3), I(0.15))

    # 로그인 인증
    e2x, e2w = I(5.0), I(1.6)
    s = rnd(slide, e2x, ey, e2w, eh, bd=BDR)
    stxt(s, "로그인 인증", 12, True, DGRAY)
    rarr(slide, e2x + e2w + I(0.05), ecy, I(0.3), I(0.15))

    # SYS01
    e3x, e3w = I(7.0), I(3.3)
    s = rnd(slide, e3x, ey, e3w, eh, f=NAVY)
    stxt(s, "SYS01 초과근무관리체계", 12, True, WHITE)

    # ── Entry → Row 1 연결선 ──
    e3cx = e3x + e3w // 2
    tree1_y = I(1.32)
    vln(slide, e3cx, ey + eh, tree1_y)

    # Row 1 배치
    r1y = I(1.55)
    r1h = I(2.85)
    r1w = I(3.0)
    r1gap = I(0.25)
    r1xs = [I(0.25) + i * (r1w + r1gap) for i in range(4)]
    r1cxs = [x + r1w // 2 for x in r1xs]

    # 트리 연결 (수평 + 수직 화살표)
    hln(slide, r1cxs[0], r1cxs[3], tree1_y)
    for cx in r1cxs:
        vln(slide, cx, tree1_y, r1y)
        adn(slide, cx, r1y - I(0.1))

    # ── Row 1: 4개 핵심 메뉴 (각 섹션 = 2 오브젝트) ──

    section(slide, r1xs[0], r1y, r1w, r1h, "신청서 관리", [
        (1, "신청서 작성", [
            "사전/사후 신청 구분 선택",
            "초과근무 종류 (평일/공휴일/훈련)",
            "근무일자, 시작/종료시간 입력",
            "당직개소 선택, IP 자동수집",
        ]),
        (2, "신청서 결재", [
            "중간결재 \u2192 최종결재 순차",
            "승인/반려/보류 처리",
        ]),
        (3, "일괄처리", ["복수 인원 일괄 신청 등록"]),
        (4, "일괄처리 승인", ["일괄 건 결재 처리"]),
        (5, "월말결산", [
            "월 단위 초과근무 마감",
            "부대별/개인별 집계 확정",
        ]),
    ])

    section(slide, r1xs[1], r1y, r1w, r1h, "현황조회", [
        (1, "나의 근무현황", ["월별/기간별 초과근무 조회"]),
        (2, "나의 부재관리", ["부재 사유/기간 등록관리"]),
        (3, "부대 근무 현황", ["부대원 전체 근무 현황"]),
        (4, "부대 근무 통계", ["차트/그래프 통계 시각화"]),
        (5, "부대 부재 현황", []),
        (6, "월말결산 현황", []),
        (7, "자료 출력", ["Excel/PDF 다운로드"]),
    ])

    section(slide, r1xs[2], r1y, r1w, r1h, "부대관리", [
        (1, "부대인원 조회", ["군번/계급/성명 검색"]),
        (2, "최대인정시간", ["계급별/일별 한도 설정"]),
        (3, "근무시간 관리", ["기본 근무시간/요일별 설정"]),
        (4, "공휴일 관리", ["공휴일/보상휴무 등록"]),
        (5, "결재선 관리", ["중간/최종 결재자 지정"]),
    ])

    section(slide, r1xs[3], r1y, r1w, r1h, "당직업무", [
        (1, "초과근무자 관리", [
            "당일 근무자 출퇴근 확인",
            "실시간 인원 현황 조회",
        ]),
        (2, "당직개소 관리", ["당직개소 목록 CRUD"]),
        (3, "당직개소 변경", ["임시 변경 신청/승인"]),
        (4, "개인별 당직개소 승인", []),
        (5, "개인별 부서 이동 승인", []),
    ])

    # ── Row 1 → Row 2 연결선 ──
    r2y = I(4.75)
    r2h = I(2.15)
    r2w = I(4.0)
    r2gap = I(0.35)
    r2_start = I(0.35)
    r2xs = [r2_start + i * (r2w + r2gap) for i in range(3)]
    r2cxs = [x + r2w // 2 for x in r2xs]

    mid_y = I(4.50)
    vln(slide, I(6.665), r1y + r1h, mid_y)
    hln(slide, r2cxs[0], r2cxs[2], mid_y)
    for cx in r2cxs:
        vln(slide, cx, mid_y, r2y)
        adn(slide, cx, r2y - I(0.1))

    # ── Row 2: 3개 보조 메뉴 (각 섹션 = 2 오브젝트) ──

    section(slide, r2xs[0], r2y, r2w, r2h, "개인설정", [
        (1, "개인설정 정보", [
            "알림 설정 (결재/승인 알림)",
            "기본 근무 설정값 지정",
        ]),
        (2, "개인별 당직개소 설정", ["기본 당직개소 지정"]),
        (3, "개인별 부서 설정", ["소속 부서 변경 신청"]),
    ])

    section(slide, r2xs[1], r2y, r2w, r2h, "게시판", [
        (1, "공지사항", [
            "목록/상세 조회",
            "등록/수정/삭제 (관리자)",
        ]),
        (2, "질의응답", [
            "질문 등록 / 답변 확인",
            "관리자 답변 작성",
        ]),
    ])

    section(slide, r2xs[2], r2y, r2w, r2h, "관리자", [
        (1, "체계담당자 관리", ["시스템 관리자 지정/해제"]),
        (2, "메뉴관리", ["메뉴 구조 편집 / 권한 설정"]),
        (3, "메시지관리", ["시스템 메시지 편집"]),
        (4, "로그조회", ["사용자 활동 로그 조회"]),
        (5, "코드관리", ["공통코드 CRUD"]),
    ])


# ════════════════════════════════════════
#  슬라이드 3: 유스케이스 다이어그램
# ════════════════════════════════════════

def uc_actor(slide, cx, top_y, label):
    """UML 액터 스틱 피규어 — 반환: 몸통 중앙 y (연결점)"""
    hr = I(0.11)
    ovl(slide, cx - hr, top_y, hr * 2, hr * 2, bd=DGRAY)
    by = top_y + hr * 2
    bh = I(0.20)
    vln(slide, cx, by, by + bh, DGRAY)
    hln(slide, cx - I(0.16), cx + I(0.16), by + I(0.06), DGRAY)
    ly = by + bh
    fline(slide, cx, ly, cx - I(0.12), ly + I(0.16), DGRAY, Pt(1.5))
    fline(slide, cx, ly, cx + I(0.12), ly + I(0.16), DGRAY, Pt(1.5))
    tb(slide, cx - I(0.55), ly + I(0.18), I(1.1), I(0.18),
       label, 10, True, DGRAY, PP_ALIGN.CENTER)
    return by + bh // 2


def uc_oval(slide, cx, cy, text):
    """유스케이스 타원 (단일 오브젝트)"""
    w, h = I(1.65), I(0.36)
    s = ovl(slide, cx - w // 2, cy - h // 2, w, h, f=LBLUE, bd=BDR)
    stxt(s, text, 10, False, DGRAY)


def create_usecase(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(slide, "\uc720\uc2a4\ucf00\uc774\uc2a4 \ub2e4\uc774\uc5b4\uadf8\ub7a8 (Use Case Diagram)")
    tb(slide, I(0.4), I(0.55), I(12.5), I(0.25),
       "5\uac1c \uc561\ud130, 18\uac1c \uc720\uc2a4\ucf00\uc774\uc2a4 \u2014 \ucd08\uacfc\uadfc\ubb34 \uc2e0\uccad/\uacb0\uc7ac/\uad00\ub9ac \ud750\ub984",
       16, False, MGRAY)
    footer(slide, 3, 4)

    # ── 시스템 경계 ──
    sb = rnd(slide, I(1.90), I(1.12), I(9.50), I(3.78), bd=BDR)
    sb.line.width = Pt(1.5)
    tb(slide, I(2.05), I(1.14), I(4.0), I(0.22),
       "SYS01 \ucd08\uacfc\uadfc\ubb34\uad00\ub9ac\uccb4\uacc4", 12, True, NAVY, PP_ALIGN.LEFT)

    # ── 유스케이스 배치 (3열 x 6행) ──
    HW = I(1.65) // 2
    C1, C2, C3 = I(3.50), I(6.65), I(9.80)
    ys = [I(1.55), I(2.15), I(2.75), I(3.35), I(3.95), I(4.55)]

    # Col 1 — 신청자 영역
    uc_oval(slide, C1, ys[0], "\ucd08\uacfc\uadfc\ubb34 \uc2e0\uccad")
    uc_oval(slide, C1, ys[1], "\uc77c\uad04\ucc98\ub9ac")
    uc_oval(slide, C1, ys[2], "\ub098\uc758 \uadfc\ubb34\ud604\ud669 \uc870\ud68c")
    uc_oval(slide, C1, ys[3], "\ubd80\uc7ac\uad00\ub9ac")
    uc_oval(slide, C1, ys[4], "\uac1c\uc778\uc124\uc815")
    uc_oval(slide, C1, ys[5], "\uac8c\uc2dc\ud310(\uacf5\uc9c0/Q&A)")

    # Col 2 — 결재/현황
    uc_oval(slide, C2, ys[0], "\uc720\ud6a8\uc131 \uac80\uc99d")
    uc_oval(slide, C2, ys[1], "\uc2e0\uccad\uc11c \uacb0\uc7ac")
    uc_oval(slide, C2, ys[2], "\uc6d4\ub9d0\uacb0\uc0b0")
    uc_oval(slide, C2, ys[3], "\ubd80\ub300 \uadfc\ubb34\ud604\ud669 \uc870\ud68c")
    uc_oval(slide, C2, ys[4], "\ubd80\ub300 \uadfc\ubb34\ud1b5\uacc4 \uc870\ud68c")
    uc_oval(slide, C2, ys[5], "\uc790\ub8cc \ucd9c\ub825")

    # Col 3 — 당직/관리
    uc_oval(slide, C3, ys[0], "\ubd80\ub300\uc778\uc6d0 \uc870\ud68c")
    uc_oval(slide, C3, ys[1], "\ucd5c\ub300\uc778\uc815\uc2dc\uac04 \uc124\uc815")
    uc_oval(slide, C3, ys[2], "\uacb0\uc7ac\uc120 \uad00\ub9ac")
    uc_oval(slide, C3, ys[3], "\uc2dc\uc2a4\ud15c \uad00\ub9ac")
    uc_oval(slide, C3, ys[4], "\ucd08\uacfc\uadfc\ubb34\uc790 \uad00\ub9ac")
    uc_oval(slide, C3, ys[5], "\ub2f9\uc9c1\uac1c\uc18c \uad00\ub9ac")

    # ── 액터 ──
    a_req = uc_actor(slide, I(0.70), I(1.70), "\uc2e0\uccad\uc790")
    a_mid = uc_actor(slide, I(0.70), I(4.80), "\uc911\uac04\uacb0\uc7ac\uc790")
    a_fin = uc_actor(slide, I(12.65), I(1.10), "\ucd5c\uc885\uacb0\uc7ac\uc790")
    a_adm = uc_actor(slide, I(12.65), I(2.90), "\uad00\ub9ac\uc790")
    a_dut = uc_actor(slide, I(12.65), I(4.70), "\ub2f9\uc9c1\uc790")

    # ── 연결선 (액터 → 유스케이스) ──
    al = I(0.70) + I(0.16)
    ar = I(12.65) - I(0.16)
    L1, L2 = C1 - HW, C2 - HW
    R2, R3 = C2 + HW, C3 + HW

    for y in ys:
        fline(slide, al, a_req, L1, y, MGRAY, Pt(0.75))
    fline(slide, al, a_mid, L2, ys[1], MGRAY, Pt(0.75))
    fline(slide, ar, a_fin, R2, ys[1], MGRAY, Pt(0.75))
    fline(slide, ar, a_fin, R2, ys[2], MGRAY, Pt(0.75))
    for y in ys[:4]:
        fline(slide, ar, a_adm, R3, y, MGRAY, Pt(0.75))
    fline(slide, ar, a_dut, R3, ys[4], MGRAY, Pt(0.75))
    fline(slide, ar, a_dut, R3, ys[5], MGRAY, Pt(0.75))

    # ── 관계 (<<포함>>/<<확장>>) ──
    R1 = C1 + HW
    fline(slide, R1, ys[0], L2, ys[0], DGRAY, Pt(1), dash=True)
    tb(slide, (R1 + L2) // 2 - I(0.40), ys[0] - I(0.24), I(0.80), I(0.16),
       "<<\ud3ec\ud568>>", 10, False, DGRAY, PP_ALIGN.CENTER)

    fline(slide, C2 - I(0.25), ys[3] + I(0.18), C2 - I(0.25), ys[5] - I(0.18),
          DGRAY, Pt(1), dash=True)
    tb(slide, C2 - I(0.85), (ys[3] + ys[5]) // 2 - I(0.08), I(0.80), I(0.16),
       "<<\ud655\uc7a5>>", 10, False, DGRAY, PP_ALIGN.CENTER)

    fline(slide, C2 + I(0.25), ys[4] + I(0.18), C2 + I(0.25), ys[5] - I(0.18),
          DGRAY, Pt(1), dash=True)
    tb(slide, C2 + I(0.10), (ys[4] + ys[5]) // 2 - I(0.08), I(0.80), I(0.16),
       "<<\ud655\uc7a5>>", 10, False, DGRAY, PP_ALIGN.CENTER)


# ════════════════════════════════════════
#  슬라이드 4: 목표 모델
# ════════════════════════════════════════

def create_goal_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(slide, "\ubaa9\ud45c \ubaa8\ub378 (Goal Model)")
    tb(slide, I(0.4), I(0.55), I(12.5), I(0.25),
       "\ucd08\uacfc\uadfc\ubb34\uad00\ub9ac\uccb4\uacc4\uc758 \ubaa9\ud45c \xb7 \ud504\ub85c\uc138\uc2a4 \xb7 \ucc38\uc5ec\uc8fc\uccb4 \uad6c\uc870",
       16, False, MGRAY)
    footer(slide, 4, 4)

    # ── 목표 박스 (상단) ──
    tb(slide, I(2.4), I(1.02), I(0.9), I(0.25),
       "\ubaa9\ud45c", 12, True, NAVY, PP_ALIGN.CENTER)

    gx, gy, gw, gh = I(3.4), I(0.90), I(7.0), I(0.78)
    gs = rnd(slide, gx, gy, gw, gh, f=LBLUE, bd=NAVY)
    gs.line.width = Pt(1.5)
    tf = gs.text_frame
    tf.word_wrap = True
    _anchor_ctr(tf)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "\ucd08\uacfc\uadfc\ubb34 \uad00\ub9ac\uc758 \ud22c\uba85\uc131\uacfc \ud6a8\uc728\uc131 \ud655\ubcf4"
    sf(r, 16, True, NAVY)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "(\uc2e0\uccad-\uacb0\uc7ac-\ud655\uc778-\uacb0\uc0b0\uc758 \uc804 \uacfc\uc815 \ub514\uc9c0\ud138\ud654 \ubc0f \uc790\ub3d9\ud654)"
    sf(r2, 12, False, DGRAY)

    # ── 화살표 (목표 → 운영체계) ──
    cd(slide, gx + gw // 2, gy + gh, I(2.15))

    # ── 운영체계 컨테이너 ──
    cont = rnd(slide, I(0.50), I(2.20), I(12.33), I(4.65), f=LGRAY, bd=BDR)

    # ── 좌측: 업무 프로세스 ──
    lpx, lpy, lpw = I(0.80), I(2.40), I(4.50)
    lph = I(4.25)
    lps = rnd(slide, lpx, lpy, lpw, lph, f=WHITE, bd=BDR)
    lhd = rct(slide, lpx, lpy, lpw, I(0.30), f=NAVY)
    stxt(lhd, "\uc5c5\ubb34 \ud504\ub85c\uc138\uc2a4", 12, True, WHITE)

    steps = [
        ("\ucd08\uacfc\uadfc\ubb34 \uc2e0\uccad", "\uc2e0\uccad\uc11c \uc791\uc131 / \uad6c\ubd84 \uc120\ud0dd"),
        ("\uacb0\uc7ac \ucc98\ub9ac", "\uc911\uac04\uacb0\uc7ac \u2192 \ucd5c\uc885\uacb0\uc7ac"),
        ("\ub2f9\uc9c1 \ud655\uc778", "\ucd9c\ud1f4\uadfc \uae30\ub85d / \uc2e4\uc2dc\uac04 \uad00\ub9ac"),
        ("\uc2dc\uac04 \ud655\uc815", "\ucd08\uacfc\uadfc\ubb34\uc2dc\uac04 \uc790\ub3d9\uacc4\uc0b0"),
        ("\uc6d4\ub9d0\uacb0\uc0b0", "\ubd80\ub300\ubcc4\xb7\uac1c\uc778\ubcc4 \uc9d1\uacc4 \ud655\uc815"),
    ]

    sw, sh = I(3.8), I(0.50)
    sx = lpx + (lpw - sw) // 2
    sy0 = lpy + I(0.50)
    sgap = I(0.75)

    for i, (ttl, desc) in enumerate(steps):
        sy = sy0 + i * sgap
        s = rnd(slide, sx, sy, sw, sh, f=LBLUE, bd=BDR)
        tf = s.text_frame
        tf.word_wrap = True
        _anchor_ctr(tf)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = ttl
        sf(r, 12, True, NAVY)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = desc
        sf(r2, 10, False, MGRAY)
        if i < len(steps) - 1:
            cd(slide, sx + sw // 2, sy + sh, sy + sgap)

    # ── 우측: 주요 참여주체 ──
    rpx, rpy, rpw, rph = I(5.60), I(2.40), I(6.93), I(4.25)
    rps = rnd(slide, rpx, rpy, rpw, rph, f=WHITE, bd=BDR)
    rhd = rct(slide, rpx, rpy, rpw, I(0.30), f=NAVY)
    stxt(rhd, "\uc8fc\uc694 \ucc38\uc5ec\uc8fc\uccb4", 12, True, WHITE)

    parts = [
        "\uc2e0\uccad\uc790\n(\uad70\uc778)", "\uc911\uac04\n\uacb0\uc7ac\uc790",
        "\ucd5c\uc885\n\uacb0\uc7ac\uc790", "\ub2f9\uc9c1\uc790",
        "\uad00\ub9ac\uc790", "\uc2dc\uc2a4\ud15c\n(\uc790\ub3d9\ucc98\ub9ac)",
    ]

    rcx = rpx + rpw // 2
    rcy = rpy + I(0.30) + (rph - I(0.30)) // 2
    rr = I(1.50)
    nr = I(0.42)

    pos = []
    for i in range(6):
        a = math.pi / 2 - i * 2 * math.pi / 6
        pos.append((rcx + int(rr * math.cos(a)), rcy - int(rr * math.sin(a))))

    for i in range(6):
        j = (i + 1) % 6
        fline(slide, pos[i][0], pos[i][1], pos[j][0], pos[j][1], CONN_GRAY, Pt(1.5))

    for i, (nx, ny) in enumerate(pos):
        ns = ovl(slide, nx - nr, ny - nr, nr * 2, nr * 2, f=LBLUE, bd=BDR)
        stxt(ns, parts[i], 10, True, DGRAY)

    # ── 하단 라벨 ──
    tb(slide, I(3.5), I(6.95), I(6.33), I(0.30),
       "\ucd08\uacfc\uadfc\ubb34\uad00\ub9ac \uc6b4\uc601\uccb4\uacc4", 16, True, NAVY, PP_ALIGN.CENTER)


# ════════════════════════════════════════
#  메인
# ════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = I(13.33)
    prs.slide_height = I(7.5)

    create_flowchart(prs)
    create_storyboard(prs)
    create_usecase(prs)
    create_goal_model(prs)

    out_dir = os.path.join(os.path.dirname(__file__), 'ppt')
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, 'SYS01_초과근무관리체계_FC_SB.pptx')
    prs.save(out_path)
    print(f"생성 완료: {out_path}")
    print(f"  슬라이드 1: 플로우차트 (신청→결재→당직확인→월말결산)")
    print(f"  슬라이드 2: 스토리보드 (7개 메뉴, 27개 하위기능)")
    print(f"  슬라이드 3: 유스케이스 다이어그램 (5개 액터, 18개 유스케이스)")
    print(f"  슬라이드 4: 목표 모델 (목표·프로세스·참여주체)")


if __name__ == '__main__':
    main()
