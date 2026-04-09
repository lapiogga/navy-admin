"""
SYS01 초과근무관리체계 - 와이어프레임 / 플로우차트 / 스토리보드 PPT 생성
각 1슬라이드, 구체적이고 자세한 내용
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 색상
NAVY = RGBColor(0x00, 0x33, 0x66)
DARK_NAVY = RGBColor(0x00, 0x1F, 0x3F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
GRAY = RGBColor(0x66, 0x66, 0x66)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
ACCENT_BLUE = RGBColor(0x1A, 0x73, 0xE8)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
BORDER_GRAY = RGBColor(0xDD, 0xDD, 0xDD)
GREEN = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
RED = RGBColor(0xE7, 0x4C, 0x3C)
LIGHT_GREEN = RGBColor(0xE8, 0xF8, 0xE8)
LIGHT_ORANGE = RGBColor(0xFD, 0xF2, 0xE9)
LIGHT_RED = RGBColor(0xFD, 0xED, 0xED)
LIGHT_PURPLE = RGBColor(0xF0, 0xE8, 0xFE)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=12, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_size=9, color=DARK_GRAY, bold_first=False, spacing=2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        if bold_first and i == 0:
            p.font.bold = True
        p.space_before = Pt(spacing)
        p.space_after = Pt(spacing)
    return txBox


def add_arrow(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    return shape


def add_down_arrow(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    return shape


def add_header(slide, title):
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.7), fill_color=NAVY)
    add_text_box(slide, Inches(0.4), Inches(0.12), Inches(12), Inches(0.5),
                 f"SYS01 초과근무관리체계 - {title}", font_size=18, bold=True, color=WHITE)


def add_footer(slide, page_num):
    add_shape(slide, Inches(0), Inches(7.0), Inches(13.33), Inches(0.5), fill_color=LIGHT_GRAY)
    add_text_box(slide, Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.3),
                 f"해병대 행정포탈 | SYS01 초과근무관리체계 | {page_num}/3 | 2026-04-09",
                 font_size=8, color=GRAY, alignment=PP_ALIGN.RIGHT)


# ============================================================
# 슬라이드 1: 와이어프레임
# ============================================================
def create_wireframe_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "와이어프레임 (Wireframe)")

    # --- 좌측: 초과근무 신청서 작성 화면 ---
    x0 = Inches(0.3)
    add_text_box(slide, x0, Inches(0.8), Inches(4.1), Inches(0.3),
                 "[화면1] 초과근무 신청서 작성", font_size=11, bold=True, color=NAVY)

    # 외곽 프레임
    add_shape(slide, x0, Inches(1.1), Inches(4.1), Inches(5.7), line_color=BORDER_GRAY)

    # 상단 바
    add_shape(slide, x0, Inches(1.1), Inches(4.1), Inches(0.35), fill_color=NAVY)
    add_text_box(slide, Inches(0.4), Inches(1.12), Inches(3.5), Inches(0.3),
                 "초과근무관리 > 신청서 관리 > 신청서 작성", font_size=8, color=WHITE)

    # 검색 영역
    add_shape(slide, Inches(0.4), Inches(1.55), Inches(3.9), Inches(0.55), fill_color=LIGHT_BLUE, line_color=BORDER_GRAY)
    add_text_box(slide, Inches(0.5), Inches(1.57), Inches(3.5), Inches(0.25),
                 "신청구분: [사전신청 v]  근무일: [2026-04-09]  부대: [해병대사령부 v]", font_size=7, color=DARK_GRAY)
    add_text_box(slide, Inches(0.5), Inches(1.82), Inches(3.5), Inches(0.2),
                 "계급: [소령 v]  성명: [홍길동]       [검색]  [초기화]", font_size=7, color=DARK_GRAY)

    # 폼 영역
    form_y = Inches(2.2)
    form_fields = [
        "신청구분     [사전신청 v]        초과근무 종류  [평일 v]",
        "근무일자     [2026-04-09]        시작시간      [18:00]",
        "종료시간     [21:00]             초과근무시간   3시간 (자동계산)",
        "신청자 IP    192.168.1.100       당직개소      [본부 당직실 v]",
        "초과근무 내용 [                                          ]",
        "             [                                          ]",
    ]
    for i, field in enumerate(form_fields):
        add_text_box(slide, Inches(0.5), form_y + Inches(i * 0.28), Inches(3.8), Inches(0.25),
                     field, font_size=7, color=DARK_GRAY)

    # 버튼
    add_rounded_shape(slide, Inches(1.5), Inches(4.1), Inches(0.8), Inches(0.25), fill_color=NAVY)
    add_text_box(slide, Inches(1.5), Inches(4.1), Inches(0.8), Inches(0.25),
                 "신청", font_size=8, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_rounded_shape(slide, Inches(2.5), Inches(4.1), Inches(0.8), Inches(0.25), fill_color=BORDER_GRAY)
    add_text_box(slide, Inches(2.5), Inches(4.1), Inches(0.8), Inches(0.25),
                 "취소", font_size=8, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # 하단 테이블 (나의 신청 목록)
    add_text_box(slide, Inches(0.5), Inches(4.5), Inches(3.8), Inches(0.25),
                 "나의 신청 목록", font_size=9, bold=True, color=NAVY)
    add_shape(slide, Inches(0.4), Inches(4.75), Inches(3.9), Inches(0.03), fill_color=NAVY)

    # 테이블 헤더
    add_shape(slide, Inches(0.4), Inches(4.8), Inches(3.9), Inches(0.25), fill_color=LIGHT_BLUE)
    add_text_box(slide, Inches(0.45), Inches(4.82), Inches(3.8), Inches(0.2),
                 "No  구분    근무일     시작  종료  시간  상태     비고", font_size=7, bold=True, color=NAVY)

    # 테이블 행
    rows = [
        "1   사전    2026-04-08  18:00 21:00 3h    승인완료  -",
        "2   사후    2026-04-07  19:00 22:00 3h    결재대기  긴급",
        "3   사전    2026-04-05  18:00 20:00 2h    반려     시간초과",
    ]
    for i, row in enumerate(rows):
        bg = LIGHT_GRAY if i % 2 == 1 else None
        if bg:
            add_shape(slide, Inches(0.4), Inches(5.08 + i * 0.22), Inches(3.9), Inches(0.22), fill_color=bg)
        add_text_box(slide, Inches(0.45), Inches(5.08 + i * 0.22), Inches(3.8), Inches(0.2),
                     row, font_size=6.5, color=DARK_GRAY)

    # 페이징
    add_text_box(slide, Inches(0.4), Inches(5.75), Inches(3.9), Inches(0.2),
                 "< 1  2  3  4  5 >                      총 47건", font_size=7, color=GRAY, alignment=PP_ALIGN.CENTER)

    # 범례
    add_multi_text(slide, Inches(0.4), Inches(6.0), Inches(3.9), Inches(0.8), [
        "* 결재 이후 수정/삭제는 관리자만 가능",
        "* 최대인정시간 초과 시 경고 표시",
        "* 당직 PC MAC 주소 확인 별도 Client 필요",
    ], font_size=6.5, color=RED)

    # --- 중앙: 결재 화면 ---
    x1 = Inches(4.6)
    add_text_box(slide, x1, Inches(0.8), Inches(4.1), Inches(0.3),
                 "[화면2] 신청서 결재", font_size=11, bold=True, color=NAVY)

    add_shape(slide, x1, Inches(1.1), Inches(4.1), Inches(5.7), line_color=BORDER_GRAY)
    add_shape(slide, x1, Inches(1.1), Inches(4.1), Inches(0.35), fill_color=NAVY)
    add_text_box(slide, Inches(4.7), Inches(1.12), Inches(3.5), Inches(0.3),
                 "초과근무관리 > 신청서 관리 > 결재", font_size=8, color=WHITE)

    # 결재 대기 목록
    add_shape(slide, Inches(4.7), Inches(1.55), Inches(3.9), Inches(0.35), fill_color=LIGHT_BLUE, line_color=BORDER_GRAY)
    add_text_box(slide, Inches(4.8), Inches(1.57), Inches(3.5), Inches(0.3),
                 "결재대기: 12건  |  기간: [2026-04-01] ~ [2026-04-09]  [조회]", font_size=7, color=DARK_GRAY)

    # 테이블
    add_shape(slide, Inches(4.7), Inches(2.0), Inches(3.9), Inches(0.25), fill_color=LIGHT_BLUE)
    add_text_box(slide, Inches(4.75), Inches(2.02), Inches(3.8), Inches(0.2),
                 "[v] No  군번        계급  성명   근무일     시간  구분", font_size=7, bold=True, color=NAVY)

    approval_rows = [
        "[v] 1   25-70001  중위  김해군  04-09  3h   사전",
        "[v] 2   25-70002  대위  이해병  04-09  2h   사후",
        "[ ] 3   25-70003  소령  박진해  04-08  4h   사전",
        "[v] 4   25-70004  상사  최해군  04-08  3h   사전",
        "[ ] 5   25-70005  중사  정해병  04-07  2h   사후",
    ]
    for i, row in enumerate(approval_rows):
        bg = LIGHT_GRAY if i % 2 == 1 else None
        if bg:
            add_shape(slide, Inches(4.7), Inches(2.28 + i * 0.22), Inches(3.9), Inches(0.22), fill_color=bg)
        add_text_box(slide, Inches(4.75), Inches(2.28 + i * 0.22), Inches(3.8), Inches(0.2),
                     row, font_size=6.5, color=DARK_GRAY)

    # 결재 버튼
    add_rounded_shape(slide, Inches(4.8), Inches(3.5), Inches(1.0), Inches(0.28), fill_color=GREEN)
    add_text_box(slide, Inches(4.8), Inches(3.5), Inches(1.0), Inches(0.28),
                 "일괄 승인", font_size=8, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_rounded_shape(slide, Inches(6.0), Inches(3.5), Inches(1.0), Inches(0.28), fill_color=RED)
    add_text_box(slide, Inches(6.0), Inches(3.5), Inches(1.0), Inches(0.28),
                 "일괄 반려", font_size=8, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_rounded_shape(slide, Inches(7.2), Inches(3.5), Inches(1.0), Inches(0.28), fill_color=BORDER_GRAY)
    add_text_box(slide, Inches(7.2), Inches(3.5), Inches(1.0), Inches(0.28),
                 "엑셀저장", font_size=8, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # 상세 팝업 영역
    add_text_box(slide, Inches(4.8), Inches(4.0), Inches(3.5), Inches(0.25),
                 "선택 신청서 상세", font_size=9, bold=True, color=NAVY)
    add_shape(slide, Inches(4.7), Inches(4.25), Inches(3.9), Inches(0.03), fill_color=NAVY)

    detail_lines = [
        "신청자: 25-70001 / 중위 / 김해군",
        "소속: 해병대사령부 > 작전처",
        "근무일: 2026-04-09 (수) 18:00~21:00 (3시간)",
        "구분: 사전신청 / 평일초과근무",
        "당직개소: 본부 당직실",
        "내용: 전투준비태세 점검 관련 야간 작전 지원",
        "결재선: 김해군(신청) > 이소령(중간) > 박대령(최종)",
    ]
    for i, line in enumerate(detail_lines):
        add_text_box(slide, Inches(4.8), Inches(4.35 + i * 0.22), Inches(3.7), Inches(0.2),
                     line, font_size=7, color=DARK_GRAY)

    # 승인/반려 개별 버튼
    add_rounded_shape(slide, Inches(4.8), Inches(5.95), Inches(0.8), Inches(0.25), fill_color=GREEN)
    add_text_box(slide, Inches(4.8), Inches(5.95), Inches(0.8), Inches(0.25),
                 "승인", font_size=8, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_rounded_shape(slide, Inches(5.8), Inches(5.95), Inches(0.8), Inches(0.25), fill_color=RED)
    add_text_box(slide, Inches(5.8), Inches(5.95), Inches(0.8), Inches(0.25),
                 "반려", font_size=8, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # 반려사유 입력
    add_text_box(slide, Inches(4.8), Inches(6.3), Inches(3.5), Inches(0.2),
                 "반려사유: [____________________________________]", font_size=7, color=GRAY)

    # --- 우측: 월말결산/현황 ---
    x2 = Inches(8.9)
    add_text_box(slide, x2, Inches(0.8), Inches(4.1), Inches(0.3),
                 "[화면3] 월말결산 / 현황조회", font_size=11, bold=True, color=NAVY)

    add_shape(slide, x2, Inches(1.1), Inches(4.1), Inches(2.6), line_color=BORDER_GRAY)
    add_shape(slide, x2, Inches(1.1), Inches(4.1), Inches(0.35), fill_color=NAVY)
    add_text_box(slide, Inches(9.0), Inches(1.12), Inches(3.5), Inches(0.3),
                 "신청서 관리 > 월말결산", font_size=8, color=WHITE)

    closing_lines = [
        "대상월: [2026년 04월 v]   부대: [해병대사령부 v]",
        "",
        "결산 현황 요약:",
        "  총 신청건수: 247건  |  승인: 231건  |  반려: 16건",
        "  총 초과근무시간: 741시간  |  인당평균: 12.4시간",
        "  마감기한: 2026-04-30",
        "",
        "  [결산실행]   [마감]   [엑셀다운로드]",
    ]
    for i, line in enumerate(closing_lines):
        add_text_box(slide, Inches(9.0), Inches(1.55 + i * 0.22), Inches(3.8), Inches(0.2),
                     line, font_size=7, color=DARK_GRAY)

    # 현황 조회 영역
    add_shape(slide, x2, Inches(3.85), Inches(4.1), Inches(2.95), line_color=BORDER_GRAY)
    add_shape(slide, x2, Inches(3.85), Inches(4.1), Inches(0.35), fill_color=NAVY)
    add_text_box(slide, Inches(9.0), Inches(3.87), Inches(3.5), Inches(0.3),
                 "현황조회 > 부대현황", font_size=8, color=WHITE)

    status_lines = [
        "조회조건: [2026-04] [해병대사령부 v]  [조회]",
        "",
        "부서명      인원  신청  승인  반려  시간합계",
        "----------- ---- ---- ---- ---- --------",
        "작전처       15   42   40    2   126h",
        "군수처       12   36   34    2   102h",
        "인사처       10   28   27    1    81h",
        "정보처        8   24   23    1    69h",
        "교육처       15   35   33    2    99h",
        "",
        "합 계       60  165  157    8   477h",
    ]
    for i, line in enumerate(status_lines):
        bold = i == 0 or i == 2 or i == 10
        add_text_box(slide, Inches(9.0), Inches(4.3 + i * 0.19), Inches(3.8), Inches(0.18),
                     line, font_size=6.5, color=NAVY if bold else DARK_GRAY, bold=bold)

    add_footer(slide, 1)


# ============================================================
# 슬라이드 2: 플로우차트
# ============================================================
def create_flowchart_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "플로우차트 (Flowchart)")

    # --- Flow 1: 초과근무 신청/결재 프로세스 ---
    add_text_box(slide, Inches(0.3), Inches(0.85), Inches(8.0), Inches(0.3),
                 "Flow 1. 초과근무 신청/결재 프로세스", font_size=12, bold=True, color=NAVY)
    add_shape(slide, Inches(0.3), Inches(1.12), Inches(12.7), Inches(0.02), fill_color=NAVY)

    # 역할 레인
    roles = ["신청자", "결재자(중간)", "결재자(최종)", "시스템"]
    lane_y = Inches(1.25)
    lane_h = Inches(0.9)
    colors = [LIGHT_BLUE, LIGHT_GREEN, LIGHT_ORANGE, LIGHT_PURPLE]

    for i, (role, clr) in enumerate(zip(roles, colors)):
        y = lane_y + Inches(i * 0.95)
        add_shape(slide, Inches(0.3), y, Inches(1.3), lane_h, fill_color=clr, line_color=BORDER_GRAY)
        add_text_box(slide, Inches(0.35), y + Inches(0.25), Inches(1.2), Inches(0.4),
                     role, font_size=8, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # 신청자 행
    steps_r1 = [
        ("신청서\n작성", Inches(1.8)),
        ("구분선택\n(사전/사후)", Inches(3.1)),
        ("근무정보\n입력", Inches(4.4)),
        ("당직개소\n선택", Inches(5.7)),
        ("결재선\n확인", Inches(7.0)),
        ("신청\n제출", Inches(8.3)),
    ]
    y1 = lane_y + Inches(0.15)
    for text, x in steps_r1:
        add_rounded_shape(slide, x, y1, Inches(1.0), Inches(0.6), fill_color=WHITE, line_color=ACCENT_BLUE)
        add_text_box(slide, x, y1 + Inches(0.08), Inches(1.0), Inches(0.5),
                     text, font_size=7, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # 화살표 (신청자 행)
    for i in range(len(steps_r1) - 1):
        x_from = steps_r1[i][1] + Inches(1.0)
        x_to = steps_r1[i + 1][1]
        mid_y = y1 + Inches(0.25)
        add_arrow(slide, x_from, mid_y, x_to - x_from, Inches(0.12))

    # 결재자(중간) 행
    y2 = lane_y + Inches(0.95) + Inches(0.15)
    mid_steps = [
        ("결재대기\n알림수신", Inches(8.3)),
        ("신청서\n내용확인", Inches(9.6)),
        ("승인\n/반려", Inches(10.9)),
    ]
    for text, x in mid_steps:
        clr = GREEN if "승인" in text else WHITE
        add_rounded_shape(slide, x, y2, Inches(1.0), Inches(0.6), fill_color=clr if clr == WHITE else WHITE, line_color=GREEN if "승인" in text else ACCENT_BLUE)
        add_text_box(slide, x, y2 + Inches(0.08), Inches(1.0), Inches(0.5),
                     text, font_size=7, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    for i in range(len(mid_steps) - 1):
        x_from = mid_steps[i][1] + Inches(1.0)
        x_to = mid_steps[i + 1][1]
        add_arrow(slide, x_from, y2 + Inches(0.25), x_to - x_from, Inches(0.12))

    # 수직 화살표 (신청→중간결재)
    add_down_arrow(slide, Inches(8.65), y1 + Inches(0.6), Inches(0.12), y2 - y1 - Inches(0.6))

    # 결재자(최종) 행
    y3 = lane_y + Inches(1.9) + Inches(0.15)
    final_steps = [
        ("결재대기\n알림수신", Inches(10.9)),
        ("최종\n승인/반려", Inches(12.0)),
    ]
    for text, x in final_steps:
        add_rounded_shape(slide, x, y3, Inches(1.0), Inches(0.6), fill_color=WHITE, line_color=GREEN)
        add_text_box(slide, x, y3 + Inches(0.08), Inches(1.0), Inches(0.5),
                     text, font_size=7, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    add_down_arrow(slide, Inches(11.25), y2 + Inches(0.6), Inches(0.12), y3 - y2 - Inches(0.6))
    add_arrow(slide, Inches(11.9), y3 + Inches(0.25), Inches(0.1), Inches(0.12))

    # 시스템 행
    y4 = lane_y + Inches(2.85) + Inches(0.15)
    sys_steps = [
        ("MAC주소\n검증", Inches(5.7)),
        ("최대시간\n검증", Inches(7.0)),
        ("결재선\n자동배정", Inches(8.3)),
        ("알림발송\n(결재자)", Inches(9.6)),
        ("상태변경\n(승인/반려)", Inches(10.9)),
        ("실적반영\n+이력저장", Inches(12.0)),
    ]
    for text, x in sys_steps:
        add_rounded_shape(slide, x, y4, Inches(1.0), Inches(0.6), fill_color=WHITE, line_color=PURPLE)
        add_text_box(slide, x, y4 + Inches(0.08), Inches(1.0), Inches(0.5),
                     text, font_size=7, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    for i in range(len(sys_steps) - 1):
        x_from = sys_steps[i][1] + Inches(1.0)
        x_to = sys_steps[i + 1][1]
        add_arrow(slide, x_from, y4 + Inches(0.25), x_to - x_from, Inches(0.12))

    # --- Flow 2: 월말결산 프로세스 ---
    flow2_y = Inches(5.15)
    add_text_box(slide, Inches(0.3), flow2_y, Inches(6.0), Inches(0.3),
                 "Flow 2. 월말결산 프로세스", font_size=12, bold=True, color=NAVY)
    add_shape(slide, Inches(0.3), flow2_y + Inches(0.25), Inches(6.0), Inches(0.02), fill_color=NAVY)

    closing_steps = [
        ("마감기한\n설정", LIGHT_BLUE),
        ("신청서\n최종확인", LIGHT_GREEN),
        ("미결재건\n일괄처리", LIGHT_ORANGE),
        ("결산\n실행", LIGHT_PURPLE),
        ("엑셀\n다운로드", LIGHT_BLUE),
        ("마감\n처리", LIGHT_GREEN),
    ]
    for i, (text, clr) in enumerate(closing_steps):
        x = Inches(0.4 + i * 1.1)
        add_rounded_shape(slide, x, flow2_y + Inches(0.4), Inches(0.9), Inches(0.55), fill_color=clr, line_color=BORDER_GRAY)
        add_text_box(slide, x, flow2_y + Inches(0.47), Inches(0.9), Inches(0.45),
                     text, font_size=7, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
        if i < len(closing_steps) - 1:
            add_arrow(slide, x + Inches(0.9), flow2_y + Inches(0.58), Inches(0.2), Inches(0.1))

    # --- Flow 3: 당직업무 프로세스 ---
    add_text_box(slide, Inches(7.0), flow2_y, Inches(6.0), Inches(0.3),
                 "Flow 3. 당직업무 프로세스", font_size=12, bold=True, color=NAVY)
    add_shape(slide, Inches(7.0), flow2_y + Inches(0.25), Inches(6.0), Inches(0.02), fill_color=NAVY)

    duty_steps = [
        ("당직개소\n등록/관리", LIGHT_BLUE),
        ("당직자\n지정", LIGHT_GREEN),
        ("초과근무자\n확인", LIGHT_ORANGE),
        ("개인별\n근무승인", LIGHT_PURPLE),
        ("부서별\n근무승인", LIGHT_BLUE),
    ]
    for i, (text, clr) in enumerate(duty_steps):
        x = Inches(7.1 + i * 1.15)
        add_rounded_shape(slide, x, flow2_y + Inches(0.4), Inches(0.95), Inches(0.55), fill_color=clr, line_color=BORDER_GRAY)
        add_text_box(slide, x, flow2_y + Inches(0.47), Inches(0.95), Inches(0.45),
                     text, font_size=7, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
        if i < len(duty_steps) - 1:
            add_arrow(slide, x + Inches(0.95), flow2_y + Inches(0.58), Inches(0.2), Inches(0.1))

    add_footer(slide, 2)


# ============================================================
# 슬라이드 3: 스토리보드
# ============================================================
def create_storyboard_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "스토리보드 (Storyboard)")

    scenarios = [
        {
            "title": "시나리오 1: 평일 초과근무 사전신청",
            "actor": "신청자 (중위 김해군)",
            "color": LIGHT_BLUE,
            "steps": [
                "1. 메인포탈 로그인 > 초과근무관리 클릭",
                "2. 좌측메뉴 '신청서 관리 > 신청서 작성' 이동",
                "3. 신청구분 '사전신청' 선택",
                "4. 근무일 2026-04-09, 시작 18:00, 종료 21:00 입력",
                "5. 시스템이 초과근무시간 '3시간' 자동계산 표시",
                "6. 당직개소 '본부 당직실' 선택",
                "7. 초과근무 내용 '전투준비태세 점검 야간작전 지원' 입력",
                "8. [신청] 클릭 > 결재선 자동 배정 확인 팝업",
                "9. 결재선: 김해군(신청)>이소령(중간)>박대령(최종)",
                "10. [확인] > 신청완료 메시지 + 목록에 '결재대기' 표시",
            ],
            "rules": "최대인정시간(4h) 이내 / MAC 주소 미확인(사전신청) / 결재선 자동",
        },
        {
            "title": "시나리오 2: 결재자 승인/반려 처리",
            "actor": "결재자 (소령 이철수)",
            "color": LIGHT_GREEN,
            "steps": [
                "1. 로그인 시 알림 뱃지 '결재대기 5건' 확인",
                "2. '신청서 관리 > 결재' 메뉴 이동",
                "3. 결재대기 목록에서 기간/부대 필터 조회",
                "4. 김해군 신청서 클릭 > 하단 상세 정보 표시",
                "5. 근무내용/시간/당직개소 확인",
                "6. 체크박스 선택 후 [승인] 클릭",
                "7. 시스템: 상태 '결재대기'>'승인완료' 변경",
                "8. 박대령(최종결재자)에게 알림 자동 발송",
                "9. 반려 시: 반려사유 입력 필수 > [반려] 클릭",
                "10. 일괄 승인: 체크박스 다건 선택 > [일괄승인]",
            ],
            "rules": "결재 이후 수정/삭제 불가(관리자만) / 반려사유 필수 / 일괄처리 가능",
        },
        {
            "title": "시나리오 3: 월말결산 및 현황 조회",
            "actor": "부대관리자 (대위 박관리)",
            "color": LIGHT_ORANGE,
            "steps": [
                "1. '신청서 관리 > 월말결산' 메뉴 이동",
                "2. 대상월 '2026년 04월', 부대 '해병대사령부' 선택",
                "3. 결산현황: 총 247건, 승인 231건, 반려 16건 표시",
                "4. 총 초과근무시간 741h, 인당평균 12.4h 확인",
                "5. 미결재건 16건 확인 > 기간 내 처리 독촉",
                "6. [결산실행] 클릭 > 결산 데이터 확정",
                "7. [엑셀다운로드] > 부서별/개인별 시간 엑셀 저장",
                "8. '현황조회 > 부대현황' 이동",
                "9. 부서별 신청/승인/반려/시간합계 테이블 확인",
                "10. [마감] > 마감기한(04-30) 이후 수정 불가 처리",
            ],
            "rules": "마감기한 설정 필수 / 결산 후 엑셀 다운로드 / 마감 후 수정불가",
        },
        {
            "title": "시나리오 4: 당직업무 초과근무자 관리",
            "actor": "당직자 (상사 최당직)",
            "color": LIGHT_PURPLE,
            "steps": [
                "1. '당직업무 > 초과근무자 현황' 메뉴 이동",
                "2. 금일 초과근무 승인자 목록 자동 표시",
                "3. 출입현황 확인: 입실 12명, 퇴실 3명, 미퇴실 9명",
                "4. 개인별 시작/종료시간, 당직개소 매칭 확인",
                "5. MAC 주소 확인 (당직PC Client 연동)",
                "6. 미승인 사후신청 건 확인 > 개인별 승인 처리",
                "7. '당직업무 > 당직개소 관리' 이동",
                "8. 당직개소 등록/수정/삭제 (개소명/위치/PC-MAC)",
                "9. '당직업무 > 당직개소 변경' 이동",
                "10. 당직개소 변경 이력 조회 / 변경 사유 입력",
            ],
            "rules": "당직 PC MAC 주소 확인 별도 Client / 개인별+부서별 승인 분리",
        },
    ]

    for i, sc in enumerate(scenarios):
        col = i
        x = Inches(0.3 + col * 3.25)
        y_start = Inches(0.9)

        # 시나리오 제목
        add_shape(slide, x, y_start, Inches(3.05), Inches(0.35), fill_color=NAVY)
        add_text_box(slide, x + Inches(0.05), y_start + Inches(0.03), Inches(2.95), Inches(0.3),
                     sc["title"], font_size=9, bold=True, color=WHITE)

        # 액터
        add_shape(slide, x, y_start + Inches(0.35), Inches(3.05), Inches(0.25), fill_color=sc["color"])
        add_text_box(slide, x + Inches(0.05), y_start + Inches(0.37), Inches(2.95), Inches(0.2),
                     f"액터: {sc['actor']}", font_size=8, bold=True, color=DARK_GRAY)

        # 단계
        add_shape(slide, x, y_start + Inches(0.6), Inches(3.05), Inches(4.3), fill_color=WHITE, line_color=BORDER_GRAY)
        for j, step in enumerate(sc["steps"]):
            add_text_box(slide, x + Inches(0.08), y_start + Inches(0.65 + j * 0.4), Inches(2.9), Inches(0.38),
                         step, font_size=7, color=DARK_GRAY)

        # 규칙/제약
        add_shape(slide, x, y_start + Inches(4.9), Inches(3.05), Inches(0.7), fill_color=LIGHT_RED, line_color=RED)
        add_text_box(slide, x + Inches(0.08), y_start + Inches(4.93), Inches(2.9), Inches(0.15),
                     "규칙/제약사항", font_size=7, bold=True, color=RED)
        add_text_box(slide, x + Inches(0.08), y_start + Inches(5.1), Inches(2.9), Inches(0.45),
                     sc["rules"], font_size=6.5, color=DARK_GRAY)

    add_footer(slide, 3)


# ============================================================
# 메인 실행
# ============================================================
if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    create_wireframe_slide(prs)
    create_flowchart_slide(prs)
    create_storyboard_slide(prs)

    output_path = "spec-doc/ppt/SYS01_초과근무관리체계_와이어프레임_플로우차트_스토리보드.pptx"
    prs.save(output_path)
    print(f"PPT 생성 완료: {output_path}")
    print(f"  슬라이드 1: 와이어프레임 (신청서작성/결재/월말결산 3개 화면)")
    print(f"  슬라이드 2: 플로우차트 (신청결재/월말결산/당직업무 3개 프로세스)")
    print(f"  슬라이드 3: 스토리보드 (4개 시나리오 x 10단계)")
