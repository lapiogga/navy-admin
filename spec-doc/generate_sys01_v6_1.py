#!/usr/bin/env python3
"""
SYS01 초과근무관리체계 기능명세서 PPTX v6.1 (portrait, 4페이지)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
system_level_spec_guide_v6.1.md 기반 핵심 5페이지:
  P1. 시스템 개념도      P2. 목표모델 (As-Is → To-Be)
  P3. 유스케이스          P4. 스토리보드 (화면흐름도)
  P5. 플로우차트 (업무처리흐름)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ═══ 색상 팔레트 ═══
NAVY  = RGBColor(0x6B, 0x8E, 0xA8)
TEAL  = RGBColor(0x8E, 0xC5, 0xC0)
SKY   = RGBColor(0xA8, 0xC8, 0xE8)
MINT  = RGBColor(0xDC, 0xF5, 0xE4)
PEACH = RGBColor(0xFC, 0xD4, 0xD0)
CREAM = RGBColor(0xFD, 0xEE, 0xC2)
ICE   = RGBColor(0xE4, 0xEE, 0xF8)
SNOW  = RGBColor(0xF7, 0xF9, 0xFC)
CHAR  = RGBColor(0x1A, 0x1A, 0x1A)
ASH   = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRN   = RGBColor(0xA4, 0xD4, 0xB4)
WARN  = RGBColor(0xF5, 0xCC, 0xA0)
ERR   = RGBColor(0xF0, 0xB2, 0xAE)
LBLUE = RGBColor(0xDC, 0xEA, 0xFE)
LSUB  = RGBColor(0xC5, 0xDC, 0xEC)  # 부제 색상
HBAR  = RGBColor(0xD4, 0xE2, 0xEC)  # 제목바 연한 배경
THBG  = RGBColor(0xD4, 0xE2, 0xEC)  # 표 헤더 연한 배경
BORD  = RGBColor(0x4A, 0x6B, 0x8C)  # 박스 테두리 (진한)
TBLB  = RGBColor(0x80, 0x80, 0x80)  # 표 테두리 (회색)

# ═══ 치수 ═══
SW = Inches(7.5)
SH = Inches(10)
ML = Inches(0.3)       # 좌우 마진
CW = Inches(6.9)       # 콘텐츠 폭 = 7.5 - 0.3*2
FONT = 'Freesentation 5 Medium'
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'ppt')


# ══════════════════════════════════════════════════
# 헬퍼 함수
# ══════════════════════════════════════════════════

def _sl(prs):
    """빈 슬라이드 + Snow 배경"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = SNOW
    return s


def _run(p, txt, sz=9, b=False, c=CHAR):
    """단락에 서식 런 추가 (최소 9pt 보장)"""
    sz = max(sz, 10)
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = b
    r.font.color.rgb = c
    r.font.name = FONT
    return r


def _tb(s, l, t, w, h, txt, sz=9, b=False, c=CHAR, al=PP_ALIGN.LEFT):
    """텍스트박스"""
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    _run(tb.text_frame.paragraphs[0], txt, sz, b, c)
    tb.text_frame.paragraphs[0].alignment = al
    return tb


def _title(s, title, sub, pg):
    """[A] 제목 영역 (portrait) — 연한 바 + 진한 테두리 + 검정 글자"""
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = HBAR
    bar.line.color.rgb = BORD
    bar.line.width = Pt(0.5)
    _tb(s, ML, Inches(0.06), Inches(5.5), Inches(0.36), title, 20, True, CHAR)
    if sub:
        _tb(s, ML, Inches(0.42), Inches(6.5), Inches(0.26), sub, 10, False, CHAR)
    _tb(s, Inches(6.5), Inches(0.18), Inches(0.6), Inches(0.3),
        f'{pg}/4', 12, True, CHAR, PP_ALIGN.RIGHT)


def _box(s, l, t, w, h, txt, fill=ICE, stk=BORD, sw=0.5,
         sz=9, b=False, fc=CHAR, al=PP_ALIGN.CENTER):
    """둥근 사각형 + 멀티라인 텍스트"""
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = stk
    sh.line.width = Pt(sw)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.margin_left = tf.margin_right = Inches(0.03)
    for i, ln in enumerate(txt.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p, ln, sz, b, fc)
        p.alignment = al
        p.space_before = p.space_after = Pt(0)
    bp = sh._element.txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', 'ctr')
    return sh


def _darr(s, cx, y, h=Inches(0.08)):
    """하향 화살표"""
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                           cx - Inches(0.06), y, Inches(0.12), h)
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _rarr(s, x, cy, w=Inches(0.18)):
    """우향 화살표"""
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                           x, cy - Inches(0.05), w, Inches(0.1))
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _cell_bg(cell, hx):
    """셀 배경색"""
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for sf in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(sf)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), val=hx)


def _cell_border(cell, color_hex='808080', width_pt=0.5):
    """셀 테두리 4면 회색 설정"""
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for border in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for existing in tcPr.findall(qn(border)):
            tcPr.remove(existing)
        ln = etree.SubElement(tcPr, qn(border))
        ln.set('w', str(int(width_pt * 12700)))
        ln.set('cap', 'flat')
        ln.set('cmpd', 'sng')
        ln.set('algn', 'ctr')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr'), val=color_hex)
        etree.SubElement(ln, qn('a:prstDash'), val='solid')
        etree.SubElement(ln, qn('a:round'))


def _tbl(s, l, t, w, h, data, cws=None):
    """테이블 (멀티라인 셀 지원, 진한 테두리 0.5pt, 연한 헤더 배경). data[0]=헤더"""
    rows, cols = len(data), len(data[0])
    sh = s.shapes.add_table(rows, cols, l, t, w, h)
    tbl = sh.table
    if cws:
        for i, cw in enumerate(cws):
            tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ''
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_top = tf.margin_bottom = Inches(0.04)
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.06)
            hdr = (r == 0)
            for i, ln in enumerate(data[r][c].split('\n')):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                _run(p, ln, 9 if not hdr else 10, hdr, CHAR)
                p.alignment = PP_ALIGN.CENTER if hdr else PP_ALIGN.LEFT
                p.space_before = p.space_after = Pt(0)
            _cell_bg(cell, 'D4E2EC' if hdr else ('FFFFFF' if r % 2 == 1 else 'F2F7FB'))
            _cell_border(cell, '808080', 0.5)
    return tbl




# ══════════════════════════════════════════════════
# Page 2: 목표모델 (As-Is → To-Be)
# ══════════════════════════════════════════════════

def p2_target(prs):
    s = _sl(prs)
    _title(s, '목표모델 (As-Is → To-Be)',
           '현행 문제점 진단 및 목표 시스템 기능 구성', 1)

    pairs = [
        ('수기 초과근무 신청\n종이양식/엑셀, 누락·오기 빈번',
         '온라인 전자 신청\n웹폼+자동검증, 실시간 유효성'),
        ('대면/유선 결재\n처리시간 평균 2일, 추적 불가',
         '전자결재 워크플로우\n즉시 승인/반려, 상태 실시간'),
        ('수기 시간 계산\n월평균 15건 오류, 정산 불일치',
         '자동 시간 계산\nTimePicker diff, 오류 95%↓'),
        ('엑셀 월말결산\n취합 3일 소요, 버전관리 불가',
         '온라인 월말결산\n마감기한 자동, 원클릭 결산'),
        ('수기 당직대장\n변경 추적 불가, 개소관리 미흡',
         '당직개소 CRUD\nMAC주소, 변경 승인 워크플로우'),
    ]

    y0 = Inches(1.0)
    RH = Inches(0.78)
    RG = Inches(0.06)
    AW = Inches(2.6)     # As-Is / To-Be 폭
    ARR_W = Inches(0.22)
    AX = ML
    ARX = ML + AW + Inches(0.15)
    TX = ARX + ARR_W + Inches(0.15)

    # 섹션 라벨
    _box(s, AX, y0 - Inches(0.06), AW, Inches(0.18),
         '현행 시스템 (As-Is)', fill=ERR, stk=ERR, sz=8, b=True, fc=WHITE)
    _box(s, TX, y0 - Inches(0.06), AW, Inches(0.18),
         '목표 시스템 (To-Be)', fill=GRN, stk=GRN, sz=8, b=True, fc=WHITE)

    for i, (asis, tobe) in enumerate(pairs):
        y = y0 + i * (RH + RG)
        _box(s, AX, y, AW, RH, asis, fill=PEACH, stk=ERR, sz=7, fc=CHAR)
        _rarr(s, ARX, y + RH // 2, ARR_W)
        _box(s, TX, y, AW, RH, tobe, fill=MINT, stk=GRN, sz=7, fc=CHAR)

    # 우측 기능 요약 박스
    FX = TX + AW + Inches(0.12)
    FW = ML + CW - FX
    _box(s, FX, y0, FW, 5 * RH + 4 * RG,
         '핵심 기능\n─────────\n'
         '신청서관리  5\n현황조회     7\n부대관리     5\n'
         '당직업무     5\n개인설정     3\n게시판        2\n관리자        5\n'
         '─────────\n총 32개 메뉴\n99 프로세스',
         fill=ICE, stk=TEAL, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)

    # [C] 서술 테이블 — 개선 항목
    dy = y0 + 5 * (RH + RG) + Inches(0.12)
    _tbl(s, ML, dy, CW, Inches(1.25), [
        ['개선항목', '현행 문제점', '목표 개선방향', '기대 효과', '정량 목표'],
        ['업무처리', '종이/엑셀 수기 처리', '온라인 전자 신청·결재', '데이터 정합성·표준화', '오류 95%↓'],
        ['시간계산', '수기 계산, 월 15건 오류', 'TimePicker 자동 계산', '계산 오류 제거', '오류 0건'],
        ['결재업무', '대면/유선, 평균 2일', '전자결재 즉시 처리', '결재 시간 단축', '2일→즉시'],
        ['월말결산', '엑셀 취합, 3일 소요', '원클릭 온라인 결산', '결산 시간 단축', '3일→1시간'],
        ['당직관리', '수기 대장, 추적 불가', 'CRUD+MAC+승인 흐름', '변경 이력 관리', '이력 100%'],
    ], cws=[int(CW * r) for r in [0.10, 0.24, 0.24, 0.22, 0.10]])


# ══════════════════════════════════════════════════
# Page 3: 유스케이스
# ══════════════════════════════════════════════════

def p3_usecase(prs):
    s = _sl(prs)
    _title(s, '유스케이스 다이어그램',
           '6개 액터 × 23개 유스케이스, 99개 단위 프로세스', 2)

    y0 = Inches(0.82)

    # ── 좌측: 액터 ──
    actors = [
        '시스템관리자\nROLE_ADMIN',
        '부대장\nROLE_COMMANDER',
        '결재자\nROLE_APPROVER',
        '초과근무담당자\nROLE_WORKER',
        '당직근무자\nROLE_DUTY',
        '일반사용자\nROLE_USER',
    ]
    AX = ML
    AW = Inches(1.0)
    AH = Inches(0.66)
    AG = Inches(0.06)
    for i, atxt in enumerate(actors):
        _box(s, AX, y0 + i * (AH + AG), AW, AH, atxt,
             fill=ICE, stk=NAVY, sz=7, fc=CHAR)

    # ── 중앙: 시스템 바운더리 ──
    BX = AX + AW + Inches(0.2)
    BW = Inches(4.3)
    BH = 6 * AH + 5 * AG + Inches(0.3)
    bnd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, BX, y0 - Inches(0.05), BW, BH)
    bnd.fill.solid()
    bnd.fill.fore_color.rgb = RGBColor(0xFB, 0xFB, 0xFB)
    bnd.line.color.rgb = NAVY
    bnd.line.width = Pt(2)
    _tb(s, BX + Inches(0.1), y0 - Inches(0.02), Inches(4), Inches(0.18),
        '초과근무관리체계 (SYS01) — 바운더리', 8, True, NAVY)

    # UC 그룹 (2열 배치)
    uc_groups = [
        ('신청서관리 (5)', 'UC-OT-001~005\n작성·결재·일괄·일괄승인·월말결산', ICE, SKY),
        ('현황조회 (7)', 'UC-OT-006~010\n나의현황·부재·부대현황·통계·출력', ICE, SKY),
        ('부대관리 (5)', 'UC-OT-011~015\n인원·최대시간·근무시간·공휴일·결재선', ICE, SKY),
        ('당직업무 (5)', 'UC-OT-016~020\n근무자·개소·변경·당직승인·부서승인', ICE, SKY),
        ('개인설정 (3)', 'UC-OT-021~023\n설정정보·당직개소·부서', LBLUE, TEAL),
        ('게시판+관리자 (7)', 'UC-OT-024~025 공지·Q&A\nUC-SYS 시스템·코드·권한·결재선·게시판', LBLUE, TEAL),
    ]
    UY0 = y0 + Inches(0.22)
    UW = Inches(2.0)
    UH = Inches(0.70)
    UG = Inches(0.06)
    for i, (grp, desc, fill, stk) in enumerate(uc_groups):
        r, c = divmod(i, 2)
        ux = BX + Inches(0.08) + c * (UW + Inches(0.1))
        uy = UY0 + r * (UH + UG)
        sh = _box(s, ux, uy, UW, UH, f'{grp}\n{desc}',
                  fill=fill, stk=stk, sz=6, b=False, fc=CHAR, al=PP_ALIGN.LEFT)
        # 그룹명 볼드 처리 — 첫 번째 단락만
        sh.text_frame.paragraphs[0].runs[0].font.bold = True

    # ── 우측: 외부 시스템 (추후) ──
    EX = BX + BW + Inches(0.12)
    EW = ML + CW - EX
    _box(s, EX, y0 + Inches(0.2), EW, Inches(0.35),
         'GPKI/SSO\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)
    _box(s, EX, y0 + Inches(0.65), EW, Inches(0.35),
         '알림 서비스\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)

    # [C] 서술 테이블 — 주요 UC 상세
    dy = Inches(3.7)
    _tbl(s, ML, dy, CW, Inches(1.65), [
        ['UC-ID', '유스케이스명', '주 액터', '기능 설명 (구체적)', '우선'],
        ['UC-OT-001', '신청서 작성', '담당자',
         '초과근무 신청 등록(7항목), TimePicker 시간자동계산, 당직개소 선택, CRUD', '상'],
        ['UC-OT-002', '신청서 결재', '결재자',
         'Steps(작성→결재대기→승인완료), 탭별 목록, 승인/반려/사유입력 Modal', '상'],
        ['UC-OT-003', '일괄처리', '담당자',
         '복수 인원 동시 신청서 생성, 부대·기간·사유 일괄 지정', '상'],
        ['UC-OT-005', '월말결산', '부대장',
         '연도/월별 결산, 마감기한 DatePicker, close/cancel, isDeadlinePassed 자동', '상'],
        ['UC-OT-008', '부대현황/통계', '결재자',
         'Column/Bar 차트 + DataTable, 부대별 근무시간/인원 통계, 엑셀 다운로드', '중'],
        ['UC-OT-012', '최대인정시간', '담당자',
         '3탭(월별설정/예외처리/예외구분), 각 탭 DataTable+Modal CRUD', '중'],
        ['UC-OT-017', '당직개소 관리', '담당자',
         '개소 CRUD (명칭/위치/MAC주소/수용인원/활성상태), 다중 부대 배정', '중'],
    ], cws=[int(CW * r) for r in [0.09, 0.10, 0.07, 0.67, 0.05]])


# ══════════════════════════════════════════════════
# Page 4: 스토리보드 (화면흐름도)
# ══════════════════════════════════════════════════

def p4_storyboard(prs):
    s = _sl(prs)
    _title(s, '스토리보드 (화면흐름도)',
           '총 32개 화면, 7개 업무영역, 핵심 네비게이션 흐름', 3)

    y0 = Inches(0.82)

    # 로그인 → 대시보드
    _box(s, ML, y0, Inches(1.1), Inches(0.35),
         '로그인\n/login', fill=WARN, stk=WARN, sz=8, b=True, fc=WHITE)
    _rarr(s, ML + Inches(1.35), y0 + Inches(0.12), Inches(0.25))
    _box(s, ML + Inches(1.65), y0, Inches(2.0), Inches(0.35),
         '대시보드 / 홈\n/sys01', fill=NAVY, stk=NAVY, sz=8, b=True, fc=WHITE)

    # 화살표: 대시보드 → 아래
    _darr(s, ML + Inches(2.65), y0 + Inches(0.37))

    # 7개 메뉴 그룹 (2열 × 4행)
    groups = [
        ('1. 신청서관리', '작성 · 결재 · 일괄 · 일괄승인 · 월말결산', '/sys01/1/*', '5', ICE, SKY),
        ('2. 현황조회', '나의현황 · 부재 · 부대현황\n부대통계 · 부재현황 · 월말 · 출력', '/sys01/2/*', '7', ICE, SKY),
        ('3. 부대관리', '인원조회 · 최대시간 · 근무시간\n공휴일 · 결재선', '/sys01/3/*', '5', ICE, SKY),
        ('4. 당직업무', '근무자 · 당직개소 · 개소변경\n당직승인 · 부서이동승인', '/sys01/4/*', '5', ICE, SKY),
        ('5. 개인설정', '설정정보 · 당직개소설정 · 부서설정', '/sys01/5/*', '3', LBLUE, TEAL),
        ('6. 게시판', '공지사항 · 질의응답', '/sys01/6/*', '2', LBLUE, TEAL),
        ('7. 관리자', '시스템 · 코드 · 권한 · 결재선 · 게시판', '/sys01/admin/*', '5', LBLUE, TEAL),
    ]

    GY0 = y0 + Inches(0.55)
    GW = Inches(3.2)
    GH = Inches(0.90)
    GG = Inches(0.06)
    GCOL_GAP = Inches(0.15)

    for i, (name, screens, path, cnt, fill, stk) in enumerate(groups):
        r, c = divmod(i, 2)
        gx = ML + c * (GW + GCOL_GAP)
        gy = GY0 + r * (GH + GG)
        sh = _box(s, gx, gy, GW, GH,
                  f'{name} ({cnt}개)\n{screens}\n{path}',
                  fill=fill, stk=stk, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        # URL 라인 색상을 ASH로
        pars = list(sh.text_frame.paragraphs)
        if len(pars) >= 3:
            for r_obj in pars[-1].runs:
                r_obj.font.color.rgb = ASH
                r_obj.font.size = Pt(8)

    # [C] 서술 — 좌: 화면 목록 / 우: 핵심 흐름
    dy = GY0 + 4 * (GH + GG) + Inches(0.12)

    # 좌측: 화면 요약 테이블
    LW = int(CW * 0.48)
    _tbl(s, ML, dy, LW, Inches(1.25), [
        ['영역', '화면 수', 'URL 패턴', '주요 컴포넌트'],
        ['신청서관리', '5', '/sys01/1/*', 'SearchForm+DataTable+Modal'],
        ['현황조회', '7', '/sys01/2/*', 'Chart(Column/Bar)+DataTable'],
        ['부대관리', '5', '/sys01/3/*', 'Tabs+DataTable+Modal+Calendar'],
        ['당직업무', '5', '/sys01/4/*', 'DataTable+Descriptions+Select'],
        ['개인설정', '3', '/sys01/5/*', 'Card+Descriptions (bordered)'],
        ['게시판', '2', '/sys01/6/*', 'SimpleBoardPage (공통 재사용)'],
        ['관리자', '5', '/sys01/admin/*', 'AdminRoutes (공통기능 연결)'],
    ], cws=[int(LW * r) for r in [0.18, 0.10, 0.25, 0.47]])

    # 우측: 핵심 흐름 시나리오
    RX = ML + LW + Inches(0.15)
    RW = CW - LW - Inches(0.15)
    _tbl(s, RX, dy, RW, Inches(1.25), [
        ['핵심 흐름 시나리오'],
        ['시나리오 1: 초과근무 신청 흐름\n'
         '① 로그인 → 대시보드\n'
         '② 좌측메뉴 [신청서관리] > [신청서 작성]\n'
         '③ SearchForm 검색 → DataTable 목록\n'
         '④ [신청서 작성] 버튼 → Modal 폼 입력\n'
         '⑤ TimePicker 시간 입력 → 자동 계산\n'
         '⑥ [저장] → 성공 토스트 + 목록 갱신'],
        ['시나리오 2: 결재 흐름\n'
         '① [신청서 결재] 메뉴 → 탭별 목록\n'
         '② 행 클릭 → 상세 확인\n'
         '③ [승인]/[반려] → 확인 Dialog\n'
         '④ 상태 변경 → 알림 발송 → 목록 갱신'],
    ])


# ══════════════════════════════════════════════════
# Page 5: 플로우차트 (업무처리흐름)
# ══════════════════════════════════════════════════

def p5_flowchart(prs):
    s = _sl(prs)
    _title(s, '업무처리 플로우차트',
           '핵심 프로세스: 초과근무 신청 → 결재 → 마감 — 4개 레인, 12단계', 4)

    y0 = Inches(0.82)

    # 4개 스윔레인 배경
    lanes = [
        ('사용자 (신청자/결재자)', ICE, NAVY),
        ('프론트엔드', LBLUE, SKY),
        ('백엔드', MINT, GRN),
        ('DB / 외부', CREAM, WARN),
    ]
    LGAP = Inches(0.14)  # 레인 사이 간격 (10pt)
    LN_W = (CW - 3 * LGAP) // 4
    LN_H = Inches(4.1)   # 다이어그램 높이 1.5배
    LN_HDR = Inches(0.22)

    for i, (name, fill, stk) in enumerate(lanes):
        lx = ML + (LN_W + LGAP) * i
        # 레인 배경
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = fill
        bg.line.color.rgb = stk
        bg.line.width = Pt(1)
        # 레인 헤더
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_HDR)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = stk
        hdr.line.fill.background()
        tf = hdr.text_frame
        _run(tf.paragraphs[0], name, 7, True, WHITE)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 프로세스 단계 (레인 인덱스 0~3, y 오프셋)
    steps = [
        (0, 0.30, '① 신청서 메뉴 접근\n좌측 네비게이션'),
        (1, 0.30, '② 신청 폼 표시\nSearchForm + DataTable'),
        (0, 0.70, '③ 정보 입력\n종류/날짜/시간/개소/사유'),
        (1, 0.70, '④ 유효성 검증\n필수값·TimePicker·Select'),
        (1, 1.10, '⑤ POST /sys01/requests\nAPI 호출 (axios)'),
        (2, 1.10, '⑥ 업무규칙 검증\n상태체크·시간유효성'),
        (3, 1.10, '⑦ OtRequest INSERT\n데이터 저장'),
        (2, 1.50, '⑧ 결재 알림 발송\n결재자에게 통보'),
        (1, 1.50, '⑨ 성공 토스트\n+ 목록 갱신'),
        (0, 1.90, '⑩ 결재자: 승인/반려\nModal 확인'),
        (2, 1.90, '⑪ PUT approve/reject\n상태 변경'),
        (3, 1.90, '⑫ 이력 저장\n결재 이력 기록'),
    ]

    SBW = LN_W - Inches(0.12)  # 단계 박스 폭
    SBH = Inches(0.56)          # 단계 박스 높이

    for lane, y_off, txt in steps:
        sx = ML + (LN_W + LGAP) * lane + Inches(0.07)
        sy = y0 + LN_HDR + Inches(y_off)
        _box(s, sx, sy, SBW, SBH, txt,
             fill=WHITE, stk=ASH, sw=1, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)

    # 판단 분기 표시 (④ 검증 실패)
    fx = ML + (LN_W + LGAP) * 1 + Inches(0.07)
    fy = y0 + LN_HDR + Inches(0.95)
    _box(s, fx, fy, SBW, Inches(0.12),
         '검증 실패 → 필드 에러 표시 → 재입력',
         fill=PEACH, stk=ERR, sw=1, sz=5, fc=ERR, al=PP_ALIGN.CENTER)

    # [C] 서술 — 좌: 단계 설명 / 우: 업무규칙+예외
    dy = y0 + LN_H + Inches(0.1)

    _tbl(s, ML, dy, CW, Inches(1.6), [
        ['단계', '동작', '조건/상세'],
        ['①②', '메뉴접근 → 폼표시', 'SearchForm(3필드) + DataTable(11컬럼)'],
        ['③', '정보입력', '종류(사전/사후), 날짜, 시간, 당직개소, 사유'],
        ['④', '프론트 유효성 검증', '필수값 확인, TimePicker 범위, Select 선택'],
        ['⑤⑥⑦', 'API→규칙검증→DB저장', 'POST /sys01/requests → OtRequest INSERT'],
        ['⑧⑨', '알림→결과표시', '결재자 통보 + 성공 토스트 + invalidateQueries'],
        ['⑩⑪⑫', '결재처리→상태변경', 'PUT approve/reject + 이력 기록'],
    ], cws=[int(CW * r) for r in [0.08, 0.28, 0.64]])

    _tbl(s, ML, dy + Inches(1.6) + Inches(0.1), CW, Inches(1.6), [
        ['업무 규칙 / 예외 처리'],
        ['[업무 규칙]\n'
         'BR-OT-001: 총근무시간 = endTime-startTime (분→시)\n'
         'BR-OT-002: 상태 4단계 (draft→pending→approved|rejected)\n'
         'BR-OT-003: 일괄처리는 복수 인원 동시 신청\n'
         'BR-OT-004: 월말결산 마감 후 수정 불가 (isDeadlinePassed)\n'
         'BR-OT-005: 최대인정시간 3탭 (월별/예외/구분)\n'
         'BR-OT-006: 당직개소 MAC 주소 기반 위치확인\n'
         'BR-OT-007: 개인 당직/부서 변경 → 승인 필요\n'
         'BR-OT-008: 부재유형 4종 (휴가/휴직/출장/파견)'],
        ['[예외 처리]\n'
         'E-001: 세션만료(401) → 로그인 리다이렉트\n'
         'E-002: 권한없음(403) → 접근 제한 알림\n'
         'E-003: 검증실패 → 필드별 에러 표시+재입력\n'
         'E-004: 서버오류(500) → 에러 메시지+재시도'],
    ])


# ══════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════

def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    print('[1/4] 목표모델...')
    p2_target(prs)
    print('[2/4] 유스케이스...')
    p3_usecase(prs)
    print('[3/4] 스토리보드...')
    p4_storyboard(prs)
    print('[4/4] 플로우차트...')
    p5_flowchart(prs)

    out_path = os.path.join(OUT, 'SYS01_초과근무관리체계_분석설계_v6.1.pptx')
    prs.save(out_path)
    print(f'\n완료! 총 {len(prs.slides)}장 슬라이드')
    print(f'파일: {out_path}')


if __name__ == '__main__':
    main()
