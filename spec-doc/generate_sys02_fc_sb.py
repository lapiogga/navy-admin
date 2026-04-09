# -*- coding: utf-8 -*-
"""
SYS02 설문종합관리체계 - 프레젠테이션 4장 PPT 생성
 - 1장: 목표 모델 (sample_goal model.png 스타일)
 - 2장: 플로우차트 (sample_flowchart.png 스타일)
 - 3장: 스토리보드 (sample_storyboard.png 스타일)
 - 4장: 유스케이스 다이어그램 (sample_usecase.png 스타일)
 - 폰트: KoPub돋움체 10/12/16/20pt
 - 투명 배경, 복잡한 질감/색상 없음
 - 하나의 오브젝트 내 설명 가능 시 단일 오브젝트로 구성
 - 3D / 그림자 사용 금지
 - Line 속성: '선' 또는 '선 화살표'
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import math
import os

# ── 색상 (투명 배경 스타일) ──
NAVY      = RGBColor(0x00, 0x33, 0x66)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY     = RGBColor(0x33, 0x33, 0x33)
MGRAY     = RGBColor(0x66, 0x66, 0x66)
LGRAY     = RGBColor(0xF2, 0xF2, 0xF2)
BDR       = RGBColor(0xBB, 0xBB, 0xBB)
ORANGE    = RGBColor(0xF5, 0xA6, 0x23)
LBLUE     = RGBColor(0xE8, 0xF0, 0xFE)
LRED      = RGBColor(0xFD, 0xED, 0xED)
RED_BD    = RGBColor(0xE7, 0x4C, 0x3C)
RED_T     = RGBColor(0xCC, 0x33, 0x33)
GREEN_T   = RGBColor(0x27, 0xAE, 0x60)
CONN_GRAY = RGBColor(0xDD, 0xDD, 0xDD)

FONT = 'KoPub돋움체 Medium'
I = Inches


# ════════════════════════════════════════
#  헬퍼 함수
# ════════════════════════════════════════

def sf(r, sz, b=False, c=DGRAY):
    """폰트 설정"""
    r.font.name = FONT
    r.font.size = Pt(sz)
    r.font.bold = b
    r.font.color.rgb = c


def _anchor_ctr(tf):
    try:
        tf._txBody.bodyPr.attrib['anchor'] = 'ctr'
    except Exception:
        pass


def _make(slide, typ, l, t, w, h, fill=None, bd=None, bw=1):
    """도형 생성 공통"""
    s = slide.shapes.add_shape(typ, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if bd:
        s.line.fill.solid()
        s.line.color.rgb = bd
        s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    return s


def rct(s, l, t, w, h, f=None, bd=None, bw=1):
    return _make(s, MSO_SHAPE.RECTANGLE, l, t, w, h, f, bd, bw)


def rnd(s, l, t, w, h, f=None, bd=None):
    return _make(s, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h, f, bd)


def ovl(s, l, t, w, h, f=None, bd=None):
    return _make(s, MSO_SHAPE.OVAL, l, t, w, h, f, bd)


def dia(s, l, t, w, h):
    return _make(s, MSO_SHAPE.FLOWCHART_DECISION, l, t, w, h, WHITE, BDR)


def stxt(sh, txt, sz=10, b=False, c=DGRAY, a=PP_ALIGN.CENTER):
    """도형 내부 텍스트 (단일 오브젝트 유지)"""
    tf = sh.text_frame
    tf.word_wrap = True
    _anchor_ctr(tf)
    p = tf.paragraphs[0]
    p.alignment = a
    r = p.add_run()
    r.text = txt
    sf(r, sz, b, c)


def tb(s, l, t, w, h, txt, sz=10, b=False, c=DGRAY, a=PP_ALIGN.LEFT):
    """독립 텍스트 박스"""
    box = s.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.alignment = a
    r = p.add_run()
    r.text = txt
    sf(r, sz, b, c)
    return box


# ── 연결선 / 화살표 ──

def vln(s, cx, y1, y2, c=DGRAY):
    w = Pt(1.5)
    rct(s, cx - w // 2, y1, w, y2 - y1, f=c)


def hln(s, x1, x2, cy, c=DGRAY):
    h = Pt(1.5)
    rct(s, x1, cy - h // 2, x2 - x1, h, f=c)


def adn(s, cx, y, c=DGRAY, sz=I(0.1)):
    sh = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, cx - sz // 2, y, sz, sz)
    sh.rotation = 180.0
    sh.fill.solid()
    sh.fill.fore_color.rgb = c
    sh.line.fill.background()


def art(s, x, cy, c=DGRAY, sz=I(0.1)):
    sh = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, x, cy - sz // 2, sz, sz)
    sh.rotation = 90.0
    sh.fill.solid()
    sh.fill.fore_color.rgb = c
    sh.line.fill.background()


def cd(s, cx, y1, y2, c=DGRAY):
    """수직 연결 (선 + 화살표)"""
    vln(s, cx, y1, y2 - I(0.08), c)
    adn(s, cx, y2 - I(0.1), c)


def cr(s, x1, x2, cy, c=DGRAY):
    """수평 연결 (선 + 화살표)"""
    hln(s, x1, x2 - I(0.08), cy, c)
    art(s, x2 - I(0.1), cy, c)


def rarr(s, x, cy, w=I(0.3), h=I(0.15)):
    """우측 화살표 도형"""
    sh = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, cy - h // 2, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = NAVY
    sh.line.fill.background()


def fline(slide, x1, y1, x2, y2, c=DGRAY, w=Pt(1), dash=False):
    """임의 각도 직선 (freeform 또는 connector)"""
    s = None
    try:
        ff = slide.shapes.build_freeform(x1, y1)
        ff.add_line_segments([(x2, y2)])
        s = ff.convert_to_shape()
        s.fill.background()
    except (AttributeError, TypeError):
        pass
    if s is None:
        try:
            s = slide.shapes.add_connector(1, x1, y1, x2, y2)
        except (AttributeError, TypeError):
            return None
    s.line.color.rgb = c
    s.line.width = w
    if dash:
        try:
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            s.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        except (ImportError, AttributeError):
            pass
    return s


# ── 공통 슬라이드 요소 ──

def title_bar(slide, title):
    """타이틀 바 (단일 오브젝트: 네이비 사각형 + 내부 텍스트)"""
    s = rct(slide, I(0), I(0), I(13.33), I(0.5), f=NAVY)
    tf = s.text_frame
    tf.margin_left = Pt(30)
    _anchor_ctr(tf)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"SYS02 설문종합관리체계 \u2014 {title}"
    sf(r, 20, True, WHITE)


def footer(slide, pg, total):
    """푸터 (단일 오브젝트)"""
    s = rct(slide, I(0), I(7.15), I(13.33), I(0.35), f=LGRAY)
    tf = s.text_frame
    tf.margin_right = Pt(20)
    _anchor_ctr(tf)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"해병대 행정포탈  |  SYS02 설문종합관리체계  |  {pg}/{total}  |  2026-04-09"
    sf(r, 10, False, MGRAY)


def annot(slide, x, y, w_inch, title, items):
    """주석 박스 — 단일 오브젝트 (둥근 사각형 + 멀티 패러그래프 텍스트)"""
    lh = 0.155
    h = 0.15 + (1 + len(items)) * lh + 0.08

    s = rnd(slide, x, y, I(w_inch), I(h), bd=BDR)
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(4)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = f"[{title}]"
    sf(r, 10, True, NAVY)
    p.space_after = Pt(3)

    for item in items:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"\u2022 {item}"
        sf(r, 10, False, DGRAY)
        p.space_before = Pt(1)
        p.space_after = Pt(1)

    return h


def section(slide, x, y, w, h, header, items):
    """스토리보드 섹션 — 2개 오브젝트 (콘텐츠 사각형 + 오렌지 헤더)"""
    s = rct(slide, x, y, w, h, bd=BDR)
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(4)
    tf.margin_top = I(0.30)
    tf.margin_bottom = Pt(4)

    first = True
    for num, text, subs in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        r = p.add_run()
        r.text = f"{num}  {text}"
        sf(r, 10, True, DGRAY)
        p.space_before = Pt(2)
        p.space_after = Pt(0)

        for sub in subs:
            p = tf.add_paragraph()
            r = p.add_run()
            r.text = f"   \u2022 {sub}"
            sf(r, 10, False, MGRAY)
            p.space_before = Pt(0)
            p.space_after = Pt(0)

    hs = rct(slide, x, y, w, I(0.26), f=ORANGE)
    stxt(hs, header, 12, True, WHITE)


# ════════════════════════════════════════
#  슬라이드 1: 목표 모델
# ════════════════════════════════════════

def create_goal_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(slide, "목표 모델 (Goal Model)")
    tb(slide, I(0.4), I(0.55), I(12.5), I(0.25),
       "설문종합관리체계의 목표 \xb7 프로세스 \xb7 참여주체 구조",
       16, False, MGRAY)
    footer(slide, 1, 4)

    # ── 목표 박스 (상단) ──
    tb(slide, I(2.4), I(1.02), I(0.9), I(0.25),
       "목표", 12, True, NAVY, PP_ALIGN.CENTER)

    gx, gy, gw, gh = I(3.4), I(0.90), I(7.0), I(0.78)
    gs = rnd(slide, gx, gy, gw, gh, f=LBLUE, bd=NAVY)
    gs.line.width = Pt(1.5)
    tf = gs.text_frame
    tf.word_wrap = True
    _anchor_ctr(tf)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "설문 운영의 체계화와 응답 데이터 활용 극대화"
    sf(r, 16, True, NAVY)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "(설문 생성·승인·배포·응답·분석의 전 과정 디지털화 및 자동화)"
    sf(r2, 12, False, DGRAY)

    # ── 화살표 (목표 → 운영체계) ──
    cd(slide, gx + gw // 2, gy + gh, I(2.15))

    # ── 운영체계 컨테이너 ──
    cont = rnd(slide, I(0.50), I(2.20), I(12.33), I(4.65), f=LGRAY, bd=BDR)

    # ── 좌측: 업무 프로세스 ──
    lpx, lpy, lpw = I(0.80), I(2.40), I(4.50)
    lph = I(4.25)
    lps = rnd(slide, lpx, lpy, lpw, lph, f=WHITE, bd=BDR)
    lhd = rct(slide, lpx, lpy, lpw, I(0.30), f=NAVY)
    stxt(lhd, "업무 프로세스", 12, True, WHITE)

    steps = [
        ("설문 생성", "설문명·기간·대상자·문항 작성"),
        ("승인 처리", "관리자 검토 → 승인/반려"),
        ("설문 배포", "대상자에게 설문 공개 (active)"),
        ("응답 수집", "문항별 응답 제출 · 참여 관리"),
        ("결과 분석", "통계·차트 시각화 · 엑셀 출력"),
    ]

    sw, sh = I(3.8), I(0.50)
    sx = lpx + (lpw - sw) // 2
    sy0 = lpy + I(0.50)
    sgap = I(0.75)

    for i, (ttl, desc) in enumerate(steps):
        sy = sy0 + i * sgap
        s = rnd(slide, sx, sy, sw, sh, f=LBLUE, bd=BDR)
        tf = s.text_frame
        tf.word_wrap = True
        _anchor_ctr(tf)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = ttl
        sf(r, 12, True, NAVY)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = desc
        sf(r2, 10, False, MGRAY)
        if i < len(steps) - 1:
            cd(slide, sx + sw // 2, sy + sh, sy + sgap)

    # ── 우측: 주요 참여주체 ──
    rpx, rpy, rpw, rph = I(5.60), I(2.40), I(6.93), I(4.25)
    rps = rnd(slide, rpx, rpy, rpw, rph, f=WHITE, bd=BDR)
    rhd = rct(slide, rpx, rpy, rpw, I(0.30), f=NAVY)
    stxt(rhd, "주요 참여주체", 12, True, WHITE)

    parts = [
        "설문\n작성자", "승인\n관리자",
        "응답자\n(군인)", "대상자\n관리",
        "카테고리\n/템플릿", "시스템\n(자동처리)",
    ]

    rcx = rpx + rpw // 2
    rcy = rpy + I(0.30) + (rph - I(0.30)) // 2
    rr = I(1.50)
    nr = I(0.42)

    pos = []
    for i in range(6):
        a = math.pi / 2 - i * 2 * math.pi / 6
        pos.append((rcx + int(rr * math.cos(a)), rcy - int(rr * math.sin(a))))

    for i in range(6):
        j = (i + 1) % 6
        fline(slide, pos[i][0], pos[i][1], pos[j][0], pos[j][1], CONN_GRAY, Pt(1.5))

    for i, (nx, ny) in enumerate(pos):
        ns = ovl(slide, nx - nr, ny - nr, nr * 2, nr * 2, f=LBLUE, bd=BDR)
        stxt(ns, parts[i], 10, True, DGRAY)

    # ── 하단 라벨 ──
    tb(slide, I(3.5), I(6.95), I(6.33), I(0.30),
       "설문종합관리 운영체계", 16, True, NAVY, PP_ALIGN.CENTER)


# ════════════════════════════════════════
#  슬라이드 2: 플로우차트
# ════════════════════════════════════════

def create_flowchart(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(slide, "플로우차트 (Flowchart)")
    tb(slide, I(0.4), I(0.55), I(12.5), I(0.25),
       "설문 생성 → 승인 → 배포 → 응답 → 결과분석 프로세스 흐름도", 16, False, MGRAY)
    footer(slide, 2, 4)

    # ── 치수 ──
    CX = I(6.3)
    PW, PH = I(2.2), I(0.28)
    DW, DH = I(1.8), I(0.45)
    TW, TH = I(1.4), I(0.28)

    def cl(w):
        return CX - w // 2

    # Y 좌표 (12 노드)
    y = [I(v) for v in [
        0.90,   # 0  시작
        1.30,   # 1  설문 생성
        1.75,   # 2  문항 편집
        2.20,   # 3  승인 요청
        2.65,   # 4  관리자 검토 (diamond)
        3.30,   # 5  설문 배포
        3.75,   # 6  응답 수집
        4.20,   # 7  응답 완료 확인 (diamond)
        4.85,   # 8  설문 마감
        5.30,   # 9  결과 분석
        5.75,   # 10 엑셀 출력
        6.15,   # 11 종료
    ]]

    # ── 0: 시작 ──
    s = ovl(slide, cl(TW), y[0], TW, TH, f=DGRAY)
    stxt(s, "시작", 12, True, WHITE)
    cd(slide, CX, y[0] + TH, y[1])

    # ── 1: 설문 생성 ──
    s = rct(slide, cl(PW), y[1], PW, PH, bd=BDR)
    stxt(s, "설문 생성 (기본정보 입력)", 10)
    cd(slide, CX, y[1] + PH, y[2])

    # ── 2: 문항 편집 ──
    s = rct(slide, cl(PW), y[2], PW, PH, bd=BDR)
    stxt(s, "문항 편집 (추가/수정/순서변경)", 10)
    cd(slide, CX, y[2] + PH, y[3])

    # ── 3: 승인 요청 ──
    s = rct(slide, cl(PW), y[3], PW, PH, bd=BDR)
    stxt(s, "승인 요청 (draft → submitted)", 10)
    cd(slide, CX, y[3] + PH, y[4])

    # ── 4: 관리자 검토 (다이아몬드) ──
    d = dia(slide, cl(DW), y[4], DW, DH)
    stxt(d, "관리자\n검토", 10, True)

    d4cy = y[4] + DH // 2
    d4rx = cl(DW) + DW

    # 반려 → 설문작성자 통보
    rjx = d4rx + I(0.3)
    rjw, rjh = I(1.6), I(0.28)
    s = rnd(slide, rjx, d4cy - rjh // 2, rjw, rjh, f=LRED, bd=RED_BD)
    stxt(s, "반려사유 → 작성자 통보", 10, False, RED_T)
    cr(slide, d4rx, rjx, d4cy, MGRAY)
    tb(slide, d4rx + I(0.02), d4cy - I(0.22), I(0.3), I(0.14), "반려", 10, False, RED_T)
    tb(slide, CX + I(0.1), y[4] + DH - I(0.02), I(0.3), I(0.14), "승인", 10, False, GREEN_T)

    cd(slide, CX, y[4] + DH, y[5])

    # ── 5: 설문 배포 ──
    s = rct(slide, cl(PW), y[5], PW, PH, bd=BDR)
    stxt(s, "설문 배포 (approved → active)", 10)
    cd(slide, CX, y[5] + PH, y[6])

    # ── 6: 응답 수집 ──
    s = rct(slide, cl(PW), y[6], PW, PH, bd=BDR)
    stxt(s, "응답자 설문 응답 수집", 10)
    cd(slide, CX, y[6] + PH, y[7])

    # ── 7: 응답 완료 확인 (다이아몬드) ──
    d = dia(slide, cl(DW), y[7], DW, DH)
    stxt(d, "마감조건\n충족?", 10, True)

    d7cy = y[7] + DH // 2
    d7rx = cl(DW) + DW

    # No → 응답 수집 계속
    nx = d7rx + I(0.3)
    nw, nh = I(1.5), I(0.25)
    s = rnd(slide, nx, d7cy - nh // 2, nw, nh, f=LBLUE, bd=BDR)
    stxt(s, "응답 기간 연장 / 미응답 알림", 10, False, DGRAY)
    cr(slide, d7rx, nx, d7cy, MGRAY)
    tb(slide, d7rx + I(0.02), d7cy - I(0.22), I(0.25), I(0.14), "No", 10, True, RED_T)
    tb(slide, CX + I(0.1), y[7] + DH - I(0.02), I(0.3), I(0.14), "Yes", 10, True, GREEN_T)

    cd(slide, CX, y[7] + DH, y[8])

    # ── 8: 설문 마감 ──
    s = rct(slide, cl(PW), y[8], PW, PH, bd=BDR)
    stxt(s, "설문 마감 (active → closed)", 10)
    cd(slide, CX, y[8] + PH, y[9])

    # ── 9: 결과 분석 ──
    s = rct(slide, cl(PW), y[9], PW, PH, bd=BDR)
    stxt(s, "결과 분석 (통계/차트 생성)", 10)
    cd(slide, CX, y[9] + PH, y[10])

    # ── 10: 엑셀 출력 ──
    s = rct(slide, cl(PW), y[10], PW, PH, bd=BDR)
    stxt(s, "결과 내보내기 (엑셀 다운로드)", 10)
    cd(slide, CX, y[10] + PH, y[11])

    # ── 11: 종료 ──
    s = ovl(slide, cl(TW), y[11], TW, TH, f=DGRAY)
    stxt(s, "종료", 12, True, WHITE)

    # ════════════════════════════════
    #  좌측 주석
    # ════════════════════════════════

    LX = I(0.2)
    LW = 3.2

    # 주석 1: 설문 기본정보 입력항목
    annot(slide, LX, I(0.95), LW, "설문 기본정보 입력항목", [
        "설문명 (텍스트, 100자)",
        "설문내용/개요 (TextArea)",
        "시작일 / 종료일 (DatePicker)",
        "대상 구분: 부대/직급/전체",
        "대상 부대·직급·직책·성별 선택",
        "결과 공개 여부 (Switch)",
        "익명 여부 (Switch)",
        "첨부: 승인공문, 설문지, 보안검토",
    ])
    hln(slide, LX + I(LW), cl(PW), y[1] + PH // 2, CONN_GRAY)

    # 주석 2: 문항 편집 세부
    annot(slide, LX, I(2.60), LW, "문항 편집 세부", [
        "문항 유형: 단일선택/복수선택/주관식/평점",
        "문항 내용 입력 (500자)",
        "선택지 추가/삭제 (radio/checkbox)",
        "필수 응답 여부 설정 (Switch)",
        "드래그앤드롭 순서 변경",
        "엑셀 양식 업로드/다운로드",
    ])
    hln(slide, LX + I(LW), cl(PW), y[2] + PH // 2, CONN_GRAY)

    # ════════════════════════════════
    #  우측 주석
    # ════════════════════════════════

    RX = I(9.5)
    RW = 3.6

    # 주석 3: 상태 전이 규칙
    annot(slide, RX, I(1.65), RW, "설문 상태 전이", [
        "draft(작성중) → submitted(제출됨)",
        "submitted → approved(승인) / rejected(반려)",
        "approved → active(진행중)",
        "active → closed(마감)",
        "반려 시 반려사유 입력 필수",
    ])

    # 주석 4: 응답 처리 규칙
    annot(slide, RX, I(3.30), RW, "응답 처리 규칙", [
        "active 상태 설문만 응답 가능",
        "필수 문항 미입력 시 제출 차단",
        "익명 설문: 응답자 정보 비표시",
        "1인 1회 응답 원칙 (중복 방지)",
        "임시 저장 기능 지원",
    ])

    # 주석 5: 결과 분석 기능
    annot(slide, RX, I(5.00), RW, "결과 분석 기능", [
        "선택형: Bar 차트 + 비율(%) 표시",
        "주관식: 텍스트 목록 (페이지네이션)",
        "평점형: 평균 평점 + 분포 차트",
        "응답률 통계: 대상자/응답자/응답률",
        "엑셀(CSV) 다운로드 지원",
    ])
    hln(slide, cl(PW) + PW, RX, y[9] + PH // 2, CONN_GRAY)


# ════════════════════════════════════════
#  슬라이드 3: 스토리보드
# ════════════════════════════════════════

def create_storyboard(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(slide, "스토리보드 (Storyboard)")
    tb(slide, I(0.4), I(0.55), I(12.5), I(0.25),
       "5개 메뉴, 6개 관리탭, 31개 단위 프로세스", 16, False, MGRAY)
    footer(slide, 3, 4)

    # ── Entry Flow ──
    ey = I(0.88)
    eh = I(0.33)
    ecy = ey + eh // 2

    e1x, e1w = I(3.0), I(1.6)
    s = rnd(slide, e1x, ey, e1w, eh, bd=BDR)
    stxt(s, "메인 포탈", 12, True, DGRAY)
    rarr(slide, e1x + e1w + I(0.05), ecy, I(0.3), I(0.15))

    e2x, e2w = I(5.0), I(1.6)
    s = rnd(slide, e2x, ey, e2w, eh, bd=BDR)
    stxt(s, "로그인 인증", 12, True, DGRAY)
    rarr(slide, e2x + e2w + I(0.05), ecy, I(0.3), I(0.15))

    e3x, e3w = I(7.0), I(3.3)
    s = rnd(slide, e3x, ey, e3w, eh, f=NAVY)
    stxt(s, "SYS02 설문종합관리체계", 12, True, WHITE)

    # ── Entry → Row 1 연결선 ──
    e3cx = e3x + e3w // 2
    tree1_y = I(1.32)
    vln(slide, e3cx, ey + eh, tree1_y)

    # Row 1 배치 — 3개 핵심 메뉴 (일반 사용자)
    r1y = I(1.55)
    r1h = I(2.65)
    r1w = I(3.85)
    r1gap = I(0.30)
    r1_start = I(0.60)
    r1xs = [r1_start + i * (r1w + r1gap) for i in range(3)]
    r1cxs = [x + r1w // 2 for x in r1xs]

    hln(slide, r1cxs[0], r1cxs[2], tree1_y)
    for cx in r1cxs:
        vln(slide, cx, tree1_y, r1y)
        adn(slide, cx, r1y - I(0.1))

    # ── Row 1: 3개 메뉴 ──

    section(slide, r1xs[0], r1y, r1w, r1h, "나의 설문관리", [
        (1, "설문 등록/수정", [
            "설문명, 기간, 대상자, 익명여부 입력",
            "첨부: 승인공문, 설문지, 보안검토",
        ]),
        (2, "문항 편집기", [
            "단일선택/복수선택/주관식/평점",
            "선택지 추가, 필수 여부, 순서 변경",
            "엑셀 양식 업로드/다운로드",
        ]),
        (3, "상태 관리", [
            "승인 요청/배포/마감 처리",
            "상태별 설문 목록 조회",
        ]),
    ])

    section(slide, r1xs[1], r1y, r1w, r1h, "설문참여", [
        (1, "진행중 설문 조회", [
            "active 상태 설문 목록",
            "설문명, 등록자, 기간, 참여 여부",
        ]),
        (2, "설문 응답 폼", [
            "문항별 입력 위젯 자동 렌더링",
            "필수 문항 유효성 검사",
            "제출 / 임시 저장",
        ]),
        (3, "응답 완료", [
            "성공 메시지 표시",
            "목록으로 돌아가기",
        ]),
    ])

    section(slide, r1xs[2], r1y, r1w, r1h, "지난 설문보기", [
        (1, "마감 설문 조회", [
            "closed 상태 설문 목록",
            "대상자 수, 응답자 수 표시",
        ]),
        (2, "결과 분석", [
            "응답률 통계 (카드 3개)",
            "선택형 Bar 차트 + 비율",
            "주관식 텍스트 목록",
            "평점형 평균 + 분포 차트",
        ]),
        (3, "엑셀 다운로드", ["CSV 형식 결과 내보내기"]),
    ])

    # ── Row 1 → Row 2 연결선 ──
    r2y = I(4.55)
    r2h = I(2.35)
    r2w = I(3.85)
    r2gap = I(0.30)
    r2_start = I(0.60)
    r2xs = [r2_start + i * (r2w + r2gap) for i in range(3)]
    r2cxs = [x + r2w // 2 for x in r2xs]

    mid_y = I(4.32)
    vln(slide, e3cx, r1y + r1h, mid_y)
    hln(slide, r2cxs[0], r2cxs[2], mid_y)
    for cx in r2cxs:
        vln(slide, cx, mid_y, r2y)
        adn(slide, cx, r2y - I(0.1))

    # ── Row 2: 3개 보조 메뉴 ──

    section(slide, r2xs[0], r2y, r2w, r2h, "체계관리 (관리자)", [
        (1, "승인 대기", ["제출 설문 승인/반려 처리"]),
        (2, "전체 설문관리", ["모든 상태 설문 조회/검색"]),
        (3, "카테고리 관리", ["카테고리 등록/수정/삭제"]),
        (4, "통계", ["월별 등록/마감/응답 차트"]),
        (5, "대상자 관리", ["설문별 대상자 목록/응답여부"]),
        (6, "설문 템플릿", ["템플릿 등록/관리/활용"]),
    ])

    section(slide, r2xs[1], r2y, r2w, r2h, "게시판", [
        (1, "공지사항", [
            "목록/상세 조회",
            "등록/수정/삭제 (관리자)",
        ]),
        (2, "질의응답", [
            "질문 등록 / 답변 확인",
            "관리자 답변 작성",
        ]),
    ])

    section(slide, r2xs[2], r2y, r2w, r2h, "관리자", [
        (1, "체계담당자 관리", ["시스템 관리자 지정/해제"]),
        (2, "메뉴관리", ["메뉴 구조 편집 / 권한 설정"]),
        (3, "메시지관리", ["시스템 메시지 편집"]),
        (4, "로그조회", ["사용자 활동 로그 조회"]),
        (5, "코드관리", ["공통코드 CRUD"]),
        (6, "권한관리", ["권한그룹 / 메뉴별 권한 설정"]),
    ])


# ════════════════════════════════════════
#  슬라이드 4: 유스케이스 다이어그램
# ════════════════════════════════════════

def uc_actor(slide, cx, top_y, label):
    """UML 액터 스틱 피규어"""
    hr = I(0.11)
    ovl(slide, cx - hr, top_y, hr * 2, hr * 2, bd=DGRAY)
    by = top_y + hr * 2
    bh = I(0.20)
    vln(slide, cx, by, by + bh, DGRAY)
    hln(slide, cx - I(0.16), cx + I(0.16), by + I(0.06), DGRAY)
    ly = by + bh
    fline(slide, cx, ly, cx - I(0.12), ly + I(0.16), DGRAY, Pt(1.5))
    fline(slide, cx, ly, cx + I(0.12), ly + I(0.16), DGRAY, Pt(1.5))
    tb(slide, cx - I(0.55), ly + I(0.18), I(1.1), I(0.18),
       label, 10, True, DGRAY, PP_ALIGN.CENTER)
    return by + bh // 2


def uc_oval(slide, cx, cy, text):
    """유스케이스 타원 (단일 오브젝트)"""
    w, h = I(1.65), I(0.36)
    s = ovl(slide, cx - w // 2, cy - h // 2, w, h, f=LBLUE, bd=BDR)
    stxt(s, text, 10, False, DGRAY)


def create_usecase(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(slide, "유스케이스 다이어그램 (Use Case Diagram)")
    tb(slide, I(0.4), I(0.55), I(12.5), I(0.25),
       "3개 액터, 18개 유스케이스 — 설문 생성/참여/관리 흐름",
       16, False, MGRAY)
    footer(slide, 4, 4)

    # ── 시스템 경계 ──
    sb = rnd(slide, I(1.90), I(1.12), I(9.50), I(3.78), bd=BDR)
    sb.line.width = Pt(1.5)
    tb(slide, I(2.05), I(1.14), I(4.0), I(0.22),
       "SYS02 설문종합관리체계", 12, True, NAVY, PP_ALIGN.LEFT)

    # ── 유스케이스 배치 (3열 x 6행) ──
    HW = I(1.65) // 2
    C1, C2, C3 = I(3.50), I(6.65), I(9.80)
    ys = [I(1.55), I(2.15), I(2.75), I(3.35), I(3.95), I(4.55)]

    # Col 1 — 설문작성자 영역
    uc_oval(slide, C1, ys[0], "설문 등록/수정")
    uc_oval(slide, C1, ys[1], "문항 편집")
    uc_oval(slide, C1, ys[2], "승인 요청")
    uc_oval(slide, C1, ys[3], "설문 배포")
    uc_oval(slide, C1, ys[4], "설문 마감")
    uc_oval(slide, C1, ys[5], "나의 설문 조회")

    # Col 2 — 응답자 영역
    uc_oval(slide, C2, ys[0], "설문 응답")
    uc_oval(slide, C2, ys[1], "임시 저장")
    uc_oval(slide, C2, ys[2], "응답 제출")
    uc_oval(slide, C2, ys[3], "지난 설문 조회")
    uc_oval(slide, C2, ys[4], "결과 분석 조회")
    uc_oval(slide, C2, ys[5], "엑셀 다운로드")

    # Col 3 — 관리자 영역
    uc_oval(slide, C3, ys[0], "설문 승인/반려")
    uc_oval(slide, C3, ys[1], "전체 설문 관리")
    uc_oval(slide, C3, ys[2], "카테고리 관리")
    uc_oval(slide, C3, ys[3], "템플릿 관리")
    uc_oval(slide, C3, ys[4], "대상자 관리")
    uc_oval(slide, C3, ys[5], "통계 조회")

    # ── 액터 ──
    a_author = uc_actor(slide, I(0.70), I(1.50), "설문작성자")
    a_resp   = uc_actor(slide, I(0.70), I(4.10), "응답자")
    a_admin  = uc_actor(slide, I(12.65), I(2.60), "관리자")

    # ── 연결선 (액터 → 유스케이스) ──
    al = I(0.70) + I(0.16)
    ar = I(12.65) - I(0.16)
    L1 = C1 - HW
    L2 = C2 - HW
    R3 = C3 + HW

    # 설문작성자 → Col 1 전체
    for y_val in ys:
        fline(slide, al, a_author, L1, y_val, MGRAY, Pt(0.75))

    # 응답자 → Col 2 전체
    for y_val in ys:
        fline(slide, al, a_resp, L2, y_val, MGRAY, Pt(0.75))

    # 관리자 → Col 3 전체
    for y_val in ys:
        fline(slide, ar, a_admin, R3, y_val, MGRAY, Pt(0.75))

    # ── 관계 (<<포함>>/<<확장>>) ──
    R1 = C1 + HW

    # 설문 등록 --<<포함>>--> 문항 편집
    fline(slide, C1 + I(0.3), ys[0] + I(0.18), C1 + I(0.3), ys[1] - I(0.18),
          DGRAY, Pt(1), dash=True)
    tb(slide, C1 + I(0.35), (ys[0] + ys[1]) // 2 - I(0.08), I(0.80), I(0.16),
       "<<포함>>", 10, False, DGRAY, PP_ALIGN.CENTER)

    # 설문 응답 --<<포함>>--> 응답 제출
    fline(slide, C2 + I(0.3), ys[0] + I(0.18), C2 + I(0.3), ys[2] - I(0.18),
          DGRAY, Pt(1), dash=True)
    tb(slide, C2 + I(0.35), (ys[0] + ys[2]) // 2 - I(0.08), I(0.80), I(0.16),
       "<<포함>>", 10, False, DGRAY, PP_ALIGN.CENTER)

    # 설문 응답 --<<확장>>--> 임시 저장
    fline(slide, C2 - I(0.3), ys[0] + I(0.18), C2 - I(0.3), ys[1] - I(0.18),
          DGRAY, Pt(1), dash=True)
    tb(slide, C2 - I(1.05), (ys[0] + ys[1]) // 2 - I(0.08), I(0.80), I(0.16),
       "<<확장>>", 10, False, DGRAY, PP_ALIGN.CENTER)

    # 결과 분석 --<<확장>>--> 엑셀 다운로드
    fline(slide, C2 + I(0.3), ys[4] + I(0.18), C2 + I(0.3), ys[5] - I(0.18),
          DGRAY, Pt(1), dash=True)
    tb(slide, C2 + I(0.35), (ys[4] + ys[5]) // 2 - I(0.08), I(0.80), I(0.16),
       "<<확장>>", 10, False, DGRAY, PP_ALIGN.CENTER)


# ════════════════════════════════════════
#  메인
# ════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = I(13.33)
    prs.slide_height = I(7.5)

    create_goal_model(prs)
    create_flowchart(prs)
    create_storyboard(prs)
    create_usecase(prs)

    out_dir = os.path.join(os.path.dirname(__file__), 'ppt')
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, 'SYS02_설문종합관리체계.pptx')
    prs.save(out_path)
    print(f"생성 완료: {out_path}")
    print(f"  슬라이드 1: 목표 모델 (목표·프로세스·참여주체)")
    print(f"  슬라이드 2: 플로우차트 (생성→승인→배포→응답→분석)")
    print(f"  슬라이드 3: 스토리보드 (5개 메뉴, 31개 프로세스)")
    print(f"  슬라이드 4: 유스케이스 다이어그램 (3개 액터, 18개 UC)")


if __name__ == '__main__':
    main()
