#!/usr/bin/env python3
"""
SYS02 설문종합관리체계 분석설계 PPTX v6.1 (portrait, 4페이지)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
system_level_spec_guide_v6.1.md 기반 핵심 5페이지:
  P1. 시스템 개념도      P2. 목표모델 (As-Is → To-Be)
  P3. 유스케이스          P4. 스토리보드 (화면흐름도)
  P5. 플로우차트 (업무처리흐름)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ═══ 색상 팔레트 ═══
NAVY  = RGBColor(0x6B, 0x8E, 0xA8)
TEAL  = RGBColor(0x8E, 0xC5, 0xC0)
SKY   = RGBColor(0xA8, 0xC8, 0xE8)
MINT  = RGBColor(0xDC, 0xF5, 0xE4)
PEACH = RGBColor(0xFC, 0xD4, 0xD0)
CREAM = RGBColor(0xFD, 0xEE, 0xC2)
ICE   = RGBColor(0xE4, 0xEE, 0xF8)
SNOW  = RGBColor(0xF7, 0xF9, 0xFC)
CHAR  = RGBColor(0x1A, 0x1A, 0x1A)
ASH   = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRN   = RGBColor(0xA4, 0xD4, 0xB4)
WARN  = RGBColor(0xF5, 0xCC, 0xA0)
ERR   = RGBColor(0xF0, 0xB2, 0xAE)
LBLUE = RGBColor(0xDC, 0xEA, 0xFE)
LSUB  = RGBColor(0xC5, 0xDC, 0xEC)
HBAR  = RGBColor(0xD4, 0xE2, 0xEC)  # 제목바 연한 배경
THBG  = RGBColor(0xD4, 0xE2, 0xEC)  # 표 헤더 연한 배경
BORD  = RGBColor(0x4A, 0x6B, 0x8C)  # 박스 테두리 (진한)
TBLB  = RGBColor(0x80, 0x80, 0x80)  # 표 테두리 (회색)

# ═══ 치수 ═══
SW = Inches(7.5)
SH = Inches(10)
ML = Inches(0.3)       # 좌우 마진
CW = Inches(6.9)       # 콘텐츠 폭 = 7.5 - 0.3*2
FONT = 'Freesentation 5 Medium'
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'ppt')


# ══════════════════════════════════════════════════
# 헬퍼 함수 (SYS01과 동일)
# ══════════════════════════════════════════════════

def _sl(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = SNOW
    return s


def _run(p, txt, sz=9, b=False, c=CHAR):
    """단락에 서식 런 추가 (최소 9pt 보장)"""
    sz = max(sz, 10)
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = b
    r.font.color.rgb = c
    r.font.name = FONT
    return r


def _tb(s, l, t, w, h, txt, sz=9, b=False, c=CHAR, al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    _run(tb.text_frame.paragraphs[0], txt, sz, b, c)
    tb.text_frame.paragraphs[0].alignment = al
    return tb


def _title(s, title, sub, pg):
    """[A] 제목 영역 (portrait) — 연한 바 + 진한 테두리 + 검정 글자"""
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = HBAR
    bar.line.color.rgb = BORD
    bar.line.width = Pt(0.5)
    _tb(s, ML, Inches(0.06), Inches(5.5), Inches(0.36), title, 20, True, CHAR)
    if sub:
        _tb(s, ML, Inches(0.42), Inches(6.5), Inches(0.26), sub, 10, False, CHAR)
    _tb(s, Inches(6.5), Inches(0.18), Inches(0.6), Inches(0.3),
        f'{pg}/4', 12, True, CHAR, PP_ALIGN.RIGHT)


def _box(s, l, t, w, h, txt, fill=ICE, stk=BORD, sw=0.5,
         sz=9, b=False, fc=CHAR, al=PP_ALIGN.CENTER):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = stk
    sh.line.width = Pt(sw)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.margin_left = tf.margin_right = Inches(0.03)
    for i, ln in enumerate(txt.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p, ln, sz, b, fc)
        p.alignment = al
        p.space_before = p.space_after = Pt(0)
    bp = sh._element.txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', 'ctr')
    return sh


def _darr(s, cx, y, h=Inches(0.08)):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                           cx - Inches(0.06), y, Inches(0.12), h)
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _rarr(s, x, cy, w=Inches(0.18)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                           x, cy - Inches(0.05), w, Inches(0.1))
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _cell_bg(cell, hx):
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for sf in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(sf)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), val=hx)


def _cell_border(cell, color_hex='808080', width_pt=0.5):
    """셀 테두리 4면 회색 설정"""
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for border in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for existing in tcPr.findall(qn(border)):
            tcPr.remove(existing)
        ln = etree.SubElement(tcPr, qn(border))
        ln.set('w', str(int(width_pt * 12700)))
        ln.set('cap', 'flat')
        ln.set('cmpd', 'sng')
        ln.set('algn', 'ctr')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr'), val=color_hex)
        etree.SubElement(ln, qn('a:prstDash'), val='solid')
        etree.SubElement(ln, qn('a:round'))


def _tbl(s, l, t, w, h, data, cws=None):
    """테이블 (멀티라인 셀 지원, 진한 테두리 0.5pt, 연한 헤더 배경). data[0]=헤더"""
    rows, cols = len(data), len(data[0])
    sh = s.shapes.add_table(rows, cols, l, t, w, h)
    tbl = sh.table
    if cws:
        for i, cw in enumerate(cws):
            tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ''
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_top = tf.margin_bottom = Inches(0.04)
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.06)
            hdr = (r == 0)
            for i, ln in enumerate(data[r][c].split('\n')):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                _run(p, ln, 9 if not hdr else 10, hdr, CHAR)
                p.alignment = PP_ALIGN.CENTER if hdr else PP_ALIGN.LEFT
                p.space_before = p.space_after = Pt(0)
            _cell_bg(cell, 'D4E2EC' if hdr else ('FFFFFF' if r % 2 == 1 else 'F2F7FB'))
            _cell_border(cell, '808080', 0.5)
    return tbl




# ══════════════════════════════════════════════════
# Page 2: 목표모델 (As-Is → To-Be)
# ══════════════════════════════════════════════════

def p2_target(prs):
    s = _sl(prs)
    _title(s, '목표모델 (As-Is → To-Be)',
           '현행 문제점 진단 및 목표 시스템 기능 구성', 1)

    pairs = [
        ('수기 설문지 작성·배포\n종이/엑셀, 인쇄·배포 비용 과다',
         '온라인 설문 생성·배포\n웹폼 설계+원클릭 배포, 문항편집기'),
        ('수동 응답 수합\n취합 1~2주, 수기 집계 오류 빈번',
         '실시간 응답 수집\n자동 집계, 응답률 실시간 모니터링'),
        ('결과 수작업 분석\n엑셀 차트, 통계 왜곡·누락 가능',
         '자동 통계·차트 분석\nColumn 차트, 응답분포, 평균평점'),
        ('대상자 지정 비효율\n부대별 수작업 배정, 추적 불가',
         '대상자 자동 지정·관리\n계급/부대/직책 필터, 미응답 알림'),
        ('설문 승인 대면 처리\n승인 지연 평균 3일, 이력 관리 부재',
         '전자 승인 워크플로우\n즉시 승인/반려, 사유입력, 6단계 상태'),
    ]

    y0 = Inches(1.0)
    RH = Inches(0.78)
    RG = Inches(0.06)
    AW = Inches(2.6)
    ARR_W = Inches(0.22)
    AX = ML
    ARX = ML + AW + Inches(0.15)
    TX = ARX + ARR_W + Inches(0.15)

    _box(s, AX, y0 - Inches(0.06), AW, Inches(0.18),
         '현행 시스템 (As-Is)', fill=ERR, stk=ERR, sz=8, b=True, fc=WHITE)
    _box(s, TX, y0 - Inches(0.06), AW, Inches(0.18),
         '목표 시스템 (To-Be)', fill=GRN, stk=GRN, sz=8, b=True, fc=WHITE)

    for i, (asis, tobe) in enumerate(pairs):
        y = y0 + i * (RH + RG)
        _box(s, AX, y, AW, RH, asis, fill=PEACH, stk=ERR, sz=7, fc=CHAR)
        _rarr(s, ARX, y + RH // 2, ARR_W)
        _box(s, TX, y, AW, RH, tobe, fill=MINT, stk=GRN, sz=7, fc=CHAR)

    # 우측 기능 요약 박스
    FX = TX + AW + Inches(0.12)
    FW = ML + CW - FX
    _box(s, FX, y0, FW, 5 * RH + 4 * RG,
         '핵심 기능\n─────────\n'
         '설문관리     4\n게시판        2\n관리자        5\n'
         '─────────\n총 11개 메뉴\n31 프로세스\n─────────\n'
         '체계관리 6탭\n승인대기/전체\n카테고리/통계\n대상자/템플릿',
         fill=ICE, stk=TEAL, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)

    # [C] 서술 테이블
    dy = y0 + 5 * (RH + RG) + Inches(0.12)
    _tbl(s, ML, dy, CW, Inches(1.25), [
        ['개선항목', '현행 문제점', '목표 개선방향', '기대 효과', '정량 목표'],
        ['설문배포', '종이/엑셀 수기 배포', '온라인 생성+원클릭 배포', '배포 시간 단축', '2주→즉시'],
        ['응답수합', '수동 취합 1~2주', '실시간 자동 집계', '집계 오류 제거', '오류 0건'],
        ['결과분석', '엑셀 수작업 차트', '자동 통계+차트 분석', '분석 정확도 향상', '왜곡 0건'],
        ['대상관리', '수작업 배정/추적불가', '필터+자동 알림', '응답률 향상', '응답률 80%↑'],
        ['승인처리', '대면 승인, 평균 3일', '전자 승인/반려', '승인 시간 단축', '3일→즉시'],
    ], cws=[int(CW * r) for r in [0.10, 0.24, 0.24, 0.22, 0.10]])


# ══════════════════════════════════════════════════
# Page 3: 유스케이스
# ══════════════════════════════════════════════════

def p3_usecase(prs):
    s = _sl(prs)
    _title(s, '유스케이스 다이어그램',
           '5개 액터 x 15개 유스케이스, 31개 단위 프로세스', 2)

    y0 = Inches(0.82)

    # 좌측: 액터
    actors = [
        '시스템관리자\nROLE_ADMIN',
        '설문관리자\nROLE_SURVEY_MGR',
        '승인자\nROLE_APPROVER',
        '설문등록자\nROLE_CREATOR',
        '응답자\nROLE_RESPONDENT',
    ]
    AX = ML
    AW = Inches(1.0)
    AH = Inches(0.42)
    AG = Inches(0.06)
    for i, atxt in enumerate(actors):
        _box(s, AX, y0 + i * (AH + AG), AW, AH, atxt,
             fill=ICE, stk=NAVY, sz=7, fc=CHAR)

    # 중앙: 시스템 바운더리
    BX = AX + AW + Inches(0.2)
    BW = Inches(4.3)
    BH = 5 * AH + 4 * AG + Inches(0.3)
    bnd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, BX, y0 - Inches(0.05), BW, BH)
    bnd.fill.solid()
    bnd.fill.fore_color.rgb = RGBColor(0xFB, 0xFB, 0xFB)
    bnd.line.color.rgb = NAVY
    bnd.line.width = Pt(2)
    _tb(s, BX + Inches(0.1), y0 - Inches(0.02), Inches(4), Inches(0.18),
        '설문종합관리체계 (SYS02) — 바운더리', 8, True, NAVY)

    # UC 그룹 (2열 배치)
    uc_groups = [
        ('나의설문관리 (SURV-02)', 'UC-SV-001~004\n설문생성·수정·삭제·제출\n문항편집·대상지정', ICE, SKY),
        ('설문참여 (SURV-01)', 'UC-SV-005~007\n참여목록·응답폼·제출\n진행중 설문 조회', ICE, SKY),
        ('지난설문보기 (SURV-03)', 'UC-SV-008~009\n마감설문 조회·결과확인\n응답분포·통계 차트', ICE, SKY),
        ('체계관리 (SURV-04)', 'UC-SV-010~015\n승인대기·전체설문·카테고리\n통계·대상자·템플릿 6탭', LBLUE, TEAL),
        ('게시판 (SURV-05)', 'UC-SV-016~017\n공지사항·질의응답\nSimpleBoardPage 재사용', LBLUE, TEAL),
        ('관리자 (공통)', 'UC-SYS 시스템·코드·권한\n결재선·게시판 관리\n공통 AdminRoutes', LBLUE, TEAL),
    ]
    UY0 = y0 + Inches(0.22)
    UW = Inches(2.0)
    UH = Inches(0.42)
    UG = Inches(0.06)
    for i, (grp, desc, fill, stk) in enumerate(uc_groups):
        r, c = divmod(i, 2)
        ux = BX + Inches(0.08) + c * (UW + Inches(0.1))
        uy = UY0 + r * (UH + UG)
        sh = _box(s, ux, uy, UW, UH, f'{grp}\n{desc}',
                  fill=fill, stk=stk, sz=6, b=False, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True

    # 우측: 외부 시스템
    EX = BX + BW + Inches(0.12)
    EW = ML + CW - EX
    _box(s, EX, y0 + Inches(0.2), EW, Inches(0.35),
         'GPKI/SSO\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)
    _box(s, EX, y0 + Inches(0.65), EW, Inches(0.35),
         '알림 서비스\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)

    # [C] 서술 테이블 — 주요 UC 상세
    dy = Inches(3.7)
    _tbl(s, ML, dy, CW, Inches(1.65), [
        ['UC-ID', '유스케이스명', '주 액터', '기능 설명 (구체적)', '우선'],
        ['UC-SV-001', '설문 생성', '등록자',
         'CrudForm 15필드 (설문명/내용/시작일/종료일/대상구분/계급/부대/직책/성별/공개/익명/3첨부)', '상'],
        ['UC-SV-002', '문항 편집', '등록자',
         'SurveyQuestionEditor: radio/checkbox/textarea/rate 4유형, 드래그 정렬, 일괄저장', '상'],
        ['UC-SV-005', '설문 참여', '응답자',
         'SurveyFormPage: 문항별 응답 입력 (선택/체크/서술/평점), 참여여부 StatusBadge', '상'],
        ['UC-SV-010', '승인 처리', '승인자',
         'PendingTab: 승인/반려 버튼, 반려 Modal+사유입력, 6단계 상태전환 워크플로우', '상'],
        ['UC-SV-013', '통계 분석', '관리자',
         'StatsTab: 4개 Statistic 카드(총/진행/마감/응답률) + Column 차트(월별 등록/마감/응답)', '중'],
        ['UC-SV-015', '설문 템플릿', '관리자',
         'TemplateTab: 템플릿 CRUD, 카테고리 연동, "이 템플릿으로 설문 생성" 기능', '중'],
    ], cws=[int(CW * r) for r in [0.09, 0.10, 0.07, 0.67, 0.05]])


# ══════════════════════════════════════════════════
# Page 4: 스토리보드 (화면흐름도)
# ══════════════════════════════════════════════════

def p4_storyboard(prs):
    s = _sl(prs)
    _title(s, '스토리보드 (화면흐름도)',
           '총 11개 화면, 3개 업무영역, 핵심 네비게이션 흐름', 3)

    y0 = Inches(0.82)

    # 로그인 → 대시보드
    _box(s, ML, y0, Inches(1.1), Inches(0.35),
         '로그인\n/login', fill=WARN, stk=WARN, sz=8, b=True, fc=WHITE)
    _rarr(s, ML + Inches(1.35), y0 + Inches(0.12), Inches(0.25))
    _box(s, ML + Inches(1.65), y0, Inches(2.0), Inches(0.35),
         '대시보드 / 홈\n/sys02', fill=NAVY, stk=NAVY, sz=8, b=True, fc=WHITE)

    _darr(s, ML + Inches(2.65), y0 + Inches(0.37))

    # 메뉴 그룹
    groups = [
        ('1. 설문관리', '나의설문 · 설문참여\n지난설문 · 체계관리', '/sys02/1/*', '4+6탭', ICE, SKY),
        ('2. 게시판', '공지사항 · 질의응답', '/sys02/board/*', '2', ICE, SKY),
        ('3. 관리자', '시스템 · 코드 · 권한\n결재선 · 게시판', '/sys02/admin/*', '5', LBLUE, TEAL),
    ]

    GY0 = y0 + Inches(0.55)
    GW = Inches(3.2)
    GH = Inches(0.90)
    GG = Inches(0.06)
    GCOL_GAP = Inches(0.15)

    for i, (name, screens, path, cnt, fill, stk) in enumerate(groups):
        r, c = divmod(i, 2)
        gx = ML + c * (GW + GCOL_GAP)
        gy = GY0 + r * (GH + GG)
        sh = _box(s, gx, gy, GW, GH,
                  f'{name} ({cnt}개)\n{screens}\n{path}',
                  fill=fill, stk=stk, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        pars = list(sh.text_frame.paragraphs)
        if len(pars) >= 3:
            for r_obj in pars[-1].runs:
                r_obj.font.color.rgb = ASH
                r_obj.font.size = Pt(8)

    # 체계관리 6탭 상세 (추가 박스)
    TY = GY0 + 2 * (GH + GG) + Inches(0.05)
    tab_items = [
        ('승인대기', 'SearchForm\n승인/반려 버튼'),
        ('전체설문', 'StatusBadge\n7컬럼 DataTable'),
        ('카테고리', 'CRUD Modal\n정렬순서 관리'),
        ('통계', '4 Statistic 카드\nColumn 차트'),
        ('대상자', 'Select 설문선택\n응답여부 Badge'),
        ('템플릿', 'CRUD Modal\n카테고리 연동'),
    ]
    TW = (CW - Inches(0.25)) // 6
    TG = Inches(0.05)
    _tb(s, ML, TY - Inches(0.14), CW, Inches(0.14),
        '체계관리(SURV-04) 6탭 구성:', 8, True, NAVY)
    for i, (tn, td) in enumerate(tab_items):
        tx = ML + i * (TW + TG)
        _box(s, tx, TY, TW, Inches(0.45),
             f'{tn}\n{td}', fill=LBLUE, stk=TEAL, sz=6, fc=CHAR)

    # [C] 서술 — 좌: 화면 목록 / 우: 핵심 흐름
    dy = TY + Inches(0.55)

    LW = int(CW * 0.48)
    _tbl(s, ML, dy, LW, Inches(1.2), [
        ['영역', '화면 수', 'URL 패턴', '주요 컴포넌트'],
        ['나의설문관리', '1+편집', '/sys02/1/2', 'DataTable+Dropdown+CrudForm'],
        ['설문참여', '1+응답', '/sys02/1/3', 'DataTable+SurveyFormPage'],
        ['지난설문보기', '1', '/sys02/1/4', 'DataTable (closed 필터)'],
        ['체계관리', '6탭', '/sys02/1/5', 'Tabs+DataTable+Chart+Modal'],
        ['게시판', '2', '/sys02/board/*', 'SimpleBoardPage (공통 재사용)'],
        ['관리자', '5', '/sys02/admin/*', 'AdminRoutes (공통기능 연결)'],
    ], cws=[int(LW * r) for r in [0.20, 0.10, 0.25, 0.45]])

    RX = ML + LW + Inches(0.15)
    RW = CW - LW - Inches(0.15)
    _tbl(s, RX, dy, RW, Inches(1.2), [
        ['핵심 흐름 시나리오'],
        ['시나리오 1: 설문 생성 → 배포 흐름\n'
         '1. 로그인 → 대시보드\n'
         '2. [설문관리] > [나의 설문관리]\n'
         '3. [신규 등록] → CrudForm 15필드 입력\n'
         '4. [저장] → [문항 편집] 이동\n'
         '5. 4유형 문항 추가 → [일괄저장]\n'
         '6. 관리 Dropdown [승인 요청] → 제출'],
        ['시나리오 2: 설문 응답 흐름\n'
         '1. [설문참여] 메뉴 → active 목록\n'
         '2. 설문명 클릭 → SurveyFormPage\n'
         '3. 문항별 응답(선택/체크/서술/평점)\n'
         '4. [제출] → 응답률 갱신 → 목록 복귀'],
    ])


# ══════════════════════════════════════════════════
# Page 5: 플로우차트 (업무처리흐름)
# ══════════════════════════════════════════════════

def p5_flowchart(prs):
    s = _sl(prs)
    _title(s, '업무처리 플로우차트',
           '핵심 프로세스: 설문 생성 → 승인 → 배포 → 응답 → 마감 — 4개 레인, 12단계', 4)

    y0 = Inches(0.82)

    # 4개 스윔레인
    lanes = [
        ('등록자 / 응답자', ICE, NAVY),
        ('프론트엔드', LBLUE, SKY),
        ('백엔드', MINT, GRN),
        ('DB / 외부', CREAM, WARN),
    ]
    LGAP = Inches(0.14)  # 레인 사이 간격 (10pt)
    LN_W = (CW - 3 * LGAP) // 4
    LN_H = Inches(4.1)   # 다이어그램 높이 1.5배
    LN_HDR = Inches(0.22)

    for i, (name, fill, stk) in enumerate(lanes):
        lx = ML + (LN_W + LGAP) * i
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = fill
        bg.line.color.rgb = stk
        bg.line.width = Pt(1)
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_HDR)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = stk
        hdr.line.fill.background()
        tf = hdr.text_frame
        _run(tf.paragraphs[0], name, 7, True, WHITE)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 프로세스 단계
    steps = [
        (0, 0.30, '1. 설문 생성 메뉴\n[나의 설문관리] 접근'),
        (1, 0.30, '2. CrudForm 표시\n15필드 입력 폼 렌더'),
        (0, 0.70, '3. 설문정보 입력\n설문명/기간/대상/첨부'),
        (1, 0.70, '4. 유효성 검증\n필수값·날짜범위 체크'),
        (1, 1.10, '5. POST /sys02/surveys\nAPI 호출 (생성)'),
        (2, 1.10, '6. 업무규칙 검증\n상태: draft 초기화'),
        (3, 1.10, '7. Survey INSERT\n데이터 저장'),
        (0, 1.50, '8. 문항편집+제출\n4유형 문항 작성→submit'),
        (2, 1.50, '9. PUT submit/approve\n6단계 상태 전환'),
        (1, 1.90, '10. 배포→응답수집\nactive 상태 설문 목록'),
        (0, 1.90, '11. 응답자: 설문참여\n문항별 응답 입력·제출'),
        (3, 1.90, '12. Response INSERT\n응답+통계 저장'),
    ]

    SBW = LN_W - Inches(0.12)
    SBH = Inches(0.56)

    for lane, y_off, txt in steps:
        sx = ML + (LN_W + LGAP) * lane + Inches(0.07)
        sy = y0 + LN_HDR + Inches(y_off)
        _box(s, sx, sy, SBW, SBH, txt,
             fill=WHITE, stk=ASH, sw=1, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)

    # 판단 분기
    fx = ML + (LN_W + LGAP) * 1 + Inches(0.07)
    fy = y0 + LN_HDR + Inches(0.95)
    _box(s, fx, fy, SBW, Inches(0.12),
         '검증 실패 → 필드 에러 표시 → 재입력',
         fill=PEACH, stk=ERR, sw=1, sz=5, fc=ERR, al=PP_ALIGN.CENTER)

    # [C] 서술 — 좌: 단계 / 우: 업무규칙+예외
    dy = y0 + LN_H + Inches(0.1)

    _tbl(s, ML, dy, CW, Inches(1.6), [
        ['단계', '동작', '조건/상세'],
        ['1~2', '메뉴접근 → 폼표시', 'CrudForm(15필드) + SearchForm(2필드)'],
        ['3', '설문정보 입력', '설문명/내용/기간/대상구분/계급/부대/직책/성별/공개/익명/첨부3'],
        ['4', '프론트 유효성 검증', '필수값, 시작일≤종료일, 대상구분별 조건부 필드'],
        ['5~7', 'API→규칙검증→DB저장', 'POST /sys02/surveys → Survey INSERT (status=draft)'],
        ['8~9', '문항편집→승인요청', 'PUT questions → PUT submit → 6단계 상태전환'],
        ['10~12', '배포→응답수집→저장', 'active 설문 → 응답 입력 → Response INSERT + 통계'],
    ], cws=[int(CW * r) for r in [0.08, 0.28, 0.64]])

    _tbl(s, ML, dy + Inches(1.6) + Inches(0.1), CW, Inches(1.6), [
        ['업무 규칙 / 예외 처리'],
        ['[업무 규칙]\n'
         'BR-SV-001: 상태 6단계 (draft→submitted→approved→active→closed, rejected)\n'
         'BR-SV-002: 승인요청은 draft 상태에서만 가능\n'
         'BR-SV-003: 배포는 approved 상태에서만 가능\n'
         'BR-SV-004: 마감은 active 상태에서만 가능\n'
         'BR-SV-005: 문항유형 4종 (radio/checkbox/textarea/rate)\n'
         'BR-SV-006: 대상지정 필터 (계급/부대/직책/성별)\n'
         'BR-SV-007: 결과공개/익명 여부 설문별 설정\n'
         'BR-SV-008: 카테고리별 설문 분류+템플릿 관리'],
        ['[예외 처리]\n'
         'E-001: 세션만료(401) → 로그인 리다이렉트\n'
         'E-002: 권한없음(403) → 접근 제한 알림\n'
         'E-003: 검증실패 → 필드별 에러 표시+재입력\n'
         'E-004: 서버오류(500) → 에러 메시지+재시도'],
    ])


# ══════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════

def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    print('[1/4] 목표모델...')
    p2_target(prs)
    print('[2/4] 유스케이스...')
    p3_usecase(prs)
    print('[3/4] 스토리보드...')
    p4_storyboard(prs)
    print('[4/4] 플로우차트...')
    p5_flowchart(prs)

    out_path = os.path.join(OUT, 'SYS02_설문종합관리체계_분석설계_v6.1.pptx')
    prs.save(out_path)
    print(f'\n완료! 총 {len(prs.slides)}장 슬라이드')
    print(f'파일: {out_path}')
    print(f'크기: {os.path.getsize(out_path) // 1024}KB')


if __name__ == '__main__':
    main()
