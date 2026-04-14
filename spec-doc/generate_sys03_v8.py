#!/usr/bin/env python3
"""
SYS03 성과관리체계 분석설계 PPTX v8 (portrait, 4페이지)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
system_level_spec_guide_v8.md 기반 핵심 5페이지:
  P1. 시스템 개념도      P2. 목표모델 (As-Is → To-Be)
  P3. 유스케이스          P4. 스토리보드 (화면흐름도)
  P5. 플로우차트 (업무처리흐름)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ═══ 색상 팔레트 ═══
NAVY  = RGBColor(0x6B, 0x8E, 0xA8)
TEAL  = RGBColor(0x8E, 0xC5, 0xC0)
SKY   = RGBColor(0xA8, 0xC8, 0xE8)
MINT  = RGBColor(0xDC, 0xF5, 0xE4)
PEACH = RGBColor(0xFC, 0xD4, 0xD0)
CREAM = RGBColor(0xFD, 0xEE, 0xC2)
ICE   = RGBColor(0xE4, 0xEE, 0xF8)
SNOW  = RGBColor(0xF7, 0xF9, 0xFC)
CHAR  = RGBColor(0x1A, 0x1A, 0x1A)
ASH   = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRN   = RGBColor(0xA4, 0xD4, 0xB4)
WARN  = RGBColor(0xF5, 0xCC, 0xA0)
ERR   = RGBColor(0xF0, 0xB2, 0xAE)
LBLUE = RGBColor(0xDC, 0xEA, 0xFE)
LSUB  = RGBColor(0xC5, 0xDC, 0xEC)
HBAR  = RGBColor(0xD4, 0xE2, 0xEC)  # 제목바 연한 배경
THBG  = RGBColor(0xD4, 0xE2, 0xEC)  # 표 헤더 연한 배경
BORD  = RGBColor(0x4A, 0x6B, 0x8C)  # 박스 테두리 (진한)
TBLB  = RGBColor(0x80, 0x80, 0x80)  # 표 테두리 (회색)

# ═══ 치수 ═══
SW = Inches(7.5)
SH = Inches(10)
ML = Inches(0.27)      # 좌우 마진
CW = Inches(6.96)      # 콘텐츠 폭 = 7.5 - 0.27*2
FONT = 'Freesentation 5 Medium'
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'ppt')


# ══════════════════════════════════════════════════
# 헬퍼 함수
# ══════════════════════════════════════════════════

def _sl(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = SNOW
    return s


def _run(p, txt, sz=9, b=False, c=CHAR):
    """단락에 서식 런 추가 (최소 9pt 보장)"""
    sz = max(sz, 10)
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = b
    r.font.color.rgb = c
    r.font.name = FONT
    return r


def _tb(s, l, t, w, h, txt, sz=9, b=False, c=CHAR, al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    _run(tb.text_frame.paragraphs[0], txt, sz, b, c)
    tb.text_frame.paragraphs[0].alignment = al
    return tb


def _title(s, title, sub, pg):
    """[A] 제목 영역 (v8) — 10pt 한 줄, 폭 CW, 높이 0.30"""
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, ML, 0, CW, Inches(0.30))
    bar.fill.solid()
    bar.fill.fore_color.rgb = HBAR
    bar.line.color.rgb = BORD
    bar.line.width = Pt(0.5)
    txt = f'{title} — {sub}' if sub else title
    _tb(s, ML + Inches(0.08), Inches(0.05), CW - Inches(0.16), Inches(0.20),
        txt, 10, True, CHAR, PP_ALIGN.LEFT)


def _box(s, l, t, w, h, txt, fill=ICE, stk=BORD, sw=0.5,
         sz=9, b=False, fc=CHAR, al=PP_ALIGN.CENTER):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = stk
    sh.line.width = Pt(sw)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.margin_left = tf.margin_right = Inches(0.03)
    for i, ln in enumerate(txt.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p, ln, sz, b, fc)
        p.alignment = al
        p.space_before = p.space_after = Pt(0)
    bp = sh._element.txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', 'ctr')
    return sh


def _darr(s, cx, y, h=Inches(0.08)):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                           cx - Inches(0.06), y, Inches(0.12), h)
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _rarr(s, x, cy, w=Inches(0.18)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                           x, cy - Inches(0.05), w, Inches(0.1))
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _cell_bg(cell, hx):
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for sf in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(sf)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), val=hx)


def _cell_border(cell, color_hex='808080', width_pt=0.5):
    """셀 테두리 4면 회색 설정"""
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for border in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for existing in tcPr.findall(qn(border)):
            tcPr.remove(existing)
        ln = etree.SubElement(tcPr, qn(border))
        ln.set('w', str(int(width_pt * 12700)))
        ln.set('cap', 'flat')
        ln.set('cmpd', 'sng')
        ln.set('algn', 'ctr')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr'), val=color_hex)
        etree.SubElement(ln, qn('a:prstDash'), val='solid')
        etree.SubElement(ln, qn('a:round'))


def _tbl(s, l, t, w, h, data, cws=None):
    """테이블 (멀티라인 셀 지원, 진한 테두리 0.5pt, 연한 헤더 배경). data[0]=헤더"""
    rows, cols = len(data), len(data[0])
    sh = s.shapes.add_table(rows, cols, l, t, w, h)
    tbl = sh.table
    if cws:
        for i, cw in enumerate(cws):
            tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ''
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_top = tf.margin_bottom = Inches(0.04)
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.06)
            hdr = (r == 0)
            for i, ln in enumerate(data[r][c].split('\n')):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                _run(p, ln, 9 if not hdr else 10, hdr, CHAR)
                p.alignment = PP_ALIGN.CENTER if hdr else PP_ALIGN.LEFT
                p.space_before = p.space_after = Pt(0)
            _cell_bg(cell, 'D4E2EC' if hdr else ('FFFFFF' if r % 2 == 1 else 'F2F7FB'))
            _cell_border(cell, '000000', 0.5)
    return tbl




# ══════════════════════════════════════════════════
# Page 2: 목표모델 (As-Is → To-Be)
# ══════════════════════════════════════════════════

def p2_target(prs):
    s = _sl(prs)
    _title(s, '목표모델 (As-Is → To-Be)',
           '현행 문제점 진단 및 목표 시스템 기능 구성', 1)

    pairs = [
        ('수기 과제 관리\n엑셀/문서, 6단계 계층 추적 불가',
         '온라인 과제 체계 관리\n방침→중점→중→소→상세→실적 6계층'),
        ('수기 실적 입력·보고\n종이 보고서, 누락·지연 빈번',
         '온라인 실적 입력·상신\nSlider 진도율, 4단계 상태 워크플로우'),
        ('대면 결재·평가\n평가 3주 소요, 기준 불명확',
         '전자 결재·5등급 평가\n즉시 승인/반려, S/A/B/C/D 평가'),
        ('엑셀 진도율 집계\n부대별 취합 1주, 오류 빈번',
         '실시간 추진진도율\n2탭(부대별/부서별) + Bar 차트'),
        ('결과 보고서 수작업\n개인별·부대별 수기 분석',
         '자동 대시보드·통계\nGauge+Column 차트, Progress 시각화'),
    ]

    y0 = Inches(0.50)
    RH = Inches(0.78)
    RG = Inches(0.06)
    AW = Inches(2.6)
    ARR_W = Inches(0.22)
    AX = ML
    ARX = ML + AW + Inches(0.15)
    TX = ARX + ARR_W + Inches(0.15)

    _box(s, AX, y0 - Inches(0.06), AW, Inches(0.18),
         '현행 시스템 (As-Is)', fill=ERR, stk=ERR, sz=8, b=True, fc=WHITE)
    _box(s, TX, y0 - Inches(0.06), AW, Inches(0.18),
         '목표 시스템 (To-Be)', fill=GRN, stk=GRN, sz=8, b=True, fc=WHITE)

    for i, (asis, tobe) in enumerate(pairs):
        y = y0 + i * (RH + RG)
        _box(s, AX, y, AW, RH, asis, fill=PEACH, stk=ERR, sz=7, fc=CHAR)
        _rarr(s, ARX, y + RH // 2, ARR_W)
        _box(s, TX, y, AW, RH, tobe, fill=MINT, stk=GRN, sz=7, fc=CHAR)

    # 우측 기능 요약
    FX = TX + AW + Inches(0.12)
    FW = ML + CW - FX
    _box(s, FX, y0, FW, 5 * RH + 4 * RG,
         '핵심 기능\n─────────\n'
         '기준정보     8\n연간과제     6\n평가결과     2\n'
         '게시판        3\n과제검색     1\n관리자        5\n'
         '─────────\n총 25개 메뉴\n76 프로세스\n─────────\n'
         '과제 6계층\nS/A/B/C/D 평가',
         fill=ICE, stk=TEAL, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)

    # [C] 서술 테이블
    dy = y0 + 5 * (RH + RG) + Inches(0.12)
    _tbl(s, ML, dy, CW, Inches(1.25), [
        ['개선항목', '현행 문제점', '목표 개선방향', '기대 효과', '정량 목표'],
        ['과제관리', '엑셀/문서, 계층 추적불가', '6계층 온라인 과제체계', '체계적 과제 관리', '누락 0건'],
        ['실적입력', '종이 보고, 누락·지연', 'Slider+상신 워크플로우', '실적 적시 관리', '지연 90%↓'],
        ['결재평가', '대면, 3주 소요, 기준불명', '전자결재+5등급 평가', '평가 시간 단축', '3주→즉시'],
        ['진도율', '엑셀 취합, 1주 소요', '실시간 Bar 차트+2탭', '현황 파악 즉시', '1주→실시간'],
        ['결과분석', '수기 보고서', 'Gauge+Column 대시보드', '의사결정 지원', '자동 100%'],
    ], cws=[int(CW * r) for r in [0.10, 0.24, 0.24, 0.22, 0.10]])


# ══════════════════════════════════════════════════
# Page 3: 유스케이스
# ══════════════════════════════════════════════════

def p3_usecase(prs):
    s = _sl(prs)
    _title(s, '유스케이스 다이어그램',
           '5개 액터 x 20개 유스케이스, 76개 단위 프로세스', 2)

    y0 = Inches(0.40)

    # 좌측: 액터
    actors = [
        '시스템관리자\nROLE_ADMIN',
        '평가관리자\nROLE_EVAL_MGR',
        '지휘관/결재자\nROLE_APPROVER',
        '과제담당자\nROLE_TASK_MGR',
        '평가대상자\nROLE_TARGET',
    ]
    AX = ML
    AW = Inches(1.0)
    AH = Inches(0.42)
    AG = Inches(0.06)
    for i, atxt in enumerate(actors):
        _box(s, AX, y0 + i * (AH + AG), AW, AH, atxt,
             fill=ICE, stk=NAVY, sz=7, fc=CHAR)

    # 중앙: 시스템 바운더리
    BX = AX + AW + Inches(0.2)
    BW = Inches(4.3)
    BH = 5 * AH + 4 * AG + Inches(0.3)
    bnd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, BX, y0 - Inches(0.05), BW, BH)
    bnd.fill.solid()
    bnd.fill.fore_color.rgb = RGBColor(0xFB, 0xFB, 0xFB)
    bnd.line.color.rgb = NAVY
    bnd.line.width = Pt(2)
    _tb(s, BX + Inches(0.1), y0 - Inches(0.02), Inches(4), Inches(0.18),
        '성과관리체계 (SYS03) — 바운더리', 8, True, NAVY)

    # UC 그룹 (2열 배치)
    uc_groups = [
        ('기준정보관리 (8)', 'UC-PF-001~008\n기준년도·평가조직·개인대상\n방침·중점·중·소·상세과제 CRUD', ICE, SKY),
        ('연간과제관리 (6)', 'UC-PF-009~014\n추진진도율(2탭)·과제등록\n실적입력·승인·과제평가·개인평가', ICE, SKY),
        ('평가결과 (2)', 'UC-PF-015~016\n부대별 평가결과+Progress\n입력현황 DataTable', ICE, SKY),
        ('과제검색 (1)', 'UC-PF-017\n전체 과제 통합검색\n키워드+부대(서) 필터', LBLUE, TEAL),
        ('게시판 (3)', 'UC-PF-018~020\n공지사항·질의응답·자료실\nSimpleBoardPage 재사용', LBLUE, TEAL),
        ('관리자 (공통)', 'UC-SYS 시스템·코드·권한\n결재선·게시판 관리\n공통 AdminRoutes', LBLUE, TEAL),
    ]
    UY0 = y0 + Inches(0.22)
    UW = Inches(2.0)
    UH = Inches(0.42)
    UG = Inches(0.06)
    for i, (grp, desc, fill, stk) in enumerate(uc_groups):
        r, c = divmod(i, 2)
        ux = BX + Inches(0.08) + c * (UW + Inches(0.1))
        uy = UY0 + r * (UH + UG)
        sh = _box(s, ux, uy, UW, UH, f'{grp}\n{desc}',
                  fill=fill, stk=stk, sz=6, b=False, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True

    # 우측: 외부 시스템
    EX = BX + BW + Inches(0.12)
    EW = ML + CW - EX
    _box(s, EX, y0 + Inches(0.2), EW, Inches(0.35),
         'GPKI/SSO\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)
    _box(s, EX, y0 + Inches(0.65), EW, Inches(0.35),
         '알림 서비스\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)

    # [C] 서술 테이블 — 주요 UC 상세
    dy = Inches(3.7)
    _tbl(s, ML, dy, CW, Inches(1.65), [
        ['UC-ID', '유스케이스명', '주 액터', '기능 설명 (구체적)', '우선'],
        ['UC-PF-001', '과제체계 관리', '관리자',
         '6계층 과제구조 CRUD: 방침→중점→중→소→상세, SearchForm+DataTable+Modal, 엑셀Export', '상'],
        ['UC-PF-009', '추진진도율 조회', '결재자',
         'Tabs 2탭(부대별/부서별), Bar 차트, SearchForm+DataTable, 엑셀저장+인쇄, 상세Modal', '상'],
        ['UC-PF-011', '업무실적 입력', '담당자',
         'DataTable+Modal, Slider(0~100%) 진도율, 4단계 상태(draft→pending→approved|rejected)', '상'],
        ['UC-PF-012', '과제실적 승인', '결재자',
         '상신된 실적 목록, 승인/반려+사유입력, StatusBadge 4색 상태', '상'],
        ['UC-PF-013', '과제실적 평가', '결재자',
         'S/A/B/C/D 5등급 평가, pending/evaluated 상태, 군번/계급/성명 컬럼', '중'],
        ['UC-PF-015', '평가결과 조회', '관리자',
         '부대별 DataTable(6컬럼)+Progress 시각화, 평균등급·평가율, 엑셀저장', '중'],
        ['UC-PF-000', '대시보드', '전체',
         '4 Statistic카드(달성률/과제수/완료/평균) + Gauge + 2 Column차트(방침별/부서별)', '상'],
    ], cws=[int(CW * r) for r in [0.09, 0.11, 0.07, 0.66, 0.05]])


# ══════════════════════════════════════════════════
# Page 4: 스토리보드 (화면흐름도)
# ══════════════════════════════════════════════════

def p4_storyboard(prs):
    s = _sl(prs)
    _title(s, '스토리보드 (화면흐름도)',
           '총 25개 화면, 6개 업무영역, 핵심 네비게이션 흐름', 3)

    y0 = Inches(0.40)

    # 로그인 → 대시보드
    _box(s, ML, y0, Inches(1.1), Inches(0.35),
         '로그인\n/login', fill=WARN, stk=WARN, sz=8, b=True, fc=WHITE)
    _rarr(s, ML + Inches(1.35), y0 + Inches(0.12), Inches(0.25))
    _box(s, ML + Inches(1.65), y0, Inches(2.0), Inches(0.35),
         '대시보드 / 홈\n/sys03', fill=NAVY, stk=NAVY, sz=8, b=True, fc=WHITE)

    _darr(s, ML + Inches(2.65), y0 + Inches(0.37))

    # 메뉴 그룹
    groups = [
        ('1. 기준정보관리', '시스템·평가조직·개인대상\n방침·중점·중·소·상세과제', '/sys03/2/*', '8', ICE, SKY),
        ('2. 연간과제관리', '진도율(2탭)·과제등록·실적\n실적승인·과제평가·개인평가', '/sys03/3/*', '6', ICE, SKY),
        ('3. 평가결과', '평가결과 + 입력현황\nProgress 시각화', '/sys03/4/*', '2', ICE, SKY),
        ('4. 게시판', '공지사항·질의응답·자료실', '/sys03/5/*', '3', LBLUE, TEAL),
        ('5. 과제검색', '통합 과제 검색', '/sys03/6/*', '1', LBLUE, TEAL),
        ('6. 관리자', '시스템·코드·권한·결재선·게시판', '/sys03/admin/*', '5', LBLUE, TEAL),
    ]

    GY0 = y0 + Inches(0.55)
    GW = Inches(3.2)
    GH = Inches(0.90)
    GG = Inches(0.06)
    GCOL_GAP = Inches(0.15)

    for i, (name, screens, path, cnt, fill, stk) in enumerate(groups):
        r, c = divmod(i, 2)
        gx = ML + c * (GW + GCOL_GAP)
        gy = GY0 + r * (GH + GG)
        sh = _box(s, gx, gy, GW, GH,
                  f'{name} ({cnt}개)\n{screens}\n{path}',
                  fill=fill, stk=stk, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        pars = list(sh.text_frame.paragraphs)
        if len(pars) >= 3:
            for r_obj in pars[-1].runs:
                r_obj.font.color.rgb = ASH
                r_obj.font.size = Pt(8)

    # 과제 6계층 구조도 (추가)
    HY = GY0 + 3 * (GH + GG) + Inches(0.08)
    _tb(s, ML, HY - Inches(0.14), CW, Inches(0.14),
        '과제 6계층 구조: (기준정보관리에서 CRUD)', 8, True, NAVY)
    hier = [
        ('지휘방침', NAVY),
        ('추진중점과제', TEAL),
        ('중과제', SKY),
        ('소과제', GRN),
        ('상세과제', WARN),
        ('업무실적', ERR),
    ]
    HW = Inches(1.35)
    HG = Inches(0.12)
    for i, (nm, clr) in enumerate(hier):
        hx = ML + i * (HW + HG)
        _box(s, hx, HY, HW, Inches(0.28), nm,
             fill=WHITE, stk=clr, sz=7, b=True, fc=clr)
        if i < len(hier) - 1:
            _rarr(s, hx + HW + Inches(0.01), HY + Inches(0.09), Inches(0.1))

    # [C] 서술
    dy = HY + Inches(0.42)

    LW = int(CW * 0.48)
    _tbl(s, ML, dy, LW, Inches(1.15), [
        ['영역', '화면 수', 'URL 패턴', '주요 컴포넌트'],
        ['대시보드', '1', '/sys03', 'Gauge+Column(2)+Statistic(4)'],
        ['기준정보', '8', '/sys03/2/*', 'DataTable+Modal CRUD'],
        ['연간과제', '6', '/sys03/3/*', 'Tabs+Slider+Bar+StatusBadge'],
        ['평가결과', '2', '/sys03/4/*', 'DataTable+Progress+엑셀'],
        ['게시판', '3', '/sys03/5/*', 'SimpleBoardPage (공통)'],
        ['과제검색+관리자', '1+5', '/sys03/6,admin/*', 'SearchForm+AdminRoutes'],
    ], cws=[int(LW * r) for r in [0.18, 0.10, 0.25, 0.47]])

    RX = ML + LW + Inches(0.15)
    RW = CW - LW - Inches(0.15)
    _tbl(s, RX, dy, RW, Inches(1.15), [
        ['핵심 흐름 시나리오'],
        ['시나리오 1: 과제등록 → 실적입력 → 평가\n'
         '1. [기준정보] > [지휘방침] CRUD\n'
         '2. 중점과제→중→소→상세과제 순차 등록\n'
         '3. [연간과제] > [업무실적 입력] Slider\n'
         '4. [상신] → 결재자 승인/반려\n'
         '5. [과제실적 평가] S/A/B/C/D 등급 부여'],
        ['시나리오 2: 진도율 조회 → 결과 확인\n'
         '1. [연간과제] > [추진진도율] 2탭 조회\n'
         '2. 부대별 Bar 차트 + 엑셀 저장\n'
         '3. [평가결과] Progress 시각화\n'
         '4. 대시보드 Gauge + Column 차트 확인'],
    ])


# ══════════════════════════════════════════════════
# Page 5: 플로우차트 (업무처리흐름)
# ══════════════════════════════════════════════════

def p5_flowchart(prs):
    s = _sl(prs)
    _title(s, '업무처리 플로우차트',
           '핵심 프로세스: 과제등록 → 실적입력 → 승인 → 평가 — 4개 레인, 12단계', 4)

    y0 = Inches(0.40)

    # 4개 스윔레인
    lanes = [
        ('담당자 / 결재자', ICE, NAVY),
        ('프론트엔드', LBLUE, SKY),
        ('백엔드', MINT, GRN),
        ('DB / 외부', CREAM, WARN),
    ]
    LGAP = Inches(0.14)  # 레인 사이 간격 (10pt)
    LN_W = (CW - 3 * LGAP) // 4
    LN_H = Inches(3.2)   # 다이어그램 높이 (겹침 방지 축소)
    LN_HDR = Inches(0.22)

    for i, (name, fill, stk) in enumerate(lanes):
        lx = ML + (LN_W + LGAP) * i
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = fill
        bg.line.color.rgb = stk
        bg.line.width = Pt(1)
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_HDR)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = stk
        hdr.line.fill.background()
        tf = hdr.text_frame
        _run(tf.paragraphs[0], name, 7, True, WHITE)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 프로세스 단계
    steps = [
        (0, 0.30, '1. 과제체계 구성\n6계층 과제 순차 등록'),
        (1, 0.30, '2. CRUD 폼 표시\nDataTable + Modal'),
        (0, 0.70, '3. 업무실적 입력\nSlider 진도율 + 내용'),
        (1, 0.70, '4. 유효성 검증\n필수값·진도율 범위'),
        (1, 1.10, '5. POST task-results\nAPI 호출 (저장)'),
        (2, 1.10, '6. 업무규칙 검증\nstatus=draft 저장'),
        (3, 1.10, '7. TaskResult INSERT\n데이터 저장'),
        (0, 1.50, '8. 상신 요청\nPopconfirm 확인'),
        (2, 1.50, '9. POST submit\nstatus→pending'),
        (0, 1.90, '10. 결재자 승인/반려\n+ S/A/B/C/D 평가'),
        (2, 1.90, '11. POST approve/reject\n+ POST task-evals'),
        (3, 1.90, '12. 이력+평가 저장\n등급·일시 기록'),
    ]

    SBW = LN_W - Inches(0.12)
    SBH = Inches(0.56)

    for lane, y_off, txt in steps:
        sx = ML + (LN_W + LGAP) * lane + Inches(0.07)
        sy = y0 + LN_HDR + Inches(y_off)
        _box(s, sx, sy, SBW, SBH, txt,
             fill=WHITE, stk=ASH, sw=1, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)

    # 판단 분기
    fx = ML + (LN_W + LGAP) * 1 + Inches(0.07)
    fy = y0 + LN_HDR + Inches(0.95)
    _box(s, fx, fy, SBW, Inches(0.12),
         '검증 실패 → 필드 에러 표시 → 재입력',
         fill=PEACH, stk=ERR, sw=1, sz=5, fc=ERR, al=PP_ALIGN.CENTER)

    # [C] 서술
    dy = y0 + LN_H + Inches(0.1)

    _tbl(s, ML, dy, CW, Inches(1.6), [
        ['단계', '동작', '조건/상세'],
        ['1~2', '과제체계 구성', '6계층 순차 CRUD: 방침→중점→중→소→상세→실적'],
        ['3', '업무실적 입력', 'Slider(0~100%), TextArea 진행내용, 반려사유 Alert 표시'],
        ['4', '프론트 유효성 검증', '필수값 확인, 진도율 0~100 범위, 상태 체크'],
        ['5~7', 'API→규칙검증→DB저장', 'POST /sys03/task-results → TaskResult INSERT'],
        ['8~9', '상신 요청→상태변경', 'Popconfirm → POST submit → status: pending'],
        ['10~12', '승인/평가→이력저장', 'approve/reject + S/A/B/C/D 5등급 평가 기록'],
    ], cws=[int(CW * r) for r in [0.08, 0.28, 0.64]])

    _tbl(s, ML, dy + Inches(1.6) + Inches(0.1), CW, Inches(1.6), [
        ['업무 규칙 / 예외 처리'],
        ['[업무 규칙]\n'
         'BR-PF-001: 과제 6계층 (방침→중점→중→소→상세→실적)\n'
         'BR-PF-002: 실적 상태 4단계 (draft→pending→approved|rejected)\n'
         'BR-PF-003: 평가 등급 5단계 (S/A/B/C/D)\n'
         'BR-PF-004: 진도율 Slider 0~100% (marks: 0/50/100)\n'
         'BR-PF-005: 반려 시 사유 입력 필수 + Alert 표시\n'
         'BR-PF-006: 추진진도율 2탭 (부대별/부서별 과제별)\n'
         'BR-PF-007: 대시보드 3차트 (Gauge+Column 2개)\n'
         'BR-PF-008: 엑셀 Export + 인쇄 기능 내장'],
        ['[예외 처리]\n'
         'E-001: 세션만료(401) → 로그인 리다이렉트\n'
         'E-002: 권한없음(403) → 접근 제한 알림\n'
         'E-003: 검증실패 → 필드별 에러 표시+재입력\n'
         'E-004: 서버오류(500) → 에러 메시지+재시도'],
    ])


# ══════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════

def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    print('[1/4] 목표모델...')
    p2_target(prs)
    print('[2/4] 유스케이스...')
    p3_usecase(prs)
    print('[3/4] 스토리보드...')
    p4_storyboard(prs)
    print('[4/4] 플로우차트...')
    p5_flowchart(prs)

    out_path = os.path.join(OUT, 'SYS03_성과관리체계_분석설계_v8.pptx')
    prs.save(out_path)
    print(f'\n완료! 총 {len(prs.slides)}장 슬라이드')
    print(f'파일: {out_path}')
    print(f'크기: {os.path.getsize(out_path) // 1024}KB')


if __name__ == '__main__':
    main()
