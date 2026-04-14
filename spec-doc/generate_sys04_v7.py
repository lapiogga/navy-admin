#!/usr/bin/env python3
"""
SYS04 인증서발급신청체계 기능명세서 PPTX v7 (portrait, 4페이지)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
system_level_spec_guide_v7.md 기반 핵심 5페이지:
  P1. 시스템 개념도      P2. 목표모델 (As-Is → To-Be)
  P3. 유스케이스          P4. 스토리보드 (화면흐름도)
  P5. 플로우차트 (업무처리흐름)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ═══ 색상 팔레트 ═══
NAVY  = RGBColor(0x6B, 0x8E, 0xA8)
TEAL  = RGBColor(0x8E, 0xC5, 0xC0)
SKY   = RGBColor(0xA8, 0xC8, 0xE8)
MINT  = RGBColor(0xDC, 0xF5, 0xE4)
PEACH = RGBColor(0xFC, 0xD4, 0xD0)
CREAM = RGBColor(0xFD, 0xEE, 0xC2)
ICE   = RGBColor(0xE4, 0xEE, 0xF8)
SNOW  = RGBColor(0xF7, 0xF9, 0xFC)
CHAR  = RGBColor(0x1A, 0x1A, 0x1A)
ASH   = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRN   = RGBColor(0xA4, 0xD4, 0xB4)
WARN  = RGBColor(0xF5, 0xCC, 0xA0)
ERR   = RGBColor(0xF0, 0xB2, 0xAE)
LBLUE = RGBColor(0xDC, 0xEA, 0xFE)
LSUB  = RGBColor(0xC5, 0xDC, 0xEC)
HBAR  = RGBColor(0xD4, 0xE2, 0xEC)  # 제목바 연한 배경
THBG  = RGBColor(0xD4, 0xE2, 0xEC)  # 표 헤더 연한 배경
BORD  = RGBColor(0x4A, 0x6B, 0x8C)  # 박스 테두리 (진한)
TBLB  = RGBColor(0x80, 0x80, 0x80)  # 표 테두리 (회색)

# ═══ 치수 ═══
SW = Inches(7.5)
SH = Inches(10)
ML = Inches(0.27)      # 좌우 마진
CW = Inches(6.96)      # 콘텐츠 폭 = 7.5 - 0.27*2
FONT = 'Freesentation 5 Medium'
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'ppt')


# ══════════════════════════════════════════════════
# 헬퍼 함수
# ══════════════════════════════════════════════════

def _sl(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = SNOW
    return s


def _run(p, txt, sz=9, b=False, c=CHAR):
    """단락에 서식 런 추가 (최소 9pt 보장)"""
    sz = max(sz, 10)
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = b
    r.font.color.rgb = c
    r.font.name = FONT
    return r


def _tb(s, l, t, w, h, txt, sz=9, b=False, c=CHAR, al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    _run(tb.text_frame.paragraphs[0], txt, sz, b, c)
    tb.text_frame.paragraphs[0].alignment = al
    return tb


def _title(s, title, sub, pg):
    """[A] 제목 영역 (portrait) — 연한 바 + 진한 테두리 + 검정 글자"""
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = HBAR
    bar.line.color.rgb = BORD
    bar.line.width = Pt(0.5)
    _tb(s, ML, Inches(0.06), Inches(5.5), Inches(0.36), title, 20, True, CHAR)
    if sub:
        _tb(s, ML, Inches(0.42), Inches(6.5), Inches(0.26), sub, 10, False, CHAR)


def _box(s, l, t, w, h, txt, fill=ICE, stk=BORD, sw=0.5,
         sz=9, b=False, fc=CHAR, al=PP_ALIGN.CENTER):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = stk
    sh.line.width = Pt(sw)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.margin_left = tf.margin_right = Inches(0.03)
    for i, ln in enumerate(txt.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p, ln, sz, b, fc)
        p.alignment = al
        p.space_before = p.space_after = Pt(0)
    bp = sh._element.txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', 'ctr')
    return sh


def _darr(s, cx, y, h=Inches(0.08)):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                           cx - Inches(0.06), y, Inches(0.12), h)
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _rarr(s, x, cy, w=Inches(0.18)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                           x, cy - Inches(0.05), w, Inches(0.1))
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _cell_bg(cell, hx):
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for sf in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(sf)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), val=hx)


def _cell_border(cell, color_hex='808080', width_pt=0.5):
    """셀 테두리 4면 회색 설정"""
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for border in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for existing in tcPr.findall(qn(border)):
            tcPr.remove(existing)
        ln = etree.SubElement(tcPr, qn(border))
        ln.set('w', str(int(width_pt * 12700)))
        ln.set('cap', 'flat')
        ln.set('cmpd', 'sng')
        ln.set('algn', 'ctr')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr'), val=color_hex)
        etree.SubElement(ln, qn('a:prstDash'), val='solid')
        etree.SubElement(ln, qn('a:round'))


def _tbl(s, l, t, w, h, data, cws=None):
    """테이블 (멀티라인 셀 지원, 진한 테두리 0.5pt, 연한 헤더 배경). data[0]=헤더"""
    rows, cols = len(data), len(data[0])
    sh = s.shapes.add_table(rows, cols, l, t, w, h)
    tbl = sh.table
    if cws:
        for i, cw in enumerate(cws):
            tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ''
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_top = tf.margin_bottom = Inches(0.04)
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.06)
            hdr = (r == 0)
            for i, ln in enumerate(data[r][c].split('\n')):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                _run(p, ln, 9 if not hdr else 10, hdr, CHAR)
                p.alignment = PP_ALIGN.CENTER if hdr else PP_ALIGN.LEFT
                p.space_before = p.space_after = Pt(0)
            _cell_bg(cell, 'D4E2EC' if hdr else ('FFFFFF' if r % 2 == 1 else 'F2F7FB'))
            _cell_border(cell, '000000', 0.5)
    return tbl




# ══════════════════════════════════════════════════
# Page 2: 목표모델 (As-Is → To-Be)
# ══════════════════════════════════════════════════

def p2_target(prs):
    s = _sl(prs)
    _title(s, '목표모델 (As-Is → To-Be)',
           '현행 문제점 진단 및 목표 시스템 기능 구성', 1)

    pairs = [
        ('수기 인증서 발급 신청\n종이양식 작성, 대면 제출 필요',
         '온라인 전자 신청\n웹폼 9항목 입력, 즉시 접수'),
        ('대면/유선 승인 처리\n처리시간 평균 3일, 추적 불가',
         '전자 승인 워크플로우\n온라인 승인/반려, 상태 실시간'),
        ('수기 대장 관리\n엑셀 등록대장, 누락·불일치 빈번',
         '온라인 등록대장\n자동 기록, 검색·통계·엑셀출력'),
        ('NDSCA 연계 부재\n국방전자서명 확인 수동 처리',
         'NDSCA 상태 표시\n인증센터 처리상태 연동 (추후)'),
        ('인증서 현황 파악 어려움\n부서별 발급현황 수기 집계',
         '실시간 통계 대시보드\n전체/승인/반려/대기 건수 표시'),
    ]

    y0 = Inches(1.0)
    RH = Inches(0.78)
    RG = Inches(0.06)
    AW = Inches(2.6)
    ARR_W = Inches(0.22)
    AX = ML
    ARX = ML + AW + Inches(0.15)
    TX = ARX + ARR_W + Inches(0.15)

    _box(s, AX, y0 - Inches(0.06), AW, Inches(0.18),
         '현행 시스템 (As-Is)', fill=ERR, stk=ERR, sz=8, b=True, fc=WHITE)
    _box(s, TX, y0 - Inches(0.06), AW, Inches(0.18),
         '목표 시스템 (To-Be)', fill=GRN, stk=GRN, sz=8, b=True, fc=WHITE)

    for i, (asis, tobe) in enumerate(pairs):
        y = y0 + i * (RH + RG)
        _box(s, AX, y, AW, RH, asis, fill=PEACH, stk=ERR, sz=7, fc=CHAR)
        _rarr(s, ARX, y + RH // 2, ARR_W)
        _box(s, TX, y, AW, RH, tobe, fill=MINT, stk=GRN, sz=7, fc=CHAR)

    # 우측 기능 요약
    FX = TX + AW + Inches(0.12)
    FW = ML + CW - FX
    _box(s, FX, y0, FW, 5 * RH + 4 * RG,
         '핵심 기능\n─────────\n'
         '인증서신청  3\n게시판        2\n관리자        5\n'
         '─────────\n총 10개 메뉴\n14 프로세스\n'
         '─────────\n인증서 3종\n재직/경력/복무',
         fill=ICE, stk=TEAL, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)

    # [C] 서술 테이블
    dy = y0 + 5 * (RH + RG) + Inches(0.12)
    _tbl(s, ML, dy, CW, Inches(1.25), [
        ['개선항목', '현행 문제점', '목표 개선방향', '기대 효과', '정량 목표'],
        ['신청처리', '종이양식 수기 작성', '온라인 전자 신청 (9필드)', '접수 누락 방지', '누락 0건'],
        ['승인업무', '대면/유선, 평균 3일', '온라인 승인/반려 즉시', '처리 시간 단축', '3일→즉시'],
        ['대장관리', '엑셀 수기 기록', '자동 등록대장+CSV출력', '데이터 정합성', '불일치 0건'],
        ['현황파악', '수기 집계, 주 1회', '실시간 Statistic 통계', '현황 실시간 조회', '즉시 확인'],
        ['인증연계', 'NDSCA 수동 확인', '상태 자동 표시 (추후)', '중복 확인 제거', '자동화 100%'],
    ], cws=[int(CW * r) for r in [0.10, 0.24, 0.24, 0.22, 0.10]])


# ══════════════════════════════════════════════════
# Page 3: 유스케이스
# ══════════════════════════════════════════════════

def p3_usecase(prs):
    s = _sl(prs)
    _title(s, '유스케이스 다이어그램',
           '3개 액터 x 10개 유스케이스, 14개 단위 프로세스', 2)

    y0 = Inches(0.82)

    # 좌측: 액터
    actors = [
        '신청자\n(일반사용자)\nROLE_USER',
        '승인자\n(관리자)\nROLE_APPROVER',
        '시스템관리자\nROLE_ADMIN',
    ]
    AX = ML
    AW = Inches(1.0)
    AH = Inches(0.55)
    AG = Inches(0.15)
    for i, atxt in enumerate(actors):
        _box(s, AX, y0 + i * (AH + AG), AW, AH, atxt,
             fill=ICE, stk=NAVY, sz=7, fc=CHAR)

    # 중앙: 시스템 바운더리
    BX = AX + AW + Inches(0.2)
    BW = Inches(4.3)
    BH = 3 * AH + 2 * AG + Inches(0.3)
    bnd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, BX, y0 - Inches(0.05), BW, BH)
    bnd.fill.solid()
    bnd.fill.fore_color.rgb = RGBColor(0xFB, 0xFB, 0xFB)
    bnd.line.color.rgb = NAVY
    bnd.line.width = Pt(2)
    _tb(s, BX + Inches(0.1), y0 - Inches(0.02), Inches(5), Inches(0.18),
        '인증서발급신청체계 (SYS04) — 바운더리', 8, True, NAVY)

    # UC 그룹
    uc_groups = [
        ('인증서 신청 (4)', 'UC-CERT-001~004\n신청등록·수정·삭제·회수', ICE, SKY),
        ('인증서 승인관리 (3)', 'UC-CERT-005~007\n목록조회·승인처리·반려처리', ICE, SKY),
        ('등록대장 관리 (3)', 'UC-CERT-008~010\n대장조회·상세보기·엑셀다운', ICE, SKY),
        ('게시판 (2)', 'UC-CERT-011~012\n공지사항·질의응답', LBLUE, TEAL),
        ('관리자 (5)', 'UC-SYS 시스템·코드·권한\n결재선·게시판 관리', LBLUE, TEAL),
    ]
    UY0 = y0 + Inches(0.22)
    UW = Inches(2.0)
    UH = Inches(0.70)
    UG = Inches(0.06)
    for i, (grp, desc, fill, stk) in enumerate(uc_groups):
        r, c = divmod(i, 2)
        ux = BX + Inches(0.08) + c * (UW + Inches(0.1))
        uy = UY0 + r * (UH + UG)
        sh = _box(s, ux, uy, UW, UH, f'{grp}\n{desc}',
                  fill=fill, stk=stk, sz=6, b=False, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True

    # 우측: 외부 시스템
    EX = BX + BW + Inches(0.12)
    EW = ML + CW - EX
    _box(s, EX, y0 + Inches(0.2), EW, Inches(0.35),
         'NDSCA\n국방전자서명인증센터', fill=SNOW, stk=ASH, sz=7, fc=ASH)
    _box(s, EX, y0 + Inches(0.65), EW, Inches(0.35),
         'GPKI/SSO\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)

    # [C] 서술 테이블
    dy = Inches(3.2)
    _tbl(s, ML, dy, CW, Inches(2.1), [
        ['UC-ID', '유스케이스명', '주 액터', '기능 설명 (구체적)', '우선'],
        ['UC-CERT-001', '인증서 신청 등록', '신청자',
         'CrudForm 9필드(인증서구분/신청구분/용도/사유/군번/소속/이메일/전화/동의), Modal 팝업', '상'],
        ['UC-CERT-002', '신청서 수정', '신청자',
         'pending 상태만 수정 가능, editTarget 로딩 → Modal 폼 수정', '상'],
        ['UC-CERT-003', '신청서 삭제/회수', '신청자',
         'pending 상태만 삭제·회수 가능, Popconfirm 확인 다이얼로그', '상'],
        ['UC-CERT-005', '승인 목록 조회', '승인자',
         'SearchForm 3필드 검색 + DataTable 11컬럼, NDSCA 상태 표시', '상'],
        ['UC-CERT-006', '승인/반려 처리', '승인자',
         '승인: Popconfirm 확인, 반려: Modal+TextArea 사유입력', '상'],
        ['UC-CERT-008', '등록대장 조회', '승인자',
         'Tabs(목록/통계), SearchForm 2필드, DataTable+행클릭 DetailModal', '중'],
        ['UC-CERT-010', '엑셀 다운로드', '승인자',
         'CSV(UTF-8 BOM) 출력, 11컬럼 데이터, 날짜기반 파일명', '중'],
    ], cws=[int(CW * r) for r in [0.10, 0.11, 0.07, 0.66, 0.05]])


# ══════════════════════════════════════════════════
# Page 4: 스토리보드 (화면흐름도)
# ══════════════════════════════════════════════════

def p4_storyboard(prs):
    s = _sl(prs)
    _title(s, '스토리보드 (화면흐름도)',
           '총 10개 화면, 3개 업무영역, 핵심 네비게이션 흐름', 3)

    y0 = Inches(0.82)

    # 로그인 → 대시보드
    _box(s, ML, y0, Inches(1.1), Inches(0.35),
         '로그인\n/login', fill=WARN, stk=WARN, sz=8, b=True, fc=WHITE)
    _rarr(s, ML + Inches(1.35), y0 + Inches(0.12), Inches(0.25))
    _box(s, ML + Inches(1.65), y0, Inches(2.0), Inches(0.35),
         '대시보드 / 홈\n/sys04', fill=NAVY, stk=NAVY, sz=8, b=True, fc=WHITE)

    _darr(s, ML + Inches(2.65), y0 + Inches(0.37))

    # 3개 메뉴 그룹
    groups = [
        ('1. 인증서발급신청', '인증서 신청 · 승인/관리 · 등록대장', '/sys04/1/*', '3', ICE, SKY),
        ('2. 게시판', '공지사항 · 질의응답', '/sys04/board/*', '2', LBLUE, TEAL),
        ('3. 관리자', '시스템 · 코드 · 권한 · 결재선 · 게시판', '/sys04/admin/*', '5', LBLUE, TEAL),
    ]

    GY0 = y0 + Inches(0.55)
    GW = Inches(3.2)
    GH = Inches(0.90)
    GG = Inches(0.06)
    GCOL_GAP = Inches(0.15)

    for i, (name, screens, path, cnt, fill, stk) in enumerate(groups):
        r, c = divmod(i, 2)
        gx = ML + c * (GW + GCOL_GAP)
        gy = GY0 + r * (GH + GG)
        sh = _box(s, gx, gy, GW, GH,
                  f'{name} ({cnt}개)\n{screens}\n{path}',
                  fill=fill, stk=stk, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        pars = list(sh.text_frame.paragraphs)
        if len(pars) >= 3:
            for r_obj in pars[-1].runs:
                r_obj.font.color.rgb = ASH
                r_obj.font.size = Pt(8)

    # 인증서 신청 화면 구성 상세 (3개 핵심 화면)
    DY = GY0 + 2 * (GH + GG) + Inches(0.08)
    det_boxes = [
        ('인증서 신청\n/sys04/1/2',
         'SearchForm 3필드\n(인증서구분/신청구분/상태)\nDataTable 9컬럼\nCrudForm 9필드 Modal\n수정/삭제/회수 버튼',
         ICE, SKY),
        ('인증서 승인/관리\n/sys04/1/3',
         'SearchForm 3필드\nDataTable 11컬럼\nNDSCA 상태 표시\n승인 Popconfirm\n반려 Modal+사유입력',
         MINT, GRN),
        ('인증서 등록대장\n/sys04/1/4',
         'Tabs (목록/통계)\nSearchForm 2필드\nDataTable+DetailModal\n엑셀 CSV 다운로드\nStatistic 4건 통계',
         CREAM, WARN),
    ]
    DW = Inches(2.95)
    DH = Inches(1.0)
    DG = Inches(0.1)
    for i, (ttl, desc, fill, stk) in enumerate(det_boxes):
        dx = ML + i * (DW + DG)
        sh = _box(s, dx, DY, DW, DH,
                  f'{ttl}\n──────────\n{desc}',
                  fill=fill, stk=stk, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        sh.text_frame.paragraphs[1].runs[0].font.bold = True

    # [C] 서술 — 좌: 화면 목록 / 우: 핵심 흐름
    dy = DY + DH + Inches(0.1)

    LW = int(CW * 0.48)
    _tbl(s, ML, dy, LW, Inches(1.1), [
        ['영역', '화면 수', 'URL 패턴', '주요 컴포넌트'],
        ['인증서신청', '1', '/sys04/1/2', 'SearchForm+DataTable+CrudForm'],
        ['승인관리', '1', '/sys04/1/3', 'SearchForm+DataTable+RejectModal'],
        ['등록대장', '1', '/sys04/1/4', 'Tabs+DataTable+DetailModal+CSV'],
        ['게시판', '2', '/sys04/board/*', 'SimpleBoardPage (공통 재사용)'],
        ['관리자', '5', '/sys04/admin/*', 'AdminRoutes (공통기능 연결)'],
    ], cws=[int(LW * r) for r in [0.18, 0.10, 0.25, 0.47]])

    RX = ML + LW + Inches(0.15)
    RW = CW - LW - Inches(0.15)
    _tbl(s, RX, dy, RW, Inches(1.1), [
        ['핵심 흐름 시나리오'],
        ['시나리오 1: 인증서 신청 흐름\n'
         '① 로그인 → 대시보드\n'
         '② 좌측메뉴 [인증서발급신청] > [인증서 신청]\n'
         '③ SearchForm 검색 → DataTable 목록\n'
         '④ [신청서 작성] 버튼 → Modal 폼 9항목\n'
         '⑤ 인증서구분/신청구분/용도/사유 입력\n'
         '⑥ [저장] → 성공 토스트 + 목록 갱신'],
        ['시나리오 2: 승인 흐름\n'
         '① [인증서 승인/관리] 메뉴 → 대기 목록\n'
         '② NDSCA 상태 확인 → 승인/반려 선택\n'
         '③ 승인: Popconfirm → 처리완료\n'
         '④ 반려: Modal 사유입력 → 반려 처리'],
    ])


# ══════════════════════════════════════════════════
# Page 5: 플로우차트 (업무처리흐름)
# ══════════════════════════════════════════════════

def p5_flowchart(prs):
    s = _sl(prs)
    _title(s, '업무처리 플로우차트',
           '핵심 프로세스: 인증서 신청 → 승인 → 등록대장 — 4개 레인, 12단계', 4)

    y0 = Inches(0.82)

    # 4개 스윔레인
    lanes = [
        ('사용자 (신청자/승인자)', ICE, NAVY),
        ('프론트엔드', LBLUE, SKY),
        ('백엔드', MINT, GRN),
        ('DB / 외부연계', CREAM, WARN),
    ]
    LGAP = Inches(0.14)  # 레인 사이 간격 (10pt)
    LN_W = (CW - 3 * LGAP) // 4
    LN_H = Inches(3.2)   # 다이어그램 높이 (겹침 방지 축소)
    LN_HDR = Inches(0.22)

    for i, (name, fill, stk) in enumerate(lanes):
        lx = ML + (LN_W + LGAP) * i
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = fill
        bg.line.color.rgb = stk
        bg.line.width = Pt(1)
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_HDR)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = stk
        hdr.line.fill.background()
        tf = hdr.text_frame
        _run(tf.paragraphs[0], name, 7, True, WHITE)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    steps = [
        (0, 0.30, '① 신청 메뉴 접근\n좌측 네비게이션'),
        (1, 0.30, '② 신청 폼 표시\nSearchForm + DataTable'),
        (0, 0.70, '③ 신청 정보 입력\n인증서구분/용도/사유 등 9항목'),
        (1, 0.70, '④ 유효성 검증\n필수값·Select·Checkbox'),
        (1, 1.10, '⑤ POST /certificates\nAPI 호출 (axios)'),
        (2, 1.10, '⑥ 업무규칙 검증\n상태·필수값·중복 체크'),
        (3, 1.10, '⑦ CertApplication INSERT\n데이터 저장 (status: pending)'),
        (2, 1.50, '⑧ 승인자 알림 발송\n(추후) 통보'),
        (1, 1.50, '⑨ 성공 토스트\n+ 목록 갱신 (invalidate)'),
        (0, 1.90, '⑩ 승인자: 승인/반려\nPopconfirm / Modal 사유'),
        (2, 1.90, '⑪ PATCH approve/reject\n상태 변경 + 처리일시'),
        (3, 1.90, '⑫ 등록대장 갱신\n상태 반영 + 이력 기록'),
    ]

    SBW = LN_W - Inches(0.12)
    SBH = Inches(0.56)

    for lane, y_off, txt in steps:
        sx = ML + (LN_W + LGAP) * lane + Inches(0.07)
        sy = y0 + LN_HDR + Inches(y_off)
        _box(s, sx, sy, SBW, SBH, txt,
             fill=WHITE, stk=ASH, sw=1, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)

    # 판단 분기 (④ 검증 실패)
    fx = ML + (LN_W + LGAP) * 1 + Inches(0.07)
    fy = y0 + LN_HDR + Inches(0.95)
    _box(s, fx, fy, SBW, Inches(0.12),
         '검증 실패 → 필드 에러 표시 → 재입력',
         fill=PEACH, stk=ERR, sw=1, sz=5, fc=ERR, al=PP_ALIGN.CENTER)

    # [C] 서술 — 좌: 단계 설명 / 우: 업무규칙+예외
    dy = y0 + LN_H + Inches(0.1)

    _tbl(s, ML, dy, CW, Inches(1.6), [
        ['단계', '동작', '조건/상세'],
        ['①②', '메뉴접근 → 폼표시', 'SearchForm(3필드) + DataTable(9컬럼)'],
        ['③', '정보입력', '인증서구분(3종), 신청구분(3종), 용도, 사유, 이메일 등'],
        ['④', '프론트 유효성 검증', '필수값 확인, Select 선택, Checkbox 동의'],
        ['⑤⑥⑦', 'API→규칙검증→DB저장', 'POST /certificates → status: pending'],
        ['⑧⑨', '알림→결과표시', '(추후) 승인자 통보 + 성공 토스트 + invalidate'],
        ['⑩⑪⑫', '승인/반려→상태변경', 'PATCH approve/reject + 처리일시 + 등록대장 갱신'],
    ], cws=[int(CW * r) for r in [0.08, 0.28, 0.64]])

    _tbl(s, ML, dy + Inches(1.6) + Inches(0.1), CW, Inches(1.6), [
        ['업무 규칙 / 예외 처리'],
        ['[업무 규칙]\n'
         'BR-CERT-001: 인증서 3종 (재직/경력/복무증명서)\n'
         'BR-CERT-002: 신청구분 3종 (신규/재발급/갱신)\n'
         'BR-CERT-003: 상태 4단계 (pending→approved|rejected|withdrawn)\n'
         'BR-CERT-004: pending 상태에서만 수정/삭제/회수 가능\n'
         'BR-CERT-005: 반려 시 사유 필수 입력 (TextArea)\n'
         'BR-CERT-006: NDSCA 상태 3단계 (처리중/승인/반려)\n'
         'BR-CERT-007: 등록대장 CSV 엑셀 출력 (UTF-8 BOM)\n'
         'BR-CERT-008: 통계 4건 (전체/승인/반려/대기)'],
        ['[예외 처리]\n'
         'E-001: 세션만료(401) → 로그인 리다이렉트\n'
         'E-002: 권한없음(403) → 접근 제한 알림\n'
         'E-003: 검증실패 → 필드별 에러 표시+재입력\n'
         'E-004: 서버오류(500) → 에러 메시지+재시도'],
    ])


# ══════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════

def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    print('[1/4] 목표모델...')
    p2_target(prs)
    print('[2/4] 유스케이스...')
    p3_usecase(prs)
    print('[3/4] 스토리보드...')
    p4_storyboard(prs)
    print('[4/4] 플로우차트...')
    p5_flowchart(prs)

    out_path = os.path.join(OUT, 'SYS04_인증서발급신청체계_분석설계_v7.pptx')
    prs.save(out_path)
    print(f'\n완료! 총 {len(prs.slides)}장 슬라이드')
    print(f'파일: {out_path}')


if __name__ == '__main__':
    main()
