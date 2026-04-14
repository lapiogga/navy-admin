#!/usr/bin/env python3
"""
SYS06 해병대규정관리체계 기능명세서 PPTX v8 (portrait, 4페이지)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
system_level_spec_guide_v8.md 기반 핵심 5페이지:
  P1. 시스템 개념도      P2. 목표모델 (As-Is → To-Be)
  P3. 유스케이스          P4. 스토리보드 (화면흐름도)
  P5. 플로우차트 (업무처리흐름)

소스 기반:
  - index.tsx: 4 lazy 페이지 + 3 공통게시판 + 권한관리 + 관리자
  - sys06.ts: Regulation 인터페이스, 4 데이터셋(49건), 17 endpoints
  - RegulationListPage.tsx: Tree 조직도(부/실/처) + SearchForm 3필드 + DataTable 9컬럼 + CrudForm 7필드 + 즐겨찾기
  - PrecedentHQPage.tsx: SearchForm 2필드 + DataTable 8컬럼 + CrudForm 5필드
  - PrecedentUnitPage.tsx: Card 그리드(4열) + 8개 부대 Tag
  - DirectiveListPage.tsx: SearchForm 2필드 + DataTable 8컬럼 + CrudForm 5필드
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ═══ 색상 팔레트 ═══
NAVY  = RGBColor(0x6B, 0x8E, 0xA8)
TEAL  = RGBColor(0x8E, 0xC5, 0xC0)
SKY   = RGBColor(0xA8, 0xC8, 0xE8)
MINT  = RGBColor(0xDC, 0xF5, 0xE4)
PEACH = RGBColor(0xFC, 0xD4, 0xD0)
CREAM = RGBColor(0xFD, 0xEE, 0xC2)
ICE   = RGBColor(0xE4, 0xEE, 0xF8)
SNOW  = RGBColor(0xF7, 0xF9, 0xFC)
CHAR  = RGBColor(0x1A, 0x1A, 0x1A)
ASH   = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRN   = RGBColor(0xA4, 0xD4, 0xB4)
WARN  = RGBColor(0xF5, 0xCC, 0xA0)
ERR   = RGBColor(0xF0, 0xB2, 0xAE)
LBLUE = RGBColor(0xDC, 0xEA, 0xFE)
LSUB  = RGBColor(0xC5, 0xDC, 0xEC)
HBAR  = RGBColor(0xD4, 0xE2, 0xEC)  # 제목바 연한 배경
THBG  = RGBColor(0xD4, 0xE2, 0xEC)  # 표 헤더 연한 배경
BORD  = RGBColor(0x4A, 0x6B, 0x8C)  # 박스 테두리 (진한)
TBLB  = RGBColor(0x80, 0x80, 0x80)  # 표 테두리 (회색)

# ═══ 치수 ═══
SW = Inches(7.5)
SH = Inches(10)
ML = Inches(0.27)      # 좌우 마진
CW = Inches(6.96)      # 콘텐츠 폭 = 7.5 - 0.27*2
FONT = 'Freesentation 5 Medium'
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'ppt')


# ══════════════════════════════════════════════════
# 헬퍼 함수
# ══════════════════════════════════════════════════

def _sl(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = SNOW
    return s


def _run(p, txt, sz=9, b=False, c=CHAR):
    """단락에 서식 런 추가 (최소 9pt 보장)"""
    sz = max(sz, 10)
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = b
    r.font.color.rgb = c
    r.font.name = FONT
    return r


def _tb(s, l, t, w, h, txt, sz=9, b=False, c=CHAR, al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    _run(tb.text_frame.paragraphs[0], txt, sz, b, c)
    tb.text_frame.paragraphs[0].alignment = al
    return tb


def _title(s, title, sub, pg):
    """[A] 제목 영역 (v8) — 10pt 한 줄, 폭 CW, 높이 0.30"""
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, ML, 0, CW, Inches(0.30))
    bar.fill.solid()
    bar.fill.fore_color.rgb = HBAR
    bar.line.color.rgb = BORD
    bar.line.width = Pt(0.5)
    txt = f'{title} — {sub}' if sub else title
    _tb(s, ML + Inches(0.08), Inches(0.05), CW - Inches(0.16), Inches(0.20),
        txt, 10, True, CHAR, PP_ALIGN.LEFT)


def _box(s, l, t, w, h, txt, fill=ICE, stk=BORD, sw=0.5,
         sz=9, b=False, fc=CHAR, al=PP_ALIGN.CENTER):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = stk
    sh.line.width = Pt(sw)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.margin_left = tf.margin_right = Inches(0.03)
    for i, ln in enumerate(txt.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p, ln, sz, b, fc)
        p.alignment = al
        p.space_before = p.space_after = Pt(0)
    bp = sh._element.txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', 'ctr')
    return sh


def _darr(s, cx, y, h=Inches(0.08)):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                           cx - Inches(0.06), y, Inches(0.12), h)
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _rarr(s, x, cy, w=Inches(0.18)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                           x, cy - Inches(0.05), w, Inches(0.1))
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _cell_bg(cell, hx):
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for sf in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(sf)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), val=hx)


def _cell_border(cell, color_hex='808080', width_pt=0.5):
    """셀 테두리 4면 회색 설정"""
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for border in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for existing in tcPr.findall(qn(border)):
            tcPr.remove(existing)
        ln = etree.SubElement(tcPr, qn(border))
        ln.set('w', str(int(width_pt * 12700)))
        ln.set('cap', 'flat')
        ln.set('cmpd', 'sng')
        ln.set('algn', 'ctr')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr'), val=color_hex)
        etree.SubElement(ln, qn('a:prstDash'), val='solid')
        etree.SubElement(ln, qn('a:round'))


def _tbl(s, l, t, w, h, data, cws=None):
    """테이블 (멀티라인 셀 지원, 진한 테두리 0.5pt, 연한 헤더 배경). data[0]=헤더"""
    rows, cols = len(data), len(data[0])
    sh = s.shapes.add_table(rows, cols, l, t, w, h)
    tbl = sh.table
    if cws:
        for i, cw in enumerate(cws):
            tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ''
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_top = tf.margin_bottom = Inches(0.04)
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.06)
            hdr = (r == 0)
            for i, ln in enumerate(data[r][c].split('\n')):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                _run(p, ln, 9 if not hdr else 10, hdr, CHAR)
                p.alignment = PP_ALIGN.CENTER if hdr else PP_ALIGN.LEFT
                p.space_before = p.space_after = Pt(0)
            _cell_bg(cell, 'D4E2EC' if hdr else ('FFFFFF' if r % 2 == 1 else 'F2F7FB'))
            _cell_border(cell, '000000', 0.5)
    return tbl




# ══════════════════════════════════════════════════
# Page 2: 목표모델 (As-Is → To-Be)
# ══════════════════════════════════════════════════

def p2_target(prs):
    s = _sl(prs)
    _title(s, '목표모델 (As-Is → To-Be)',
           '현행 문제점 진단 및 목표 시스템 기능 구성', 1)

    pairs = [
        ('규정 분산 관리\n부서별 개별 파일·문서함, 최신본 혼란',
         '통합 규정 포탈\n조직도 Tree 필터 + 단일 DB, 문서번호 검색'),
        ('수기 예규 관리\n사령부/예하부대 예규 별도 관리',
         '예규 2계층 통합\n사령부+예하부대 8개 부대 포탈'),
        ('지시문서 추적 어려움\n발령일·담당부서·이행여부 불명확',
         '지시문서 CRUD\n발령부서·발령일 검색, 상세 조회'),
        ('조직별 규정 찾기 어려움\n부서별 규정 분류 체계 부재',
         '조직도 Tree 필터\n참모부/직할부서/처 트리 선택 조회'),
        ('열람 이력 관리 부재\n즐겨찾기, 다운로드 이력 없음',
         '즐겨찾기 + 다운로드\nStar 토글, PDF 다운로드/프린트'),
    ]

    y0 = Inches(0.50)
    RH = Inches(0.78)
    RG = Inches(0.06)
    AW = Inches(2.6)
    ARR_W = Inches(0.22)
    AX = ML
    ARX = ML + AW + Inches(0.15)
    TX = ARX + ARR_W + Inches(0.15)

    _box(s, AX, y0 - Inches(0.06), AW, Inches(0.18),
         '현행 시스템 (As-Is)', fill=ERR, stk=ERR, sz=8, b=True, fc=WHITE)
    _box(s, TX, y0 - Inches(0.06), AW, Inches(0.18),
         '목표 시스템 (To-Be)', fill=GRN, stk=GRN, sz=8, b=True, fc=WHITE)

    for i, (asis, tobe) in enumerate(pairs):
        y = y0 + i * (RH + RG)
        _box(s, AX, y, AW, RH, asis, fill=PEACH, stk=ERR, sz=7, fc=CHAR)
        _rarr(s, ARX, y + RH // 2, ARR_W)
        _box(s, TX, y, AW, RH, tobe, fill=MINT, stk=GRN, sz=7, fc=CHAR)

    # 우측 기능 요약
    FX = TX + AW + Inches(0.12)
    FW = ML + CW - FX
    _box(s, FX, y0, FW, 5 * RH + 4 * RG,
         '핵심 기능\n─────────\n'
         '현행규정     1\n예규           2\n지시           1\n'
         '게시판        3\n관리자        5\n'
         '─────────\n총 12개 메뉴\n30 프로세스\n'
         '─────────\n문서분류 4종\n법령/훈령/예규/고시',
         fill=ICE, stk=TEAL, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)

    # [C] 서술 테이블
    dy = y0 + 5 * (RH + RG) + Inches(0.12)
    _tbl(s, ML, dy, CW, Inches(1.25), [
        ['개선항목', '현행 문제점', '목표 개선방향', '기대 효과', '정량 목표'],
        ['규정관리', '부서별 분산, 최신본 혼란', 'Tree 조직도 + 단일 DB 관리', '최신 규정 즉시 접근', '검색 1초'],
        ['예규관리', '본부/부대 별도 관리', '2계층 통합 (사령부+부대 8개)', '예규 통합 조회', '통합 100%'],
        ['지시문서', '추적 불가, 이행 미확인', 'CRUD + 발령부서·일자 검색', '지시 추적 가능', '이력 100%'],
        ['조직필터', '부서별 분류 체계 없음', 'Tree 조직도 부서 선택 조회', '조직별 규정 즉시 접근', '3뎁스 트리'],
        ['열람기능', '즐겨찾기·이력 없음', 'Star 즐겨찾기 + PDF 다운로드', '자주 쓰는 규정 관리', '즐겨찾기 ON'],
    ], cws=[int(CW * r) for r in [0.10, 0.24, 0.24, 0.22, 0.10]])


# ══════════════════════════════════════════════════
# Page 3: 유스케이스
# ══════════════════════════════════════════════════

def p3_usecase(prs):
    s = _sl(prs)
    _title(s, '유스케이스 다이어그램',
           '3개 액터 x 12개 유스케이스, 30개 단위 프로세스', 2)

    y0 = Inches(0.40)

    # 좌측: 액터
    actors = [
        '규정관리자\n규정/예규/지시 등록\nROLE_MANAGER',
        '업무담당자\n열람·검색·즐겨찾기\nROLE_USER',
        '시스템관리자\nROLE_ADMIN',
    ]
    AX = ML
    AW = Inches(1.0)
    AH = Inches(0.55)
    AG = Inches(0.15)
    for i, atxt in enumerate(actors):
        _box(s, AX, y0 + i * (AH + AG), AW, AH, atxt,
             fill=ICE, stk=NAVY, sz=7, fc=CHAR)

    # 중앙: 시스템 바운더리
    BX = AX + AW + Inches(0.2)
    BW = Inches(4.3)
    BH = 3 * AH + 2 * AG + Inches(0.3)
    bnd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, BX, y0 - Inches(0.05), BW, BH)
    bnd.fill.solid()
    bnd.fill.fore_color.rgb = RGBColor(0xFB, 0xFB, 0xFB)
    bnd.line.color.rgb = NAVY
    bnd.line.width = Pt(2)
    _tb(s, BX + Inches(0.1), y0 - Inches(0.02), Inches(5), Inches(0.18),
        '해병대규정관리체계 (SYS06) — 바운더리', 8, True, NAVY)

    # UC 그룹
    uc_groups = [
        ('현행규정 관리 (4)', 'UC-REG-001~004\nCRUD+검색+조직도필터+즐겨찾기', ICE, SKY),
        ('사령부 예규 (3)', 'UC-REG-005~007\n예규 CRUD+검색+상세조회', ICE, SKY),
        ('예하부대 예규 (1)', 'UC-REG-008\n8개 부대 Card 그리드 조회', ICE, SKY),
        ('지시문서 관리 (3)', 'UC-REG-009~011\n지시 CRUD+검색+상세조회', ICE, SKY),
        ('게시판 (3)', 'UC-REG-012~014\n공지사항·규정예고·자료실', LBLUE, TEAL),
        ('관리자 (5)', 'UC-SYS 시스템·코드·권한\n결재선·게시판 관리', LBLUE, TEAL),
    ]
    UY0 = y0 + Inches(0.22)
    UW = Inches(2.0)
    UH = Inches(0.70)
    UG = Inches(0.06)
    for i, (grp, desc, fill, stk) in enumerate(uc_groups):
        r, c = divmod(i, 2)
        ux = BX + Inches(0.08) + c * (UW + Inches(0.1))
        uy = UY0 + r * (UH + UG)
        sh = _box(s, ux, uy, UW, UH, f'{grp}\n{desc}',
                  fill=fill, stk=stk, sz=6, b=False, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True

    # 우측: 외부 시스템
    EX = BX + BW + Inches(0.12)
    EW = ML + CW - EX
    _box(s, EX, y0 + Inches(0.2), EW, Inches(0.35),
         '국방법령정보시스템\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)
    _box(s, EX, y0 + Inches(0.65), EW, Inches(0.35),
         'GPKI/SSO\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)

    # [C] 서술 테이블
    dy = Inches(3.2)
    _tbl(s, ML, dy, CW, Inches(2.1), [
        ['UC-ID', '유스케이스명', '주 액터', '기능 설명 (구체적)', '우선'],
        ['UC-REG-001', '현행규정 조회', '사용자',
         'Row 6:18 레이아웃 — 좌: Tree 조직도(3뎁스), 우: SearchForm(3필드)+DataTable(9컬럼)', '상'],
        ['UC-REG-002', '규정 등록/수정', '관리자',
         'CrudForm 7필드 Modal(규정명/문서번호/분류/소관부서/시행일/내용/첨부파일)', '상'],
        ['UC-REG-003', '규정 상세/즐겨찾기', '사용자',
         'Descriptions(bordered) 상세 Modal + Star 즐겨찾기 토글 + 다운로드/프린트 버튼', '상'],
        ['UC-REG-005', '사령부 예규 관리', '관리자',
         'SearchForm(2필드)+DataTable(8컬럼)+CrudForm(5필드) Modal, CRUD 완전 지원', '상'],
        ['UC-REG-008', '예하부대 예규 조회', '사용자',
         'Card 그리드(4열), 8개 부대 카드, 부대유형 Tag(사단/여단/단), 예규 건수 표시', '중'],
        ['UC-REG-009', '지시문서 관리', '관리자',
         'SearchForm(2필드)+DataTable(8컬럼)+CrudForm(5필드) Modal, 발령부서·발령일', '중'],
        ['UC-REG-012', '게시판', '사용자',
         'SimpleBoardPage 재사용: 공지사항(/sys06/board/1), 규정예고(/sys06/board/2), 자료실(/sys06/board/3)', '하'],
    ], cws=[int(CW * r) for r in [0.10, 0.12, 0.06, 0.66, 0.05]])


# ══════════════════════════════════════════════════
# Page 4: 스토리보드 (화면흐름도)
# ══════════════════════════════════════════════════

def p4_storyboard(prs):
    s = _sl(prs)
    _title(s, '스토리보드 (화면흐름도)',
           '총 12개 화면, 5개 업무영역, 핵심 네비게이션 흐름', 3)

    y0 = Inches(0.40)

    # 로그인 → 대시보드
    _box(s, ML, y0, Inches(1.1), Inches(0.35),
         '로그인\n/login', fill=WARN, stk=WARN, sz=8, b=True, fc=WHITE)
    _rarr(s, ML + Inches(1.35), y0 + Inches(0.12), Inches(0.25))
    _box(s, ML + Inches(1.65), y0, Inches(2.0), Inches(0.35),
         '대시보드 / 홈\n/sys06', fill=NAVY, stk=NAVY, sz=8, b=True, fc=WHITE)

    _darr(s, ML + Inches(2.65), y0 + Inches(0.37))

    # 5개 메뉴 그룹
    groups = [
        ('1. 해병대규정', '현행규정 (Tree+DataTable)', '/sys06/1/*', '1', ICE, SKY),
        ('2. 예규', '사령부 · 예하부대', '/sys06/2/*', '2', ICE, SKY),
        ('3. 지시', '지시문서', '/sys06/3/*', '1', ICE, SKY),
        ('4. 게시판', '공지사항 · 규정예고 · 자료실', '/sys06/board/*', '3', LBLUE, TEAL),
        ('5. 관리자', '시스템 · 코드 · 권한 · 결재선 · 게시판', '/sys06/admin/*', '5', LBLUE, TEAL),
    ]

    GY0 = y0 + Inches(0.55)
    GW = Inches(3.2)
    GH = Inches(0.90)
    GG = Inches(0.06)
    GCOL_GAP = Inches(0.15)

    for i, (name, screens, path, cnt, fill, stk) in enumerate(groups):
        r, c = divmod(i, 2)
        gx = ML + c * (GW + GCOL_GAP)
        gy = GY0 + r * (GH + GG)
        sh = _box(s, gx, gy, GW, GH,
                  f'{name} ({cnt}개)\n{screens}\n{path}',
                  fill=fill, stk=stk, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        pars = list(sh.text_frame.paragraphs)
        if len(pars) >= 3:
            for r_obj in pars[-1].runs:
                r_obj.font.color.rgb = ASH
                r_obj.font.size = Pt(8)

    # 핵심 화면 상세 (3개 고유 페이지)
    DY = GY0 + 3 * (GH + GG) + Inches(0.08)
    det_boxes = [
        ('현행규정\n/sys06/1/1',
         'Row 6:18 레이아웃\n좌: Tree 조직도(3뎁스)\n우: SearchForm(3필드)\n'
         'DataTable(9컬럼)\n즐겨찾기 Star 토글\n상세Modal+다운로드/프린트',
         ICE, SKY),
        ('사령부 예규\n/sys06/2/1',
         'SearchForm(2필드)\nDataTable(8컬럼)\nCrudForm(5필드) Modal\n'
         'Descriptions 상세\n다운로드 버튼',
         MINT, GRN),
        ('예하부대 예규\n/sys06/2/2',
         'Card 그리드(4열 반응형)\n8개 부대 카드\n부대유형 Tag 색상\n'
         '(사단/여단/단)\n예규 건수 표시\nBankOutlined 아이콘',
         CREAM, WARN),
    ]
    DW = Inches(2.95)
    DH = Inches(1.1)
    DG = Inches(0.1)
    for i, (ttl, desc, fill, stk) in enumerate(det_boxes):
        dx = ML + i * (DW + DG)
        sh = _box(s, dx, DY, DW, DH,
                  f'{ttl}\n──────────\n{desc}',
                  fill=fill, stk=stk, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        sh.text_frame.paragraphs[1].runs[0].font.bold = True

    # [C] 서술 테이블 — 2개 좌우 배치
    dy = DY + DH + Inches(0.08)

    LW = int(CW * 0.48)
    _tbl(s, ML, dy, LW, Inches(0.95), [
        ['영역', '화면 수', 'URL 패턴', '주요 컴포넌트'],
        ['현행규정', '1', '/sys06/1/1', 'Tree+SearchForm+DataTable+Modal'],
        ['예규(사령부)', '1', '/sys06/2/1', 'SearchForm+DataTable+CrudForm'],
        ['예규(부대)', '1', '/sys06/2/2', 'Card Grid(4열) + Tag'],
        ['지시문서', '1', '/sys06/3/1', 'SearchForm+DataTable+CrudForm'],
        ['게시판+관리자', '8', '/sys06/board,admin/*', 'SimpleBoardPage+AdminRoutes'],
    ], cws=[int(LW * r) for r in [0.20, 0.10, 0.28, 0.42]])

    RX = ML + LW + Inches(0.15)
    RW = CW - LW - Inches(0.15)
    _tbl(s, RX, dy, RW, Inches(0.95), [
        ['핵심 흐름 시나리오'],
        ['시나리오 1: 규정 검색·열람 흐름\n'
         '① 로그인 → 대시보드\n'
         '② [해병대규정] > [현행규정] 메뉴 선택\n'
         '③ 좌측 조직도 Tree 부서 선택 → 필터링\n'
         '④ SearchForm 규정명/문서번호 검색\n'
         '⑤ 행 클릭 → 상세 Modal (다운로드/프린트)\n'
         '⑥ Star 클릭 → 즐겨찾기 토글'],
        ['시나리오 2: 규정 등록 흐름\n'
         '① [규정 등록] 버튼 → CrudForm Modal\n'
         '② 7항목 입력 (규정명/번호/분류/부서/시행일)\n'
         '③ [저장] → 성공 토스트 + 목록 갱신'],
    ])


# ══════════════════════════════════════════════════
# Page 5: 플로우차트 (업무처리흐름)
# ══════════════════════════════════════════════════

def p5_flowchart(prs):
    s = _sl(prs)
    _title(s, '업무처리 플로우차트',
           '핵심 프로세스: 규정 검색 → 열람 → 등록 — 4개 레인, 12단계', 4)

    y0 = Inches(0.40)

    # 4개 스윔레인
    lanes = [
        ('사용자 (조회자/관리자)', ICE, NAVY),
        ('프론트엔드', LBLUE, SKY),
        ('백엔드', MINT, GRN),
        ('DB / 외부', CREAM, WARN),
    ]
    LGAP = Inches(0.14)  # 레인 사이 간격 (10pt)
    LN_W = (CW - 3 * LGAP) // 4
    LN_H = Inches(3.2)   # 다이어그램 높이 (겹침 방지 축소)
    LN_HDR = Inches(0.22)

    for i, (name, fill, stk) in enumerate(lanes):
        lx = ML + (LN_W + LGAP) * i
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = fill
        bg.line.color.rgb = stk
        bg.line.width = Pt(1)
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_HDR)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = stk
        hdr.line.fill.background()
        tf = hdr.text_frame
        _run(tf.paragraphs[0], name, 7, True, WHITE)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    steps = [
        (0, 0.30, '① 규정 메뉴 접근\n좌측 네비게이션'),
        (1, 0.30, '② 조직도 Tree + 목록\nRow(6:18) 레이아웃 렌더'),
        (0, 0.70, '③ 부서 선택 + 검색\nTree 클릭 → SearchForm'),
        (1, 0.70, '④ API 호출 (필터)\nGET /regulations?dept=X'),
        (2, 0.70, '⑤ 필터링 + 페이지네이션\nkeyword/dept/category'),
        (3, 0.70, '⑥ Regulation SELECT\n조건별 조회 (Mock 16건)'),
        (1, 1.10, '⑦ DataTable 렌더\n9컬럼 + 즐겨찾기 Star'),
        (0, 1.10, '⑧ 행 클릭 → 상세 열람\nDescriptions Modal'),
        (0, 1.50, '⑨ 관리자: 규정 등록\n[규정 등록] 버튼 클릭'),
        (1, 1.50, '⑩ CrudForm Modal\n7필드 입력 폼'),
        (2, 1.90, '⑪ POST /regulations\n신규 규정 저장'),
        (3, 1.90, '⑫ Regulation INSERT\n데이터 저장 + 이력'),
    ]

    SBW = LN_W - Inches(0.12)
    SBH = Inches(0.56)

    for lane, y_off, txt in steps:
        sx = ML + (LN_W + LGAP) * lane + Inches(0.07)
        sy = y0 + LN_HDR + Inches(y_off)
        _box(s, sx, sy, SBW, SBH, txt,
             fill=WHITE, stk=ASH, sw=1, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)

    # 판단 분기
    fx = ML + (LN_W + LGAP) * 1 + Inches(0.07)
    fy = y0 + LN_HDR + Inches(1.35)
    _box(s, fx, fy, SBW, Inches(0.12),
         '검증 실패 → 필드 에러 표시 → 재입력',
         fill=PEACH, stk=ERR, sw=1, sz=5, fc=ERR, al=PP_ALIGN.CENTER)

    # [C] 서술 테이블 — 2개 좌우 배치
    dy = y0 + LN_H + Inches(0.1)

    _tbl(s, ML, dy, CW, Inches(1.6), [
        ['단계', '동작', '조건/상세'],
        ['①②', '메뉴접근 → 레이아웃', 'Row(6:18), Tree 조직도 3뎁스 + DataTable 9컬럼'],
        ['③④', '부서선택 → API호출', 'Tree 클릭 → ?dept=X, SearchForm 3필드 필터'],
        ['⑤⑥', '서버 필터링 → DB조회', '규정명/문서번호/분류/소관부서 4개 필터 조합'],
        ['⑦⑧', '목록렌더 → 상세열람', 'DataTable + Star 즐겨찾기 + Descriptions Modal'],
        ['⑨⑩', '등록버튼 → CrudForm', 'Modal 7필드 (규정명/번호/분류/부서/시행일/내용/첨부)'],
        ['⑪⑫', 'API등록 → DB저장', 'POST /regulations + Regulation INSERT'],
    ], cws=[int(CW * r) for r in [0.08, 0.28, 0.64]])

    _tbl(s, ML, dy + Inches(1.6) + Inches(0.1), CW, Inches(1.6), [
        ['업무 규칙 / 예외 처리'],
        ['[업무 규칙]\n'
         'BR-REG-001: 문서분류 4종 (법령/훈령/예규/고시)\n'
         'BR-REG-002: 조직도 3뎁스 (사령부→부/실/처→부서)\n'
         'BR-REG-003: 데이터 4종 (규정16+사령부예규13+부대예규11+지시9)\n'
         'BR-REG-004: 예하부대 8개 (1사단/2사단/교육훈련단/상륙기동단/항공단/군수단/6여단/9여단)\n'
         'BR-REG-005: 부대유형 3종 Tag (사단/여단/단)\n'
         'BR-REG-006: CrudForm 7필드 CRUD (규정), 5필드 (예규/지시)\n'
         'BR-REG-007: 즐겨찾기 useFavorites 훅 (LocalStorage)\n'
         'BR-REG-008: 상세 Modal에 다운로드/프린트 버튼'],
        ['[예외 처리]\n'
         'E-001: 세션만료(401) → 로그인 리다이렉트\n'
         'E-002: 권한없음(403) → 접근 제한 알림\n'
         'E-003: 검증실패 → 필드별 에러 표시+재입력\n'
         'E-004: 서버오류(500) → 에러 메시지+재시도'],
    ])


# ══════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════

def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    print('[1/4] 목표모델...')
    p2_target(prs)
    print('[2/4] 유스케이스...')
    p3_usecase(prs)
    print('[3/4] 스토리보드...')
    p4_storyboard(prs)
    print('[4/4] 플로우차트...')
    p5_flowchart(prs)

    out_path = os.path.join(OUT, 'SYS06_해병대규정관리체계_분석설계_v8.pptx')
    prs.save(out_path)
    print(f'\n완료! 총 {len(prs.slides)}장 슬라이드')
    print(f'파일: {out_path}')


if __name__ == '__main__':
    main()
