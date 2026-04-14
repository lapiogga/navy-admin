#!/usr/bin/env python3
"""
SYS07 군사자료관리체계 기능명세서 PPTX v5
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
system_level_spec_guide_v5.md 기반 핵심 5페이지:
  P1. 시스템 개념도      P2. 목표모델 (As-Is → To-Be)
  P3. 유스케이스          P4. 스토리보드 (화면흐름도)
  P5. 플로우차트 (업무처리흐름)

SYS05 포맷 완전 동일 구조로 재작성.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ═══ 색상 팔레트 ═══
NAVY  = RGBColor(0x6B, 0x8E, 0xA8)
TEAL  = RGBColor(0x8E, 0xC5, 0xC0)
SKY   = RGBColor(0xA8, 0xC8, 0xE8)
MINT  = RGBColor(0xDC, 0xF5, 0xE4)
PEACH = RGBColor(0xFC, 0xD4, 0xD0)
CREAM = RGBColor(0xFD, 0xEE, 0xC2)
ICE   = RGBColor(0xE4, 0xEE, 0xF8)
SNOW  = RGBColor(0xF7, 0xF9, 0xFC)
CHAR  = RGBColor(0x1A, 0x1A, 0x1A)
ASH   = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRN   = RGBColor(0xA4, 0xD4, 0xB4)
WARN  = RGBColor(0xF5, 0xCC, 0xA0)
ERR   = RGBColor(0xF0, 0xB2, 0xAE)
LBLUE = RGBColor(0xDC, 0xEA, 0xFE)
LSUB  = RGBColor(0xC5, 0xDC, 0xEC)
HBAR  = RGBColor(0xD4, 0xE2, 0xEC)  # 제목바 연한 배경
THBG  = RGBColor(0xD4, 0xE2, 0xEC)  # 표 헤더 연한 배경
BORD  = RGBColor(0x4A, 0x6B, 0x8C)  # 박스 테두리 (진한)
TBLB  = RGBColor(0x80, 0x80, 0x80)  # 표 테두리 (회색)

# ═══ 치수 ═══
SW = Inches(10)
SH = Inches(5.625)
ML = Inches(0.4)
CW = Inches(9.2)
FONT = 'Freesentation 5 Medium'
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'ppt')


# ══════════════════════════════════════════════════
# 헬퍼 함수
# ══════════════════════════════════════════════════

def _sl(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = SNOW
    return s


def _run(p, txt, sz=9, b=False, c=CHAR):
    """단락에 서식 런 추가 (최소 9pt 보장)"""
    sz = max(sz, 9)
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = b
    r.font.color.rgb = c
    r.font.name = FONT
    return r


def _tb(s, l, t, w, h, txt, sz=9, b=False, c=CHAR, al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    _run(tb.text_frame.paragraphs[0], txt, sz, b, c)
    tb.text_frame.paragraphs[0].alignment = al
    return tb


def _title(s, title, sub, pg):
    """[A] 제목 영역 — 연한 바 + 진한 테두리 + 검정 글자"""
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = HBAR
    bar.line.color.rgb = BORD
    bar.line.width = Pt(0.5)
    _tb(s, ML, Inches(0.06), Inches(7.5), Inches(0.36), title, 22, True, CHAR)
    if sub:
        _tb(s, ML, Inches(0.42), Inches(8.5), Inches(0.26), sub, 11, False, CHAR)
    _tb(s, Inches(9.0), Inches(0.18), Inches(0.6), Inches(0.3),
        f'{pg}/5', 12, True, CHAR, PP_ALIGN.RIGHT)


def _box(s, l, t, w, h, txt, fill=ICE, stk=BORD, sw=0.5,
         sz=9, b=False, fc=CHAR, al=PP_ALIGN.CENTER):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = stk
    sh.line.width = Pt(sw)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.margin_left = tf.margin_right = Inches(0.03)
    for i, ln in enumerate(txt.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p, ln, sz, b, fc)
        p.alignment = al
        p.space_before = p.space_after = Pt(0)
    bp = sh._element.txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', 'ctr')
    return sh


def _darr(s, cx, y, h=Inches(0.08)):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                           cx - Inches(0.06), y, Inches(0.12), h)
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _rarr(s, x, cy, w=Inches(0.18)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                           x, cy - Inches(0.05), w, Inches(0.1))
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _cell_bg(cell, hx):
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for sf in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(sf)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), val=hx)


def _cell_border(cell, color_hex='808080', width_pt=0.5):
    """셀 테두리 4면 회색 설정"""
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for border in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for existing in tcPr.findall(qn(border)):
            tcPr.remove(existing)
        ln = etree.SubElement(tcPr, qn(border))
        ln.set('w', str(int(width_pt * 12700)))
        ln.set('cap', 'flat')
        ln.set('cmpd', 'sng')
        ln.set('algn', 'ctr')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr'), val=color_hex)
        etree.SubElement(ln, qn('a:prstDash'), val='solid')
        etree.SubElement(ln, qn('a:round'))


def _tbl(s, l, t, w, h, data, cws=None):
    """테이블 (멀티라인 셀 지원, 진한 테두리 0.5pt, 연한 헤더 배경). data[0]=헤더"""
    rows, cols = len(data), len(data[0])
    sh = s.shapes.add_table(rows, cols, l, t, w, h)
    tbl = sh.table
    if cws:
        for i, cw in enumerate(cws):
            tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ''
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_top = tf.margin_bottom = Inches(0.04)
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.06)
            hdr = (r == 0)
            for i, ln in enumerate(data[r][c].split('\n')):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                _run(p, ln, 9 if not hdr else 10, hdr, CHAR)
                p.alignment = PP_ALIGN.CENTER if hdr else PP_ALIGN.LEFT
                p.space_before = p.space_after = Pt(0)
            _cell_bg(cell, 'D4E2EC' if hdr else ('FFFFFF' if r % 2 == 1 else 'F2F7FB'))
            _cell_border(cell, '808080', 0.5)
    return tbl


# ══════════════════════════════════════════════════
# Page 1: 시스템 개념도
# ══════════════════════════════════════════════════

def p1_concept(prs):
    s = _sl(prs)
    _title(s, '시스템 개념도',
           '군사자료관리체계의 전체 아키텍처 및 기술 구성 — 6개 계층', 1)

    layers = [
        ('사용자\n계층', [
            ('자료관리관\n자료 등록·수정·삭제·일괄등록', ICE, NAVY),
            ('대출·열람자\n대출/열람/지출 신청', ICE, NAVY),
            ('승인권자\n승인·반려·대출·반납처리', ICE, NAVY),
        ]),
        ('채널\n계층', [
            ('웹 브라우저\nChrome/Edge 1920×1080', ICE, TEAL),
            ('반응형 모바일 웹\n360px~', ICE, TEAL),
            ('API Gateway\nRESTful ~25 endpoints', ICE, TEAL),
        ]),
        ('프론트\n엔드', [
            ('React 18 + TS 5\nVite 5 · Router 7', LBLUE, SKY),
            ('Ant Design 5.24\nTabs·Steps·DataTable', LBLUE, SKY),
            ('Zustand 5 + TQ 5\n상태관리 · 서버캐싱', LBLUE, SKY),
        ]),
        ('백엔드\n계층', [
            ('Spring Boot 3.x\n(추후 개발) Java 17', MINT, GRN),
            ('MSW 2.12 Mock\nFaker.js 9 데이터', MINT, GRN),
            ('REST API 설계\nJSON/HTTPS ~25개', MINT, GRN),
        ]),
        ('데이터\n계층', [
            ('PostgreSQL 15\n(추후) 정규화 설계', CREAM, WARN),
            ('Mock 데이터 70건\n문서30+활용20+해기단20', CREAM, WARN),
            ('파일 스토리지\n첨부파일 PDF (추후)', CREAM, WARN),
        ]),
        ('외부\n연계', [
            ('Phase 1 공통\n게시판·코드·권한', SNOW, ASH),
            ('SYS09 공유\nPrintableReport 인쇄', SNOW, ASH),
            ('알림톡/SMS\n(추후 연계)', SNOW, ASH),
        ]),
    ]

    y0 = Inches(0.82)
    LH = Inches(0.44)
    GAP = Inches(0.05)
    LBL_W = Inches(0.85)
    LBL_G = Inches(0.05)

    for idx, (lbl, items) in enumerate(layers):
        y = y0 + idx * (LH + GAP)
        _box(s, ML, y, LBL_W, LH, lbl, fill=NAVY, stk=NAVY, sz=7, b=True, fc=WHITE)
        n = len(items)
        avail = CW - LBL_W - LBL_G
        igap = Inches(0.04)
        iw = (avail - igap * (n - 1)) // n
        x = ML + LBL_W + LBL_G
        for txt, fill, stk in items:
            _box(s, x, y, iw, LH, txt, fill=fill, stk=stk, sz=7, fc=CHAR)
            x += iw + igap
        if idx < len(layers) - 1:
            _darr(s, ML + CW // 2, y + LH + Inches(0.01))

    # [C] 서술 테이블
    dy = y0 + 6 * (LH + GAP) + Inches(0.05)
    _tbl(s, ML, dy, CW, Inches(1.45), [
        ['프론트엔드 기술스택', '백엔드 기술스택', '인프라 / 보안 / 연계'],
        ['• React 18.x + TypeScript 5.x\n'
         '• Ant Design 5.24 (Tabs/Steps/Table)\n'
         '• Zustand 5 + TanStack Query 5\n'
         '• React Router 7 + Vite 5\n'
         '• @ant-design/charts (Column/Pie/Line/Bar)\n'
         '• 4개 고유 페이지, 12개 메뉴',
         '• Spring Boot 3.x (추후 개발)\n'
         '• MSW 2.12 Mock API (현재)\n'
         '• Faker.js 9 한국어 데이터\n'
         '• REST API ~25개 endpoint\n'
         '• JWT + Redis 인증 (추후)\n'
         '• MyBatis 3.5 ORM (추후)',
         '• PostgreSQL 15 (추후 개발)\n'
         '• Linux / Tomcat 10.x (추후)\n'
         '• SSL/TLS 1.3, GPKI (추후)\n'
         '• Phase 1 공통(게시판/코드/권한)\n'
         '• SYS09 PrintableReport 공유\n'
         '• 사용자: 4역할, 40 프로세스'],
    ], cws=[CW // 3] * 3)


# ══════════════════════════════════════════════════
# Page 2: 목표모델 (As-Is → To-Be)
# ══════════════════════════════════════════════════

def p2_target(prs):
    s = _sl(prs)
    _title(s, '목표모델 (As-Is → To-Be)',
           '현행 문제점 진단 및 목표 시스템 기능 구성', 2)

    pairs = [
        ('수기 대장(종이) 기록\n개별 엑셀, 12+ 검색 불가',
         'SearchForm 14필드 디지털 검색\nDataTable 14컬럼 + 일괄등록(엑셀)'),
        ('대출대장 수기 기입\n반납 추적·비밀등급 통제 미흡',
         'Steps 워크플로우(신청→승인→대출→반납)\n비밀/대외비 경고 Alert'),
        ('보존기간 만료 수동 확인\n평가결과 개별 입력·이력 없음',
         '보존만료 자동필터 + 결과 일괄업로드\n변경이력(CHANGE_HISTORY) 추적'),
        ('월/분기 수동 집계\n차트 없이 숫자만 보고',
         '6탭 대시보드(Column/Pie/Line/Bar)\n접수용/활용지원 기록부 + 엑셀/인쇄'),
        ('해기단 자료 별도 관리\n캐비넷 위치·관리자 정보 없음',
         'HaegidanListPage 통합관리\n캐비넷보관위치 + 관리자 군번/계급/성명'),
    ]

    y0 = Inches(1.0)
    RH = Inches(0.52)
    RG = Inches(0.03)
    AW = Inches(3.5)
    ARR_W = Inches(0.22)
    AX = ML
    ARX = ML + AW + Inches(0.15)
    TX = ARX + ARR_W + Inches(0.15)

    _box(s, AX, y0 - Inches(0.06), AW, Inches(0.18),
         '현행 시스템 (As-Is)', fill=ERR, stk=ERR, sz=8, b=True, fc=WHITE)
    _box(s, TX, y0 - Inches(0.06), AW, Inches(0.18),
         '목표 시스템 (To-Be)', fill=GRN, stk=GRN, sz=8, b=True, fc=WHITE)

    for i, (asis, tobe) in enumerate(pairs):
        y = y0 + i * (RH + RG)
        _box(s, AX, y, AW, RH, asis, fill=PEACH, stk=ERR, sz=7, fc=CHAR)
        _rarr(s, ARX, y + RH // 2, ARR_W)
        _box(s, TX, y, AW, RH, tobe, fill=MINT, stk=GRN, sz=7, fc=CHAR)

    # 우측 기능 요약
    FX = TX + AW + Inches(0.12)
    FW = ML + CW - FX
    _box(s, FX, y0, FW, 5 * RH + 4 * RG,
         '핵심 기능\n─────────\n'
         '자료관리     1\n자료활용     1\n통계자료     1\n'
         '해기단        1\n게시판        2\n관리자        5\n'
         '─────────\n총 12개 메뉴\n40 프로세스\n'
         '─────────\n비밀등급 3종\n비밀/대외비/일반',
         fill=ICE, stk=TEAL, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)

    # [C] 서술 테이블
    dy = y0 + 5 * (RH + RG) + Inches(0.12)
    _tbl(s, ML, dy, CW, Inches(1.25), [
        ['개선항목', '현행 문제점', '목표 개선방향', '기대 효과', '정량 목표'],
        ['자료관리', '수기 대장, 엑셀 개별 관리', 'SearchForm 14필드 + DataTable 14컬럼', '디지털 검색 즉시 접근', '검색 1초'],
        ['대출·활용', '대출대장 수기, 비밀통제 미흡', 'Steps 워크플로우 + 비밀경고 Alert', '대출·반납 실시간 추적', '추적 100%'],
        ['평가심의', '만료 수동 확인, 이력 없음', '보존만료 자동필터 + 일괄업로드', '평가 자동화', '자동필터'],
        ['통계현황', '수동 집계, 차트 없음', '6탭 차트 + 기록부 2종 + 엑셀/인쇄', '실시간 통계 대시보드', '6탭 차트'],
        ['해기단자료', '별도 관리, 위치·관리자 누락', '통합관리 + 캐비넷위치 + 관리자 R6', '자료 통합 조회', '통합 100%'],
    ], cws=[int(CW * r) for r in [0.10, 0.24, 0.24, 0.22, 0.10]])


# ══════════════════════════════════════════════════
# Page 3: 유스케이스
# ══════════════════════════════════════════════════

def p3_usecase(prs):
    s = _sl(prs)
    _title(s, '유스케이스 다이어그램',
           '4개 액터 x 7개 UC 그룹, 40개 단위 프로세스', 3)

    y0 = Inches(0.82)

    # 좌측: 액터
    actors = [
        '자료관리관\n자료 등록·수정·삭제\nROLE_MANAGER',
        '대출·열람자\n대출/열람/지출 신청\nROLE_USER',
        '승인권자\n승인·반려·대출·반납\nROLE_APPROVER',
        '시스템관리자\nROLE_ADMIN',
    ]
    AX = ML
    AW = Inches(1.15)
    AH = Inches(0.55)
    AG = Inches(0.15)
    for i, atxt in enumerate(actors):
        _box(s, AX, y0 + i * (AH + AG), AW, AH, atxt,
             fill=ICE, stk=NAVY, sz=7, fc=CHAR)

    # 중앙: 시스템 바운더리
    BX = AX + AW + Inches(0.2)
    BW = Inches(6.2)
    BH = 3 * AH + 2 * AG + Inches(0.3)
    bnd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, BX, y0 - Inches(0.05), BW, BH)
    bnd.fill.solid()
    bnd.fill.fore_color.rgb = RGBColor(0xFB, 0xFB, 0xFB)
    bnd.line.color.rgb = NAVY
    bnd.line.width = Pt(2)
    _tb(s, BX + Inches(0.1), y0 - Inches(0.02), Inches(5), Inches(0.18),
        '군사자료관리체계 (SYS07) — 바운더리', 8, True, NAVY)

    # UC 그룹
    uc_groups = [
        ('군사자료 관리 (8)', 'UC-MDATA-001~008\nCRUD+일괄등록+14필드검색+14컬럼', ICE, SKY),
        ('평가심의 (4)', 'UC-MDATA-009~012\n보존만료 자동필터+결과 일괄업로드', ICE, SKY),
        ('자료 활용 신청 (6)', 'UC-MDATA-013~018\n대출/열람/지출 신청+비밀경고', ICE, SKY),
        ('활용 승인/처리 (6)', 'UC-MDATA-019~024\n승인/반려/대출처리/반납처리', ICE, SKY),
        ('통계자료 (6)', 'UC-MDATA-025~030\n6탭 차트+기록부2종+엑셀/인쇄', LBLUE, TEAL),
        ('해기단자료 (5)', 'UC-MDATA-031~035\n4필드검색+12컬럼+deleteReason', LBLUE, TEAL),
        ('게시판+관리 (5)', 'UC-SYS 시스템·코드·권한\n결재선·게시판 관리', LBLUE, TEAL),
    ]
    UY0 = y0 + Inches(0.22)
    UW = Inches(2.95)
    UH = Inches(0.48)
    UG = Inches(0.03)
    for i, (grp, desc, fill, stk) in enumerate(uc_groups):
        r, c = divmod(i, 2)
        ux = BX + Inches(0.08) + c * (UW + Inches(0.1))
        uy = UY0 + r * (UH + UG)
        sh = _box(s, ux, uy, UW, UH, f'{grp}\n{desc}',
                  fill=fill, stk=stk, sz=6, b=False, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True

    # 우측: 외부 시스템
    EX = BX + BW + Inches(0.12)
    EW = ML + CW - EX
    _box(s, EX, y0 + Inches(0.2), EW, Inches(0.35),
         'Phase 1 공통\n게시판·코드·권한', fill=SNOW, stk=ASH, sz=7, fc=ASH)
    _box(s, EX, y0 + Inches(0.65), EW, Inches(0.35),
         'SYS09 공유\nPrintableReport', fill=SNOW, stk=ASH, sz=7, fc=ASH)

    # [C] 서술 테이블
    dy = Inches(3.2)
    _tbl(s, ML, dy, CW, Inches(2.1), [
        ['UC-ID', '유스케이스명', '주 액터', '기능 설명 (구체적)', '우선'],
        ['UC-MDATA-001', '군사자료 조회/검색', '자료관리관',
         'Tabs(목록/평가심의) — SearchForm 14필드 + DataTable 14컬럼 + StatusBadge 8종 + 비밀등급 Tag 3종', '상'],
        ['UC-MDATA-002', '군사자료 등록/수정', '자료관리관',
         'CrudForm Modal + 일괄등록(Upload.Dragger → bulk-validate → 오류Table → bulk-save)', '상'],
        ['UC-MDATA-013', '자료 활용 신청', '대출·열람자',
         '신청Form 10필드(2열grid) + 비밀/대외비 경고Alert + 활용목적 필수', '상'],
        ['UC-MDATA-019', '활용 승인/처리', '승인권자',
         'Steps(신청→승인→대출→반납) + 승인/반려/대출처리/반납처리 버튼 + 반려사유 Modal', '상'],
        ['UC-MDATA-025', '통계자료', '전체',
         '6탭: Column(문서별)/Pie(등급별)/Line(활용추이)/Bar(교차)/접수용기록부/활용지원기록부', '중'],
        ['UC-MDATA-031', '해기단자료 관리', '자료관리관',
         'SearchForm 4필드 + DataTable 12컬럼 + deleteReason 필수 + 관리자 군번/계급/성명(R6)', '중'],
        ['UC-MDATA-036', '게시판+관리자', '관리자',
         'SimpleBoardPage 재사용: 공지사항·질의응답 + 코드/권한/시스템/결재선/게시판 관리', '하'],
    ], cws=[int(CW * r) for r in [0.10, 0.12, 0.08, 0.64, 0.05]])


# ══════════════════════════════════════════════════
# Page 4: 스토리보드 (화면흐름도)
# ══════════════════════════════════════════════════

def p4_storyboard(prs):
    s = _sl(prs)
    _title(s, '스토리보드 (화면흐름도)',
           '총 12개 메뉴, 4개 전용 화면 + 2 게시판 + 관리자, 핵심 네비게이션 흐름', 4)

    y0 = Inches(0.82)

    # 로그인 → 대시보드
    _box(s, ML, y0, Inches(1.3), Inches(0.35),
         '로그인\n/login', fill=WARN, stk=WARN, sz=8, b=True, fc=WHITE)
    _rarr(s, ML + Inches(1.35), y0 + Inches(0.12), Inches(0.25))
    _box(s, ML + Inches(1.65), y0, Inches(2.0), Inches(0.35),
         '대시보드 / 홈\n/sys07', fill=NAVY, stk=NAVY, sz=8, b=True, fc=WHITE)

    _darr(s, ML + Inches(2.65), y0 + Inches(0.37))

    # 5개 메뉴 그룹
    groups = [
        ('1. 군사자료', '자료관리(Tabs) · 자료활용(Steps)', '/sys07/1/*', '2', ICE, SKY),
        ('2. 통계자료', '6탭 차트 + 기록부 2종', '/sys07/1/3', '1', ICE, SKY),
        ('3. 해기단자료', '검색·목록·등록·수정', '/sys07/2/*', '1', ICE, SKY),
        ('4. 게시판', '공지사항 · 질의응답', '/sys07/board/*', '2', LBLUE, TEAL),
        ('5. 관리자', '시스템·코드·권한·결재선·게시판', '/sys07/admin/*', '5', LBLUE, TEAL),
    ]

    GY0 = y0 + Inches(0.55)
    GW = Inches(4.5)
    GH = Inches(0.58)
    GG = Inches(0.03)
    GCOL_GAP = Inches(0.15)

    for i, (name, screens, path, cnt, fill, stk) in enumerate(groups):
        r, c = divmod(i, 2)
        gx = ML + c * (GW + GCOL_GAP)
        gy = GY0 + r * (GH + GG)
        sh = _box(s, gx, gy, GW, GH,
                  f'{name} ({cnt}개)\n{screens}\n{path}',
                  fill=fill, stk=stk, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        pars = list(sh.text_frame.paragraphs)
        if len(pars) >= 3:
            for r_obj in pars[-1].runs:
                r_obj.font.color.rgb = ASH
                r_obj.font.size = Pt(8)

    # 핵심 화면 상세 (3개 고유 페이지)
    DY = GY0 + 3 * (GH + GG) + Inches(0.08)
    det_boxes = [
        ('군사자료 관리\n/sys07/1/1',
         'Tabs(목록/평가심의)\nSearchForm 14필드\nDataTable 14컬럼\n'
         'StatusBadge 8종 + 비밀등급 Tag\n일괄등록(Upload.Dragger)\ndeleteReason 필수',
         ICE, SKY),
        ('군사자료 활용\n/sys07/1/2',
         'Steps(신청→승인→대출→반납)\nSearchForm 3필드\nDataTable 11컬럼\n'
         '신청Form 10필드(2열grid)\n비밀/대외비 경고 Alert\n승인/반려/대출/반납 버튼',
         MINT, GRN),
        ('해기단자료\n/sys07/2/1',
         'SearchForm 4필드\nDataTable 12컬럼\nHaegidanFormPage 등록/수정\n'
         '캐비넷보관위치 필드\n관리자 군번/계급/성명(R6)\ndeleteReason 필수',
         CREAM, WARN),
    ]
    DW = Inches(2.95)
    DH = Inches(1.1)
    DG = Inches(0.1)
    for i, (ttl, desc, fill, stk) in enumerate(det_boxes):
        dx = ML + i * (DW + DG)
        sh = _box(s, dx, DY, DW, DH,
                  f'{ttl}\n──────────\n{desc}',
                  fill=fill, stk=stk, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        sh.text_frame.paragraphs[1].runs[0].font.bold = True

    # [C] 서술 테이블 (2개 좌우 배치)
    dy = DY + DH + Inches(0.08)

    LW = int(CW * 0.48)
    _tbl(s, ML, dy, LW, Inches(0.95), [
        ['영역', '화면 수', 'URL 패턴', '주요 컴포넌트'],
        ['군사자료관리', '1', '/sys07/1/1', 'Tabs+SearchForm14+DataTable14'],
        ['군사자료활용', '1', '/sys07/1/2', 'Steps+SearchForm3+DataTable11'],
        ['통계자료', '1', '/sys07/1/3', 'Tabs6탭+Chart4종+DataTable2종'],
        ['해기단자료', '1', '/sys07/2/1', 'SearchForm4+DataTable12+Form'],
        ['게시판+관리자', '7', '/sys07/board,admin/*', 'SimpleBoardPage+AdminRoutes'],
    ], cws=[int(LW * r) for r in [0.20, 0.10, 0.28, 0.42]])

    RX = ML + LW + Inches(0.15)
    RW = CW - LW - Inches(0.15)
    _tbl(s, RX, dy, RW, Inches(0.95), [
        ['핵심 흐름 시나리오'],
        ['시나리오 1: 군사자료 검색·조회 흐름\n'
         '① 로그인 → 대시보드\n'
         '② [군사자료] > [군사자료관리] 메뉴 선택\n'
         '③ SearchForm 14필드로 조건 입력\n'
         '④ DataTable 14컬럼 목록 + StatusBadge\n'
         '⑤ 행 클릭 → 상세 Modal (비밀등급 Tag)\n'
         '⑥ 일괄등록: 엑셀 업로드→검증→저장'],
        ['시나리오 2: 자료 활용(대출) 흐름\n'
         '① [군사자료활용] 메뉴 → 신청 버튼\n'
         '② 10필드 입력 (비밀등급 경고 Alert)\n'
         '③ Steps: 신청→승인→대출→반납 진행'],
    ])


# ══════════════════════════════════════════════════
# Page 5: 플로우차트 (업무처리흐름)
# ══════════════════════════════════════════════════

def p5_flowchart(prs):
    s = _sl(prs)
    _title(s, '업무처리 플로우차트',
           '핵심 프로세스: 등록 → 활용(대출) → 평가심의 → 통계 — 4개 레인, 12단계', 5)

    y0 = Inches(0.82)

    # 4개 세로 스윔레인 컬럼
    lanes = [
        ('자료관리관 / 승인권자', ICE, NAVY),
        ('프론트엔드 (React 18)', LBLUE, SKY),
        ('백엔드 (API/Mock)', MINT, GRN),
        ('DB / 외부', CREAM, WARN),
    ]
    LN_W = CW // 4
    LN_H = Inches(2.75)
    LN_HDR = Inches(0.22)

    for i, (name, fill, stk) in enumerate(lanes):
        lx = ML + LN_W * i
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = fill
        bg.line.color.rgb = stk
        bg.line.width = Pt(1)
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_HDR)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = stk
        hdr.line.fill.background()
        tf = hdr.text_frame
        _run(tf.paragraphs[0], name, 7, True, WHITE)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    steps = [
        (0, 0.30, '① 군사자료 등록\n개별 CrudForm / 일괄엑셀'),
        (1, 0.30, '② MilDataFormPage\nPOST /documents 호출'),
        (2, 0.30, '③ bulk-validate\n→ bulk-save 처리'),
        (3, 0.30, '④ Document INSERT\n데이터 저장 (Mock 30건)'),
        (0, 0.70, '⑤ 대출·열람 신청\n10필드 입력 + 비밀경고'),
        (1, 0.70, '⑥ Steps 워크플로우\n신청→승인→대출→반납'),
        (2, 0.70, '⑦ POST /usages\nstatus=pending'),
        (0, 1.10, '⑧ 승인/반려 처리\n반려사유 Modal 입력'),
        (1, 1.10, '⑨ approve/reject API\nPUT /usages/:id'),
        (2, 1.50, '⑩ 평가심의 자동필터\nbulk-upload (파기/연장)'),
        (3, 1.50, '⑪ Evaluation UPDATE\n보존기간 변경 이력'),
        (1, 1.90, '⑫ 통계 6탭 차트\nColumn/Pie/Line/Bar'),
    ]

    SBW = LN_W - Inches(0.12)
    SBH = Inches(0.40)

    for lane, y_off, txt in steps:
        sx = ML + LN_W * lane + Inches(0.06)
        sy = y0 + LN_HDR + Inches(y_off)
        _box(s, sx, sy, SBW, SBH, txt,
             fill=WHITE, stk=ASH, sw=1, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)

    # 판단 분기
    fx = ML + LN_W * 1 + Inches(0.06)
    fy = y0 + LN_HDR + Inches(0.95)
    _box(s, fx, fy, SBW, Inches(0.12),
         '비밀/대외비 → 경고 Alert + 활용목적 필수',
         fill=PEACH, stk=ERR, sw=1, sz=5, fc=ERR, al=PP_ALIGN.CENTER)

    # [C] 서술 테이블 (2개 좌우 배치)
    dy = y0 + LN_H + Inches(0.1)

    LW = int(CW * 0.52)
    _tbl(s, ML, dy, LW, Inches(1.42), [
        ['단계', '동작', '조건/상세'],
        ['①②', '자료등록 → API호출', 'CrudForm(개별) 또는 Upload.Dragger(일괄엑셀)'],
        ['③④', '검증·저장 → DB', 'bulk-validate → 오류Table → bulk-save → INSERT'],
        ['⑤⑥', '활용신청 → Steps', '10필드 입력, 비밀/대외비 경고Alert, Steps 4단계'],
        ['⑦', 'POST /usages', 'status=pending, 활용유형(대출/열람/지출)'],
        ['⑧⑨', '승인/반려 → API', 'approve/reject + 반려사유Modal, PUT /usages/:id'],
        ['⑩⑪', '평가심의 → DB', '보존만료 자동필터 + 일괄업로드(파기/연장)'],
        ['⑫', '통계 6탭 차트', 'Column/Pie/Line/Bar + 접수용/활용지원 기록부'],
    ], cws=[int(LW * r) for r in [0.08, 0.28, 0.64]])

    RX = ML + LW + Inches(0.12)
    RW = CW - LW - Inches(0.12)
    _tbl(s, RX, dy, RW, Inches(1.42), [
        ['업무 규칙 / 예외 처리'],
        ['[업무 규칙]\n'
         'BR-MDATA-001: 비밀등급 3종 Tag (비밀=red, 대외비=orange, 일반=blue)\n'
         'BR-MDATA-002: 상태 8종 StatusBadge (정상/대출/열람/파기/재분류/합철/지출/평가심의)\n'
         'BR-MDATA-003: 삭제 시 deleteReason 필수 (400 BAD_REQUEST if empty)\n'
         'BR-MDATA-004: Steps 4단계 (신청→승인→대출→반납), rejected→status="error"\n'
         'BR-MDATA-005: 일괄등록 (Upload.Dragger→bulk-validate→오류Table→bulk-save)\n'
         'BR-MDATA-006: 인쇄 PrintableReport 공유(SYS09), window.print()'],
        ['[예외 처리]\n'
         'E-001: 세션만료(401) → 로그인 리다이렉트\n'
         'E-002: 권한없음(403) → 접근 제한 알림\n'
         'E-003: 비밀자료 경고 → Alert + 활용목적 필수 입력\n'
         'E-004: 서버오류(500) → 에러 메시지+재시도'],
    ])


# ══════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════

def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    print('[1/5] 시스템 개념도...')
    p1_concept(prs)
    print('[2/5] 목표모델...')
    p2_target(prs)
    print('[3/5] 유스케이스...')
    p3_usecase(prs)
    print('[4/5] 스토리보드...')
    p4_storyboard(prs)
    print('[5/5] 플로우차트...')
    p5_flowchart(prs)

    out_path = os.path.join(OUT, 'SYS07_군사자료관리체계_분석설계_v5.pptx')
    prs.save(out_path)
    print(f'\n완료! 총 {len(prs.slides)}장 슬라이드')
    print(f'파일: {out_path}')


if __name__ == '__main__':
    main()
