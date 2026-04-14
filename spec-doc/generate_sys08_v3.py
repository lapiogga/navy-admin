#!/usr/bin/env python3
"""
SYS08 부대계보관리체계 기능명세서 PPTX v3
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
system_level_spec_guide_v3.md 기반 핵심 5페이지:
  P1. 시스템 개념도      P2. 목표모델 (As-Is → To-Be)
  P3. 유스케이스          P4. 스토리보드 (화면흐름도)
  P5. 플로우차트 (업무처리흐름)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ═══ 색상 팔레트 ═══
NAVY  = RGBColor(0x1B, 0x3A, 0x5C)
TEAL  = RGBColor(0x2E, 0x7D, 0x9B)
SKY   = RGBColor(0x4E, 0xAB, 0xE6)
MINT  = RGBColor(0xD1, 0xFA, 0xE5)
PEACH = RGBColor(0xFE, 0xE2, 0xE2)
CREAM = RGBColor(0xFE, 0xF3, 0xC7)
ICE   = RGBColor(0xE8, 0xF4, 0xFD)
SNOW  = RGBColor(0xF5, 0xF7, 0xFA)
CHAR  = RGBColor(0x21, 0x21, 0x21)
ASH   = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRN   = RGBColor(0x2E, 0x7D, 0x32)
WARN  = RGBColor(0xF5, 0x7C, 0x00)
ERR   = RGBColor(0xC6, 0x28, 0x28)
LBLUE = RGBColor(0xDB, 0xEA, 0xFE)
LSUB  = RGBColor(0xA0, 0xD0, 0xE8)

# ═══ 치수 ═══
SW = Inches(10)
SH = Inches(5.625)
ML = Inches(0.4)
CW = Inches(9.2)
FONT = '맑은 고딕'
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'ppt')


# ══════════════════════════════════════════════════
# 헬퍼 함수
# ══════════════════════════════════════════════════

def _sl(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = SNOW
    return s


def _run(p, txt, sz=9, b=False, c=CHAR):
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = b
    r.font.color.rgb = c
    r.font.name = FONT
    return r


def _tb(s, l, t, w, h, txt, sz=9, b=False, c=CHAR, al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    _run(tb.text_frame.paragraphs[0], txt, sz, b, c)
    tb.text_frame.paragraphs[0].alignment = al
    return tb


def _title(s, title, sub, pg):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.68))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    _tb(s, ML, Inches(0.06), Inches(7.5), Inches(0.34), title, 22, True, WHITE)
    if sub:
        _tb(s, ML, Inches(0.38), Inches(8.5), Inches(0.26), sub, 11, False, LSUB)
    _tb(s, Inches(9.0), Inches(0.15), Inches(0.6), Inches(0.3),
        f'{pg}/5', 12, True, WHITE, PP_ALIGN.RIGHT)


def _box(s, l, t, w, h, txt, fill=ICE, stk=NAVY, sw=1,
         sz=9, b=False, fc=CHAR, al=PP_ALIGN.CENTER):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = stk
    sh.line.width = Pt(sw)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.margin_left = tf.margin_right = Inches(0.03)
    for i, ln in enumerate(txt.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p, ln, sz, b, fc)
        p.alignment = al
        p.space_before = p.space_after = Pt(0)
    bp = sh._element.txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', 'ctr')
    return sh


def _darr(s, cx, y, h=Inches(0.08)):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                           cx - Inches(0.06), y, Inches(0.12), h)
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _rarr(s, x, cy, w=Inches(0.18)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                           x, cy - Inches(0.05), w, Inches(0.1))
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _cell_bg(cell, hx):
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for sf in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(sf)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), val=hx)


def _tbl(s, l, t, w, h, data, cws=None):
    rows, cols = len(data), len(data[0])
    sh = s.shapes.add_table(rows, cols, l, t, w, h)
    tbl = sh.table
    if cws:
        for i, cw in enumerate(cws):
            tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ''
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_top = tf.margin_bottom = Inches(0.01)
            tf.margin_left = Inches(0.04)
            tf.margin_right = Inches(0.03)
            hdr = (r == 0)
            for i, ln in enumerate(data[r][c].split('\n')):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                _run(p, ln, 7 if not hdr else 8, hdr, WHITE if hdr else CHAR)
                p.alignment = PP_ALIGN.CENTER if hdr else PP_ALIGN.LEFT
                p.space_before = p.space_after = Pt(0)
            _cell_bg(cell, '1B3A5C' if hdr else ('FFFFFF' if r % 2 == 1 else 'F0F4F8'))
    return tbl


# ══════════════════════════════════════════════════
# Page 1: 시스템 개념도
# ══════════════════════════════════════════════════

def p1_concept(prs):
    s = _sl(prs)
    _title(s, '시스템 개념도',
           '부대계보관리체계의 전체 아키텍처 및 기술 구성 — 6개 계층', 1)

    layers = [
        ('사용자\n계층', [
            ('계보담당\n부대기록·계승·직위자·활동 등록', ICE, NAVY),
            ('결재권자\n주요활동 검토·승인 (2단계)', ICE, NAVY),
            ('시스템관리자\n코드·권한·시스템 관리', ICE, NAVY),
        ]),
        ('채널\n계층', [
            ('웹 브라우저\nChrome/Edge 1920×1080', ICE, TEAL),
            ('반응형 모바일 웹\n360px~', ICE, TEAL),
            ('API Gateway\nRESTful ~33 endpoints', ICE, TEAL),
        ]),
        ('프론트\n엔드', [
            ('React 18 + TS 5\nVite 5 · Router 7', LBLUE, SKY),
            ('Ant Design 5.24\nTree · Steps · Upload', LBLUE, SKY),
            ('Zustand 5 + TQ 5\n상태관리 · 서버캐싱', LBLUE, SKY),
        ]),
        ('백엔드\n계층', [
            ('Spring Boot 3.x\n(추후 개발) Java 17', MINT, GRN),
            ('MSW 2.12 Mock\nFaker.js 9 데이터', MINT, GRN),
            ('REST API 설계\nJSON/HTTPS ~33개', MINT, GRN),
        ]),
        ('데이터\n계층', [
            ('PostgreSQL 15\n(추후) 정규화 설계', CREAM, WARN),
            ('Mock 데이터 156건\n7 데이터셋 94필드 합계', CREAM, WARN),
            ('파일 스토리지\n사진·첨부파일 (추후)', CREAM, WARN),
        ]),
        ('외부\n연계', [
            ('Phase 1 공통\n게시판·코드·권한', SNOW, ASH),
            ('SYS09 PrintableReport\n인쇄 미리보기 공유', SNOW, ASH),
            ('GPKI/SSO 인증\n(추후 연계)', SNOW, ASH),
        ]),
    ]

    y0 = Inches(0.82)
    LH = Inches(0.37)
    GAP = Inches(0.1)
    LBL_W = Inches(0.85)
    LBL_G = Inches(0.05)

    for idx, (lbl, items) in enumerate(layers):
        y = y0 + idx * (LH + GAP)
        _box(s, ML, y, LBL_W, LH, lbl, fill=NAVY, stk=NAVY, sz=7, b=True, fc=WHITE)
        n = len(items)
        avail = CW - LBL_W - LBL_G
        igap = Inches(0.04)
        iw = (avail - igap * (n - 1)) // n
        x = ML + LBL_W + LBL_G
        for txt, fill, stk in items:
            _box(s, x, y, iw, LH, txt, fill=fill, stk=stk, sz=7, fc=CHAR)
            x += iw + igap
        if idx < len(layers) - 1:
            _darr(s, ML + CW // 2, y + LH + Inches(0.01))

    # [C] 서술 테이블
    dy = y0 + 6 * (LH + GAP) + Inches(0.05)
    _tbl(s, ML, dy, CW, Inches(1.45), [
        ['프론트엔드 기술스택', '백엔드 기술스택', '인프라 / 보안 / 연계'],
        ['• React 18.x + TypeScript 5.x\n'
         '• Ant Design 5.24 (Tree/Steps/Upload)\n'
         '• Zustand 5 + TanStack Query 5\n'
         '• React Router 7 + Vite 5\n'
         '• React Hook Form + Zod\n'
         '• 10개 전용 페이지, 17개 메뉴',
         '• Spring Boot 3.x (추후 개발)\n'
         '• MSW 2.12 Mock API (현재)\n'
         '• Faker.js 9 한국어 데이터\n'
         '• REST API ~33개 endpoint\n'
         '• JWT + Redis 인증 (추후)\n'
         '• MyBatis 3.5 ORM (추후)',
         '• PostgreSQL 15 (추후 개발)\n'
         '• Linux / Tomcat 10.x (추후)\n'
         '• SSL/TLS 1.3, GPKI (추후)\n'
         '• Phase 1 공통 (게시판·권한)\n'
         '• SYS09 PrintableReport 공유\n'
         '• 사용자: 4역할, 59 프로세스'],
    ], cws=[CW // 3] * 3)


# ══════════════════════════════════════════════════
# Page 2: 목표모델 (As-Is → To-Be)
# ══════════════════════════════════════════════════

def p2_target(prs):
    s = _sl(prs)
    _title(s, '목표모델 (As-Is → To-Be)',
           '현행 문제점 진단 및 목표 시스템 기능 구성', 2)

    pairs = [
        ('부대 이력 수기 관리\n종이 대장, 창설/해체 이력 분실, 검색 불가',
         '디지털 부대기록부\nSearchForm 2필드 + DataTable 7컬럼 CRUD'),
        ('계승관계 파악 불가\n수기 관리, 계층 미표현, 사진 미보관',
         'Tree 3단 계층 관리\n함대→전단→함정 + CrudForm 11필드 + Upload'),
        ('보직이력 수기 기록\n역대 직위자 추적 곤란, 직무대리 미관리',
         '주요직위자 통합 관리\nCrudForm 16필드 + Timeline 이력 + R6'),
        ('주요활동 엑셀 관리\n결재 미비, 비밀등급 미관리, 일괄등록 불가',
         '활동 관리 + 3단계 결재\nSteps 결재 + BulkUpload + 비밀등급 Switch'),
        ('통계/권한 수동 관리\n수동 집계, 무분별 접근, 인쇄 불편',
         '10종 통계 + 권한 워크플로우\nColumn/Pie 차트 + 3p 권한 + PrintableReport'),
    ]

    y0 = Inches(0.82)
    RH = Inches(0.44)
    RG = Inches(0.06)
    AW = Inches(3.5)
    ARR_W = Inches(0.22)
    AX = ML
    ARX = ML + AW + Inches(0.15)
    TX = ARX + ARR_W + Inches(0.15)

    _box(s, AX, y0 - Inches(0.22), AW, Inches(0.18),
         '현행 시스템 (As-Is)', fill=ERR, stk=ERR, sz=8, b=True, fc=WHITE)
    _box(s, TX, y0 - Inches(0.22), AW, Inches(0.18),
         '목표 시스템 (To-Be)', fill=GRN, stk=GRN, sz=8, b=True, fc=WHITE)

    for i, (asis, tobe) in enumerate(pairs):
        y = y0 + i * (RH + RG)
        _box(s, AX, y, AW, RH, asis, fill=PEACH, stk=ERR, sz=7, fc=CHAR)
        _rarr(s, ARX, y + RH // 2, ARR_W)
        _box(s, TX, y, AW, RH, tobe, fill=MINT, stk=GRN, sz=7, fc=CHAR)

    # 우측 기능 요약
    FX = TX + AW + Inches(0.12)
    FW = ML + CW - FX
    _box(s, FX, y0, FW, 5 * RH + 4 * RG,
         '핵심 기능\n─────────\n'
         '부대기록부   1\n제원/계승     1\n주요직위자   1\n'
         '주요활동     1\n활동결재     1\n부대기/마크  1\n'
         '통계출력     1\n권한관리     3\n'
         '─────────\n총 17개 메뉴\n59 프로세스\n'
         '─────────\nTree 3단계\nSteps 3단계\n4종 인쇄',
         fill=ICE, stk=TEAL, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)

    # [C] 서술 테이블
    dy = y0 + 5 * (RH + RG) + Inches(0.12)
    _tbl(s, ML, dy, CW, Inches(1.25), [
        ['개선항목', '현행 문제점', '목표 개선방향', '기대 효과', '정량 목표'],
        ['부대기록', '종이 대장, 검색 불가', 'SearchForm+DataTable CRUD', '즉시 검색·추적', '검색 1초'],
        ['계승관계', '수기, 계층 미표현', 'Tree 3단+CrudForm 11필드', '계층 시각화', '3뎁스 트리'],
        ['직위자', '이력 추적 곤란', '16필드 CRUD+Timeline', '역대 이력 관리', '이력 100%'],
        ['주요활동', '엑셀, 결재 없음', 'Steps 3단계+BulkUpload', '전자결재 구현', '3단계 결재'],
        ['통계/권한', '수동 집계, 무분별', '10종 차트+3p 권한관리', '현황 즉시 파악', '완료율 표시'],
    ], cws=[int(CW * r) for r in [0.10, 0.24, 0.24, 0.22, 0.10]])


# ══════════════════════════════════════════════════
# Page 3: 유스케이스
# ══════════════════════════════════════════════════

def p3_usecase(prs):
    s = _sl(prs)
    _title(s, '유스케이스 다이어그램',
           '4개 액터 x 8개 유스케이스 그룹, 59개 단위 프로세스', 3)

    y0 = Inches(0.82)

    # 좌측: 액터
    actors = [
        '계보담당\n부대기록·계승·직위자\nROLE_LINEAGE',
        '중간결재자\n주요활동 1차 검토\nROLE_REVIEWER',
        '확인관\n주요활동 최종 승인\nROLE_APPROVER',
        '시스템관리자\nROLE_ADMIN',
    ]
    AX = ML
    AW = Inches(1.15)
    AH = Inches(0.48)
    AG = Inches(0.10)
    for i, atxt in enumerate(actors):
        _box(s, AX, y0 + i * (AH + AG), AW, AH, atxt,
             fill=ICE, stk=NAVY, sz=7, fc=CHAR)

    # 중앙: 시스템 바운더리
    BX = AX + AW + Inches(0.2)
    BW = Inches(6.2)
    BH = 4 * Inches(0.4) + 3 * Inches(0.06) + Inches(0.35)
    bnd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, BX, y0 - Inches(0.05), BW, BH)
    bnd.fill.solid()
    bnd.fill.fore_color.rgb = RGBColor(0xFB, 0xFB, 0xFB)
    bnd.line.color.rgb = NAVY
    bnd.line.width = Pt(2)
    _tb(s, BX + Inches(0.1), y0 - Inches(0.02), Inches(5), Inches(0.18),
        '부대계보관리체계 (SYS08) — 바운더리', 8, True, NAVY)

    # UC 그룹
    uc_groups = [
        ('부대기록부 (6)', 'UC-LIN-001~006\nSearchForm 2+DataTable 7+CRUD 6필드', ICE, SKY),
        ('제원/계승부대 (11)', 'UC-LIN-007~017\nTree 3단+DataTable 5+CrudForm 11', ICE, SKY),
        ('주요직위자 (5)', 'UC-LIN-018~022\n16필드+Timeline 이력+R6 군번/계급', ICE, SKY),
        ('주요활동 관리 (10)', 'UC-LIN-023~032\nBulkUpload 엑셀+비밀등급 Switch', ICE, SKY),
        ('활동결재 (6)', 'UC-LIN-033~038\nSteps 3단계 승인/반려', ICE, SKY),
        ('부대기/마크 (4)', 'UC-LIN-039~042\nUpload.Dragger+Image 미리보기', ICE, SKY),
        ('권한관리 (8)', 'UC-LIN-043~050\n신청/관리/조회 3p 워크플로우', LBLUE, TEAL),
        ('통계+게시판 (9)', 'UC-LIN-051~059\n10종 차트+게시판 2+관리자 5', LBLUE, TEAL),
    ]
    UY0 = y0 + Inches(0.22)
    UW = Inches(2.95)
    UH = Inches(0.4)
    UG = Inches(0.06)
    for i, (grp, desc, fill, stk) in enumerate(uc_groups):
        r, c = divmod(i, 2)
        ux = BX + Inches(0.08) + c * (UW + Inches(0.1))
        uy = UY0 + r * (UH + UG)
        sh = _box(s, ux, uy, UW, UH, f'{grp}\n{desc}',
                  fill=fill, stk=stk, sz=6, b=False, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True

    # 우측: 외부 시스템
    EX = BX + BW + Inches(0.12)
    EW = ML + CW - EX
    _box(s, EX, y0 + Inches(0.2), EW, Inches(0.35),
         'Phase 1 공통\n게시판·코드·권한', fill=SNOW, stk=ASH, sz=7, fc=ASH)
    _box(s, EX, y0 + Inches(0.65), EW, Inches(0.35),
         'SYS09 PrintableReport\n인쇄 공유', fill=SNOW, stk=ASH, sz=7, fc=ASH)

    # [C] 서술 테이블
    dy = Inches(3.2)
    _tbl(s, ML, dy, CW, Inches(2.1), [
        ['UC-ID', '유스케이스명', '주 액터', 'API/화면 (구체적)', '우선'],
        ['UC-LIN-001', '부대기록부 관리', '계보담당',
         'UnitRecordPage — SearchForm(2필드)+DataTable(7컬럼)+CrudForm(6필드) Modal', '상'],
        ['UC-LIN-007', '제원/계승부대', '계보담당',
         'UnitLineageTreePage — Tree 3단(함대→전단→함정)+DataTable(5컬럼)+CrudForm(11필드)+Upload', '상'],
        ['UC-LIN-018', '주요직위자 관리', '계보담당',
         'UnitKeyPersonPage — DataTable(9컬럼)+CrudForm(16필드)+Timeline 이력+militaryPersonColumn', '상'],
        ['UC-LIN-023', '주요활동 관리', '계보담당',
         'UnitActivityPage — DataTable(8컬럼)+CrudForm(10필드)+BulkUpload(.xlsx)+비밀등급 Switch', '상'],
        ['UC-LIN-033', '활동결재', '결재권자',
         'ApprovalPage — Tabs(대기/전체)+Steps 3단계(계보담당→중간→확인관)+승인/반려 Modal', '상'],
        ['UC-LIN-039', '부대기/마크', '계보담당',
         'UnitFlagPage — DataTable(7컬럼, 썸네일)+CrudForm(8필드)+Upload.Dragger 이미지+미리보기', '중'],
        ['UC-LIN-043', '권한관리', '전체',
         'AuthRequest(Card 3열 폼)+AuthMgmt(승인/반려)+AuthView(Tabs 내권한/부대별)', '중'],
    ], cws=[int(CW * r) for r in [0.10, 0.12, 0.06, 0.66, 0.05]])


# ══════════════════════════════════════════════════
# Page 4: 스토리보드 (화면흐름도)
# ══════════════════════════════════════════════════

def p4_storyboard(prs):
    s = _sl(prs)
    _title(s, '스토리보드 (화면흐름도)',
           '총 17개 메뉴, 10개 전용 화면 + 게시판 2 + 관리자 5, 핵심 네비게이션 흐름', 4)

    y0 = Inches(0.82)

    # 로그인 → 대시보드
    _box(s, ML, y0, Inches(1.3), Inches(0.35),
         '로그인\n/login', fill=WARN, stk=WARN, sz=8, b=True, fc=WHITE)
    _rarr(s, ML + Inches(1.35), y0 + Inches(0.12), Inches(0.25))
    _box(s, ML + Inches(1.65), y0, Inches(2.0), Inches(0.35),
         '대시보드 / 홈\n/sys08', fill=NAVY, stk=NAVY, sz=8, b=True, fc=WHITE)

    _darr(s, ML + Inches(2.65), y0 + Inches(0.37))

    # 6개 메뉴 그룹
    groups = [
        ('1. 부대기록부', '부대기록부 CRUD', '/sys08/8/1', '1', ICE, SKY),
        ('2. 제원/계승', 'Tree 3단 · 계승관계', '/sys08/5/1', '1', ICE, SKY),
        ('3. 주요직위자', 'CrudForm 16필드 · Timeline', '/sys08/4/1', '1', ICE, SKY),
        ('4. 주요활동/결재', '활동관리 · 결재(Steps)', '/sys08/3/*', '2', ICE, SKY),
        ('5. 부대기/통계', '부대기마크 · 10종 통계', '/sys08/6,7/*', '2', LBLUE, TEAL),
        ('6. 권한+게시+관리', '권한3p · 게시판2 · 관리자5', '/sys08/1,9/*', '10', LBLUE, TEAL),
    ]

    GY0 = y0 + Inches(0.55)
    GW = Inches(4.5)
    GH = Inches(0.5)
    GG = Inches(0.06)
    GCOL_GAP = Inches(0.15)

    for i, (name, screens, path, cnt, fill, stk) in enumerate(groups):
        r, c = divmod(i, 2)
        gx = ML + c * (GW + GCOL_GAP)
        gy = GY0 + r * (GH + GG)
        sh = _box(s, gx, gy, GW, GH,
                  f'{name} ({cnt}개)\n{screens}\n{path}',
                  fill=fill, stk=stk, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        pars = list(sh.text_frame.paragraphs)
        if len(pars) >= 3:
            for r_obj in pars[-1].runs:
                r_obj.font.color.rgb = ASH
                r_obj.font.size = Pt(6)

    # 핵심 화면 상세 (3개)
    DY = GY0 + 3 * (GH + GG) + Inches(0.08)
    det_boxes = [
        ('제원/계승부대\n/sys08/5/1',
         'Row 8:16 레이아웃\n좌: Card Tree(3뎁스)\n'
         '  함대→전단→함정\n  defaultExpandAll\n'
         '우: DataTable(5컬럼)\nCrudForm(11필드) Modal\n'
         '사진 Upload 첨부',
         ICE, SKY),
        ('주요활동+결재\n/sys08/3/1, 3/2',
         'SearchForm(2필드)\nDataTable(8컬럼)\n'
         'BulkUpload(.xlsx 일괄등록)\n비밀등급 Switch+Select\n'
         'Steps 3단계 결재\n승인(finish)/반려(error)\n'
         'ApprovalDetailModal',
         MINT, GRN),
        ('부대기/마크+직위자\n/sys08/6/1, 4/1',
         'DataTable(7컬럼, 썸네일 50×50)\n'
         'Upload.Dragger → Image 미리보기\n'
         '부대기/마크 2종 분류\n──────────\n'
         'CrudForm(16필드) 직위자\nTimeline 역대 이력 Modal\n'
         'militaryPersonColumn(R6)',
         CREAM, WARN),
    ]
    DW = Inches(2.95)
    DH = Inches(1.1)
    DG = Inches(0.1)
    for i, (ttl, desc, fill, stk) in enumerate(det_boxes):
        dx = ML + i * (DW + DG)
        sh = _box(s, dx, DY, DW, DH,
                  f'{ttl}\n──────────\n{desc}',
                  fill=fill, stk=stk, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        sh.text_frame.paragraphs[1].runs[0].font.bold = True

    # [C] 서술 테이블 (좌+우)
    dy = DY + DH + Inches(0.08)

    LW = int(CW * 0.48)
    _tbl(s, ML, dy, LW, Inches(0.95), [
        ['영역', '화면 수', 'URL 패턴', '주요 컴포넌트'],
        ['부대기록부', '1', '/sys08/8/1', 'SearchForm+DataTable+CrudForm'],
        ['계승부대', '1', '/sys08/5/1', 'Tree+DataTable+CrudForm 11필드'],
        ['주요직위자', '1', '/sys08/4/1', 'DataTable+CrudForm 16+Timeline'],
        ['활동/결재', '2', '/sys08/3/1,2', 'DataTable+Steps+BulkUpload'],
        ['기마크/통계', '2', '/sys08/6,7/1', 'Upload.Dragger+Charts+Progress'],
        ['권한+게시+관리', '10', '/sys08/1,9/*', 'Card+Tabs+AdminRoutes'],
    ], cws=[int(LW * r) for r in [0.20, 0.10, 0.28, 0.42]])

    RX = ML + LW + Inches(0.15)
    RW = CW - LW - Inches(0.15)
    _tbl(s, RX, dy, RW, Inches(0.95), [
        ['핵심 흐름 시나리오'],
        ['시나리오 1: 계보 등록·결재 흐름\n'
         '① 로그인 → 대시보드\n'
         '② [주요활동] > [활동관리] 메뉴\n'
         '③ 주요활동 등록 (개별/일괄 엑셀)\n'
         '④ 비밀등급 Switch → 등급 Select\n'
         '⑤ 결재 요청 → Steps 3단계 진행\n'
         '⑥ 승인/반려 → 상태 업데이트'],
        ['시나리오 2: 부대 계승 관리 흐름\n'
         '① [제원/계승부대] 메뉴 선택\n'
         '② Tree 3단 부대 선택 → 필터링\n'
         '③ CrudForm 11필드 등록/수정\n'
         '④ 사진 Upload → 저장'],
    ])


# ══════════════════════════════════════════════════
# Page 5: 플로우차트 (업무처리흐름)
# ══════════════════════════════════════════════════

def p5_flowchart(prs):
    s = _sl(prs)
    _title(s, '업무처리 플로우차트',
           '핵심 프로세스: 부대기록 → 활동등록 → 결재 — 4개 레인, 12단계', 5)

    y0 = Inches(0.82)

    # 4개 스윔레인
    lanes = [
        ('계보담당 / 결재권자', ICE, NAVY),
        ('프론트엔드 (React 18)', LBLUE, SKY),
        ('백엔드 (API/Mock)', MINT, GRN),
        ('DB / 외부', CREAM, WARN),
    ]
    LN_W = CW // 4
    LN_H = Inches(2.75)
    LN_HDR = Inches(0.22)

    for i, (name, fill, stk) in enumerate(lanes):
        lx = ML + LN_W * i
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = fill
        bg.line.color.rgb = stk
        bg.line.width = Pt(1)
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_HDR)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = stk
        hdr.line.fill.background()
        tf = hdr.text_frame
        _run(tf.paragraphs[0], name, 7, True, WHITE)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    steps = [
        (0, 0.30, '① 부대기록부 접근\n좌측 메뉴 선택'),
        (1, 0.30, '② UnitRecordPage\nSearchForm+DataTable 렌더'),
        (0, 0.70, '③ Tree 부대 선택\n계승관계 메뉴 클릭'),
        (1, 0.70, '④ Tree 렌더+필터\nGET /unit-tree 호출'),
        (2, 0.70, '⑤ 계승관계 조회\nGET /lineage?unit=X'),
        (3, 0.70, '⑥ Lineage SELECT\n조건별 조회 (Mock 30건)'),
        (1, 1.10, '⑦ CrudForm 11필드\n등록/수정 Modal'),
        (0, 1.10, '⑧ 주요활동 등록\n개별/일괄(BulkUpload)'),
        (0, 1.50, '⑨ 결재 요청 → 결재권자\nSteps 3단계 진입'),
        (1, 1.50, '⑩ ApprovalPage 렌더\nSteps+승인/반려 Modal'),
        (2, 1.90, '⑪ PUT /approve·reject\n결재상태 업데이트'),
        (3, 1.90, '⑫ Activity UPDATE\n결재상태 저장+이력'),
    ]

    SBW = LN_W - Inches(0.12)
    SBH = Inches(0.32)

    for lane, y_off, txt in steps:
        sx = ML + LN_W * lane + Inches(0.06)
        sy = y0 + LN_HDR + Inches(y_off)
        _box(s, sx, sy, SBW, SBH, txt,
             fill=WHITE, stk=ASH, sw=1, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)

    # 판단 분기
    fx = ML + LN_W * 1 + Inches(0.06)
    fy = y0 + LN_HDR + Inches(1.35)
    _box(s, fx, fy, SBW, Inches(0.12),
         '반려 → 반려사유 Modal → 수정 재등록',
         fill=PEACH, stk=ERR, sw=1, sz=5, fc=ERR, al=PP_ALIGN.CENTER)

    # [C] 서술 테이블
    dy = y0 + LN_H + Inches(0.1)

    LW = int(CW * 0.52)
    _tbl(s, ML, dy, LW, Inches(1.42), [
        ['단계', '동작', '조건/상세'],
        ['①②', '기록부접근 → 목록', 'SearchForm 2필드 + DataTable 7컬럼 + CrudForm 6필드'],
        ['③④', 'Tree선택 → 계승조회', 'Tree 3단(함대→전단→함정) → GET /unit-tree'],
        ['⑤⑥', '서버조회 → DB검색', 'lineage?unit=X + Mock 30건, 12필드 반환'],
        ['⑦⑧', 'CrudForm → 활동등록', '11필드 등록/수정 + BulkUpload(.xlsx 검증→저장)'],
        ['⑨⑩', '결재요청 → 처리', 'Steps 3단계 (계보담당→중간결재→확인관)'],
        ['⑪⑫', '결재저장 → DB', 'PUT /approve·reject + 결재상태 저장'],
    ], cws=[int(LW * r) for r in [0.08, 0.28, 0.64]])

    RX = ML + LW + Inches(0.12)
    RW = CW - LW - Inches(0.12)
    _tbl(s, RX, dy, RW, Inches(1.42), [
        ['업무 규칙 / 예외 처리'],
        ['[업무 규칙]\n'
         'BR-LIN-001: Tree 3단계 (함대→전단→함정)\n'
         'BR-LIN-002: Steps 3단계 결재 (계보담당→중간→확인관)\n'
         'BR-LIN-003: 비밀등급 3종 Switch (비밀/대외비/일반)\n'
         'BR-LIN-004: BulkUpload 2단계 (검증→저장)\n'
         'BR-LIN-005: 부대기/마크 2종 분류\n'
         'BR-LIN-006: 권한 3p 워크플로우 (신청→관리→조회)\n'
         'BR-LIN-007: militaryPersonColumn(R6) 군번/계급/성명\n'
         'BR-LIN-008: PrintableReport 4종 인쇄'],
        ['[예외 처리]\n'
         'E-001: 세션만료(401) → 로그인 리다이렉트\n'
         'E-002: 권한없음(403) → 접근 제한 알림\n'
         'E-003: 반려 시 → 반려사유 Modal + 수정 재등록\n'
         'E-004: BulkUpload 검증실패 → 오류 Table 표시'],
    ])


# ══════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════

def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    print('[1/5] 시스템 개념도...')
    p1_concept(prs)
    print('[2/5] 목표모델...')
    p2_target(prs)
    print('[3/5] 유스케이스...')
    p3_usecase(prs)
    print('[4/5] 스토리보드...')
    p4_storyboard(prs)
    print('[5/5] 플로우차트...')
    p5_flowchart(prs)

    out_path = os.path.join(OUT, 'SYS08_부대계보관리체계_분석설계_v3.pptx')
    prs.save(out_path)
    print(f'\n완료! 총 {len(prs.slides)}장 슬라이드')
    print(f'파일: {out_path}')


if __name__ == '__main__':
    main()
