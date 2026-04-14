#!/usr/bin/env python3
"""
SYS09 영현보훈체계 기능명세서 PPTX v4
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
system_level_spec_guide_v4.md 기반 핵심 5페이지:
  P1. 시스템 개념도      P2. 목표모델 (As-Is → To-Be)
  P3. 유스케이스          P4. 스토리보드 (화면흐름도)
  P5. 플로우차트 (업무처리흐름)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ═══ 색상 팔레트 ═══
NAVY  = RGBColor(0x6B, 0x8E, 0xA8)
TEAL  = RGBColor(0x8E, 0xC5, 0xC0)
SKY   = RGBColor(0xA8, 0xC8, 0xE8)
MINT  = RGBColor(0xDC, 0xF5, 0xE4)
PEACH = RGBColor(0xFC, 0xD4, 0xD0)
CREAM = RGBColor(0xFD, 0xEE, 0xC2)
ICE   = RGBColor(0xE4, 0xEE, 0xF8)
SNOW  = RGBColor(0xF7, 0xF9, 0xFC)
CHAR  = RGBColor(0x4A, 0x4A, 0x4A)
ASH   = RGBColor(0x99, 0x99, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRN   = RGBColor(0xA4, 0xD4, 0xB4)
WARN  = RGBColor(0xF5, 0xCC, 0xA0)
ERR   = RGBColor(0xF0, 0xB2, 0xAE)
LBLUE = RGBColor(0xDC, 0xEA, 0xFE)
LSUB  = RGBColor(0xC5, 0xDC, 0xEC)

# ═══ 치수 ═══
SW = Inches(10)
SH = Inches(5.625)
ML = Inches(0.4)
CW = Inches(9.2)
FONT = 'Freesentation 5 Medium'
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'ppt')


# ══════════════════════════════════════════════════
# 헬퍼 함수
# ══════════════════════════════════════════════════

def _sl(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = SNOW
    return s


def _run(p, txt, sz=9, b=False, c=CHAR):
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = b
    r.font.color.rgb = c
    r.font.name = FONT
    return r


def _tb(s, l, t, w, h, txt, sz=9, b=False, c=CHAR, al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    _run(tb.text_frame.paragraphs[0], txt, sz, b, c)
    tb.text_frame.paragraphs[0].alignment = al
    return tb


def _title(s, title, sub, pg):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.68))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    _tb(s, ML, Inches(0.06), Inches(7.5), Inches(0.34), title, 22, True, WHITE)
    if sub:
        _tb(s, ML, Inches(0.38), Inches(8.5), Inches(0.26), sub, 11, False, LSUB)
    _tb(s, Inches(9.0), Inches(0.15), Inches(0.6), Inches(0.3),
        f'{pg}/5', 12, True, WHITE, PP_ALIGN.RIGHT)


def _box(s, l, t, w, h, txt, fill=ICE, stk=NAVY, sw=1,
         sz=9, b=False, fc=CHAR, al=PP_ALIGN.CENTER):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = stk
    sh.line.width = Pt(sw)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.margin_left = tf.margin_right = Inches(0.03)
    for i, ln in enumerate(txt.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p, ln, sz, b, fc)
        p.alignment = al
        p.space_before = p.space_after = Pt(0)
    bp = sh._element.txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', 'ctr')
    return sh


def _darr(s, cx, y, h=Inches(0.08)):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                           cx - Inches(0.06), y, Inches(0.12), h)
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _rarr(s, x, cy, w=Inches(0.18)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                           x, cy - Inches(0.05), w, Inches(0.1))
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _cell_bg(cell, hx):
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for sf in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(sf)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), val=hx)


def _tbl(s, l, t, w, h, data, cws=None):
    rows, cols = len(data), len(data[0])
    sh = s.shapes.add_table(rows, cols, l, t, w, h)
    tbl = sh.table
    if cws:
        for i, cw in enumerate(cws):
            tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ''
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_top = tf.margin_bottom = Inches(0.01)
            tf.margin_left = Inches(0.04)
            tf.margin_right = Inches(0.03)
            hdr = (r == 0)
            for i, ln in enumerate(data[r][c].split('\n')):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                _run(p, ln, 7 if not hdr else 8, hdr, WHITE if hdr else CHAR)
                p.alignment = PP_ALIGN.CENTER if hdr else PP_ALIGN.LEFT
                p.space_before = p.space_after = Pt(0)
            _cell_bg(cell, '6B8EA8' if hdr else ('FFFFFF' if r % 2 == 1 else 'E4EEF8'))
    return tbl


# ══════════════════════════════════════════════════
# Page 1: 시스템 개념도
# ══════════════════════════════════════════════════

def p1_concept(prs):
    s = _sl(prs)
    _title(s, '시스템 개념도',
           '영현보훈체계의 전체 아키텍처 및 기술 구성 — 6개 계층', 1)

    layers = [
        ('사용자\n계층', [
            ('보훈담당관\n사망자/상이자 정보 등록·관리', ICE, NAVY),
            ('심사위원\n전공사상심사 등록·판정', ICE, NAVY),
            ('시스템관리자\n코드·권한·시스템 관리', ICE, NAVY),
        ]),
        ('채널\n계층', [
            ('웹 브라우저\nChrome/Edge 1920x1080', ICE, TEAL),
            ('반응형 모바일 웹\n360px~', ICE, TEAL),
            ('API Gateway\nRESTful 30+ endpoints', ICE, TEAL),
        ]),
        ('프론트\n엔드', [
            ('React 18 + TS 5\nVite 5 · Router 7', LBLUE, SKY),
            ('Ant Design 5.24\nTable · Charts · Modal', LBLUE, SKY),
            ('Zustand 5 + TQ 5\n상태관리 · 서버캐싱', LBLUE, SKY),
        ]),
        ('백엔드\n계층', [
            ('Spring Boot 3.x\n(추후 개발) Java 17', MINT, GRN),
            ('MSW 2.12 Mock\nFaker.js 9 데이터', MINT, GRN),
            ('REST API 설계\nJSON/HTTPS 30+개', MINT, GRN),
        ]),
        ('데이터\n계층', [
            ('PostgreSQL 15\n(추후) 정규화 설계', CREAM, WARN),
            ('Mock 데이터\n사망자·상이자·심사건', CREAM, WARN),
            ('파일 스토리지\n확인서 PDF (추후)', CREAM, WARN),
        ]),
        ('외부\n연계', [
            ('GPKI/SSO 인증\n(추후 연계)', SNOW, ASH),
            ('국가보훈처 시스템\n(추후 연계)', SNOW, ASH),
            ('알림톡/SMS\n(추후 연계)', SNOW, ASH),
        ]),
    ]

    y0 = Inches(0.82)
    LH = Inches(0.37)
    GAP = Inches(0.1)
    LBL_W = Inches(0.85)
    LBL_G = Inches(0.05)

    for idx, (lbl, items) in enumerate(layers):
        y = y0 + idx * (LH + GAP)
        _box(s, ML, y, LBL_W, LH, lbl, fill=NAVY, stk=NAVY, sz=7, b=True, fc=WHITE)
        n = len(items)
        avail = CW - LBL_W - LBL_G
        igap = Inches(0.04)
        iw = (avail - igap * (n - 1)) // n
        x = ML + LBL_W + LBL_G
        for txt, fill, stk in items:
            _box(s, x, y, iw, LH, txt, fill=fill, stk=stk, sz=7, fc=CHAR)
            x += iw + igap
        if idx < len(layers) - 1:
            _darr(s, ML + CW // 2, y + LH + Inches(0.01))

    # [C] 서술 테이블
    dy = y0 + 6 * (LH + GAP) + Inches(0.05)
    _tbl(s, ML, dy, CW, Inches(1.45), [
        ['프론트엔드 기술스택', '백엔드 기술스택', '인프라 / 보안 / 연계'],
        ['• React 18.x + TypeScript 5.x\n'
         '• Ant Design 5.24 (Table/Modal/Charts)\n'
         '• Zustand 5 + TanStack Query 5\n'
         '• React Router 7 + Vite 5\n'
         '• @ant-design/charts (Bar)\n'
         '• 18개 고유 페이지, 20+ 화면',
         '• Spring Boot 3.x (추후 개발)\n'
         '• MSW 2.12 Mock API (현재)\n'
         '• Faker.js 9 한국어 데이터\n'
         '• REST API 30+ endpoint\n'
         '• JWT + Redis 인증 (추후)\n'
         '• MyBatis 3.5 ORM (추후)',
         '• PostgreSQL 15 (추후 개발)\n'
         '• Linux / Tomcat 10.x (추후)\n'
         '• SSL/TLS 1.3, GPKI (추후)\n'
         '• 국가보훈처 시스템 (추후)\n'
         '• 알림톡/SMS/이메일 (추후)\n'
         '• 사용자: 3역할, 35 프로세스'],
    ], cws=[CW // 3] * 3)


# ══════════════════════════════════════════════════
# Page 2: 목표모델 (As-Is → To-Be)
# ══════════════════════════════════════════════════

def p2_target(prs):
    s = _sl(prs)
    _title(s, '목표모델 (As-Is → To-Be)',
           '현행 문제점 진단 및 목표 시스템 기능 구성', 2)

    pairs = [
        ('수기 사망자 관리\n엑셀/수기장부, 최신 데이터 불일치',
         '통합 사망자 DB\nCRUD+검색(7필드)+StatusBadge'),
        ('상이자 별도 관리\n보훈청별 개별 관리, 이력 추적 불가',
         '상이자 통합 관리\nCRUD+등급관리+전공상여부 추적'),
        ('심사 이력 수기 기록\n전공사상심사 결과 문서 분산',
         '심사 통합 CRUD\n심사차수/결과/분류 DB화'),
        ('확인서 수기 발급\n순직/사망/유공자 확인서 수작업',
         '확인서 자동 발급\nModal+PrintableReport 4종'),
        ('통계 수작업 집계\n부대별/연도별/월별 통계 부재',
         '자동 통계 생성\nBar차트+Table 6종 통계'),
    ]

    y0 = Inches(0.82)
    RH = Inches(0.44)
    RG = Inches(0.06)
    AW = Inches(3.5)
    ARR_W = Inches(0.22)
    AX = ML
    ARX = ML + AW + Inches(0.15)
    TX = ARX + ARR_W + Inches(0.15)

    _box(s, AX, y0 - Inches(0.22), AW, Inches(0.18),
         '현행 시스템 (As-Is)', fill=ERR, stk=ERR, sz=8, b=True, fc=WHITE)
    _box(s, TX, y0 - Inches(0.22), AW, Inches(0.18),
         '목표 시스템 (To-Be)', fill=GRN, stk=GRN, sz=8, b=True, fc=WHITE)

    for i, (asis, tobe) in enumerate(pairs):
        y = y0 + i * (RH + RG)
        _box(s, AX, y, AW, RH, asis, fill=PEACH, stk=ERR, sz=7, fc=CHAR)
        _rarr(s, ARX, y + RH // 2, ARR_W)
        _box(s, TX, y, AW, RH, tobe, fill=MINT, stk=GRN, sz=7, fc=CHAR)

    # 우측 기능 요약
    FX = TX + AW + Inches(0.12)
    FW = ML + CW - FX
    _box(s, FX, y0, FW, 5 * RH + 4 * RG,
         '핵심 기능\n─────────\n'
         '사망자관리    1\n상이자관리    1\n심사관리       1\n'
         '확인서4종    4\n보고서2종    2\n통계6종       6\n'
         '─────────\n총 20개 메뉴\n35 프로세스\n'
         '─────────\n발급대장 1종\n게시판 2개',
         fill=ICE, stk=TEAL, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)

    # [C] 서술 테이블
    dy = y0 + 5 * (RH + RG) + Inches(0.12)
    _tbl(s, ML, dy, CW, Inches(1.25), [
        ['개선항목', '현행 문제점', '목표 개선방향', '기대 효과', '정량 목표'],
        ['사망자관리', '수기장부, 데이터 불일치', 'SearchForm(7)+DataTable+CrudForm(17)', '즉시 검색/등록', '검색 1초'],
        ['상이자관리', '보훈청별 개별 관리', 'CRUD+등급(7급)+보훈청 선택', '이력 통합 관리', '통합 100%'],
        ['심사관리', '심사결과 분산 관리', '심사차수+결과(해당/비해당)+분류', '심사 이력 DB화', '이력 100%'],
        ['확인서발급', '수작업 발급, 양식 불일치', 'Modal+PrintableReport 4종 자동', '발급 시간 단축', '발급 3초'],
        ['통계현황', '수작업 집계, 오류 빈발', 'Bar차트+Table 6종 자동 집계', '실시간 통계', '자동 100%'],
    ], cws=[int(CW * r) for r in [0.10, 0.24, 0.24, 0.22, 0.10]])


# ══════════════════════════════════════════════════
# Page 3: 유스케이스
# ══════════════════════════════════════════════════

def p3_usecase(prs):
    s = _sl(prs)
    _title(s, '유스케이스 다이어그램',
           '3개 액터 x 18개 유스케이스, 35개 단위 프로세스', 3)

    y0 = Inches(0.82)

    # 좌측: 액터
    actors = [
        '보훈담당관\n사망자/상이자 CRUD\nROLE_HONOR_ADMIN',
        '심사위원\n전공사상심사 등록/판정\nROLE_REVIEWER',
        '시스템관리자\nROLE_ADMIN',
    ]
    AX = ML
    AW = Inches(1.15)
    AH = Inches(0.55)
    AG = Inches(0.15)
    for i, atxt in enumerate(actors):
        _box(s, AX, y0 + i * (AH + AG), AW, AH, atxt,
             fill=ICE, stk=NAVY, sz=7, fc=CHAR)

    # 중앙: 시스템 바운더리
    BX = AX + AW + Inches(0.2)
    BW = Inches(6.2)
    BH = 3 * AH + 2 * AG + Inches(0.3)
    bnd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, BX, y0 - Inches(0.05), BW, BH)
    bnd.fill.solid()
    bnd.fill.fore_color.rgb = RGBColor(0xFB, 0xFB, 0xFB)
    bnd.line.color.rgb = NAVY
    bnd.line.width = Pt(2)
    _tb(s, BX + Inches(0.1), y0 - Inches(0.02), Inches(5), Inches(0.18),
        '영현보훈체계 (SYS09) — 바운더리', 8, True, NAVY)

    # UC 그룹
    uc_groups = [
        ('자료입력 (3)', 'UC-HONOR-001~003\n사망자/상이자/심사 CRUD', ICE, SKY),
        ('확인서 발급 (4)', 'UC-HONOR-012~015\n사망·유공자(사망/상이)·심사결과', ICE, SKY),
        ('보고서 (2)', 'UC-HONOR-010~011\n사망자/상이자 현황 보고서', ICE, SKY),
        ('통계 (6)', 'UC-HONOR-004~009\n부대별/신분별/연도별/월별/명부', ICE, SKY),
        ('게시판 (2)', 'UC-HONOR-017~018\n공지사항·질의응답', LBLUE, TEAL),
        ('관리자+발급대장 (2)', 'UC-SYS 시스템·코드·권한\n확인증 발급대장', LBLUE, TEAL),
    ]
    UY0 = y0 + Inches(0.22)
    UW = Inches(2.95)
    UH = Inches(0.4)
    UG = Inches(0.06)
    for i, (grp, desc, fill, stk) in enumerate(uc_groups):
        r, c = divmod(i, 2)
        ux = BX + Inches(0.08) + c * (UW + Inches(0.1))
        uy = UY0 + r * (UH + UG)
        sh = _box(s, ux, uy, UW, UH, f'{grp}\n{desc}',
                  fill=fill, stk=stk, sz=6, b=False, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True

    # 우측: 외부 시스템
    EX = BX + BW + Inches(0.12)
    EW = ML + CW - EX
    _box(s, EX, y0 + Inches(0.2), EW, Inches(0.35),
         '국가보훈처\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)
    _box(s, EX, y0 + Inches(0.65), EW, Inches(0.35),
         'GPKI/SSO\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)

    # [C] 서술 테이블
    dy = Inches(3.2)
    _tbl(s, ML, dy, CW, Inches(2.1), [
        ['UC-ID', '유스케이스명', '주 액터', '기능 설명 (구체적)', '우선'],
        ['UC-HONOR-001', '사망자 관리', '담당관',
         'SearchForm(7필드)+DataTable(7컬럼)+CrudForm(17필드) Modal, CRUD+엑셀', '상'],
        ['UC-HONOR-002', '상이자 관리', '담당관',
         'SearchForm(6필드)+DataTable(8컬럼)+CrudForm(19필드) Modal, 등급/보훈청', '상'],
        ['UC-HONOR-003', '전공사상심사', '심사위원',
         'SearchForm(6필드)+DataTable(8컬럼)+CrudForm(16필드), 심사차수/결과/분류', '상'],
        ['UC-HONOR-012', '순직/사망확인서', '담당관',
         'DataTable+행클릭→Modal+PrintableReport, Descriptions(bordered)+확인문구', '상'],
        ['UC-HONOR-004', '부대별 사망자 현황', '담당관',
         'Bar차트+Table(6컬럼) 합계행, 전사/순직/병사/사고사 분류', '중'],
        ['UC-HONOR-009', '전사망자 명부', '담당관',
         'DataTable 전체목록, 엑셀 다운로드', '중'],
        ['UC-HONOR-016', '발급대장', '담당관',
         '확인증 발급 이력 DataTable, 발급번호/일자/대상자 추적', '중'],
    ], cws=[int(CW * r) for r in [0.10, 0.12, 0.06, 0.66, 0.05]])


# ══════════════════════════════════════════════════
# Page 4: 스토리보드 (화면흐름도)
# ══════════════════════════════════════════════════

def p4_storyboard(prs):
    s = _sl(prs)
    _title(s, '스토리보드 (화면흐름도)',
           '총 20개 화면, 5개 업무영역, 핵심 네비게이션 흐름', 4)

    y0 = Inches(0.82)

    # 로그인 → 대시보드
    _box(s, ML, y0, Inches(1.3), Inches(0.35),
         '로그인\n/login', fill=WARN, stk=WARN, sz=8, b=True, fc=WHITE)
    _rarr(s, ML + Inches(1.35), y0 + Inches(0.12), Inches(0.25))
    _box(s, ML + Inches(1.65), y0, Inches(2.0), Inches(0.35),
         '대시보드 / 홈\n/sys09', fill=NAVY, stk=NAVY, sz=8, b=True, fc=WHITE)

    _darr(s, ML + Inches(2.65), y0 + Inches(0.37))

    # 5개 메뉴 그룹
    groups = [
        ('1. 자료입력', '사망자·상이자·심사', '/sys09/2/*', '3', ICE, SKY),
        ('2. 확인서 발급', '사망·유공자(사/상)·심사결과', '/sys09/3/1~4', '4', ICE, SKY),
        ('3. 보고서', '사망자·상이자 현황', '/sys09/3/5~6', '2', ICE, SKY),
        ('4. 통계현황', '부대별/신분별/연도별/월별/명부', '/sys09/3/7~13', '7', LBLUE, TEAL),
        ('5. 게시판+관리자', '공지사항·질의응답·시스템관리', '/sys09/board,admin/*', '7+', LBLUE, TEAL),
    ]

    GY0 = y0 + Inches(0.55)
    GW = Inches(4.5)
    GH = Inches(0.5)
    GG = Inches(0.06)
    GCOL_GAP = Inches(0.15)

    for i, (name, screens, path, cnt, fill, stk) in enumerate(groups):
        r, c = divmod(i, 2)
        gx = ML + c * (GW + GCOL_GAP)
        gy = GY0 + r * (GH + GG)
        sh = _box(s, gx, gy, GW, GH,
                  f'{name} ({cnt}개)\n{screens}\n{path}',
                  fill=fill, stk=stk, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        pars = list(sh.text_frame.paragraphs)
        if len(pars) >= 3:
            for r_obj in pars[-1].runs:
                r_obj.font.color.rgb = ASH
                r_obj.font.size = Pt(6)

    # 핵심 화면 상세 (3개 대표 페이지)
    DY = GY0 + 3 * (GH + GG) + Inches(0.08)
    det_boxes = [
        ('사망자 관리\n/sys09/2/1',
         'SearchForm(7필드)\n군구분/군번/성명/주민번호\n계급/소속/사망구분\n'
         'DataTable(7컬럼)\nCrudForm(17필드) Modal\n엑셀 다운로드',
         ICE, SKY),
        ('순직/사망확인서\n/sys09/3/4',
         'SearchForm(4필드)\nDataTable 사망자 목록\n행 클릭 → Modal 팝업\n'
         'PrintableReport 확인서\nDescriptions(bordered)\n문서번호 자동 채번',
         MINT, GRN),
        ('부대별 사망자 현황\n/sys09/3/10',
         'Bar 차트 (가로)\n부대별 사망자 합계\nTable 6컬럼\n'
         '전사/순직/병사/사고사\n합계행(Summary)\n엑셀 다운로드',
         CREAM, WARN),
    ]
    DW = Inches(2.95)
    DH = Inches(1.1)
    DG = Inches(0.1)
    for i, (ttl, desc, fill, stk) in enumerate(det_boxes):
        dx = ML + i * (DW + DG)
        sh = _box(s, dx, DY, DW, DH,
                  f'{ttl}\n──────────\n{desc}',
                  fill=fill, stk=stk, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        sh.text_frame.paragraphs[1].runs[0].font.bold = True

    # [C] 서술 테이블 좌우 2개
    dy = DY + DH + Inches(0.08)

    LW = int(CW * 0.48)
    _tbl(s, ML, dy, LW, Inches(0.95), [
        ['영역', '화면 수', 'URL 패턴', '주요 컴포넌트'],
        ['자료입력', '3', '/sys09/2/1~3', 'SearchForm+DataTable+CrudForm'],
        ['확인서', '4', '/sys09/3/1~4', 'DataTable+Modal+PrintableReport'],
        ['보고서', '2', '/sys09/3/5~6', 'DataTable+PrintableReport'],
        ['통계', '7', '/sys09/3/7~13', 'Bar Charts+Table+Summary'],
        ['게시판+관리자', '4+', '/sys09/board,admin/*', 'SimpleBoardPage+AdminRoutes'],
    ], cws=[int(LW * r) for r in [0.20, 0.10, 0.28, 0.42]])

    RX = ML + LW + Inches(0.15)
    RW = CW - LW - Inches(0.15)
    _tbl(s, RX, dy, RW, Inches(0.95), [
        ['핵심 흐름 시나리오'],
        ['시나리오 1: 사망자 등록 흐름\n'
         '  로그인 -> 대시보드\n'
         '  [자료입력] > [사망자 관리] 메뉴\n'
         '  SearchForm 군구분/군번/성명 검색\n'
         '  [신규 등록] -> CrudForm Modal(17필드)\n'
         '  [저장] -> 성공 토스트 + 목록 갱신'],
        ['시나리오 2: 확인서 발급 흐름\n'
         '  [확인서] > [순직/사망확인서] 메뉴\n'
         '  DataTable에서 대상자 행 클릭\n'
         '  Modal: PrintableReport 확인서 미리보기\n'
         '  Descriptions(bordered) 상세정보\n'
         '  인쇄 버튼 -> window.print()'],
    ])


# ══════════════════════════════════════════════════
# Page 5: 플로우차트 (업무처리흐름)
# ══════════════════════════════════════════════════

def p5_flowchart(prs):
    s = _sl(prs)
    _title(s, '업무처리 플로우차트',
           '핵심 프로세스: 사망자 등록 -> 심사 -> 확인서 발급 — 4개 레인, 12단계', 5)

    y0 = Inches(0.82)

    # 4개 스윔레인
    lanes = [
        ('보훈담당관/심사위원', ICE, NAVY),
        ('프론트엔드 (React 18)', LBLUE, SKY),
        ('백엔드 (API/Mock)', MINT, GRN),
        ('DB / 외부', CREAM, WARN),
    ]
    LN_W = CW // 4
    LN_H = Inches(2.75)
    LN_HDR = Inches(0.22)

    for i, (name, fill, stk) in enumerate(lanes):
        lx = ML + LN_W * i
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = fill
        bg.line.color.rgb = stk
        bg.line.width = Pt(1)
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_HDR)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = stk
        hdr.line.fill.background()
        tf = hdr.text_frame
        _run(tf.paragraphs[0], name, 7, True, WHITE)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    steps = [
        (0, 0.30, '  사망자 관리 접근\n[자료입력] 메뉴'),
        (1, 0.30, '  SearchForm+DataTable\n7필드 검색 + 7컬럼 목록'),
        (0, 0.70, '  [신규 등록] 버튼\nCrudForm Modal 17필드'),
        (1, 0.70, '  폼 검증 + API 호출\nPOST /sys09/deceased'),
        (2, 0.70, '  사망자 정보 저장\n유효성 검증 + 응답'),
        (3, 0.70, '  Deceased INSERT\n사망자 데이터 저장'),
        (0, 1.10, '  전공사상심사 등록\n심사차수/결과/분류 입력'),
        (1, 1.10, '  CrudForm Modal\n16필드 심사 정보'),
        (2, 1.10, '  POST /sys09/reviews\n심사 결과 저장'),
        (0, 1.50, '  확인서 발급 요청\n행 클릭 -> Modal'),
        (1, 1.50, '  PrintableReport\nDescriptions+확인문구'),
        (1, 1.90, '  window.print()\n확인서 PDF 인쇄'),
    ]

    SBW = LN_W - Inches(0.12)
    SBH = Inches(0.32)

    for lane, y_off, txt in steps:
        sx = ML + LN_W * lane + Inches(0.06)
        sy = y0 + LN_HDR + Inches(y_off)
        _box(s, sx, sy, SBW, SBH, txt,
             fill=WHITE, stk=ASH, sw=1, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)

    # 판단 분기
    fx = ML + LN_W * 1 + Inches(0.06)
    fy = y0 + LN_HDR + Inches(1.35)
    _box(s, fx, fy, SBW, Inches(0.12),
         '검증 실패 -> 필드 에러 표시 -> 재입력',
         fill=PEACH, stk=ERR, sw=1, sz=5, fc=ERR, al=PP_ALIGN.CENTER)

    # [C] 서술 테이블 좌우 2개
    dy = y0 + LN_H + Inches(0.1)

    LW = int(CW * 0.52)
    _tbl(s, ML, dy, LW, Inches(1.42), [
        ['단계', '동작', '조건/상세'],
        ['1~2', '메뉴접근 -> 목록', 'SearchForm 7필드(군구분/군번/성명/주민번호/계급/소속/사망구분)'],
        ['3~4', '신규등록 -> API', 'CrudForm 17필드 Modal, POST /sys09/deceased'],
        ['5~6', '서버저장 -> DB', '유효성검증 + Deceased INSERT + 이력'],
        ['7~9', '심사등록 -> 저장', '심사차수/결과(해당/비해당)/전공상분류 4종'],
        ['10~11', '확인서 발급', '행클릭->Modal, PrintableReport+Descriptions(bordered)'],
        ['12', '인쇄 출력', 'window.print() -> PDF 확인서 발급'],
    ], cws=[int(LW * r) for r in [0.08, 0.28, 0.64]])

    RX = ML + LW + Inches(0.12)
    RW = CW - LW - Inches(0.12)
    _tbl(s, RX, dy, RW, Inches(1.42), [
        ['업무 규칙 / 예외 처리'],
        ['[업무 규칙]\n'
         'BR-HONOR-001: 사망구분 4종 (전사/순직/병사/사고사)\n'
         'BR-HONOR-002: 상이구분 2종 (전상/공상), 등급 7단계\n'
         'BR-HONOR-003: 전공상분류 4종 (전투/공무/훈련/직무상이)\n'
         'BR-HONOR-004: 심사결과 2종 (해당/비해당)\n'
         'BR-HONOR-005: 확인서 4종 (사망/유공자사망/유공자상이/심사결과)\n'
         'BR-HONOR-006: 통계 6종 (부대별/신분별/연도별/월별/명부2)\n'
         'BR-HONOR-007: 군번/계급/성명 통합 표시 (militaryPersonColumn)'],
        ['[예외 처리]\n'
         'E-001: 세션만료(401) -> 로그인 리다이렉트\n'
         'E-002: 권한없음(403) -> 접근 제한 알림\n'
         'E-003: 검증실패 -> 필드별 에러 표시+재입력\n'
         'E-004: 서버오류(500) -> 에러 메시지+재시도'],
    ])


# ══════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════

def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    print('[1/5] 시스템 개념도...')
    p1_concept(prs)
    print('[2/5] 목표모델...')
    p2_target(prs)
    print('[3/5] 유스케이스...')
    p3_usecase(prs)
    print('[4/5] 스토리보드...')
    p4_storyboard(prs)
    print('[5/5] 플로우차트...')
    p5_flowchart(prs)

    out_path = os.path.join(OUT, 'SYS09_영현보훈체계_분석설계_v4.pptx')
    prs.save(out_path)
    print(f'\n완료! 총 {len(prs.slides)}장 슬라이드')
    print(f'파일: {out_path}')


if __name__ == '__main__':
    main()
