#!/usr/bin/env python3
"""
SYS10 주말버스예약관리체계 기능명세서 PPTX v6.1 (portrait, 4페이지)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
system_level_spec_guide_v6.1.md 기반 핵심 5페이지:
  P1. 시스템 개념도      P2. 목표모델 (As-Is → To-Be)
  P3. 유스케이스          P4. 스토리보드 (화면흐름도)
  P5. 플로우차트 (업무처리흐름)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ═══ 색상 팔레트 ═══
NAVY  = RGBColor(0x6B, 0x8E, 0xA8)
TEAL  = RGBColor(0x8E, 0xC5, 0xC0)
SKY   = RGBColor(0xA8, 0xC8, 0xE8)
MINT  = RGBColor(0xDC, 0xF5, 0xE4)
PEACH = RGBColor(0xFC, 0xD4, 0xD0)
CREAM = RGBColor(0xFD, 0xEE, 0xC2)
ICE   = RGBColor(0xE4, 0xEE, 0xF8)
SNOW  = RGBColor(0xF7, 0xF9, 0xFC)
CHAR  = RGBColor(0x1A, 0x1A, 0x1A)
ASH   = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRN   = RGBColor(0xA4, 0xD4, 0xB4)
WARN  = RGBColor(0xF5, 0xCC, 0xA0)
ERR   = RGBColor(0xF0, 0xB2, 0xAE)
LBLUE = RGBColor(0xDC, 0xEA, 0xFE)
LSUB  = RGBColor(0xC5, 0xDC, 0xEC)
HBAR  = RGBColor(0xD4, 0xE2, 0xEC)  # 제목바 연한 배경
THBG  = RGBColor(0xD4, 0xE2, 0xEC)  # 표 헤더 연한 배경
BORD  = RGBColor(0x4A, 0x6B, 0x8C)  # 박스 테두리 (진한)
TBLB  = RGBColor(0x80, 0x80, 0x80)  # 표 테두리 (회색)

# ═══ 치수 ═══
SW = Inches(7.5)
SH = Inches(10)
ML = Inches(0.3)       # 좌우 마진
CW = Inches(6.9)       # 콘텐츠 폭 = 7.5 - 0.3*2
FONT = 'Freesentation 5 Medium'
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'ppt')


# ══════════════════════════════════════════════════
# 헬퍼 함수
# ══════════════════════════════════════════════════

def _sl(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = SNOW
    return s


def _run(p, txt, sz=9, b=False, c=CHAR):
    """단락에 서식 런 추가 (최소 9pt 보장)"""
    sz = max(sz, 10)
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = b
    r.font.color.rgb = c
    r.font.name = FONT
    return r


def _tb(s, l, t, w, h, txt, sz=9, b=False, c=CHAR, al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    _run(tb.text_frame.paragraphs[0], txt, sz, b, c)
    tb.text_frame.paragraphs[0].alignment = al
    return tb


def _title(s, title, sub, pg):
    """[A] 제목 영역 (portrait) — 연한 바 + 진한 테두리 + 검정 글자"""
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = HBAR
    bar.line.color.rgb = BORD
    bar.line.width = Pt(0.5)
    _tb(s, ML, Inches(0.06), Inches(5.5), Inches(0.36), title, 20, True, CHAR)
    if sub:
        _tb(s, ML, Inches(0.42), Inches(6.5), Inches(0.26), sub, 10, False, CHAR)
    _tb(s, Inches(6.5), Inches(0.18), Inches(0.6), Inches(0.3),
        f'{pg}/4', 12, True, CHAR, PP_ALIGN.RIGHT)


def _box(s, l, t, w, h, txt, fill=ICE, stk=BORD, sw=0.5,
         sz=9, b=False, fc=CHAR, al=PP_ALIGN.CENTER):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = stk
    sh.line.width = Pt(sw)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.margin_left = tf.margin_right = Inches(0.03)
    for i, ln in enumerate(txt.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p, ln, sz, b, fc)
        p.alignment = al
        p.space_before = p.space_after = Pt(0)
    bp = sh._element.txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', 'ctr')
    return sh


def _darr(s, cx, y, h=Inches(0.08)):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                           cx - Inches(0.06), y, Inches(0.12), h)
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _rarr(s, x, cy, w=Inches(0.18)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                           x, cy - Inches(0.05), w, Inches(0.1))
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _cell_bg(cell, hx):
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for sf in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(sf)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), val=hx)


def _cell_border(cell, color_hex='808080', width_pt=0.5):
    """셀 테두리 4면 회색 설정"""
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for border in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for existing in tcPr.findall(qn(border)):
            tcPr.remove(existing)
        ln = etree.SubElement(tcPr, qn(border))
        ln.set('w', str(int(width_pt * 12700)))
        ln.set('cap', 'flat')
        ln.set('cmpd', 'sng')
        ln.set('algn', 'ctr')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr'), val=color_hex)
        etree.SubElement(ln, qn('a:prstDash'), val='solid')
        etree.SubElement(ln, qn('a:round'))


def _tbl(s, l, t, w, h, data, cws=None):
    """테이블 (멀티라인 셀 지원, 진한 테두리 0.5pt, 연한 헤더 배경). data[0]=헤더"""
    rows, cols = len(data), len(data[0])
    sh = s.shapes.add_table(rows, cols, l, t, w, h)
    tbl = sh.table
    if cws:
        for i, cw in enumerate(cws):
            tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ''
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_top = tf.margin_bottom = Inches(0.04)
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.06)
            hdr = (r == 0)
            for i, ln in enumerate(data[r][c].split('\n')):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                _run(p, ln, 9 if not hdr else 10, hdr, CHAR)
                p.alignment = PP_ALIGN.CENTER if hdr else PP_ALIGN.LEFT
                p.space_before = p.space_after = Pt(0)
            _cell_bg(cell, 'D4E2EC' if hdr else ('FFFFFF' if r % 2 == 1 else 'F2F7FB'))
            _cell_border(cell, '808080', 0.5)
    return tbl




# ══════════════════════════════════════════════════
# Page 2: 목표모델 (As-Is → To-Be)
# ══════════════════════════════════════════════════

def p2_target(prs):
    s = _sl(prs)
    _title(s, '목표모델 (As-Is → To-Be)',
           '현행 문제점 진단 및 목표 시스템 기능 구성', 1)

    pairs = [
        ('수기 예약 접수\n전화/대면 예약, 중복 발생 빈번',
         '온라인 좌석 예약\nSeatGrid 실시간 좌석 선택+예약'),
        ('배차 수작업 관리\n엑셀 배차표, 변경 추적 어려움',
         '배차 통합 CRUD\nSearchForm+DataTable+배정상태 관리'),
        ('대기자 순번 수기 관리\n순번 기록 누락, 배정 지연',
         '자동/수동 배정\n대기 순번 자동관리+빈좌석 배정'),
        ('예약시간 개별 통보\n계급별 예약 오픈/마감 혼란',
         '예약시간 체계 관리\n노선별·계급별 오픈/마감 CRUD'),
        ('위규자 관리 부재\n무단취소/탑승거부 제재 불가',
         '위규자+타군 관리\n제재기간 관리+타군 승인/반려'),
    ]

    y0 = Inches(1.0)
    RH = Inches(0.78)
    RG = Inches(0.06)
    AW = Inches(2.6)
    ARR_W = Inches(0.22)
    AX = ML
    ARX = ML + AW + Inches(0.15)
    TX = ARX + ARR_W + Inches(0.15)

    _box(s, AX, y0 - Inches(0.06), AW, Inches(0.18),
         '현행 시스템 (As-Is)', fill=ERR, stk=ERR, sz=8, b=True, fc=WHITE)
    _box(s, TX, y0 - Inches(0.06), AW, Inches(0.18),
         '목표 시스템 (To-Be)', fill=GRN, stk=GRN, sz=8, b=True, fc=WHITE)

    for i, (asis, tobe) in enumerate(pairs):
        y = y0 + i * (RH + RG)
        _box(s, AX, y, AW, RH, asis, fill=PEACH, stk=ERR, sz=7, fc=CHAR)
        _rarr(s, ARX, y + RH // 2, ARR_W)
        _box(s, TX, y, AW, RH, tobe, fill=MINT, stk=GRN, sz=7, fc=CHAR)

    # 우측 기능 요약
    FX = TX + AW + Inches(0.12)
    FW = ML + CW - FX
    _box(s, FX, y0, FW, 5 * RH + 4 * RG,
         '핵심 기능\n─────────\n'
         '예약           1\n대기자관리  1\n예약현황     1\n'
         '배차관리     1\n예약시간     1\n사용현황     1\n'
         '─────────\n총 13개 메뉴\n44 프로세스\n'
         '─────────\n위규자 1\n타군관리 1\n게시판 2',
         fill=ICE, stk=TEAL, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)

    # [C] 서술 테이블
    dy = y0 + 5 * (RH + RG) + Inches(0.12)
    _tbl(s, ML, dy, CW, Inches(1.25), [
        ['개선항목', '현행 문제점', '목표 개선방향', '기대 효과', '정량 목표'],
        ['예약관리', '전화/대면, 중복 발생', 'SeatGrid 좌석 선택+실시간 예약', '중복 예약 방지', '예약 30초'],
        ['배차관리', '엑셀 수작업', 'SearchForm+DataTable CRUD+배정상태', '배차 이력 추적', '배차 100%'],
        ['대기자', '순번 수기, 배정 지연', '자동배정+수동배정, 순번 자동관리', '배정 시간 단축', '자동 100%'],
        ['예약시간', '개별 통보, 혼란', '노선별·계급별 오픈/마감 CRUD', '체계적 시간관리', '일괄 설정'],
        ['위규/타군', '관리 부재', '위규자 제재+타군 승인/반려', '질서 유지', '제재 자동화'],
    ], cws=[int(CW * r) for r in [0.10, 0.24, 0.24, 0.22, 0.10]])


# ══════════════════════════════════════════════════
# Page 3: 유스케이스
# ══════════════════════════════════════════════════

def p3_usecase(prs):
    s = _sl(prs)
    _title(s, '유스케이스 다이어그램',
           '3개 액터 x 12개 유스케이스, 44개 단위 프로세스', 2)

    y0 = Inches(0.82)

    # 좌측: 액터
    actors = [
        '예약자 (장병)\n좌석 선택·예약·승차권\nROLE_USER',
        '배차관리자\n배차·배정·시간관리\nROLE_BUS_ADMIN',
        '시스템관리자\nROLE_ADMIN',
    ]
    AX = ML
    AW = Inches(1.0)
    AH = Inches(0.55)
    AG = Inches(0.15)
    for i, atxt in enumerate(actors):
        _box(s, AX, y0 + i * (AH + AG), AW, AH, atxt,
             fill=ICE, stk=NAVY, sz=7, fc=CHAR)

    # 중앙: 시스템 바운더리
    BX = AX + AW + Inches(0.2)
    BW = Inches(4.3)
    BH = 3 * AH + 2 * AG + Inches(0.3)
    bnd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, BX, y0 - Inches(0.05), BW, BH)
    bnd.fill.solid()
    bnd.fill.fore_color.rgb = RGBColor(0xFB, 0xFB, 0xFB)
    bnd.line.color.rgb = NAVY
    bnd.line.width = Pt(2)
    _tb(s, BX + Inches(0.1), y0 - Inches(0.02), Inches(5), Inches(0.18),
        '주말버스예약관리체계 (SYS10) — 바운더리', 8, True, NAVY)

    # UC 그룹
    uc_groups = [
        ('예약 관리 (3)', 'UC-BUS-001~003\n좌석선택예약·대기자·예약현황', ICE, SKY),
        ('배차 관리 (3)', 'UC-BUS-004~006\n배차등록·예약시간·사용현황', ICE, SKY),
        ('기타 관리 (2)', 'UC-BUS-007~008\n위규자관리·타군사용자관리', ICE, SKY),
        ('승차권 (1)', 'UC-BUS-009\nTicketPrint 승차권 발급/인쇄', ICE, SKY),
        ('게시판 (2)', 'UC-BUS-010~011\n공지사항·질의응답', LBLUE, TEAL),
        ('관리자 (2)', 'UC-SYS 코드관리·권한그룹\n시스템 관리', LBLUE, TEAL),
    ]
    UY0 = y0 + Inches(0.22)
    UW = Inches(2.0)
    UH = Inches(0.70)
    UG = Inches(0.06)
    for i, (grp, desc, fill, stk) in enumerate(uc_groups):
        r, c = divmod(i, 2)
        ux = BX + Inches(0.08) + c * (UW + Inches(0.1))
        uy = UY0 + r * (UH + UG)
        sh = _box(s, ux, uy, UW, UH, f'{grp}\n{desc}',
                  fill=fill, stk=stk, sz=6, b=False, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True

    # 우측: 외부 시스템
    EX = BX + BW + Inches(0.12)
    EW = ML + CW - EX
    _box(s, EX, y0 + Inches(0.2), EW, Inches(0.35),
         '타군 인증시스템\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)
    _box(s, EX, y0 + Inches(0.65), EW, Inches(0.35),
         'GPKI/SSO\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)

    # [C] 서술 테이블
    dy = Inches(3.2)
    _tbl(s, ML, dy, CW, Inches(2.1), [
        ['UC-ID', '유스케이스명', '주 액터', '기능 설명 (구체적)', '우선'],
        ['UC-BUS-001', '주말버스 예약', '예약자',
         'Form(노선/일자/시간)+SeatGrid 좌석선택+예약신청, TicketPrint 승차권 발급', '상'],
        ['UC-BUS-002', '대기자관리', '관리자',
         'SearchForm(4필드)+DataTable(10컬럼)+자동배정+수동배정(Modal 빈좌석 Select)', '상'],
        ['UC-BUS-003', '예약현황', '관리자',
         'SearchForm(4필드)+DataTable(10컬럼)+승차권발급+예약취소, 인쇄 미리보기', '상'],
        ['UC-BUS-004', '배차관리', '관리자',
         'SearchForm(5필드)+DataTable(9컬럼)+CRUD Modal(8필드)+좌석배치확인 SeatGrid', '상'],
        ['UC-BUS-005', '예약시간관리', '관리자',
         'SearchForm(3필드)+DataTable(7컬럼)+CRUD Modal(6필드), 노선별·계급별 관리', '중'],
        ['UC-BUS-007', '위규자관리', '관리자',
         'SearchForm(4필드)+DataTable(10컬럼)+CRUD Modal(7필드)+제재기간+인쇄', '중'],
        ['UC-BUS-008', '타군사용자관리', '관리자',
         'SearchForm(6필드)+DataTable(11컬럼)+승인/반려+패스워드초기화', '중'],
    ], cws=[int(CW * r) for r in [0.10, 0.12, 0.06, 0.66, 0.05]])


# ══════════════════════════════════════════════════
# Page 4: 스토리보드 (화면흐름도)
# ══════════════════════════════════════════════════

def p4_storyboard(prs):
    s = _sl(prs)
    _title(s, '스토리보드 (화면흐름도)',
           '총 13개 화면, 5개 업무영역, 핵심 네비게이션 흐름', 3)

    y0 = Inches(0.82)

    # 로그인 → 대시보드
    _box(s, ML, y0, Inches(1.1), Inches(0.35),
         '로그인\n/login', fill=WARN, stk=WARN, sz=8, b=True, fc=WHITE)
    _rarr(s, ML + Inches(1.35), y0 + Inches(0.12), Inches(0.25))
    _box(s, ML + Inches(1.65), y0, Inches(2.0), Inches(0.35),
         '대시보드 / 홈\n/sys10', fill=NAVY, stk=NAVY, sz=8, b=True, fc=WHITE)

    _darr(s, ML + Inches(2.65), y0 + Inches(0.37))

    # 5개 메뉴 그룹
    groups = [
        ('1. 예약', '예약·대기자·예약현황', '/sys10/1/2~4', '3', ICE, SKY),
        ('2. 배차', '배차관리·예약시간·사용현황', '/sys10/1/5~7', '3', ICE, SKY),
        ('3. 기타관리', '위규자·타군사용자', '/sys10/1/8~9', '2', ICE, SKY),
        ('4. 게시판', '공지사항·질의응답', '/sys10/board/*', '2', LBLUE, TEAL),
        ('5. 관리자', '코드관리·권한그룹·시스템', '/sys10/admin/*', '3+', LBLUE, TEAL),
    ]

    GY0 = y0 + Inches(0.55)
    GW = Inches(3.2)
    GH = Inches(0.90)
    GG = Inches(0.06)
    GCOL_GAP = Inches(0.15)

    for i, (name, screens, path, cnt, fill, stk) in enumerate(groups):
        r, c = divmod(i, 2)
        gx = ML + c * (GW + GCOL_GAP)
        gy = GY0 + r * (GH + GG)
        sh = _box(s, gx, gy, GW, GH,
                  f'{name} ({cnt}개)\n{screens}\n{path}',
                  fill=fill, stk=stk, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        pars = list(sh.text_frame.paragraphs)
        if len(pars) >= 3:
            for r_obj in pars[-1].runs:
                r_obj.font.color.rgb = ASH
                r_obj.font.size = Pt(8)

    # 핵심 화면 상세 (3개 대표 페이지)
    DY = GY0 + 3 * (GH + GG) + Inches(0.08)
    det_boxes = [
        ('주말버스 예약\n/sys10/1/2',
         'Row(12:12) 2열 레이아웃\n좌: Form(노선/일자/시간)\n경유지 표시\n'
         '우: SeatGrid 좌석 선택\n예약신청 버튼\nTicketPrint 승차권 Modal',
         ICE, SKY),
        ('배차관리\n/sys10/1/5',
         'SearchForm(5필드)\n방향/일자/출발지/도착지/배정\n'
         'DataTable(9컬럼)\nCRUD Modal(8필드)\n배정/배정취소 버튼\nSeatGrid 좌석배치 확인',
         MINT, GRN),
        ('위규자관리\n/sys10/1/8',
         'SearchForm(4필드)\n성명/군번/유형/부대\nDataTable(10컬럼)\n'
         'CRUD Modal(7필드)\n제재기간 RangePicker\n인쇄 미리보기 Modal',
         CREAM, WARN),
    ]
    DW = Inches(2.95)
    DH = Inches(1.1)
    DG = Inches(0.1)
    for i, (ttl, desc, fill, stk) in enumerate(det_boxes):
        dx = ML + i * (DW + DG)
        sh = _box(s, dx, DY, DW, DH,
                  f'{ttl}\n──────────\n{desc}',
                  fill=fill, stk=stk, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        sh.text_frame.paragraphs[1].runs[0].font.bold = True

    # [C] 서술 테이블 좌우 2개
    dy = DY + DH + Inches(0.08)

    LW = int(CW * 0.48)
    _tbl(s, ML, dy, LW, Inches(0.95), [
        ['영역', '화면 수', 'URL 패턴', '주요 컴포넌트'],
        ['예약', '3', '/sys10/1/2~4', 'Form+SeatGrid+DataTable+TicketPrint'],
        ['배차', '3', '/sys10/1/5~7', 'SearchForm+DataTable+CRUD Modal'],
        ['기타관리', '2', '/sys10/1/8~9', 'DataTable+Form Modal+StatusBadge'],
        ['게시판', '2', '/sys10/board/*', 'SimpleBoardPage 공지/질의응답'],
        ['관리자', '3+', '/sys10/admin/*', 'AdminRoutes 코드/권한/시스템'],
    ], cws=[int(LW * r) for r in [0.20, 0.10, 0.28, 0.42]])

    RX = ML + LW + Inches(0.15)
    RW = CW - LW - Inches(0.15)
    _tbl(s, RX, dy, RW, Inches(0.95), [
        ['핵심 흐름 시나리오'],
        ['시나리오 1: 주말버스 예약 흐름\n'
         '  로그인 -> 대시보드\n'
         '  [예약] > [주말버스 예약] 메뉴\n'
         '  Form: 운행일자/노선/출발시간 선택\n'
         '  SeatGrid: 빈좌석 클릭 -> 선택\n'
         '  [예약 신청] -> 성공 -> 승차권 발급'],
        ['시나리오 2: 배차 + 대기자 배정 흐름\n'
         '  [배차] > [배차관리] 메뉴\n'
         '  [신규 등록] -> 배차 정보 입력\n'
         '  [배정] 버튼 -> 배차 배정 완료\n'
         '  [대기자관리] -> [자동 배정] 실행\n'
         '  빈좌석에 대기 순번순 자동 배정'],
    ])


# ══════════════════════════════════════════════════
# Page 5: 플로우차트 (업무처리흐름)
# ══════════════════════════════════════════════════

def p5_flowchart(prs):
    s = _sl(prs)
    _title(s, '업무처리 플로우차트',
           '핵심 프로세스: 예약 -> 배차 배정 -> 대기자 관리 — 4개 레인, 12단계', 4)

    y0 = Inches(0.82)

    # 4개 스윔레인
    lanes = [
        ('예약자/배차관리자', ICE, NAVY),
        ('프론트엔드', LBLUE, SKY),
        ('백엔드', MINT, GRN),
        ('DB / 외부', CREAM, WARN),
    ]
    LGAP = Inches(0.14)  # 레인 사이 간격 (10pt)
    LN_W = (CW - 3 * LGAP) // 4
    LN_H = Inches(4.1)   # 다이어그램 높이 1.5배
    LN_HDR = Inches(0.22)

    for i, (name, fill, stk) in enumerate(lanes):
        lx = ML + (LN_W + LGAP) * i
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = fill
        bg.line.color.rgb = stk
        bg.line.width = Pt(1)
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_HDR)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = stk
        hdr.line.fill.background()
        tf = hdr.text_frame
        _run(tf.paragraphs[0], name, 7, True, WHITE)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    steps = [
        (0, 0.30, '  예약 메뉴 접근\n노선/일자/시간 선택'),
        (1, 0.30, '  SeatGrid 렌더\nGET /routes/{id}/seats'),
        (0, 0.70, '  빈좌석 선택\n[예약 신청] 버튼 클릭'),
        (1, 0.70, '  POST /reservations\n예약 데이터 전송'),
        (2, 0.70, '  예약 생성 + 좌석상태\nreserved로 변경'),
        (3, 0.70, '  Reservation INSERT\n좌석 상태 UPDATE'),
        (1, 1.10, '  TicketPrint 승차권\n승차권 발급 Modal'),
        (0, 1.10, '  관리자: 배차 등록\n[배차관리] 메뉴 접근'),
        (0, 1.50, '  [배정] 버튼 클릭\n배차 배정 처리'),
        (2, 1.50, '  PUT /dispatches/{id}\nassignStatus=assigned'),
        (0, 1.90, '  [자동 배정] 실행\n대기자 빈좌석 배정'),
        (2, 1.90, '  POST /waitlist/auto\n순번순 자동 배정'),
    ]

    SBW = LN_W - Inches(0.12)
    SBH = Inches(0.56)

    for lane, y_off, txt in steps:
        sx = ML + (LN_W + LGAP) * lane + Inches(0.07)
        sy = y0 + LN_HDR + Inches(y_off)
        _box(s, sx, sy, SBW, SBH, txt,
             fill=WHITE, stk=ASH, sw=1, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)

    # 판단 분기
    fx = ML + (LN_W + LGAP) * 1 + Inches(0.07)
    fy = y0 + LN_HDR + Inches(1.35)
    _box(s, fx, fy, SBW, Inches(0.12),
         '좌석 만석 -> 대기자 등록 -> 순번 부여',
         fill=PEACH, stk=ERR, sw=1, sz=5, fc=ERR, al=PP_ALIGN.CENTER)

    # [C] 서술 테이블 좌우 2개
    dy = y0 + LN_H + Inches(0.1)

    _tbl(s, ML, dy, CW, Inches(1.6), [
        ['단계', '동작', '조건/상세'],
        ['1~2', '예약접근 -> 좌석조회', 'Form(노선/일자/시간) -> SeatGrid 좌석 렌더'],
        ['3~4', '좌석선택 -> 예약신청', 'SeatGrid 클릭 -> POST /reservations'],
        ['5~6', '예약생성 -> DB저장', '좌석 status=reserved, Reservation INSERT'],
        ['7~8', '승차권발급 -> 배차관리', 'TicketPrint Modal, 관리자 배차 CRUD'],
        ['9~10', '배정처리 -> API', 'PUT /dispatches/{id} assignStatus=assigned'],
        ['11~12', '자동배정 -> 대기자', 'POST /waitlist/auto-assign, 순번순 빈좌석 배정'],
    ], cws=[int(CW * r) for r in [0.08, 0.28, 0.64]])

    _tbl(s, ML, dy + Inches(1.6) + Inches(0.1), CW, Inches(1.6), [
        ['업무 규칙 / 예외 처리'],
        ['[업무 규칙]\n'
         'BR-BUS-001: 노선 5개 (서울-포항/청주/대전/광주/부산)\n'
         'BR-BUS-002: 좌석상태 3종 (available/reserved/selected)\n'
         'BR-BUS-003: 예약상태 3종 (reserved/cancelled/waiting)\n'
         'BR-BUS-004: 배정상태 2종 (assigned/unassigned)\n'
         'BR-BUS-005: 위규유형 3종 (무단취소/탑승거부/기타)\n'
         'BR-BUS-006: 제재상태 3종 (제재중/만료/예정)\n'
         'BR-BUS-007: 계급별 예약오픈시간 차등 적용\n'
         'BR-BUS-008: 타군사용자 상태 3종 (신청대기/승인/반려)'],
        ['[예외 처리]\n'
         'E-001: 세션만료(401) -> 로그인 리다이렉트\n'
         'E-002: 좌석 만석 -> 대기자 자동 등록\n'
         'E-003: 중복예약 -> 경고 메시지\n'
         'E-004: 서버오류(500) -> 에러 메시지+재시도'],
    ])


# ══════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════

def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    print('[1/4] 목표모델...')
    p2_target(prs)
    print('[2/4] 유스케이스...')
    p3_usecase(prs)
    print('[3/4] 스토리보드...')
    p4_storyboard(prs)
    print('[4/4] 플로우차트...')
    p5_flowchart(prs)

    out_path = os.path.join(OUT, 'SYS10_주말버스예약관리체계_분석설계_v6.1.pptx')
    prs.save(out_path)
    print(f'\n완료! 총 {len(prs.slides)}장 슬라이드')
    print(f'파일: {out_path}')


if __name__ == '__main__':
    main()
