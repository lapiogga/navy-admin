#!/usr/bin/env python3
"""
SYS11 연구자료종합관리체계 기능명세서 PPTX v5
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
system_level_spec_guide_v5.md 기반 핵심 5페이지:
  P1. 시스템 개념도      P2. 목표모델 (As-Is -> To-Be)
  P3. 유스케이스          P4. 스토리보드 (화면흐름도)
  P5. 플로우차트 (업무처리흐름)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ═══ 색상 팔레트 ═══
NAVY  = RGBColor(0x6B, 0x8E, 0xA8)
TEAL  = RGBColor(0x8E, 0xC5, 0xC0)
SKY   = RGBColor(0xA8, 0xC8, 0xE8)
MINT  = RGBColor(0xDC, 0xF5, 0xE4)
PEACH = RGBColor(0xFC, 0xD4, 0xD0)
CREAM = RGBColor(0xFD, 0xEE, 0xC2)
ICE   = RGBColor(0xE4, 0xEE, 0xF8)
SNOW  = RGBColor(0xF7, 0xF9, 0xFC)
CHAR  = RGBColor(0x1A, 0x1A, 0x1A)
ASH   = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRN   = RGBColor(0xA4, 0xD4, 0xB4)
WARN  = RGBColor(0xF5, 0xCC, 0xA0)
ERR   = RGBColor(0xF0, 0xB2, 0xAE)
LBLUE = RGBColor(0xDC, 0xEA, 0xFE)
LSUB  = RGBColor(0xC5, 0xDC, 0xEC)
HBAR  = RGBColor(0xD4, 0xE2, 0xEC)  # 제목바 연한 배경
THBG  = RGBColor(0xD4, 0xE2, 0xEC)  # 표 헤더 연한 배경
BORD  = RGBColor(0x4A, 0x6B, 0x8C)  # 박스 테두리 (진한)
TBLB  = RGBColor(0x80, 0x80, 0x80)  # 표 테두리 (회색)

# ═══ 치수 ═══
SW = Inches(10)
SH = Inches(5.625)
ML = Inches(0.4)
CW = Inches(9.2)
FONT = 'Freesentation 5 Medium'
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'ppt')


# ══════════════════════════════════════════════════
# 헬퍼 함수
# ══════════════════════════════════════════════════

def _sl(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = SNOW
    return s


def _run(p, txt, sz=9, b=False, c=CHAR):
    """단락에 서식 런 추가 (최소 9pt 보장)"""
    sz = max(sz, 9)
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = b
    r.font.color.rgb = c
    r.font.name = FONT
    return r


def _tb(s, l, t, w, h, txt, sz=9, b=False, c=CHAR, al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    _run(tb.text_frame.paragraphs[0], txt, sz, b, c)
    tb.text_frame.paragraphs[0].alignment = al
    return tb


def _title(s, title, sub, pg):
    """[A] 제목 영역 — 연한 바 + 진한 테두리 + 검정 글자"""
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = HBAR
    bar.line.color.rgb = BORD
    bar.line.width = Pt(0.5)
    _tb(s, ML, Inches(0.06), Inches(7.5), Inches(0.36), title, 22, True, CHAR)
    if sub:
        _tb(s, ML, Inches(0.42), Inches(8.5), Inches(0.26), sub, 11, False, CHAR)
    _tb(s, Inches(9.0), Inches(0.18), Inches(0.6), Inches(0.3),
        f'{pg}/5', 12, True, CHAR, PP_ALIGN.RIGHT)


def _box(s, l, t, w, h, txt, fill=ICE, stk=BORD, sw=0.5,
         sz=9, b=False, fc=CHAR, al=PP_ALIGN.CENTER):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = stk
    sh.line.width = Pt(sw)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.margin_left = tf.margin_right = Inches(0.03)
    for i, ln in enumerate(txt.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p, ln, sz, b, fc)
        p.alignment = al
        p.space_before = p.space_after = Pt(0)
    bp = sh._element.txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', 'ctr')
    return sh


def _darr(s, cx, y, h=Inches(0.08)):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                           cx - Inches(0.06), y, Inches(0.12), h)
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _rarr(s, x, cy, w=Inches(0.18)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                           x, cy - Inches(0.05), w, Inches(0.1))
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _cell_bg(cell, hx):
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for sf in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(sf)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), val=hx)


def _cell_border(cell, color_hex='808080', width_pt=0.5):
    """셀 테두리 4면 회색 설정"""
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for border in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for existing in tcPr.findall(qn(border)):
            tcPr.remove(existing)
        ln = etree.SubElement(tcPr, qn(border))
        ln.set('w', str(int(width_pt * 12700)))
        ln.set('cap', 'flat')
        ln.set('cmpd', 'sng')
        ln.set('algn', 'ctr')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr'), val=color_hex)
        etree.SubElement(ln, qn('a:prstDash'), val='solid')
        etree.SubElement(ln, qn('a:round'))


def _tbl(s, l, t, w, h, data, cws=None):
    """테이블 (멀티라인 셀 지원, 진한 테두리 0.5pt, 연한 헤더 배경). data[0]=헤더"""
    rows, cols = len(data), len(data[0])
    sh = s.shapes.add_table(rows, cols, l, t, w, h)
    tbl = sh.table
    if cws:
        for i, cw in enumerate(cws):
            tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ''
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_top = tf.margin_bottom = Inches(0.04)
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.06)
            hdr = (r == 0)
            for i, ln in enumerate(data[r][c].split('\n')):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                _run(p, ln, 9 if not hdr else 10, hdr, CHAR)
                p.alignment = PP_ALIGN.CENTER if hdr else PP_ALIGN.LEFT
                p.space_before = p.space_after = Pt(0)
            _cell_bg(cell, 'D4E2EC' if hdr else ('FFFFFF' if r % 2 == 1 else 'F2F7FB'))
            _cell_border(cell, '808080', 0.5)
    return tbl


# ══════════════════════════════════════════════════
# Page 1: 시스템 개념도
# ══════════════════════════════════════════════════

def p1_concept(prs):
    s = _sl(prs)
    _title(s, '시스템 개념도',
           '연구자료종합관리체계의 전체 아키텍처 및 기술 구성 — 6개 계층', 1)

    layers = [
        ('사용자\n계층', [
            ('연구관리자\n연구자료 등록·관리·카테고리', ICE, NAVY),
            ('연구자/일반사용자\n자료 검색·열람·다운로드', ICE, NAVY),
            ('시스템관리자\n코드·권한·시스템 관리', ICE, NAVY),
        ]),
        ('채널\n계층', [
            ('웹 브라우저\nChrome/Edge 1920x1080', ICE, TEAL),
            ('반응형 모바일 웹\n360px~', ICE, TEAL),
            ('API Gateway\nRESTful 12 endpoints', ICE, TEAL),
        ]),
        ('프론트\n엔드', [
            ('React 18 + TS 5\nVite 5 · Router 7', LBLUE, SKY),
            ('Ant Design 5.24\nTable · Card · Statistic', LBLUE, SKY),
            ('Zustand 5 + TQ 5\n상태관리 · 서버캐싱', LBLUE, SKY),
        ]),
        ('백엔드\n계층', [
            ('Spring Boot 3.x\n(추후 개발) Java 17', MINT, GRN),
            ('MSW 2.12 Mock\nFaker.js 9 데이터', MINT, GRN),
            ('REST API 설계\nJSON/HTTPS 12개', MINT, GRN),
        ]),
        ('데이터\n계층', [
            ('PostgreSQL 15\n(추후) 정규화 설계', CREAM, WARN),
            ('Mock 데이터 50건\n연구자료+카테고리10종', CREAM, WARN),
            ('파일 스토리지\n연구계획서/보고서 PDF', CREAM, WARN),
        ]),
        ('외부\n연계', [
            ('GPKI/SSO 인증\n(추후 연계)', SNOW, ASH),
            ('국방연구정보시스템\n(추후 연계)', SNOW, ASH),
            ('알림톡/SMS\n(추후 연계)', SNOW, ASH),
        ]),
    ]

    y0 = Inches(0.82)
    LH = Inches(0.44)
    GAP = Inches(0.05)
    LBL_W = Inches(0.85)
    LBL_G = Inches(0.05)

    for idx, (lbl, items) in enumerate(layers):
        y = y0 + idx * (LH + GAP)
        _box(s, ML, y, LBL_W, LH, lbl, fill=NAVY, stk=NAVY, sz=7, b=True, fc=WHITE)
        n = len(items)
        avail = CW - LBL_W - LBL_G
        igap = Inches(0.04)
        iw = (avail - igap * (n - 1)) // n
        x = ML + LBL_W + LBL_G
        for txt, fill, stk in items:
            _box(s, x, y, iw, LH, txt, fill=fill, stk=stk, sz=7, fc=CHAR)
            x += iw + igap
        if idx < len(layers) - 1:
            _darr(s, ML + CW // 2, y + LH + Inches(0.01))

    # [C] 서술 테이블
    dy = y0 + 6 * (LH + GAP) + Inches(0.05)
    _tbl(s, ML, dy, CW, Inches(1.45), [
        ['프론트엔드 기술스택', '백엔드 기술스택', '인프라 / 보안 / 연계'],
        ['• React 18.x + TypeScript 5.x\n'
         '• Ant Design 5.24 (Table/Card/Statistic)\n'
         '• Zustand 5 + TanStack Query 5\n'
         '• React Router 7 + Vite 5\n'
         '• React Hook Form + Zod\n'
         '• 4개 고유 페이지, 10개 화면',
         '• Spring Boot 3.x (추후 개발)\n'
         '• MSW 2.12 Mock API (현재)\n'
         '• Faker.js 9 한국어 데이터\n'
         '• REST API 12개 endpoint\n'
         '• JWT + Redis 인증 (추후)\n'
         '• MyBatis 3.5 ORM (추후)',
         '• PostgreSQL 15 (추후 개발)\n'
         '• Linux / Tomcat 10.x (추후)\n'
         '• SSL/TLS 1.3, GPKI (추후)\n'
         '• 국방연구정보시스템 (추후)\n'
         '• 알림톡/SMS/이메일 (추후)\n'
         '• 사용자: 3역할, 19 프로세스'],
    ], cws=[CW // 3] * 3)


# ══════════════════════════════════════════════════
# Page 2: 목표모델 (As-Is -> To-Be)
# ══════════════════════════════════════════════════

def p2_target(prs):
    s = _sl(prs)
    _title(s, '목표모델 (As-Is -> To-Be)',
           '현행 문제점 진단 및 목표 시스템 기능 구성', 2)

    pairs = [
        ('연구자료 분산 관리\n부서별 개별 파일서버, 검색 불가',
         '통합 연구자료 포탈\n10개 카테고리 분류, 통합 검색'),
        ('진행현황 파악 곤란\n최초/중간/최종평가 추적 불가',
         '진행상황 관리\n최초/중간/최종평가 단계별 추적'),
        ('다운로드 이력 부재\n누가 언제 어떤 자료 열람했는지 불명',
         '다운로드 이력 추적\n군번/IP/파일명별 이력 조회'),
        ('카테고리 관리 미흡\n분류 체계 없이 무분별 등록',
         '카테고리 관리\n10개 연구분야 체계적 분류'),
        ('통계/현황 부재\n연도별/분야별 현황 파악 불가',
         '통계 대시보드\n연도별 카테고리 현황+예산+다운로드'),
    ]

    y0 = Inches(1.0)
    RH = Inches(0.52)
    RG = Inches(0.03)
    AW = Inches(3.5)
    ARR_W = Inches(0.22)
    AX = ML
    ARX = ML + AW + Inches(0.15)
    TX = ARX + ARR_W + Inches(0.15)

    _box(s, AX, y0 - Inches(0.06), AW, Inches(0.18),
         '현행 시스템 (As-Is)', fill=ERR, stk=ERR, sz=8, b=True, fc=WHITE)
    _box(s, TX, y0 - Inches(0.06), AW, Inches(0.18),
         '목표 시스템 (To-Be)', fill=GRN, stk=GRN, sz=8, b=True, fc=WHITE)

    for i, (asis, tobe) in enumerate(pairs):
        y = y0 + i * (RH + RG)
        _box(s, AX, y, AW, RH, asis, fill=PEACH, stk=ERR, sz=7, fc=CHAR)
        _rarr(s, ARX, y + RH // 2, ARR_W)
        _box(s, TX, y, AW, RH, tobe, fill=MINT, stk=GRN, sz=7, fc=CHAR)

    # 우측 기능 요약
    FX = TX + AW + Inches(0.12)
    FW = ML + CW - FX
    _box(s, FX, y0, FW, 5 * RH + 4 * RG,
         '핵심 기능\n─────────\n'
         '메인화면     1\n연구자료     1\n자료실        1\n'
         '게시판        2\n관리자        7탭\n'
         '─────────\n총 10개 메뉴\n19 프로세스\n'
         '─────────\n카테고리 10종\n연구분야별 분류',
         fill=ICE, stk=TEAL, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)

    # [C] 서술 테이블
    dy = y0 + 5 * (RH + RG) + Inches(0.12)
    _tbl(s, ML, dy, CW, Inches(1.25), [
        ['개선항목', '현행 문제점', '목표 개선방향', '기대 효과', '정량 목표'],
        ['자료관리', '부서별 분산, 검색 불가', '통합 DB, 10개 카테고리 분류', '자료 즉시 검색', '검색 1초'],
        ['진행관리', '평가 단계 추적 불가', '최초/중간/최종 단계별 관리', '진행상황 투명화', '추적 100%'],
        ['이력추적', '다운로드 이력 없음', '군번/IP/파일별 이력 조회', '보안 감사 가능', '이력 100%'],
        ['카테고리', '무분별 등록, 분류 없음', '10개 연구분야 체계 분류', '자료 분류 체계화', '10종 분류'],
        ['통계현황', '현황 파악 불가', '연도별 카테고리 통계 대시보드', '현황 즉시 파악', '실시간 통계'],
    ], cws=[int(CW * r) for r in [0.10, 0.24, 0.24, 0.22, 0.10]])


# ══════════════════════════════════════════════════
# Page 3: 유스케이스
# ══════════════════════════════════════════════════

def p3_usecase(prs):
    s = _sl(prs)
    _title(s, '유스케이스 다이어그램',
           '3개 액터 x 10개 유스케이스, 19개 단위 프로세스', 3)

    y0 = Inches(0.82)

    # 좌측: 액터
    actors = [
        '연구관리자\n자료등록·카테고리관리\nROLE_MANAGER',
        '연구자/일반사용자\n검색·열람·다운로드\nROLE_USER',
        '시스템관리자\nROLE_ADMIN',
    ]
    AX = ML
    AW = Inches(1.15)
    AH = Inches(0.55)
    AG = Inches(0.15)
    for i, atxt in enumerate(actors):
        _box(s, AX, y0 + i * (AH + AG), AW, AH, atxt,
             fill=ICE, stk=NAVY, sz=7, fc=CHAR)

    # 중앙: 시스템 바운더리
    BX = AX + AW + Inches(0.2)
    BW = Inches(6.2)
    BH = 3 * AH + 2 * AG + Inches(0.3)
    bnd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, BX, y0 - Inches(0.05), BW, BH)
    bnd.fill.solid()
    bnd.fill.fore_color.rgb = RGBColor(0xFB, 0xFB, 0xFB)
    bnd.line.color.rgb = NAVY
    bnd.line.width = Pt(2)
    _tb(s, BX + Inches(0.1), y0 - Inches(0.02), Inches(5), Inches(0.18),
        '연구자료종합관리체계 (SYS11) — 바운더리', 8, True, NAVY)

    # UC 그룹
    uc_groups = [
        ('메인화면 (4)', 'UC-RSRC-001~004\n통계카드+카테고리현황+최신/인기자료', ICE, SKY),
        ('연구자료 CRUD (5)', 'UC-RSRC-005~009\nSearchForm(5필드)+DataTable(10컬럼)+CrudForm(11필드)', ICE, SKY),
        ('자료실 (2)', 'UC-RSRC-010~011\nDataTable(7컬럼)+다운로드', ICE, SKY),
        ('관리자 (7탭)', 'UC-RSRC-012~018\n자료/카테고리/통계/사용자/삭제/권한/이력', ICE, SKY),
        ('게시판 (2)', 'UC-RSRC-019~020\n공지사항·질의응답', LBLUE, TEAL),
        ('공통 관리자 (5)', 'UC-SYS 시스템·코드·권한\n결재선·게시판 관리', LBLUE, TEAL),
    ]
    UY0 = y0 + Inches(0.22)
    UW = Inches(2.95)
    UH = Inches(0.48)
    UG = Inches(0.03)
    for i, (grp, desc, fill, stk) in enumerate(uc_groups):
        r, c = divmod(i, 2)
        ux = BX + Inches(0.08) + c * (UW + Inches(0.1))
        uy = UY0 + r * (UH + UG)
        sh = _box(s, ux, uy, UW, UH, f'{grp}\n{desc}',
                  fill=fill, stk=stk, sz=6, b=False, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True

    # 우측: 외부 시스템
    EX = BX + BW + Inches(0.12)
    EW = ML + CW - EX
    _box(s, EX, y0 + Inches(0.2), EW, Inches(0.35),
         '국방연구정보시스템\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)
    _box(s, EX, y0 + Inches(0.65), EW, Inches(0.35),
         'GPKI/SSO\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)

    # [C] 서술 테이블
    dy = Inches(3.2)
    _tbl(s, ML, dy, CW, Inches(2.1), [
        ['UC-ID', '유스케이스명', '주 액터', '기능 설명 (구체적)', '우선'],
        ['UC-RSRC-001', '메인 대시보드', '사용자',
         'Statistic 카드4개(전체자료/이번달/인기분야/총다운로드)+연도별 카테고리 현황 Table', '상'],
        ['UC-RSRC-002', '최신/인기자료 조회', '사용자',
         'Row 12:12 — 좌: 최신자료 List, 우: 인기자료 List(다운로드수)+전체보기 링크', '상'],
        ['UC-RSRC-003', '소개자료 카드', '사용자',
         '10개 카테고리 Card 그리드(5열x2행), 카테고리별 설명+자료보기 링크', '중'],
        ['UC-RSRC-005', '연구자료 CRUD', '관리자',
         'SearchForm(5필드:진행상황/제목/연구자/년도/담당관)+DataTable(10컬럼)+CrudForm(11필드) Modal', '상'],
        ['UC-RSRC-010', '자료실', '사용자',
         'DataTable(7컬럼: 번호/제목/파일명/파일크기/다운로드수/등록일/다운로드버튼)', '중'],
        ['UC-RSRC-012', '관리자 자료관리', '관리자',
         'Tabs 7개(자료/카테고리/통계/사용자/삭제/권한/다운로드이력), 자료 CRUD+카테고리 등록', '상'],
        ['UC-RSRC-019', '게시판', '사용자',
         'SimpleBoardPage 재사용: 공지사항(/sys11/board/1), 질의응답(/sys11/board/2)', '하'],
    ], cws=[int(CW * r) for r in [0.10, 0.12, 0.06, 0.66, 0.05]])


# ══════════════════════════════════════════════════
# Page 4: 스토리보드 (화면흐름도)
# ══════════════════════════════════════════════════

def p4_storyboard(prs):
    s = _sl(prs)
    _title(s, '스토리보드 (화면흐름도)',
           '총 10개 화면, 4개 업무영역, 핵심 네비게이션 흐름', 4)

    y0 = Inches(0.82)

    # 로그인 -> 대시보드
    _box(s, ML, y0, Inches(1.3), Inches(0.35),
         '로그인\n/login', fill=WARN, stk=WARN, sz=8, b=True, fc=WHITE)
    _rarr(s, ML + Inches(1.35), y0 + Inches(0.12), Inches(0.25))
    _box(s, ML + Inches(1.65), y0, Inches(2.0), Inches(0.35),
         '대시보드 / 홈\n/sys11', fill=NAVY, stk=NAVY, sz=8, b=True, fc=WHITE)

    _darr(s, ML + Inches(2.65), y0 + Inches(0.37))

    # 메뉴 그룹
    groups = [
        ('1. 연구자료관리', '메인화면 · 연구자료 · 자료실 · 공지사항', '/sys11/1/*', '5', ICE, SKY),
        ('2. 관리자', '자료·카테고리·통계·사용자·삭제·권한·이력', '/sys11/1/5', '7탭', ICE, SKY),
        ('3. 게시판', '공지사항 · 질의응답', '/sys11/board/*', '2', LBLUE, TEAL),
        ('4. 공통관리자', '시스템 · 코드 · 권한 · 결재선 · 게시판', '/sys11/admin/*', '5', LBLUE, TEAL),
    ]

    GY0 = y0 + Inches(0.55)
    GW = Inches(4.5)
    GH = Inches(0.58)
    GG = Inches(0.03)
    GCOL_GAP = Inches(0.15)

    for i, (name, screens, path, cnt, fill, stk) in enumerate(groups):
        r, c = divmod(i, 2)
        gx = ML + c * (GW + GCOL_GAP)
        gy = GY0 + r * (GH + GG)
        sh = _box(s, gx, gy, GW, GH,
                  f'{name} ({cnt}개)\n{screens}\n{path}',
                  fill=fill, stk=stk, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        pars = list(sh.text_frame.paragraphs)
        if len(pars) >= 3:
            for r_obj in pars[-1].runs:
                r_obj.font.color.rgb = ASH
                r_obj.font.size = Pt(8)

    # 핵심 화면 상세
    DY = GY0 + 3 * (GH + GG) + Inches(0.08)
    det_boxes = [
        ('메인화면\n/sys11/1/1',
         'Statistic 카드4개\n연도별 카테고리 Table\n최신자료 List(좌)\n'
         '인기자료 List(우)\n소개자료 Card 10개\n(5열x2행 그리드)',
         ICE, SKY),
        ('연구자료\n/sys11/1/2',
         'SearchForm(5필드)\nDataTable(10컬럼)\nCrudForm(11필드) Modal\n'
         '연구계획서 4종 Upload\n연구보고서 Upload\n엑셀 다운로드',
         MINT, GRN),
        ('관리자\n/sys11/1/5',
         'Tabs 7개 구성\n자료관리 DataTable\n카테고리 CRUD\n'
         '통계 Statistic 4개\n다운로드이력 DataTable\n(게시판/군번/IP 검색)',
         CREAM, WARN),
    ]
    DW = Inches(2.95)
    DH = Inches(1.1)
    DG = Inches(0.1)
    for i, (ttl, desc, fill, stk) in enumerate(det_boxes):
        dx = ML + i * (DW + DG)
        sh = _box(s, dx, DY, DW, DH,
                  f'{ttl}\n──────────\n{desc}',
                  fill=fill, stk=stk, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        sh.text_frame.paragraphs[1].runs[0].font.bold = True

    # [C] 서술 테이블
    dy = DY + DH + Inches(0.08)

    LW = int(CW * 0.48)
    _tbl(s, ML, dy, LW, Inches(0.95), [
        ['영역', '화면 수', 'URL 패턴', '주요 컴포넌트'],
        ['메인화면', '1', '/sys11/1/1', 'Statistic+Table+List+Card'],
        ['연구자료', '1', '/sys11/1/2', 'SearchForm+DataTable+CrudForm'],
        ['자료실', '1', '/sys11/1/3', 'DataTable+Download Button'],
        ['관리자', '1(7탭)', '/sys11/1/5', 'Tabs+DataTable+CrudForm+Statistic'],
        ['게시판+공통관리', '6', '/sys11/board,admin/*', 'SimpleBoardPage+AdminRoutes'],
    ], cws=[int(LW * r) for r in [0.20, 0.10, 0.28, 0.42]])

    RX = ML + LW + Inches(0.15)
    RW = CW - LW - Inches(0.15)
    _tbl(s, RX, dy, RW, Inches(0.95), [
        ['핵심 흐름 시나리오'],
        ['시나리오 1: 연구자료 검색·열람 흐름\n'
         '1. 로그인 -> 대시보드\n'
         '2. [연구자료관리] > [연구자료] 메뉴 선택\n'
         '3. SearchForm 진행상황/제목/연구자 검색\n'
         '4. DataTable 행 클릭 -> 상세 Modal\n'
         '5. Descriptions 12항목 확인\n'
         '6. 다운로드 버튼 클릭'],
        ['시나리오 2: 연구자료 등록 흐름\n'
         '1. [자료 등록] 버튼 -> CrudForm Modal\n'
         '2. 11항목 입력 (제목/분야/년도/예산/연구자 등)\n'
         '3. 연구계획서(최초/중간/최종/기타) Upload\n'
         '4. [저장] -> 성공 토스트 + 목록 갱신'],
    ])


# ══════════════════════════════════════════════════
# Page 5: 플로우차트 (업무처리흐름)
# ══════════════════════════════════════════════════

def p5_flowchart(prs):
    s = _sl(prs)
    _title(s, '업무처리 플로우차트',
           '핵심 프로세스: 자료 검색 -> 열람 -> 등록 — 4개 레인, 12단계', 5)

    y0 = Inches(0.82)

    # 4개 스윔레인
    lanes = [
        ('사용자 (연구자/관리자)', ICE, NAVY),
        ('프론트엔드 (React 18)', LBLUE, SKY),
        ('백엔드 (API/Mock)', MINT, GRN),
        ('DB / 외부', CREAM, WARN),
    ]
    LN_W = CW // 4
    LN_H = Inches(2.75)
    LN_HDR = Inches(0.22)

    for i, (name, fill, stk) in enumerate(lanes):
        lx = ML + LN_W * i
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = fill
        bg.line.color.rgb = stk
        bg.line.width = Pt(1)
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_HDR)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = stk
        hdr.line.fill.background()
        tf = hdr.text_frame
        _run(tf.paragraphs[0], name, 7, True, WHITE)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    steps = [
        (0, 0.30, '1. 메인화면 접근\n대시보드 통계 확인'),
        (1, 0.30, '2. 대시보드 렌더링\nStatistic+Table+List+Card'),
        (0, 0.70, '3. 연구자료 메뉴 선택\nSearchForm 5필드 입력'),
        (1, 0.70, '4. API 호출 (필터)\nGET /sys11/research?...'),
        (2, 0.70, '5. 필터링 + 페이지네이션\nprogress/keyword/year'),
        (3, 0.70, '6. Research SELECT\n조건별 조회 (Mock 50건)'),
        (1, 1.10, '7. DataTable 렌더\n10컬럼 + 군번/계급/성명'),
        (0, 1.10, '8. 행 클릭 -> 상세 열람\nDescriptions Modal 12항목'),
        (0, 1.50, '9. 관리자: 자료 등록\n[자료 등록] 버튼 클릭'),
        (1, 1.50, '10. CrudForm Modal\n11필드 + 계획서/보고서 Upload'),
        (2, 1.90, '11. POST /sys11/research\n신규 자료 저장'),
        (3, 1.90, '12. Research INSERT\n데이터 저장 + 이력'),
    ]

    SBW = LN_W - Inches(0.12)
    SBH = Inches(0.40)

    for lane, y_off, txt in steps:
        sx = ML + LN_W * lane + Inches(0.06)
        sy = y0 + LN_HDR + Inches(y_off)
        _box(s, sx, sy, SBW, SBH, txt,
             fill=WHITE, stk=ASH, sw=1, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)

    # 판단 분기
    fx = ML + LN_W * 1 + Inches(0.06)
    fy = y0 + LN_HDR + Inches(1.35)
    _box(s, fx, fy, SBW, Inches(0.12),
         '검증 실패 -> 필드 에러 표시 -> 재입력',
         fill=PEACH, stk=ERR, sw=1, sz=5, fc=ERR, al=PP_ALIGN.CENTER)

    # [C] 서술 테이블
    dy = y0 + LN_H + Inches(0.1)

    LW = int(CW * 0.52)
    _tbl(s, ML, dy, LW, Inches(1.42), [
        ['단계', '동작', '조건/상세'],
        ['1-2', '대시보드 접근 -> 렌더링', 'Statistic 4카드 + 연도별 카테고리 Table + 최신/인기 List'],
        ['3-4', '검색조건 입력 -> API호출', 'SearchForm 5필드(진행상황/제목/연구자/년도/담당관)'],
        ['5-6', '서버 필터링 -> DB조회', 'progressStatus/keyword/researcher/year 필터 조합'],
        ['7-8', '목록렌더 -> 상세열람', 'DataTable 10컬럼 + Descriptions Modal 12항목'],
        ['9-10', '등록버튼 -> CrudForm', 'Modal 11필드 + 연구계획서 4종 + 보고서 Upload'],
        ['11-12', 'API등록 -> DB저장', 'POST /sys11/research + Research INSERT'],
    ], cws=[int(LW * r) for r in [0.08, 0.28, 0.64]])

    RX = ML + LW + Inches(0.12)
    RW = CW - LW - Inches(0.12)
    _tbl(s, RX, dy, RW, Inches(1.42), [
        ['업무 규칙 / 예외 처리'],
        ['[업무 규칙]\n'
         'BR-RSRC-001: 카테고리 10종 (국방정책/개별사업/해사학술/군사학술/해군발전/함정정비/전투발전/전투실험/함정기술/장성정책연구)\n'
         'BR-RSRC-002: 진행상황 3단계 (최초평가/중간평가/최종평가)\n'
         'BR-RSRC-003: 연구자 구분 2종 (군내/외부)\n'
         'BR-RSRC-004: 연구계획서 4종 Upload (최초/중간/최종/기타)\n'
         'BR-RSRC-005: 다운로드 이력 9컬럼 (순번/게시판/글번호/군번/계급/성명/IP/일시/파일명)'],
        ['[예외 처리]\n'
         'E-001: 세션만료(401) -> 로그인 리다이렉트\n'
         'E-002: 권한없음(403) -> 접근 제한 알림\n'
         'E-003: 검증실패 -> 필드별 에러 표시+재입력\n'
         'E-004: 서버오류(500) -> 에러 메시지+재시도'],
    ])


# ══════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════

def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    print('[1/5] 시스템 개념도...')
    p1_concept(prs)
    print('[2/5] 목표모델...')
    p2_target(prs)
    print('[3/5] 유스케이스...')
    p3_usecase(prs)
    print('[4/5] 스토리보드...')
    p4_storyboard(prs)
    print('[5/5] 플로우차트...')
    p5_flowchart(prs)

    out_path = os.path.join(OUT, 'SYS11_연구자료종합관리체계_분석설계_v5.pptx')
    prs.save(out_path)
    print(f'\n완료! 총 {len(prs.slides)}장 슬라이드')
    print(f'파일: {out_path}')


if __name__ == '__main__':
    main()
