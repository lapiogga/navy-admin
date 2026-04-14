#!/usr/bin/env python3
"""
SYS12 지시건의사항관리체계 기능명세서 PPTX v6 (portrait)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
system_level_spec_guide_v6.md 기반 핵심 5페이지:
  P1. 시스템 개념도      P2. 목표모델 (As-Is -> To-Be)
  P3. 유스케이스          P4. 스토리보드 (화면흐름도)
  P5. 플로우차트 (업무처리흐름)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ═══ 색상 팔레트 ═══
NAVY  = RGBColor(0x6B, 0x8E, 0xA8)
TEAL  = RGBColor(0x8E, 0xC5, 0xC0)
SKY   = RGBColor(0xA8, 0xC8, 0xE8)
MINT  = RGBColor(0xDC, 0xF5, 0xE4)
PEACH = RGBColor(0xFC, 0xD4, 0xD0)
CREAM = RGBColor(0xFD, 0xEE, 0xC2)
ICE   = RGBColor(0xE4, 0xEE, 0xF8)
SNOW  = RGBColor(0xF7, 0xF9, 0xFC)
CHAR  = RGBColor(0x1A, 0x1A, 0x1A)
ASH   = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRN   = RGBColor(0xA4, 0xD4, 0xB4)
WARN  = RGBColor(0xF5, 0xCC, 0xA0)
ERR   = RGBColor(0xF0, 0xB2, 0xAE)
LBLUE = RGBColor(0xDC, 0xEA, 0xFE)
LSUB  = RGBColor(0xC5, 0xDC, 0xEC)
HBAR  = RGBColor(0xD4, 0xE2, 0xEC)  # 제목바 연한 배경
THBG  = RGBColor(0xD4, 0xE2, 0xEC)  # 표 헤더 연한 배경
BORD  = RGBColor(0x4A, 0x6B, 0x8C)  # 박스 테두리 (진한)
TBLB  = RGBColor(0x80, 0x80, 0x80)  # 표 테두리 (회색)

# ═══ 치수 ═══
SW = Inches(7.5)
SH = Inches(10)
ML = Inches(0.4)
CW = Inches(6.7)       # 콘텐츠 폭 = 7.5 - 0.4*2
FONT = 'Freesentation 5 Medium'
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'ppt')


# ══════════════════════════════════════════════════
# 헬퍼 함수
# ══════════════════════════════════════════════════

def _sl(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = SNOW
    return s


def _run(p, txt, sz=9, b=False, c=CHAR):
    """단락에 서식 런 추가 (최소 9pt 보장)"""
    sz = max(sz, 9)
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = b
    r.font.color.rgb = c
    r.font.name = FONT
    return r


def _tb(s, l, t, w, h, txt, sz=9, b=False, c=CHAR, al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    _run(tb.text_frame.paragraphs[0], txt, sz, b, c)
    tb.text_frame.paragraphs[0].alignment = al
    return tb


def _title(s, title, sub, pg):
    """[A] 제목 영역 (portrait) — 연한 바 + 진한 테두리 + 검정 글자"""
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = HBAR
    bar.line.color.rgb = BORD
    bar.line.width = Pt(0.5)
    _tb(s, ML, Inches(0.06), Inches(5.5), Inches(0.36), title, 20, True, CHAR)
    if sub:
        _tb(s, ML, Inches(0.42), Inches(6.5), Inches(0.26), sub, 10, False, CHAR)
    _tb(s, Inches(6.5), Inches(0.18), Inches(0.6), Inches(0.3),
        f'{pg}/5', 12, True, CHAR, PP_ALIGN.RIGHT)


def _box(s, l, t, w, h, txt, fill=ICE, stk=BORD, sw=0.5,
         sz=9, b=False, fc=CHAR, al=PP_ALIGN.CENTER):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = stk
    sh.line.width = Pt(sw)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.margin_left = tf.margin_right = Inches(0.03)
    for i, ln in enumerate(txt.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p, ln, sz, b, fc)
        p.alignment = al
        p.space_before = p.space_after = Pt(0)
    bp = sh._element.txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', 'ctr')
    return sh


def _darr(s, cx, y, h=Inches(0.08)):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                           cx - Inches(0.06), y, Inches(0.12), h)
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _rarr(s, x, cy, w=Inches(0.18)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                           x, cy - Inches(0.05), w, Inches(0.1))
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _cell_bg(cell, hx):
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for sf in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(sf)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), val=hx)


def _cell_border(cell, color_hex='808080', width_pt=0.5):
    """셀 테두리 4면 회색 설정"""
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for border in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for existing in tcPr.findall(qn(border)):
            tcPr.remove(existing)
        ln = etree.SubElement(tcPr, qn(border))
        ln.set('w', str(int(width_pt * 12700)))
        ln.set('cap', 'flat')
        ln.set('cmpd', 'sng')
        ln.set('algn', 'ctr')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr'), val=color_hex)
        etree.SubElement(ln, qn('a:prstDash'), val='solid')
        etree.SubElement(ln, qn('a:round'))


def _tbl(s, l, t, w, h, data, cws=None):
    """테이블 (멀티라인 셀 지원, 진한 테두리 0.5pt, 연한 헤더 배경). data[0]=헤더"""
    rows, cols = len(data), len(data[0])
    sh = s.shapes.add_table(rows, cols, l, t, w, h)
    tbl = sh.table
    if cws:
        for i, cw in enumerate(cws):
            tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ''
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_top = tf.margin_bottom = Inches(0.04)
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.06)
            hdr = (r == 0)
            for i, ln in enumerate(data[r][c].split('\n')):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                _run(p, ln, 9 if not hdr else 10, hdr, CHAR)
                p.alignment = PP_ALIGN.CENTER if hdr else PP_ALIGN.LEFT
                p.space_before = p.space_after = Pt(0)
            _cell_bg(cell, 'D4E2EC' if hdr else ('FFFFFF' if r % 2 == 1 else 'F2F7FB'))
            _cell_border(cell, '808080', 0.5)
    return tbl


# ══════════════════════════════════════════════════
# Page 1: 시스템 개념도
# ══════════════════════════════════════════════════

def p1_concept(prs):
    s = _sl(prs)
    _title(s, '시스템 개념도',
           '지시건의사항관리체계의 전체 아키텍처 및 기술 구성 — 6개 계층', 1)

    layers = [
        ('사용자\n계층', [
            ('지시관리자\n지시사항 등록·수정·삭제', ICE, NAVY),
            ('건의자/일반사용자\n건의사항 등록·조회·이력', ICE, NAVY),
            ('시스템관리자\n통계·카테고리·부대·권한', ICE, NAVY),
        ]),
        ('채널\n계층', [
            ('웹 브라우저\nChrome/Edge 1920x1080', ICE, TEAL),
            ('반응형 모바일 웹\n360px~', ICE, TEAL),
            ('API Gateway\nRESTful 20+ endpoints', ICE, TEAL),
        ]),
        ('프론트\n엔드', [
            ('React 18 + TS 5\nVite 5 · Router 7', LBLUE, SKY),
            ('Ant Design 5.24\nProgress · Timeline · Tabs', LBLUE, SKY),
            ('Zustand 5 + TQ 5\n상태관리 · 서버캐싱', LBLUE, SKY),
        ]),
        ('백엔드\n계층', [
            ('Spring Boot 3.x\n(추후 개발) Java 17', MINT, GRN),
            ('MSW 2.12 Mock\nFaker.js 9 데이터', MINT, GRN),
            ('REST API 설계\nJSON/HTTPS 20+개', MINT, GRN),
        ]),
        ('데이터\n계층', [
            ('PostgreSQL 15\n(추후) 정규화 설계', CREAM, WARN),
            ('Mock 데이터 60건\n지시30+건의30+이력', CREAM, WARN),
            ('파일 스토리지\n첨부파일 (추후)', CREAM, WARN),
        ]),
        ('외부\n연계', [
            ('GPKI/SSO 인증\n(추후 연계)', SNOW, ASH),
            ('국방통합정보시스템\n(추후 연계)', SNOW, ASH),
            ('알림톡/SMS\n(추후 연계)', SNOW, ASH),
        ]),
    ]

    y0 = Inches(0.82)
    LH = Inches(0.62)
    GAP = Inches(0.08)
    LBL_W = Inches(0.85)
    LBL_G = Inches(0.05)

    for idx, (lbl, items) in enumerate(layers):
        y = y0 + idx * (LH + GAP)
        _box(s, ML, y, LBL_W, LH, lbl, fill=NAVY, stk=NAVY, sz=7, b=True, fc=WHITE)
        n = len(items)
        avail = CW - LBL_W - LBL_G
        igap = Inches(0.04)
        iw = (avail - igap * (n - 1)) // n
        x = ML + LBL_W + LBL_G
        for txt, fill, stk in items:
            _box(s, x, y, iw, LH, txt, fill=fill, stk=stk, sz=7, fc=CHAR)
            x += iw + igap
        if idx < len(layers) - 1:
            _darr(s, ML + CW // 2, y + LH + Inches(0.01))

    # [C] 서술 테이블
    dy = y0 + 6 * (LH + GAP) + Inches(0.05)
    _tbl(s, ML, dy, CW, Inches(1.45), [
        ['프론트엔드 기술스택', '백엔드 기술스택', '인프라 / 보안 / 연계'],
        ['• React 18.x + TypeScript 5.x\n'
         '• Ant Design 5.24 (Progress/Timeline/Tabs)\n'
         '• Zustand 5 + TanStack Query 5\n'
         '• React Router 7 + Vite 5\n'
         '• React Hook Form + Zod\n'
         '• 5개 고유 페이지, 13개 화면',
         '• Spring Boot 3.x (추후 개발)\n'
         '• MSW 2.12 Mock API (현재)\n'
         '• Faker.js 9 한국어 데이터\n'
         '• REST API 20+개 endpoint\n'
         '• JWT + Redis 인증 (추후)\n'
         '• MyBatis 3.5 ORM (추후)',
         '• PostgreSQL 15 (추후 개발)\n'
         '• Linux / Tomcat 10.x (추후)\n'
         '• SSL/TLS 1.3, GPKI (추후)\n'
         '• 국방통합정보시스템 (추후)\n'
         '• 알림톡/SMS/이메일 (추후)\n'
         '• 사용자: 3역할, 32 프로세스'],
    ], cws=[CW // 3] * 3)


# ══════════════════════════════════════════════════
# Page 2: 목표모델 (As-Is -> To-Be)
# ══════════════════════════════════════════════════

def p2_target(prs):
    s = _sl(prs)
    _title(s, '목표모델 (As-Is -> To-Be)',
           '현행 문제점 진단 및 목표 시스템 기능 구성', 2)

    pairs = [
        ('지시사항 분산 관리\n대통령/장관/지휘관 지시 각각 별도 대장',
         '통합 지시관리 포탈\n3개 카테고리 통합, 추진현황 대시보드'),
        ('추진현황 파악 곤란\n완료/진행/미착수/지연 상태 추적 불가',
         '실시간 추진현황\nStatistic+Progress+매트릭스 테이블'),
        ('건의사항 처리 미흡\n건의 접수->처리 프로세스 부재',
         '건의사항 CRUD+이력\n건의등록·조치사항·처리이력 Timeline'),
        ('이행이력 관리 부재\n누가 언제 어떤 조치를 했는지 불명',
         '이행이력 Timeline\n조치자/일자/상태/내용 이력 추적'),
        ('통계/알림 부재\n카테고리별 현황 파악·알림 불가',
         '통계+알림 관리\nBar차트+카테고리/부대/기간/알림 설정'),
    ]

    y0 = Inches(1.0)
    RH = Inches(0.78)
    RG = Inches(0.06)
    AW = Inches(2.6)
    ARR_W = Inches(0.22)
    AX = ML
    ARX = ML + AW + Inches(0.15)
    TX = ARX + ARR_W + Inches(0.15)

    _box(s, AX, y0 - Inches(0.06), AW, Inches(0.18),
         '현행 시스템 (As-Is)', fill=ERR, stk=ERR, sz=8, b=True, fc=WHITE)
    _box(s, TX, y0 - Inches(0.06), AW, Inches(0.18),
         '목표 시스템 (To-Be)', fill=GRN, stk=GRN, sz=8, b=True, fc=WHITE)

    for i, (asis, tobe) in enumerate(pairs):
        y = y0 + i * (RH + RG)
        _box(s, AX, y, AW, RH, asis, fill=PEACH, stk=ERR, sz=7, fc=CHAR)
        _rarr(s, ARX, y + RH // 2, ARR_W)
        _box(s, TX, y, AW, RH, tobe, fill=MINT, stk=GRN, sz=7, fc=CHAR)

    # 우측 기능 요약
    FX = TX + AW + Inches(0.12)
    FW = ML + CW - FX
    _box(s, FX, y0, FW, 5 * RH + 4 * RG,
         '핵심 기능\n─────────\n'
         '지시사항     3종\n(대통령/장관/지휘관)\n건의사항     1\n'
         '게시판        2\n관리자        9탭\n'
         '─────────\n총 13개 메뉴\n32 프로세스\n'
         '─────────\n상태 4종\n완료/진행/미착수/지연',
         fill=ICE, stk=TEAL, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)

    # [C] 서술 테이블
    dy = y0 + 5 * (RH + RG) + Inches(0.12)
    _tbl(s, ML, dy, CW, Inches(1.25), [
        ['개선항목', '현행 문제점', '목표 개선방향', '기대 효과', '정량 목표'],
        ['지시통합', '3종 지시 별도 대장', '통합 포탈, 카테고리별 관리', '지시현황 일원화', '통합 100%'],
        ['추진현황', '상태 추적 불가', 'Statistic+Progress 대시보드', '실시간 현황 파악', '실시간 갱신'],
        ['건의관리', '처리 프로세스 부재', 'CRUD+조치사항+이력 관리', '건의 처리 투명화', '이력 100%'],
        ['이행이력', '조치 이력 없음', 'Timeline 이행이력 추적', '책임 소재 명확', '이력 완전'],
        ['통계알림', '현황/알림 불가', 'Bar차트+카테고리+알림설정', '선제적 관리 가능', '주기 알림'],
    ], cws=[int(CW * r) for r in [0.10, 0.24, 0.24, 0.22, 0.10]])


# ══════════════════════════════════════════════════
# Page 3: 유스케이스
# ══════════════════════════════════════════════════

def p3_usecase(prs):
    s = _sl(prs)
    _title(s, '유스케이스 다이어그램',
           '3개 액터 x 14개 유스케이스, 32개 단위 프로세스', 3)

    y0 = Inches(0.82)

    # 좌측: 액터
    actors = [
        '지시관리자\n지시등록·수정·삭제\nROLE_MANAGER',
        '건의자/사용자\n건의등록·조회·이력\nROLE_USER',
        '시스템관리자\nROLE_ADMIN',
    ]
    AX = ML
    AW = Inches(1.0)
    AH = Inches(0.55)
    AG = Inches(0.15)
    for i, atxt in enumerate(actors):
        _box(s, AX, y0 + i * (AH + AG), AW, AH, atxt,
             fill=ICE, stk=NAVY, sz=7, fc=CHAR)

    # 중앙: 시스템 바운더리
    BX = AX + AW + Inches(0.2)
    BW = Inches(4.3)
    BH = 3 * AH + 2 * AG + Inches(0.3)
    bnd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, BX, y0 - Inches(0.05), BW, BH)
    bnd.fill.solid()
    bnd.fill.fore_color.rgb = RGBColor(0xFB, 0xFB, 0xFB)
    bnd.line.color.rgb = NAVY
    bnd.line.width = Pt(2)
    _tb(s, BX + Inches(0.1), y0 - Inches(0.02), Inches(5), Inches(0.18),
        '지시건의사항관리체계 (SYS12) — 바운더리', 8, True, NAVY)

    # UC 그룹
    uc_groups = [
        ('지시사항 3종 (6)', 'UC-DRCT-001~006\n추진현황+CRUD+상세+이력+조치사항', ICE, SKY),
        ('건의사항 (4)', 'UC-DRCT-007~010\n추진현황+CRUD+상세+이력+조치사항', ICE, SKY),
        ('게시판 (2)', 'UC-DRCT-011~012\n공지사항·질의응답', LBLUE, TEAL),
        ('관리자 9탭 (9)', 'UC-DRCT-013~021\n지시/건의/카테고리/통계/사용자/부대/기간/알림/권한', ICE, SKY),
        ('공통 관리자 (5)', 'UC-SYS 시스템·코드·권한\n결재선·게시판 관리', LBLUE, TEAL),
        ('검색기능 (5)', 'UC-DRCT-022~026\n지시일/지시자/수명부대/상태/내용', LBLUE, TEAL),
    ]
    UY0 = y0 + Inches(0.22)
    UW = Inches(2.0)
    UH = Inches(0.70)
    UG = Inches(0.06)
    for i, (grp, desc, fill, stk) in enumerate(uc_groups):
        r, c = divmod(i, 2)
        ux = BX + Inches(0.08) + c * (UW + Inches(0.1))
        uy = UY0 + r * (UH + UG)
        sh = _box(s, ux, uy, UW, UH, f'{grp}\n{desc}',
                  fill=fill, stk=stk, sz=6, b=False, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True

    # 우측: 외부 시스템
    EX = BX + BW + Inches(0.12)
    EW = ML + CW - EX
    _box(s, EX, y0 + Inches(0.2), EW, Inches(0.35),
         '국방통합정보시스템\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)
    _box(s, EX, y0 + Inches(0.65), EW, Inches(0.35),
         'GPKI/SSO\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)

    # [C] 서술 테이블
    dy = Inches(3.2)
    _tbl(s, ML, dy, CW, Inches(2.1), [
        ['UC-ID', '유스케이스명', '주 액터', '기능 설명 (구체적)', '우선'],
        ['UC-DRCT-001', '지시 추진현황', '관리자',
         'Statistic 카드5개(총/완료/진행/미착수+추진율) + Progress + 카테고리별 매트릭스 DataTable(7컬럼)', '상'],
        ['UC-DRCT-002', '지시사항 CRUD', '관리자',
         'SearchForm(5필드:지시일/자/부대/상태/내용)+DataTable(8컬럼)+CrudForm(8필드)+StatusBadge', '상'],
        ['UC-DRCT-003', '지시 상세+이력', '사용자',
         '상세Modal(기본정보Grid+StatusBadge) + 이행이력 Timeline(일자/담당자/상태/내용)', '상'],
        ['UC-DRCT-004', '조치사항 등록', '관리자',
         'CrudForm(7필드: 담당자군번/계급/성명+추진상태+추진계획+내용+첨부파일) Modal', '상'],
        ['UC-DRCT-007', '건의 추진현황', '사용자',
         'Statistic 카드5개 + Progress + 카테고리별 처리현황 DataTable(7컬럼)', '상'],
        ['UC-DRCT-008', '건의사항 CRUD', '사용자',
         'SearchForm(5필드:건의일/자/부대/상태/내용)+DataTable(7컬럼)+CrudForm(7필드)+StatusBadge', '상'],
        ['UC-DRCT-013', '관리자 (9탭)', '관리자',
         'Tabs: 지시/건의관리+카테고리+통계Bar차트+사용자+부대+기간설정+알림설정+권한관리', '중'],
    ], cws=[int(CW * r) for r in [0.10, 0.12, 0.06, 0.66, 0.05]])


# ══════════════════════════════════════════════════
# Page 4: 스토리보드 (화면흐름도)
# ══════════════════════════════════════════════════

def p4_storyboard(prs):
    s = _sl(prs)
    _title(s, '스토리보드 (화면흐름도)',
           '총 13개 화면, 4개 업무영역, 핵심 네비게이션 흐름', 4)

    y0 = Inches(0.82)

    # 로그인 -> 대시보드
    _box(s, ML, y0, Inches(1.1), Inches(0.35),
         '로그인\n/login', fill=WARN, stk=WARN, sz=8, b=True, fc=WHITE)
    _rarr(s, ML + Inches(1.35), y0 + Inches(0.12), Inches(0.25))
    _box(s, ML + Inches(1.65), y0, Inches(2.0), Inches(0.35),
         '대시보드 / 홈\n/sys12', fill=NAVY, stk=NAVY, sz=8, b=True, fc=WHITE)

    _darr(s, ML + Inches(2.65), y0 + Inches(0.37))

    # 메뉴 그룹
    groups = [
        ('1. 지시사항', '대통령 · 국방부장관 · 지휘관 (추진현황+목록)', '/sys12/2/*', '6', ICE, SKY),
        ('2. 건의사항', '건의사항 (추진현황+목록)', '/sys12/3/*', '2', ICE, SKY),
        ('3. 게시판', '공지사항 · 질의응답', '/sys12/1/*', '2', LBLUE, TEAL),
        ('4. 관리자', '지시/건의/카테고리/통계/사용자/부대/기간/알림/권한', '/sys12/4/1', '9탭', LBLUE, TEAL),
    ]

    GY0 = y0 + Inches(0.55)
    GW = Inches(3.2)
    GH = Inches(0.90)
    GG = Inches(0.06)
    GCOL_GAP = Inches(0.15)

    for i, (name, screens, path, cnt, fill, stk) in enumerate(groups):
        r, c = divmod(i, 2)
        gx = ML + c * (GW + GCOL_GAP)
        gy = GY0 + r * (GH + GG)
        sh = _box(s, gx, gy, GW, GH,
                  f'{name} ({cnt}개)\n{screens}\n{path}',
                  fill=fill, stk=stk, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        pars = list(sh.text_frame.paragraphs)
        if len(pars) >= 3:
            for r_obj in pars[-1].runs:
                r_obj.font.color.rgb = ASH
                r_obj.font.size = Pt(8)

    # 핵심 화면 상세
    DY = GY0 + 3 * (GH + GG) + Inches(0.08)
    det_boxes = [
        ('지시 추진현황\n/sys12/2/3',
         'Statistic 카드5개\n(총/완료/진행/미착수+추진율)\nProgress 게이지\n'
         '카테고리별 매트릭스\nDataTable 7컬럼\n엑셀 다운로드',
         ICE, SKY),
        ('지시사항 목록\n/sys12/2/3',
         'SearchForm(5필드)\nDataTable(8컬럼)\nCrudForm(8필드) Modal\n'
         '상세Modal+이행이력\nTimeline+조치사항등록\nStatusBadge 4종',
         MINT, GRN),
        ('관리자\n/sys12/4/1',
         'Tabs 9개 구성\n지시/건의사항 관리\n카테고리 CRUD\n'
         '통계 Bar 차트\n사용자/부대 관리\n기간/알림/권한 설정',
         CREAM, WARN),
    ]
    DW = Inches(2.95)
    DH = Inches(1.1)
    DG = Inches(0.1)
    for i, (ttl, desc, fill, stk) in enumerate(det_boxes):
        dx = ML + i * (DW + DG)
        sh = _box(s, dx, DY, DW, DH,
                  f'{ttl}\n──────────\n{desc}',
                  fill=fill, stk=stk, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        sh.text_frame.paragraphs[1].runs[0].font.bold = True

    # [C] 서술 테이블
    dy = DY + DH + Inches(0.08)

    LW = int(CW * 0.48)
    _tbl(s, ML, dy, LW, Inches(0.95), [
        ['영역', '화면 수', 'URL 패턴', '주요 컴포넌트'],
        ['지시(대통령)', '2', '/sys12/2/1', 'Progress+DataTable+CrudForm'],
        ['지시(장관)', '2', '/sys12/2/2', 'Progress+DataTable+CrudForm'],
        ['지시(지휘관)', '2', '/sys12/2/3', 'Progress+DataTable+CrudForm'],
        ['건의사항', '2', '/sys12/3/1', 'Progress+DataTable+CrudForm'],
        ['게시판+관리자', '5', '/sys12/1,4,admin/*', 'SimpleBoardPage+Tabs+AdminRoutes'],
    ], cws=[int(LW * r) for r in [0.20, 0.10, 0.28, 0.42]])

    RX = ML + LW + Inches(0.15)
    RW = CW - LW - Inches(0.15)
    _tbl(s, RX, dy, RW, Inches(0.95), [
        ['핵심 흐름 시나리오'],
        ['시나리오 1: 지시사항 등록+추적 흐름\n'
         '1. 로그인 -> 대시보드\n'
         '2. [지시사항] > [지휘관 지시사항] 메뉴\n'
         '3. 추진현황 Statistic+Progress 확인\n'
         '4. [지시사항 등록] -> CrudForm(8필드)\n'
         '5. 목록 DataTable 행 -> 상세Modal\n'
         '6. 이행이력 Timeline 확인'],
        ['시나리오 2: 건의사항 등록+조치 흐름\n'
         '1. [건의사항] 메뉴 -> 추진현황 확인\n'
         '2. [건의사항 등록] -> CrudForm(7필드)\n'
         '3. 상세Modal -> [조치사항 등록] 클릭\n'
         '4. 조치 CrudForm(7필드:담당자+계획+내용)'],
    ])


# ══════════════════════════════════════════════════
# Page 5: 플로우차트 (업무처리흐름)
# ══════════════════════════════════════════════════

def p5_flowchart(prs):
    s = _sl(prs)
    _title(s, '업무처리 플로우차트',
           '핵심 프로세스: 지시 등록 -> 추적 -> 조치 — 4개 레인, 12단계', 5)

    y0 = Inches(0.82)

    # 4개 스윔레인
    lanes = [
        ('사용자 (관리자/건의자)', ICE, NAVY),
        ('프론트엔드 (React 18)', LBLUE, SKY),
        ('백엔드 (API/Mock)', MINT, GRN),
        ('DB / 외부', CREAM, WARN),
    ]
    LN_W = CW // 4
    LN_H = Inches(2.75)
    LN_HDR = Inches(0.22)

    for i, (name, fill, stk) in enumerate(lanes):
        lx = ML + LN_W * i
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = fill
        bg.line.color.rgb = stk
        bg.line.width = Pt(1)
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_HDR)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = stk
        hdr.line.fill.background()
        tf = hdr.text_frame
        _run(tf.paragraphs[0], name, 7, True, WHITE)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    steps = [
        (0, 0.30, '1. 지시사항 메뉴 접근\n추진현황 확인'),
        (1, 0.30, '2. 추진현황 렌더링\nStatistic+Progress+Matrix'),
        (0, 0.70, '3. 지시사항 등록 클릭\nCrudForm 8필드 입력'),
        (1, 0.70, '4. API 호출\nPOST /sys12/directives'),
        (2, 0.70, '5. 데이터 검증+저장\n지시번호 자동생성'),
        (3, 0.70, '6. Directive INSERT\n데이터 저장'),
        (1, 1.10, '7. DataTable 갱신\n8컬럼 + StatusBadge'),
        (0, 1.10, '8. 행 클릭 -> 상세Modal\n기본정보 + 이행이력'),
        (0, 1.50, '9. 조치사항 등록 클릭\nCrudForm 7필드 입력'),
        (1, 1.50, '10. API 호출\nPOST /.../actions'),
        (2, 1.90, '11. 조치사항 저장\n이력 추가'),
        (3, 1.90, '12. ActionHistory INSERT\n이행이력 저장'),
    ]

    SBW = LN_W - Inches(0.12)
    SBH = Inches(0.56)

    for lane, y_off, txt in steps:
        sx = ML + LN_W * lane + Inches(0.06)
        sy = y0 + LN_HDR + Inches(y_off)
        _box(s, sx, sy, SBW, SBH, txt,
             fill=WHITE, stk=ASH, sw=1, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)

    # 판단 분기
    fx = ML + LN_W * 1 + Inches(0.06)
    fy = y0 + LN_HDR + Inches(1.35)
    _box(s, fx, fy, SBW, Inches(0.12),
         '검증 실패 -> 필드 에러 표시 -> 재입력',
         fill=PEACH, stk=ERR, sw=1, sz=5, fc=ERR, al=PP_ALIGN.CENTER)

    # [C] 서술 테이블
    dy = y0 + LN_H + Inches(0.1)

    LW = int(CW * 0.52)
    _tbl(s, ML, dy, LW, Inches(1.42), [
        ['단계', '동작', '조건/상세'],
        ['1-2', '추진현황 접근 -> 렌더링', 'Statistic 5카드(총/완료/진행/미착수+추진율) + 매트릭스 Table'],
        ['3-4', '지시등록 -> API호출', 'CrudForm 8필드(지시자군번/계급/성명/일자/부대/상태/내용/종류)'],
        ['5-6', '서버 검증 -> DB저장', '지시번호 자동생성, 상태 기본값 미착수'],
        ['7-8', '목록갱신 -> 상세열람', 'DataTable 8컬럼 + StatusBadge + 상세Modal + Timeline'],
        ['9-10', '조치등록 -> API호출', 'CrudForm 7필드(담당자군번/계급/성명/상태/계획/내용/첨부)'],
        ['11-12', '조치저장 -> 이력추가', 'ActionHistory INSERT + Timeline 갱신'],
    ], cws=[int(LW * r) for r in [0.08, 0.28, 0.64]])

    RX = ML + LW + Inches(0.12)
    RW = CW - LW - Inches(0.12)
    _tbl(s, RX, dy, RW, Inches(1.42), [
        ['업무 규칙 / 예외 처리'],
        ['[업무 규칙]\n'
         'BR-DRCT-001: 지시사항 3종 (대통령/국방부장관/지휘관)\n'
         'BR-DRCT-002: 추진상태 4종 (미착수/진행중/완료/지연)\n'
         'BR-DRCT-003: 지시종류 2종 (문서/구두)\n'
         'BR-DRCT-004: 수명부대 5개 (1사단/2사단/사령부/교육훈련단/상륙기동단)\n'
         'BR-DRCT-005: 카테고리 5종 (작전/교육/군수/인사/정보화)\n'
         'BR-DRCT-006: 조치사항 7필드 (담당자군번/계급/성명/상태/계획/내용/첨부)'],
        ['[예외 처리]\n'
         'E-001: 세션만료(401) -> 로그인 리다이렉트\n'
         'E-002: 권한없음(403) -> 접근 제한 알림\n'
         'E-003: 검증실패 -> 필드별 에러 표시+재입력\n'
         'E-004: 서버오류(500) -> 에러 메시지+재시도'],
    ])


# ══════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════

def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    print('[1/5] 시스템 개념도...')
    p1_concept(prs)
    print('[2/5] 목표모델...')
    p2_target(prs)
    print('[3/5] 유스케이스...')
    p3_usecase(prs)
    print('[4/5] 스토리보드...')
    p4_storyboard(prs)
    print('[5/5] 플로우차트...')
    p5_flowchart(prs)

    out_path = os.path.join(OUT, 'SYS12_지시건의사항관리체계_분석설계_v6.pptx')
    prs.save(out_path)
    print(f'\n완료! 총 {len(prs.slides)}장 슬라이드')
    print(f'파일: {out_path}')


if __name__ == '__main__':
    main()
