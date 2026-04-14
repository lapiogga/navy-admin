#!/usr/bin/env python3
"""
SYS14 나의 제언 기능명세서 PPTX v5
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
system_level_spec_guide_v5.md 기반 핵심 5페이지:
  P1. 시스템 개념도      P2. 목표모델 (As-Is → To-Be)
  P3. 유스케이스          P4. 스토리보드 (화면흐름도)
  P5. 플로우차트 (업무처리흐름)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ═══ 색상 팔레트 ═══
NAVY  = RGBColor(0x6B, 0x8E, 0xA8)
TEAL  = RGBColor(0x8E, 0xC5, 0xC0)
SKY   = RGBColor(0xA8, 0xC8, 0xE8)
MINT  = RGBColor(0xDC, 0xF5, 0xE4)
PEACH = RGBColor(0xFC, 0xD4, 0xD0)
CREAM = RGBColor(0xFD, 0xEE, 0xC2)
ICE   = RGBColor(0xE4, 0xEE, 0xF8)
SNOW  = RGBColor(0xF7, 0xF9, 0xFC)
CHAR  = RGBColor(0x1A, 0x1A, 0x1A)
ASH   = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRN   = RGBColor(0xA4, 0xD4, 0xB4)
WARN  = RGBColor(0xF5, 0xCC, 0xA0)
ERR   = RGBColor(0xF0, 0xB2, 0xAE)
LBLUE = RGBColor(0xDC, 0xEA, 0xFE)
LSUB  = RGBColor(0xC5, 0xDC, 0xEC)
HBAR  = RGBColor(0xD4, 0xE2, 0xEC)  # 제목바 연한 배경
THBG  = RGBColor(0xD4, 0xE2, 0xEC)  # 표 헤더 연한 배경
BORD  = RGBColor(0x4A, 0x6B, 0x8C)  # 박스 테두리 (진한)
TBLB  = RGBColor(0x80, 0x80, 0x80)  # 표 테두리 (회색)

# ═══ 치수 ═══
SW = Inches(10)
SH = Inches(5.625)
ML = Inches(0.4)
CW = Inches(9.2)
FONT = 'Freesentation 5 Medium'
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'ppt')


# ══════════════════════════════════════════════════
# 헬퍼 함수
# ══════════════════════════════════════════════════

def _sl(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = SNOW
    return s


def _run(p, txt, sz=9, b=False, c=CHAR):
    """단락에 서식 런 추가 (최소 9pt 보장)"""
    sz = max(sz, 9)
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = b
    r.font.color.rgb = c
    r.font.name = FONT
    return r


def _tb(s, l, t, w, h, txt, sz=9, b=False, c=CHAR, al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    _run(tb.text_frame.paragraphs[0], txt, sz, b, c)
    tb.text_frame.paragraphs[0].alignment = al
    return tb


def _title(s, title, sub, pg):
    """[A] 제목 영역 — 연한 바 + 진한 테두리 + 검정 글자"""
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = HBAR
    bar.line.color.rgb = BORD
    bar.line.width = Pt(0.5)
    _tb(s, ML, Inches(0.06), Inches(7.5), Inches(0.36), title, 22, True, CHAR)
    if sub:
        _tb(s, ML, Inches(0.42), Inches(8.5), Inches(0.26), sub, 11, False, CHAR)
    _tb(s, Inches(9.0), Inches(0.18), Inches(0.6), Inches(0.3),
        f'{pg}/5', 12, True, CHAR, PP_ALIGN.RIGHT)


def _box(s, l, t, w, h, txt, fill=ICE, stk=BORD, sw=0.5,
         sz=9, b=False, fc=CHAR, al=PP_ALIGN.CENTER):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = stk
    sh.line.width = Pt(sw)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.margin_left = tf.margin_right = Inches(0.03)
    for i, ln in enumerate(txt.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p, ln, sz, b, fc)
        p.alignment = al
        p.space_before = p.space_after = Pt(0)
    bp = sh._element.txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', 'ctr')
    return sh


def _darr(s, cx, y, h=Inches(0.08)):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                           cx - Inches(0.06), y, Inches(0.12), h)
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _rarr(s, x, cy, w=Inches(0.18)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                           x, cy - Inches(0.05), w, Inches(0.1))
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _cell_bg(cell, hx):
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for sf in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(sf)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), val=hx)


def _cell_border(cell, color_hex='808080', width_pt=0.5):
    """셀 테두리 4면 회색 설정"""
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for border in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for existing in tcPr.findall(qn(border)):
            tcPr.remove(existing)
        ln = etree.SubElement(tcPr, qn(border))
        ln.set('w', str(int(width_pt * 12700)))
        ln.set('cap', 'flat')
        ln.set('cmpd', 'sng')
        ln.set('algn', 'ctr')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr'), val=color_hex)
        etree.SubElement(ln, qn('a:prstDash'), val='solid')
        etree.SubElement(ln, qn('a:round'))


def _tbl(s, l, t, w, h, data, cws=None):
    """테이블 (멀티라인 셀 지원, 진한 테두리 0.5pt, 연한 헤더 배경). data[0]=헤더"""
    rows, cols = len(data), len(data[0])
    sh = s.shapes.add_table(rows, cols, l, t, w, h)
    tbl = sh.table
    if cws:
        for i, cw in enumerate(cws):
            tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ''
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_top = tf.margin_bottom = Inches(0.04)
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.06)
            hdr = (r == 0)
            for i, ln in enumerate(data[r][c].split('\n')):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                _run(p, ln, 9 if not hdr else 10, hdr, CHAR)
                p.alignment = PP_ALIGN.CENTER if hdr else PP_ALIGN.LEFT
                p.space_before = p.space_after = Pt(0)
            _cell_bg(cell, 'D4E2EC' if hdr else ('FFFFFF' if r % 2 == 1 else 'F2F7FB'))
            _cell_border(cell, '808080', 0.5)
    return tbl


# ══════════════════════════════════════════════════
# Page 1: 시스템 개념도
# ══════════════════════════════════════════════════

def p1_concept(prs):
    s = _sl(prs)
    _title(s, '시스템 개념도',
           '나의 제언 체계의 전체 아키텍처 및 기술 구성 — 6개 계층', 1)

    layers = [
        ('사용자\n계층', [
            ('제언자 (일반사용자)\n제언 작성·추천·댓글', ICE, NAVY),
            ('관리자\n답변·상태변경·비공개', ICE, NAVY),
            ('시스템관리자\n서식·권한·코드 관리', ICE, NAVY),
        ]),
        ('채널\n계층', [
            ('웹 브라우저\nChrome/Edge 1920×1080', ICE, TEAL),
            ('반응형 모바일 웹\n360px~', ICE, TEAL),
            ('API Gateway\nRESTful 15 endpoints', ICE, TEAL),
        ]),
        ('프론트\n엔드', [
            ('React 18 + TS 5\nVite 5 · Router 7', LBLUE, SKY),
            ('Ant Design 5.24\nStatistic · Upload · List', LBLUE, SKY),
            ('Zustand 5 + TQ 5\n상태관리 · 서버캐싱', LBLUE, SKY),
        ]),
        ('백엔드\n계층', [
            ('Spring Boot 3.x\n(추후 개발) Java 17', MINT, GRN),
            ('MSW 2.12 Mock\nFaker.js 9 데이터', MINT, GRN),
            ('REST API 설계\nJSON/HTTPS 15개', MINT, GRN),
        ]),
        ('데이터\n계층', [
            ('PostgreSQL 15\n(추후) 정규화 설계', CREAM, WARN),
            ('Mock 데이터 25건\n제언+댓글+서식', CREAM, WARN),
            ('파일 스토리지\n첨부파일 PDF (추후)', CREAM, WARN),
        ]),
        ('외부\n연계', [
            ('GPKI/SSO 인증\n(추후 연계)', SNOW, ASH),
            ('국방민원포탈\n(추후 연계)', SNOW, ASH),
            ('알림톡/SMS\n(추후 연계)', SNOW, ASH),
        ]),
    ]

    y0 = Inches(0.82)
    LH = Inches(0.44)
    GAP = Inches(0.05)
    LBL_W = Inches(0.85)
    LBL_G = Inches(0.05)

    for idx, (lbl, items) in enumerate(layers):
        y = y0 + idx * (LH + GAP)
        _box(s, ML, y, LBL_W, LH, lbl, fill=NAVY, stk=NAVY, sz=7, b=True, fc=WHITE)
        n = len(items)
        avail = CW - LBL_W - LBL_G
        igap = Inches(0.04)
        iw = (avail - igap * (n - 1)) // n
        x = ML + LBL_W + LBL_G
        for txt, fill, stk in items:
            _box(s, x, y, iw, LH, txt, fill=fill, stk=stk, sz=7, fc=CHAR)
            x += iw + igap
        if idx < len(layers) - 1:
            _darr(s, ML + CW // 2, y + LH + Inches(0.01))

    # [C] 서술 테이블
    dy = y0 + 6 * (LH + GAP) + Inches(0.05)
    _tbl(s, ML, dy, CW, Inches(1.45), [
        ['프론트엔드 기술스택', '백엔드 기술스택', '인프라 / 보안 / 연계'],
        ['• React 18.x + TypeScript 5.x\n'
         '• Ant Design 5.24 (Statistic/Upload)\n'
         '• Zustand 5 + TanStack Query 5\n'
         '• React Router 7 + Vite 5\n'
         '• StatusBadge + CrudForm\n'
         '• 3개 고유 페이지, 7개 화면',
         '• Spring Boot 3.x (추후 개발)\n'
         '• MSW 2.12 Mock API (현재)\n'
         '• Faker.js 9 한국어 데이터\n'
         '• REST API 15개 endpoint\n'
         '• JWT + Redis 인증 (추후)\n'
         '• MyBatis 3.5 ORM (추후)',
         '• PostgreSQL 15 (추후 개발)\n'
         '• Linux / Tomcat 10.x (추후)\n'
         '• SSL/TLS 1.3, GPKI (추후)\n'
         '• 국방민원포탈 연계 (추후)\n'
         '• 알림톡/SMS/이메일 (추후)\n'
         '• 사용자: 3역할, 16 프로세스'],
    ], cws=[CW // 3] * 3)


# ══════════════════════════════════════════════════
# Page 2: 목표모델 (As-Is → To-Be)
# ══════════════════════════════════════════════════

def p2_target(prs):
    s = _sl(prs)
    _title(s, '목표모델 (As-Is → To-Be)',
           '현행 문제점 진단 및 목표 시스템 기능 구성', 2)

    pairs = [
        ('건의 문서 수기 관리\n오프라인 문서함, 접수·추적 불가',
         '온라인 제언 포탈\n등록·접수·진행·완료 4단계 추적'),
        ('진행상태 파악 불가\n접수 여부·조치 현황 확인 어려움',
         '5단계 상태 관리\n등록→접수→진행→완료/반려 추적'),
        ('익명성·비공개 미지원\n민감 제언 공개 우려로 활용 저조',
         '비공개 처리 + 추천/신고\n관리자 비공개 전환 + 추천·신고'),
        ('답변 체계 부재\n관리자 답변·조치유형 관리 불가',
         '답변·조치유형 관리\n관리자 답변+조치유형(4종)+조치일'),
        ('서식·엑셀 출력 없음\n표준 서식 없음, 데이터 활용 불가',
         '서식관리 + CSV 출력\n서식 CRUD + 전체 데이터 엑셀 내보내기'),
    ]

    y0 = Inches(1.0)
    RH = Inches(0.52)
    RG = Inches(0.03)
    AW = Inches(3.5)
    ARR_W = Inches(0.22)
    AX = ML
    ARX = ML + AW + Inches(0.15)
    TX = ARX + ARR_W + Inches(0.15)

    _box(s, AX, y0 - Inches(0.06), AW, Inches(0.18),
         '현행 시스템 (As-Is)', fill=ERR, stk=ERR, sz=8, b=True, fc=WHITE)
    _box(s, TX, y0 - Inches(0.06), AW, Inches(0.18),
         '목표 시스템 (To-Be)', fill=GRN, stk=GRN, sz=8, b=True, fc=WHITE)

    for i, (asis, tobe) in enumerate(pairs):
        y = y0 + i * (RH + RG)
        _box(s, AX, y, AW, RH, asis, fill=PEACH, stk=ERR, sz=7, fc=CHAR)
        _rarr(s, ARX, y + RH // 2, ARR_W)
        _box(s, TX, y, AW, RH, tobe, fill=MINT, stk=GRN, sz=7, fc=CHAR)

    # 우측 기능 요약
    FX = TX + AW + Inches(0.12)
    FW = ML + CW - FX
    _box(s, FX, y0, FW, 5 * RH + 4 * RG,
         '핵심 기능\n─────────\n'
         '제언메인     1\n제언확인     1\n관리자        1\n'
         '게시판        2\n권한관리     1\n'
         '─────────\n총 7개 메뉴\n16 프로세스\n'
         '─────────\n상태 5종\n등록/접수/진행/완료/반려',
         fill=ICE, stk=TEAL, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)

    # [C] 서술 테이블
    dy = y0 + 5 * (RH + RG) + Inches(0.12)
    _tbl(s, ML, dy, CW, Inches(1.25), [
        ['개선항목', '현행 문제점', '목표 개선방향', '기대 효과', '정량 목표'],
        ['접수관리', '수기 접수, 추적 불가', '온라인 등록+4단계 추적', '제언 100% 추적', '추적률 100%'],
        ['상태관리', '진행 여부 확인 불가', '5단계 상태(등록~반려)', '실시간 상태 확인', '5단계'],
        ['비공개', '민감 제언 공개 우려', '비공개 전환+추천/신고', '안전한 제언 환경', '비공개 ON'],
        ['답변체계', '답변·조치 관리 불가', '답변+조치유형(4종)+조치일', '체계적 조치 관리', '답변률 90%'],
        ['서식출력', '서식 없음, 활용 불가', '서식 CRUD+CSV 엑셀 출력', '데이터 활용 극대화', '출력 즉시'],
    ], cws=[int(CW * r) for r in [0.10, 0.24, 0.24, 0.22, 0.10]])


# ══════════════════════════════════════════════════
# Page 3: 유스케이스
# ══════════════════════════════════════════════════

def p3_usecase(prs):
    s = _sl(prs)
    _title(s, '유스케이스 다이어그램',
           '3개 액터 x 8개 유스케이스, 16개 단위 프로세스', 3)

    y0 = Inches(0.82)

    # 좌측: 액터
    actors = [
        '제언자 (일반사용자)\n제언 작성·추천·댓글\nROLE_USER',
        '관리자\n답변·상태·비공개\nROLE_MANAGER',
        '시스템관리자\n서식·권한 관리\nROLE_ADMIN',
    ]
    AX = ML
    AW = Inches(1.15)
    AH = Inches(0.55)
    AG = Inches(0.15)
    for i, atxt in enumerate(actors):
        _box(s, AX, y0 + i * (AH + AG), AW, AH, atxt,
             fill=ICE, stk=NAVY, sz=7, fc=CHAR)

    # 중앙: 시스템 바운더리
    BX = AX + AW + Inches(0.2)
    BW = Inches(6.2)
    BH = 3 * AH + 2 * AG + Inches(0.3)
    bnd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, BX, y0 - Inches(0.05), BW, BH)
    bnd.fill.solid()
    bnd.fill.fore_color.rgb = RGBColor(0xFB, 0xFB, 0xFB)
    bnd.line.color.rgb = NAVY
    bnd.line.width = Pt(2)
    _tb(s, BX + Inches(0.1), y0 - Inches(0.02), Inches(5), Inches(0.18),
        '나의 제언 (SYS14) — 바운더리', 8, True, NAVY)

    # UC 그룹
    uc_groups = [
        ('제언 메인 (3)', 'UC-SUGST-001~003\n통계카드4종+공지+최신제언', ICE, SKY),
        ('제언확인 (6)', 'UC-SUGST-004~009\n검색·등록·상세·추천·댓글·엑셀', ICE, SKY),
        ('관리자-제언관리 (4)', 'UC-SUGST-010~013\n답변·상태변경·비공개·반려', ICE, SKY),
        ('관리자-서식관리 (2)', 'UC-SUGST-014~015\n서식 CRUD + 목록', LBLUE, TEAL),
        ('게시판 (2)', 'UC-SUGST-016~017\n공지사항·질의응답', LBLUE, TEAL),
        ('권한관리 (1)', 'UC-SYS\n사용자별 권한 등록', LBLUE, TEAL),
    ]
    UY0 = y0 + Inches(0.22)
    UW = Inches(2.95)
    UH = Inches(0.48)
    UG = Inches(0.03)
    for i, (grp, desc, fill, stk) in enumerate(uc_groups):
        r, c = divmod(i, 2)
        ux = BX + Inches(0.08) + c * (UW + Inches(0.1))
        uy = UY0 + r * (UH + UG)
        sh = _box(s, ux, uy, UW, UH, f'{grp}\n{desc}',
                  fill=fill, stk=stk, sz=6, b=False, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True

    # 우측: 외부 시스템
    EX = BX + BW + Inches(0.12)
    EW = ML + CW - EX
    _box(s, EX, y0 + Inches(0.2), EW, Inches(0.35),
         '국방민원포탈\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)
    _box(s, EX, y0 + Inches(0.65), EW, Inches(0.35),
         'GPKI/SSO\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)

    # [C] 서술 테이블
    dy = Inches(3.2)
    _tbl(s, ML, dy, CW, Inches(2.1), [
        ['UC-ID', '유스케이스명', '주 액터', '기능 설명 (구체적)', '우선'],
        ['UC-SUGST-001', '제언 메인 대시보드', '사용자',
         'Statistic 카드4종(전체/이번달/답변완료/대기) + 공지사항 List + 최신제언 List', '상'],
        ['UC-SUGST-004', '제언 목록/검색', '사용자',
         'SearchForm(4필드: 제목/작성자/진행상태/조치유형) + DataTable(8컬럼: 순번~조치유형)', '상'],
        ['UC-SUGST-005', '제언 작성', '사용자',
         'CrudForm 6필드 Modal (소속/제언자/직책/전화/제목/내용) + Upload 첨부 최대5개', '상'],
        ['UC-SUGST-006', '제언 상세', '사용자',
         'Descriptions(bordered) 12항목 + 답변영역 + 반려사유 + 댓글 CRUD + 상태관리', '상'],
        ['UC-SUGST-010', '제언 관리(관리자)', '관리자',
         'DataTable(9컬럼) + 상태변경Modal(진행상태/조치유형/조치일/담당부서/반려사유)', '상'],
        ['UC-SUGST-014', '서식관리', '관리자',
         'Tabs 2탭(제언관리/서식관리) + DataTable + CrudForm(서식명/내용) CRUD', '중'],
        ['UC-SUGST-016', '게시판', '사용자',
         'SimpleBoardPage 재사용: 공지사항(/sys14/board/1), 질의응답(/sys14/board/2)', '하'],
    ], cws=[int(CW * r) for r in [0.10, 0.14, 0.06, 0.64, 0.05]])


# ══════════════════════════════════════════════════
# Page 4: 스토리보드 (화면흐름도)
# ══════════════════════════════════════════════════

def p4_storyboard(prs):
    s = _sl(prs)
    _title(s, '스토리보드 (화면흐름도)',
           '총 7개 화면, 4개 업무영역, 핵심 네비게이션 흐름', 4)

    y0 = Inches(0.82)

    # 로그인 → 대시보드
    _box(s, ML, y0, Inches(1.3), Inches(0.35),
         '로그인\n/login', fill=WARN, stk=WARN, sz=8, b=True, fc=WHITE)
    _rarr(s, ML + Inches(1.35), y0 + Inches(0.12), Inches(0.25))
    _box(s, ML + Inches(1.65), y0, Inches(2.0), Inches(0.35),
         '대시보드 / 홈\n/sys14', fill=NAVY, stk=NAVY, sz=8, b=True, fc=WHITE)

    _darr(s, ML + Inches(2.65), y0 + Inches(0.37))

    # 4개 메뉴 그룹
    groups = [
        ('1. 나의제언', '메인 · 공지 · 제언확인 · 관리자', '/sys14/1/*', '4', ICE, SKY),
        ('2. 게시판', '공지사항 · 질의응답', '/sys14/board/*', '2', LBLUE, TEAL),
        ('3. 관리자', '권한관리', '/sys14/admin/*', '1', LBLUE, TEAL),
    ]

    GY0 = y0 + Inches(0.55)
    GW = Inches(4.5)
    GH = Inches(0.58)
    GG = Inches(0.03)
    GCOL_GAP = Inches(0.15)

    for i, (name, screens, path, cnt, fill, stk) in enumerate(groups):
        r, c = divmod(i, 2)
        gx = ML + c * (GW + GCOL_GAP)
        gy = GY0 + r * (GH + GG)
        sh = _box(s, gx, gy, GW, GH,
                  f'{name} ({cnt}개)\n{screens}\n{path}',
                  fill=fill, stk=stk, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        pars = list(sh.text_frame.paragraphs)
        if len(pars) >= 3:
            for r_obj in pars[-1].runs:
                r_obj.font.color.rgb = ASH
                r_obj.font.size = Pt(8)

    # 핵심 화면 상세 (3개 고유 페이지)
    DY = GY0 + 2 * (GH + GG) + Inches(0.08)
    det_boxes = [
        ('제언 메인\n/sys14/1/1',
         'Statistic 카드 4종\n전체/이번달/답변완료/대기\n'
         'Row(12:12) 2열 배치\n좌: 공지사항 List\n우: 최신제언 List\n'
         'StatusBadge 상태 표시',
         ICE, SKY),
        ('제언확인\n/sys14/1/3',
         'SearchForm(4필드)\n제목/작성자/상태/조치유형\n'
         'DataTable(8컬럼)\n순번(연도별)/제목/제언자\n'
         '상세Modal+댓글CRUD\n엑셀 출력 + Upload 첨부',
         MINT, GRN),
        ('관리자\n/sys14/1/4',
         'Tabs 2탭 구성\n① 제언관리: DataTable(9컬럼)\n'
         '상태변경Modal 5항목\n답변등록 Modal\n'
         '② 서식관리: DataTable\nCrudForm(서식명/내용) CRUD',
         CREAM, WARN),
    ]
    DW = Inches(2.95)
    DH = Inches(1.1)
    DG = Inches(0.1)
    for i, (ttl, desc, fill, stk) in enumerate(det_boxes):
        dx = ML + i * (DW + DG)
        sh = _box(s, dx, DY, DW, DH,
                  f'{ttl}\n──────────\n{desc}',
                  fill=fill, stk=stk, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        sh.text_frame.paragraphs[1].runs[0].font.bold = True

    # [C] 서술 테이블 (좌우 2개)
    dy = DY + DH + Inches(0.08)

    LW = int(CW * 0.48)
    _tbl(s, ML, dy, LW, Inches(0.95), [
        ['영역', '화면 수', 'URL 패턴', '주요 컴포넌트'],
        ['제언메인', '1', '/sys14/1/1', 'Statistic+Card+List'],
        ['제언확인', '1', '/sys14/1/3', 'SearchForm+DataTable+Modal'],
        ['관리자', '2탭', '/sys14/1/4', 'Tabs+DataTable+CrudForm'],
        ['공지사항', '1', '/sys14/1/2', 'SimpleBoardPage'],
        ['게시판+권한', '3', '/sys14/board,admin/*', 'SimpleBoardPage+AuthGroup'],
    ], cws=[int(LW * r) for r in [0.20, 0.10, 0.28, 0.42]])

    RX = ML + LW + Inches(0.15)
    RW = CW - LW - Inches(0.15)
    _tbl(s, RX, dy, RW, Inches(0.95), [
        ['핵심 흐름 시나리오'],
        ['시나리오 1: 제언 작성·확인 흐름\n'
         '① 로그인 → 메인 대시보드 (통계카드 4종)\n'
         '② [제언확인] 메뉴 → 목록 + SearchForm\n'
         '③ [제언 작성] 버튼 → CrudForm+Upload Modal\n'
         '④ 6항목 입력 + 첨부파일 → [저장]\n'
         '⑤ 상태: 등록 → 관리자 접수 대기\n'
         '⑥ 행 클릭 → 상세 Modal (추천/신고/댓글)'],
        ['시나리오 2: 관리자 처리 흐름\n'
         '① [관리자] > [제언관리] 탭 → DataTable\n'
         '② [상태변경] → Modal (상태/조치유형/조치일/부서)\n'
         '③ [답변] 버튼 → 답변 등록 Modal\n'
         '④ [비공개] → Popconfirm 확인 → 비공개 전환'],
    ])


# ══════════════════════════════════════════════════
# Page 5: 플로우차트 (업무처리흐름)
# ══════════════════════════════════════════════════

def p5_flowchart(prs):
    s = _sl(prs)
    _title(s, '업무처리 플로우차트',
           '핵심 프로세스: 제언 작성 → 접수 → 답변 — 4개 레인, 12단계', 5)

    y0 = Inches(0.82)

    # 4개 스윔레인
    lanes = [
        ('제언자 (일반사용자)', ICE, NAVY),
        ('프론트엔드 (React 18)', LBLUE, SKY),
        ('백엔드 (API/Mock)', MINT, GRN),
        ('DB / 외부', CREAM, WARN),
    ]
    LN_W = CW // 4
    LN_H = Inches(2.75)
    LN_HDR = Inches(0.22)

    for i, (name, fill, stk) in enumerate(lanes):
        lx = ML + LN_W * i
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = fill
        bg.line.color.rgb = stk
        bg.line.width = Pt(1)
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_HDR)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = stk
        hdr.line.fill.background()
        tf = hdr.text_frame
        _run(tf.paragraphs[0], name, 7, True, WHITE)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    steps = [
        (0, 0.30, '① 제언 작성 버튼\n제언확인 페이지'),
        (1, 0.30, '② CrudForm + Upload\n6필드+첨부파일 Modal'),
        (0, 0.70, '③ 6항목 입력 + 저장\n소속/제언자/직책/전화/제목/내용'),
        (1, 0.70, '④ POST /suggestions\nAPI 호출 (JSON body)'),
        (2, 0.70, '⑤ 유효성 검증\n필수항목 + 첨부파일 체크'),
        (3, 0.70, '⑥ Suggestion INSERT\nstatus=registered 저장'),
        (1, 1.10, '⑦ 관리자: 목록 조회\n제언관리 탭 DataTable'),
        (0, 1.10, '⑧ 상태변경 버튼\nModal(상태/조치/부서/반려)'),
        (2, 1.10, '⑨ PATCH /status\n상태+조치유형+조치일 변경'),
        (0, 1.50, '⑩ 답변 버튼 클릭\n답변 등록 Modal'),
        (1, 1.50, '⑪ PATCH /answer\n답변 내용 API 호출'),
        (3, 1.90, '⑫ Suggestion UPDATE\nstatus+answer 갱신'),
    ]

    SBW = LN_W - Inches(0.12)
    SBH = Inches(0.40)

    for lane, y_off, txt in steps:
        sx = ML + LN_W * lane + Inches(0.06)
        sy = y0 + LN_HDR + Inches(y_off)
        _box(s, sx, sy, SBW, SBH, txt,
             fill=WHITE, stk=ASH, sw=1, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)

    # 판단 분기
    fx = ML + LN_W * 1 + Inches(0.06)
    fy = y0 + LN_HDR + Inches(1.35)
    _box(s, fx, fy, SBW, Inches(0.12),
         '반려 선택 → 반려사유 Modal → 사유 필수 입력',
         fill=PEACH, stk=ERR, sw=1, sz=5, fc=ERR, al=PP_ALIGN.CENTER)

    # [C] 서술 테이블 (좌우 2개)
    dy = y0 + LN_H + Inches(0.1)

    LW = int(CW * 0.52)
    _tbl(s, ML, dy, LW, Inches(1.42), [
        ['단계', '동작', '조건/상세'],
        ['①②', '작성버튼 → Modal', 'CrudForm 6필드 + Upload 첨부 최대5개'],
        ['③④', '입력완료 → API호출', 'POST /api/sys14/suggestions (JSON body)'],
        ['⑤⑥', '서버 검증 → DB저장', 'status=registered, 관리자 접수 대기'],
        ['⑦⑧', '관리자 목록 → 상태변경', 'Modal: 진행상태/조치유형/조치일/담당부서/반려사유'],
        ['⑨', '상태변경 API', 'PATCH /status (5단계: 등록→접수→진행→완료/반려)'],
        ['⑩⑪⑫', '답변등록 → DB갱신', 'PATCH /answer → Suggestion UPDATE'],
    ], cws=[int(LW * r) for r in [0.08, 0.28, 0.64]])

    RX = ML + LW + Inches(0.12)
    RW = CW - LW - Inches(0.12)
    _tbl(s, RX, dy, RW, Inches(1.42), [
        ['업무 규칙 / 예외 처리'],
        ['[업무 규칙]\n'
         'BR-SUGST-001: 진행상태 5종 (등록/접수/진행/완료/반려)\n'
         'BR-SUGST-002: 조치유형 4종 (정책반영/업무추진/기추진/업무참고)\n'
         'BR-SUGST-003: 순번 연도별 형식 (ex: 2026-1, 2026-2)\n'
         'BR-SUGST-004: 반려 시 반려사유 필수 입력\n'
         'BR-SUGST-005: 첨부파일 최대 5개 제한\n'
         'BR-SUGST-006: 로그인 정보(소속/성명/직책/전화) 자동 채움\n'
         'BR-SUGST-007: 비공개 전환 시 Popconfirm 확인\n'
         'BR-SUGST-008: 엑셀(CSV) 출력 — 전체 데이터 기반'],
        ['[예외 처리]\n'
         'E-001: 세션만료(401) → 로그인 리다이렉트\n'
         'E-002: 권한없음(403) → 접근 제한 알림\n'
         'E-003: 검증실패 → 필드별 에러 표시+재입력\n'
         'E-004: 서버오류(500) → 에러 메시지+재시도'],
    ])


# ══════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════

def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    print('[1/5] 시스템 개념도...')
    p1_concept(prs)
    print('[2/5] 목표모델...')
    p2_target(prs)
    print('[3/5] 유스케이스...')
    p3_usecase(prs)
    print('[4/5] 스토리보드...')
    p4_storyboard(prs)
    print('[5/5] 플로우차트...')
    p5_flowchart(prs)

    out_path = os.path.join(OUT, 'SYS14_나의제언_분석설계_v5.pptx')
    prs.save(out_path)
    print(f'\n완료! 총 {len(prs.slides)}장 슬라이드')
    print(f'파일: {out_path}')


if __name__ == '__main__':
    main()
