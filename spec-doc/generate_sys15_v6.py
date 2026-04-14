#!/usr/bin/env python3
"""
SYS15 보안일일결산체계 기능명세서 PPTX v6 (portrait)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
system_level_spec_guide_v6.md 기반 핵심 5페이지:
  P1. 시스템 개념도      P2. 목표모델 (As-Is → To-Be)
  P3. 유스케이스          P4. 스토리보드 (화면흐름도)
  P5. 플로우차트 (업무처리흐름)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ═══ 색상 팔레트 ═══
NAVY  = RGBColor(0x6B, 0x8E, 0xA8)
TEAL  = RGBColor(0x8E, 0xC5, 0xC0)
SKY   = RGBColor(0xA8, 0xC8, 0xE8)
MINT  = RGBColor(0xDC, 0xF5, 0xE4)
PEACH = RGBColor(0xFC, 0xD4, 0xD0)
CREAM = RGBColor(0xFD, 0xEE, 0xC2)
ICE   = RGBColor(0xE4, 0xEE, 0xF8)
SNOW  = RGBColor(0xF7, 0xF9, 0xFC)
CHAR  = RGBColor(0x1A, 0x1A, 0x1A)
ASH   = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRN   = RGBColor(0xA4, 0xD4, 0xB4)
WARN  = RGBColor(0xF5, 0xCC, 0xA0)
ERR   = RGBColor(0xF0, 0xB2, 0xAE)
LBLUE = RGBColor(0xDC, 0xEA, 0xFE)
LSUB  = RGBColor(0xC5, 0xDC, 0xEC)
HBAR  = RGBColor(0xD4, 0xE2, 0xEC)  # 제목바 연한 배경
THBG  = RGBColor(0xD4, 0xE2, 0xEC)  # 표 헤더 연한 배경
BORD  = RGBColor(0x4A, 0x6B, 0x8C)  # 박스 테두리 (진한)
TBLB  = RGBColor(0x80, 0x80, 0x80)  # 표 테두리 (회색)

# ═══ 치수 ═══
SW = Inches(7.5)
SH = Inches(10)
ML = Inches(0.4)
CW = Inches(6.7)       # 콘텐츠 폭 = 7.5 - 0.4*2
FONT = 'Freesentation 5 Medium'
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'ppt')


# ══════════════════════════════════════════════════
# 헬퍼 함수
# ══════════════════════════════════════════════════

def _sl(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = SNOW
    return s


def _run(p, txt, sz=9, b=False, c=CHAR):
    """단락에 서식 런 추가 (최소 9pt 보장)"""
    sz = max(sz, 9)
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = b
    r.font.color.rgb = c
    r.font.name = FONT
    return r


def _tb(s, l, t, w, h, txt, sz=9, b=False, c=CHAR, al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    _run(tb.text_frame.paragraphs[0], txt, sz, b, c)
    tb.text_frame.paragraphs[0].alignment = al
    return tb


def _title(s, title, sub, pg):
    """[A] 제목 영역 (portrait) — 연한 바 + 진한 테두리 + 검정 글자"""
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = HBAR
    bar.line.color.rgb = BORD
    bar.line.width = Pt(0.5)
    _tb(s, ML, Inches(0.06), Inches(5.5), Inches(0.36), title, 20, True, CHAR)
    if sub:
        _tb(s, ML, Inches(0.42), Inches(6.5), Inches(0.26), sub, 10, False, CHAR)
    _tb(s, Inches(6.5), Inches(0.18), Inches(0.6), Inches(0.3),
        f'{pg}/5', 12, True, CHAR, PP_ALIGN.RIGHT)


def _box(s, l, t, w, h, txt, fill=ICE, stk=BORD, sw=0.5,
         sz=9, b=False, fc=CHAR, al=PP_ALIGN.CENTER):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = stk
    sh.line.width = Pt(sw)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.margin_left = tf.margin_right = Inches(0.03)
    for i, ln in enumerate(txt.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p, ln, sz, b, fc)
        p.alignment = al
        p.space_before = p.space_after = Pt(0)
    bp = sh._element.txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', 'ctr')
    return sh


def _darr(s, cx, y, h=Inches(0.08)):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                           cx - Inches(0.06), y, Inches(0.12), h)
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _rarr(s, x, cy, w=Inches(0.18)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                           x, cy - Inches(0.05), w, Inches(0.1))
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _cell_bg(cell, hx):
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for sf in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(sf)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), val=hx)


def _cell_border(cell, color_hex='808080', width_pt=0.5):
    """셀 테두리 4면 회색 설정"""
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for border in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for existing in tcPr.findall(qn(border)):
            tcPr.remove(existing)
        ln = etree.SubElement(tcPr, qn(border))
        ln.set('w', str(int(width_pt * 12700)))
        ln.set('cap', 'flat')
        ln.set('cmpd', 'sng')
        ln.set('algn', 'ctr')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr'), val=color_hex)
        etree.SubElement(ln, qn('a:prstDash'), val='solid')
        etree.SubElement(ln, qn('a:round'))


def _tbl(s, l, t, w, h, data, cws=None):
    """테이블 (멀티라인 셀 지원, 진한 테두리 0.5pt, 연한 헤더 배경). data[0]=헤더"""
    rows, cols = len(data), len(data[0])
    sh = s.shapes.add_table(rows, cols, l, t, w, h)
    tbl = sh.table
    if cws:
        for i, cw in enumerate(cws):
            tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ''
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_top = tf.margin_bottom = Inches(0.04)
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.06)
            hdr = (r == 0)
            for i, ln in enumerate(data[r][c].split('\n')):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                _run(p, ln, 9 if not hdr else 10, hdr, CHAR)
                p.alignment = PP_ALIGN.CENTER if hdr else PP_ALIGN.LEFT
                p.space_before = p.space_after = Pt(0)
            _cell_bg(cell, 'D4E2EC' if hdr else ('FFFFFF' if r % 2 == 1 else 'F2F7FB'))
            _cell_border(cell, '808080', 0.5)
    return tbl


# ══════════════════════════════════════════════════
# Page 1: 시스템 개념도
# ══════════════════════════════════════════════════

def p1_concept(prs):
    s = _sl(prs)
    _title(s, '시스템 개념도',
           '보안일일결산체계의 전체 아키텍처 및 기술 구성 — 6개 계층', 1)

    layers = [
        ('사용자\n계층', [
            ('보안담당관\n비밀/매체/자재 관리·결산', ICE, NAVY),
            ('일반사용자\n개인보안 일일결산 실시', ICE, NAVY),
            ('시스템관리자\n점검항목·예외·알림 관리', ICE, NAVY),
        ]),
        ('채널\n계층', [
            ('웹 브라우저\nChrome/Edge 1920x1080', ICE, TEAL),
            ('반응형 모바일 웹\n360px~', ICE, TEAL),
            ('API Gateway\nRESTful 50+ endpoints', ICE, TEAL),
        ]),
        ('프론트\n엔드', [
            ('React 18 + TS 5\nVite 5 · Router 7', LBLUE, SKY),
            ('Ant Design 5.24\nCalendar·Checkbox·Form', LBLUE, SKY),
            ('Zustand 5 + TQ 5\n상태관리 · 서버캐싱', LBLUE, SKY),
        ]),
        ('백엔드\n계층', [
            ('Spring Boot 3.x\n(추후 개발) Java 17', MINT, GRN),
            ('MSW 2.12 Mock\nFaker.js 9 데이터', MINT, GRN),
            ('REST API 설계\nJSON/HTTPS 50+개', MINT, GRN),
        ]),
        ('데이터\n계층', [
            ('PostgreSQL 15\n(추후) 정규화 설계', CREAM, WARN),
            ('Mock 데이터 100+건\n비밀·매체·결산이력', CREAM, WARN),
            ('파일 스토리지\n비밀문서·예고문 (추후)', CREAM, WARN),
        ]),
        ('외부\n연계', [
            ('GPKI/SSO 인증\n(추후 연계)', SNOW, ASH),
            ('국방보안관리시스템\n(추후 연계)', SNOW, ASH),
            ('알림톡/SMS\n결산미실시 알림 (추후)', SNOW, ASH),
        ]),
    ]

    y0 = Inches(0.82)
    LH = Inches(0.62)
    GAP = Inches(0.08)
    LBL_W = Inches(0.85)
    LBL_G = Inches(0.05)

    for idx, (lbl, items) in enumerate(layers):
        y = y0 + idx * (LH + GAP)
        _box(s, ML, y, LBL_W, LH, lbl, fill=NAVY, stk=NAVY, sz=7, b=True, fc=WHITE)
        n = len(items)
        avail = CW - LBL_W - LBL_G
        igap = Inches(0.04)
        iw = (avail - igap * (n - 1)) // n
        x = ML + LBL_W + LBL_G
        for txt, fill, stk in items:
            _box(s, x, y, iw, LH, txt, fill=fill, stk=stk, sz=7, fc=CHAR)
            x += iw + igap
        if idx < len(layers) - 1:
            _darr(s, ML + CW // 2, y + LH + Inches(0.01))

    # [C] 서술 테이블
    dy = y0 + 6 * (LH + GAP) + Inches(0.05)
    _tbl(s, ML, dy, CW, Inches(1.45), [
        ['프론트엔드 기술스택', '백엔드 기술스택', '인프라 / 보안 / 연계'],
        ['• React 18.x + TypeScript 5.x\n'
         '• Ant Design 5.24 (Calendar/Checkbox/Form)\n'
         '• Zustand 5 + TanStack Query 5\n'
         '• React Router 7 + Vite 5\n'
         '• React Hook Form + Zod\n'
         '• 27개 고유 페이지, 138 프로세스',
         '• Spring Boot 3.x (추후 개발)\n'
         '• MSW 2.12 Mock API (현재)\n'
         '• Faker.js 9 한국어 데이터\n'
         '• REST API 50+ endpoint\n'
         '• JWT + Redis 인증 (추후)\n'
         '• MyBatis 3.5 ORM (추후)',
         '• PostgreSQL 15 (추후 개발)\n'
         '• Linux / Tomcat 10.x (추후)\n'
         '• SSL/TLS 1.3, GPKI (추후)\n'
         '• 국방보안관리시스템 (추후)\n'
         '• 알림톡/SMS/이메일 (추후)\n'
         '• 사용자: 3역할, 138 프로세스'],
    ], cws=[CW // 3] * 3)


# ══════════════════════════════════════════════════
# Page 2: 목표모델 (As-Is → To-Be)
# ══════════════════════════════════════════════════

def p2_target(prs):
    s = _sl(prs)
    _title(s, '목표모델 (As-Is → To-Be)',
           '현행 문제점 진단 및 목표 시스템 기능 구성', 2)

    pairs = [
        ('수기 보안일일결산\n종이 체크리스트로 매일 점검, 누락 빈번',
         '전자 보안일일결산\n개인/사무실 Checkbox+사유 자동 기록'),
        ('비밀/매체 대장 수기관리\nExcel 대장, 현용/파기/이관 추적 불가',
         '비밀/매체 통합 CRUD\n등급별 CRUD+인계인수+예고문 자동연계'),
        ('결산현황 파악 지연\n부대별 결산 실시여부 수합 지연',
         '결산종합현황 실시간\n비밀/개인/사무실/부재 4종 대시보드'),
        ('결재 프로세스 수기\n서면 결재로 승인 지연·분실 위험',
         '전자결재 워크플로우\n3단계 결재선(작성→검토→승인)'),
        ('보안수준 평가 부재\n개인별 보안수준 평가 체계 없음',
         '개인보안수준평가\n정기/수시 평가 + 등급산출(A~D)'),
    ]

    y0 = Inches(1.0)
    RH = Inches(0.78)
    RG = Inches(0.06)
    AW = Inches(2.6)
    ARR_W = Inches(0.22)
    AX = ML
    ARX = ML + AW + Inches(0.15)
    TX = ARX + ARR_W + Inches(0.15)

    _box(s, AX, y0 - Inches(0.06), AW, Inches(0.18),
         '현행 시스템 (As-Is)', fill=ERR, stk=ERR, sz=8, b=True, fc=WHITE)
    _box(s, TX, y0 - Inches(0.06), AW, Inches(0.18),
         '목표 시스템 (To-Be)', fill=GRN, stk=GRN, sz=8, b=True, fc=WHITE)

    for i, (asis, tobe) in enumerate(pairs):
        y = y0 + i * (RH + RG)
        _box(s, AX, y, AW, RH, asis, fill=PEACH, stk=ERR, sz=7, fc=CHAR)
        _rarr(s, ARX, y + RH // 2, ARR_W)
        _box(s, TX, y, AW, RH, tobe, fill=MINT, stk=GRN, sz=7, fc=CHAR)

    # 우측 기능 요약
    FX = TX + AW + Inches(0.12)
    FW = ML + CW - FX
    _box(s, FX, y0, FW, 5 * RH + 4 * RG,
         '핵심 기능\n─────────\n'
         '비밀/매체관리 5\n보안일일결산 6\n결재          2\n'
         '결산종합현황  4\n결산체계관리  6\n'
         '─────────\n총 27개 메뉴\n138 프로세스\n'
         '─────────\n8개 업무영역\n게시판2+관리자',
         fill=ICE, stk=TEAL, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)

    # [C] 서술 테이블
    dy = y0 + 5 * (RH + RG) + Inches(0.12)
    _tbl(s, ML, dy, CW, Inches(1.25), [
        ['개선항목', '현행 문제점', '목표 개선방향', '기대 효과', '정량 목표'],
        ['보안일일결산', '수기 체크리스트, 누락', '전자 Checkbox+미실시사유', '일일결산 자동화', '100% 전자화'],
        ['비밀/매체관리', 'Excel 대장, 추적 불가', '통합 CRUD+인계인수', '비밀 이력추적', '실시간 현황'],
        ['결산종합현황', '수합 지연, 통계 부재', '4종 대시보드 실시간', '현황 즉시 파악', '지연 0일'],
        ['결재프로세스', '서면 결재, 분실 위험', '3단계 전자결재 워크플로우', '결재 시간 단축', '당일 처리'],
        ['보안수준평가', '평가 체계 없음', '정기/수시 평가+등급산출', '개인 보안의식 향상', '평가율 100%'],
    ], cws=[int(CW * r) for r in [0.10, 0.24, 0.24, 0.22, 0.10]])


# ══════════════════════════════════════════════════
# Page 3: 유스케이스
# ══════════════════════════════════════════════════

def p3_usecase(prs):
    s = _sl(prs)
    _title(s, '유스케이스 다이어그램',
           '3개 액터 x 18개 유스케이스 그룹, 138개 단위 프로세스', 3)

    y0 = Inches(0.82)

    # 좌측: 액터
    actors = [
        '보안담당관\n비밀/매체/결산 관리\nROLE_SECURITY',
        '일반사용자\n개인결산 실시\nROLE_USER',
        '시스템관리자\n점검항목/알림 관리\nROLE_ADMIN',
    ]
    AX = ML
    AW = Inches(1.0)
    AH = Inches(0.55)
    AG = Inches(0.15)
    for i, atxt in enumerate(actors):
        _box(s, AX, y0 + i * (AH + AG), AW, AH, atxt,
             fill=ICE, stk=NAVY, sz=7, fc=CHAR)

    # 중앙: 시스템 바운더리
    BX = AX + AW + Inches(0.2)
    BW = Inches(4.3)
    BH = 3 * AH + 2 * AG + Inches(0.3)
    bnd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, BX, y0 - Inches(0.05), BW, BH)
    bnd.fill.solid()
    bnd.fill.fore_color.rgb = RGBColor(0xFB, 0xFB, 0xFB)
    bnd.line.color.rgb = NAVY
    bnd.line.width = Pt(2)
    _tb(s, BX + Inches(0.1), y0 - Inches(0.02), Inches(5), Inches(0.18),
        '보안일일결산체계 (SYS15) — 바운더리', 8, True, NAVY)

    # UC 그룹
    uc_groups = [
        ('비밀/매체관리 (5)', 'UC-SEC-001~005\n비밀/매체/자재/예고문/인계인수', ICE, SKY),
        ('보안일일결산 (6)', 'UC-SEC-006~011\n개인/사무실/점검관/수준평가/부재/교육', ICE, SKY),
        ('결재 (2)', 'UC-SEC-012~013\n결재대기 3단계, 결재완료이력', ICE, SKY),
        ('결산종합현황 (4)', 'UC-SEC-014~017\n비밀/개인/사무실/부재 종합', LBLUE, TEAL),
        ('결산체계관리 (6)', 'UC-SEC-018~023\n점검항목/휴무일/알림/로그/예외/개인설정', LBLUE, TEAL),
        ('게시판+관리자 (7)', 'UC-SYS 공지/질의/시스템관리\n코드·권한·결재선·게시판 관리', LBLUE, TEAL),
    ]
    UY0 = y0 + Inches(0.22)
    UW = Inches(2.0)
    UH = Inches(0.70)
    UG = Inches(0.06)
    for i, (grp, desc, fill, stk) in enumerate(uc_groups):
        r, c = divmod(i, 2)
        ux = BX + Inches(0.08) + c * (UW + Inches(0.1))
        uy = UY0 + r * (UH + UG)
        sh = _box(s, ux, uy, UW, UH, f'{grp}\n{desc}',
                  fill=fill, stk=stk, sz=6, b=False, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True

    # 우측: 외부 시스템
    EX = BX + BW + Inches(0.12)
    EW = ML + CW - EX
    _box(s, EX, y0 + Inches(0.2), EW, Inches(0.35),
         '국방보안관리시스템\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)
    _box(s, EX, y0 + Inches(0.65), EW, Inches(0.35),
         'GPKI/SSO\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)

    # [C] 서술 테이블
    dy = Inches(3.2)
    _tbl(s, ML, dy, CW, Inches(2.1), [
        ['UC-ID', '유스케이스명', '주 액터', '기능 설명 (구체적)', '우선'],
        ['UC-SEC-001', '비밀 관리', '보안담당관',
         'SecretMediaPage 공유 — 등급별(1/2/3/II/III급) CRUD + SearchForm(3필드) + DataTable(10컬럼)', '상'],
        ['UC-SEC-002', '저장매체 관리', '보안담당관',
         'SecretMediaPage type=media — 매체유형(USB/HDD/SSD/CD) CRUD + 용량·시리얼 관리', '상'],
        ['UC-SEC-006', '개인보안일일결산', '일반사용자',
         '필수5항목+선택3항목 Checkbox + 미체크시 사유입력 TextArea + 결재요청 제출', '상'],
        ['UC-SEC-007', '사무실보안일일결산', '보안담당관',
         '필수6항목+선택3항목 Checkbox + 미실시자/부재자 관리 Form + 결재요청', '상'],
        ['UC-SEC-012', '결재대기/승인', '보안담당관',
         '3단계 Steps(작성→검토→승인) + SearchForm(3필드) + 승인/반려 Modal + 사유입력', '상'],
        ['UC-SEC-014', '비밀/매체 종합현황', '보안담당관',
         'SearchForm(부대 Select) + DataTable(9컬럼: 비밀전체/현용/만료, 매체, 자재) + 엑셀저장', '중'],
        ['UC-SEC-018', '점검항목 관리', '시스템관리자',
         'Tabs 5종(개인필수/사무실필수/개인선택/사무실선택/보안수준) + DataTable CRUD + 가중치·필수여부', '중'],
    ], cws=[int(CW * r) for r in [0.10, 0.12, 0.08, 0.64, 0.05]])


# ══════════════════════════════════════════════════
# Page 4: 스토리보드 (화면흐름도)
# ══════════════════════════════════════════════════

def p4_storyboard(prs):
    s = _sl(prs)
    _title(s, '스토리보드 (화면흐름도)',
           '총 27개 화면, 8개 업무영역, 핵심 네비게이션 흐름', 4)

    y0 = Inches(0.82)

    # 로그인 → 대시보드
    _box(s, ML, y0, Inches(1.1), Inches(0.35),
         '로그인\n/login', fill=WARN, stk=WARN, sz=8, b=True, fc=WHITE)
    _rarr(s, ML + Inches(1.35), y0 + Inches(0.12), Inches(0.25))
    _box(s, ML + Inches(1.65), y0, Inches(2.0), Inches(0.35),
         '대시보드 / 결산캘린더\n/sys15', fill=NAVY, stk=NAVY, sz=8, b=True, fc=WHITE)

    _darr(s, ML + Inches(2.65), y0 + Inches(0.37))

    # 8개 메뉴 그룹
    groups = [
        ('1. 비밀/매체관리', '매체/비밀/예고문/자재/인계', '/sys15/2/*', '5', ICE, SKY),
        ('2. 보안일일결산', '개인/사무실/점검관/수준/부재/교육', '/sys15/3/*', '6', ICE, SKY),
        ('3. 결재', '결재대기 · 결재완료', '/sys15/4/*', '2', ICE, SKY),
        ('4. 결산종합현황', '비밀/개인/사무실/부재', '/sys15/5/*', '4', LBLUE, TEAL),
        ('5. 결산체계관리', '개인설정/점검항목/휴무일/알림/로그/예외', '/sys15/6,8/*', '6', LBLUE, TEAL),
        ('6. 게시판+관리자', '공지·질의·시스템관리', '/sys15/7,admin/*', '4+', LBLUE, TEAL),
    ]

    GY0 = y0 + Inches(0.55)
    GW = Inches(3.2)
    GH = Inches(0.90)
    GG = Inches(0.06)
    GCOL_GAP = Inches(0.15)

    for i, (name, screens, path, cnt, fill, stk) in enumerate(groups):
        r, c = divmod(i, 2)
        gx = ML + c * (GW + GCOL_GAP)
        gy = GY0 + r * (GH + GG)
        sh = _box(s, gx, gy, GW, GH,
                  f'{name} ({cnt}개)\n{screens}\n{path}',
                  fill=fill, stk=stk, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        pars = list(sh.text_frame.paragraphs)
        if len(pars) >= 3:
            for r_obj in pars[-1].runs:
                r_obj.font.color.rgb = ASH
                r_obj.font.size = Pt(8)

    # 핵심 화면 상세 (3개)
    DY = GY0 + 3 * (GH + GG) + Inches(0.08)
    det_boxes = [
        ('대시보드 (SecMainPage)\n/sys15/1/1',
         'Tabs(개인/사무실)\nCalendar Badge 월간현황\n완료(green)/미실시(red)\n'
         '날짜 클릭→상세 Modal\n미실시 Alert 자동 표시',
         ICE, SKY),
        ('개인보안일일결산\n/sys15/3/1',
         'Checkbox.Group 필수5+선택3\n미체크→사유 TextArea\n'
         '임시저장/제출(결재요청)\nRangePicker 이력 조회\nDataTable 이력 목록',
         MINT, GRN),
        ('비밀 관리\n/sys15/2/2',
         'SearchForm(3필드)\nDataTable(10컬럼)\n등급(1/2/3/II/III급)\n'
         'CrudForm Modal CRUD\n등록 시 예고문 자동 Modal',
         CREAM, WARN),
    ]
    DW = Inches(2.95)
    DH = Inches(1.1)
    DG = Inches(0.1)
    for i, (ttl, desc, fill, stk) in enumerate(det_boxes):
        dx = ML + i * (DW + DG)
        sh = _box(s, dx, DY, DW, DH,
                  f'{ttl}\n──────────\n{desc}',
                  fill=fill, stk=stk, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        sh.text_frame.paragraphs[1].runs[0].font.bold = True

    # [C] 서술 테이블 2개 좌우 배치
    dy = DY + DH + Inches(0.08)

    LW = int(CW * 0.48)
    _tbl(s, ML, dy, LW, Inches(0.95), [
        ['영역', '화면 수', 'URL 패턴', '주요 컴포넌트'],
        ['비밀/매체관리', '5', '/sys15/2/*', 'SecretMediaPage+CrudForm'],
        ['보안일일결산', '6', '/sys15/3/*', 'Checkbox+Form+DataTable'],
        ['결재', '2', '/sys15/4/*', 'Steps+SearchForm+Modal'],
        ['결산종합현황', '4', '/sys15/5/*', 'SearchForm+DataTable+엑셀'],
        ['체계관리+기타', '10+', '/sys15/6,7,8,admin', 'Tabs+DataTable+CRUD'],
    ], cws=[int(LW * r) for r in [0.20, 0.10, 0.28, 0.42]])

    RX = ML + LW + Inches(0.15)
    RW = CW - LW - Inches(0.15)
    _tbl(s, RX, dy, RW, Inches(0.95), [
        ['핵심 흐름 시나리오'],
        ['시나리오 1: 개인보안일일결산 흐름\n'
         '1. 로그인 -> 대시보드 캘린더 확인\n'
         '2. [보안일일결산] > [개인보안] 메뉴 선택\n'
         '3. 필수5항목 Checkbox 체크\n'
         '4. 미체크 항목 사유 입력\n'
         '5. [제출(결재요청)] 클릭 -> 결재선 전달\n'
         '6. 캘린더 Badge "완료"(green) 표시'],
        ['시나리오 2: 비밀 등록+예고문 흐름\n'
         '1. [비밀/매체관리] > [비밀 관리] 메뉴\n'
         '2. [등록] -> CrudForm Modal 입력\n'
         '3. 저장 -> 예고문 자동 Modal 오픈\n'
         '4. 예고문 제목/내용/예고일/수신부대 입력'],
    ])


# ══════════════════════════════════════════════════
# Page 5: 플로우차트 (업무처리흐름)
# ══════════════════════════════════════════════════

def p5_flowchart(prs):
    s = _sl(prs)
    _title(s, '업무처리 플로우차트',
           '핵심 프로세스: 일일결산 실시 -> 결재 -> 현황 — 4개 레인, 12단계', 5)

    y0 = Inches(0.82)

    # 4개 스윔레인
    lanes = [
        ('사용자 (보안담당관/일반)', ICE, NAVY),
        ('프론트엔드 (React 18)', LBLUE, SKY),
        ('백엔드 (API/Mock)', MINT, GRN),
        ('DB / 외부', CREAM, WARN),
    ]
    LN_W = CW // 4
    LN_H = Inches(2.75)
    LN_HDR = Inches(0.22)

    for i, (name, fill, stk) in enumerate(lanes):
        lx = ML + LN_W * i
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = fill
        bg.line.color.rgb = stk
        bg.line.width = Pt(1)
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_HDR)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = stk
        hdr.line.fill.background()
        tf = hdr.text_frame
        _run(tf.paragraphs[0], name, 7, True, WHITE)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    steps = [
        (0, 0.30, '1. 결산 메뉴 접근\n보안일일결산 선택'),
        (1, 0.30, '2. 체크리스트 렌더\nCheckbox.Group 필수+선택'),
        (0, 0.70, '3. 점검항목 체크\n미체크 시 사유 입력'),
        (1, 0.70, '4. 제출 API 호출\nPOST /personal-daily'),
        (2, 0.70, '5. 결산 데이터 저장\n상태: submitted'),
        (3, 0.70, '6. Daily INSERT\n결산이력 저장 (Mock)'),
        (1, 1.10, '7. 결재대기 목록 렌더\nSteps 3단계 표시'),
        (0, 1.10, '8. 결재자 승인/반려\nModal 사유입력'),
        (2, 1.10, '9. PATCH 결재상태\n/approvals/:id/approve'),
        (0, 1.50, '10. 종합현황 조회\n결산종합현황 메뉴'),
        (1, 1.50, '11. 대시보드 렌더\n4종 DataTable+통계'),
        (2, 1.90, '12. 현황 집계 API\nGET /summary/*'),
    ]

    SBW = LN_W - Inches(0.12)
    SBH = Inches(0.56)

    for lane, y_off, txt in steps:
        sx = ML + LN_W * lane + Inches(0.06)
        sy = y0 + LN_HDR + Inches(y_off)
        _box(s, sx, sy, SBW, SBH, txt,
             fill=WHITE, stk=ASH, sw=1, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)

    # 판단 분기
    fx = ML + LN_W * 1 + Inches(0.06)
    fy = y0 + LN_HDR + Inches(1.35)
    _box(s, fx, fy, SBW, Inches(0.12),
         '미체크 사유 누락 -> 경고 메시지 -> 재입력',
         fill=PEACH, stk=ERR, sw=1, sz=5, fc=ERR, al=PP_ALIGN.CENTER)

    # [C] 서술 테이블 2개 좌우 배치
    dy = y0 + LN_H + Inches(0.1)

    LW = int(CW * 0.52)
    _tbl(s, ML, dy, LW, Inches(1.42), [
        ['단계', '동작', '조건/상세'],
        ['1~2', '메뉴접근 -> 체크리스트', 'Checkbox.Group 필수5+선택3, 미체크시 TextArea 사유'],
        ['3~4', '점검체크 -> 제출', 'POST /personal-daily {checkedItems, uncheckedReasons, status}'],
        ['5~6', '서버저장 -> DB기록', 'status=submitted, 결재선 자동 연결'],
        ['7~8', '결재대기 -> 승인/반려', 'Steps 3단계(작성->검토->승인), Modal 사유입력'],
        ['9', '결재상태 갱신', 'PATCH /approvals/:id/approve or reject'],
        ['10~12', '종합현황 조회', 'GET /summary/secret, /summary/personal, /summary/office, /summary/absence'],
    ], cws=[int(LW * r) for r in [0.08, 0.28, 0.64]])

    RX = ML + LW + Inches(0.12)
    RW = CW - LW - Inches(0.12)
    _tbl(s, RX, dy, RW, Inches(1.42), [
        ['업무 규칙 / 예외 처리'],
        ['[업무 규칙]\n'
         'BR-SEC-001: 필수 점검항목 미체크 시 사유 필수 입력\n'
         'BR-SEC-002: 사무실결산 미실시자/부재자 사유 필수\n'
         'BR-SEC-003: 비밀등급 5종 (1급/2급/3급/II급/III급)\n'
         'BR-SEC-004: 비밀등록 시 예고문 자동 Modal 연계\n'
         'BR-SEC-005: 결재 3단계 워크플로우 (작성->검토->승인)\n'
         'BR-SEC-006: 개인보안수준 등급 A/B/C/D 산출\n'
         'BR-SEC-007: 점검항목 5종 탭 관리 (개인필수/사무실필수/개인선택/사무실선택/보안수준)'],
        ['[예외 처리]\n'
         'E-001: 세션만료(401) -> 로그인 리다이렉트\n'
         'E-002: 권한없음(403) -> 접근 제한 알림\n'
         'E-003: 미체크 사유 누락 -> message.warning 표시\n'
         'E-004: 서버오류(500) -> 에러 메시지+재시도'],
    ])


# ══════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════

def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    print('[1/5] 시스템 개념도...')
    p1_concept(prs)
    print('[2/5] 목표모델...')
    p2_target(prs)
    print('[3/5] 유스케이스...')
    p3_usecase(prs)
    print('[4/5] 스토리보드...')
    p4_storyboard(prs)
    print('[5/5] 플로우차트...')
    p5_flowchart(prs)

    out_path = os.path.join(OUT, 'SYS15_보안일일결산체계_분석설계_v6.pptx')
    prs.save(out_path)
    print(f'\n완료! 총 {len(prs.slides)}장 슬라이드')
    print(f'파일: {out_path}')


if __name__ == '__main__':
    main()
