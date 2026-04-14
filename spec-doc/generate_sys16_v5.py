#!/usr/bin/env python3
"""
SYS16 회의실예약관리체계 기능명세서 PPTX v5
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
system_level_spec_guide_v5.md 기반 핵심 5페이지:
  P1. 시스템 개념도      P2. 목표모델 (As-Is → To-Be)
  P3. 유스케이스          P4. 스토리보드 (화면흐름도)
  P5. 플로우차트 (업무처리흐름)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ═══ 색상 팔레트 ═══
NAVY  = RGBColor(0x6B, 0x8E, 0xA8)
TEAL  = RGBColor(0x8E, 0xC5, 0xC0)
SKY   = RGBColor(0xA8, 0xC8, 0xE8)
MINT  = RGBColor(0xDC, 0xF5, 0xE4)
PEACH = RGBColor(0xFC, 0xD4, 0xD0)
CREAM = RGBColor(0xFD, 0xEE, 0xC2)
ICE   = RGBColor(0xE4, 0xEE, 0xF8)
SNOW  = RGBColor(0xF7, 0xF9, 0xFC)
CHAR  = RGBColor(0x1A, 0x1A, 0x1A)
ASH   = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRN   = RGBColor(0xA4, 0xD4, 0xB4)
WARN  = RGBColor(0xF5, 0xCC, 0xA0)
ERR   = RGBColor(0xF0, 0xB2, 0xAE)
LBLUE = RGBColor(0xDC, 0xEA, 0xFE)
LSUB  = RGBColor(0xC5, 0xDC, 0xEC)
HBAR  = RGBColor(0xD4, 0xE2, 0xEC)  # 제목바 연한 배경
THBG  = RGBColor(0xD4, 0xE2, 0xEC)  # 표 헤더 연한 배경
BORD  = RGBColor(0x4A, 0x6B, 0x8C)  # 박스 테두리 (진한)
TBLB  = RGBColor(0x80, 0x80, 0x80)  # 표 테두리 (회색)

# ═══ 치수 ═══
SW = Inches(10)
SH = Inches(5.625)
ML = Inches(0.4)
CW = Inches(9.2)
FONT = 'Freesentation 5 Medium'
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'ppt')


# ══════════════════════════════════════════════════
# 헬퍼 함수
# ══════════════════════════════════════════════════

def _sl(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = SNOW
    return s


def _run(p, txt, sz=9, b=False, c=CHAR):
    """단락에 서식 런 추가 (최소 9pt 보장)"""
    sz = max(sz, 9)
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = b
    r.font.color.rgb = c
    r.font.name = FONT
    return r


def _tb(s, l, t, w, h, txt, sz=9, b=False, c=CHAR, al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    _run(tb.text_frame.paragraphs[0], txt, sz, b, c)
    tb.text_frame.paragraphs[0].alignment = al
    return tb


def _title(s, title, sub, pg):
    """[A] 제목 영역 — 연한 바 + 진한 테두리 + 검정 글자"""
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = HBAR
    bar.line.color.rgb = BORD
    bar.line.width = Pt(0.5)
    _tb(s, ML, Inches(0.06), Inches(7.5), Inches(0.36), title, 22, True, CHAR)
    if sub:
        _tb(s, ML, Inches(0.42), Inches(8.5), Inches(0.26), sub, 11, False, CHAR)
    _tb(s, Inches(9.0), Inches(0.18), Inches(0.6), Inches(0.3),
        f'{pg}/5', 12, True, CHAR, PP_ALIGN.RIGHT)


def _box(s, l, t, w, h, txt, fill=ICE, stk=BORD, sw=0.5,
         sz=9, b=False, fc=CHAR, al=PP_ALIGN.CENTER):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = stk
    sh.line.width = Pt(sw)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.margin_left = tf.margin_right = Inches(0.03)
    for i, ln in enumerate(txt.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p, ln, sz, b, fc)
        p.alignment = al
        p.space_before = p.space_after = Pt(0)
    bp = sh._element.txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', 'ctr')
    return sh


def _darr(s, cx, y, h=Inches(0.08)):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                           cx - Inches(0.06), y, Inches(0.12), h)
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _rarr(s, x, cy, w=Inches(0.18)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                           x, cy - Inches(0.05), w, Inches(0.1))
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _cell_bg(cell, hx):
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for sf in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(sf)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), val=hx)


def _cell_border(cell, color_hex='808080', width_pt=0.5):
    """셀 테두리 4면 회색 설정"""
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for border in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for existing in tcPr.findall(qn(border)):
            tcPr.remove(existing)
        ln = etree.SubElement(tcPr, qn(border))
        ln.set('w', str(int(width_pt * 12700)))
        ln.set('cap', 'flat')
        ln.set('cmpd', 'sng')
        ln.set('algn', 'ctr')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr'), val=color_hex)
        etree.SubElement(ln, qn('a:prstDash'), val='solid')
        etree.SubElement(ln, qn('a:round'))


def _tbl(s, l, t, w, h, data, cws=None):
    """테이블 (멀티라인 셀 지원, 진한 테두리 0.5pt, 연한 헤더 배경). data[0]=헤더"""
    rows, cols = len(data), len(data[0])
    sh = s.shapes.add_table(rows, cols, l, t, w, h)
    tbl = sh.table
    if cws:
        for i, cw in enumerate(cws):
            tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ''
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_top = tf.margin_bottom = Inches(0.04)
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.06)
            hdr = (r == 0)
            for i, ln in enumerate(data[r][c].split('\n')):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                _run(p, ln, 9 if not hdr else 10, hdr, CHAR)
                p.alignment = PP_ALIGN.CENTER if hdr else PP_ALIGN.LEFT
                p.space_before = p.space_after = Pt(0)
            _cell_bg(cell, 'D4E2EC' if hdr else ('FFFFFF' if r % 2 == 1 else 'F2F7FB'))
            _cell_border(cell, '808080', 0.5)
    return tbl


# ══════════════════════════════════════════════════
# Page 1: 시스템 개념도
# ══════════════════════════════════════════════════

def p1_concept(prs):
    s = _sl(prs)
    _title(s, '시스템 개념도',
           '회의실예약관리체계의 전체 아키텍처 및 기술 구성 — 6개 계층', 1)

    layers = [
        ('사용자\n계층', [
            ('일반사용자\n회의예약 신청·확인·현황', ICE, NAVY),
            ('예약관리자\n예약 승인·반려·관리', ICE, NAVY),
            ('시스템관리자\n회의실·코드·권한 관리', ICE, NAVY),
        ]),
        ('채널\n계층', [
            ('웹 브라우저\nChrome/Edge 1920x1080', ICE, TEAL),
            ('반응형 모바일 웹\n360px~', ICE, TEAL),
            ('API Gateway\nRESTful 15+ endpoints', ICE, TEAL),
        ]),
        ('프론트\n엔드', [
            ('React 18 + TS 5\nVite 5 · Router 7', LBLUE, SKY),
            ('Ant Design 5.24\nForm·Select·TimePicker', LBLUE, SKY),
            ('Zustand 5 + TQ 5\n상태관리 · 서버캐싱', LBLUE, SKY),
        ]),
        ('백엔드\n계층', [
            ('Spring Boot 3.x\n(추후 개발) Java 17', MINT, GRN),
            ('MSW 2.12 Mock\nFaker.js 9 데이터', MINT, GRN),
            ('REST API 설계\nJSON/HTTPS 15+개', MINT, GRN),
        ]),
        ('데이터\n계층', [
            ('PostgreSQL 15\n(추후) 정규화 설계', CREAM, WARN),
            ('Mock 데이터 30+건\n회의실5+예약20+', CREAM, WARN),
            ('파일 스토리지\n회의실 사진 (추후)', CREAM, WARN),
        ]),
        ('외부\n연계', [
            ('GPKI/SSO 인증\n(추후 연계)', SNOW, ASH),
            ('알림톡/SMS\n예약확인 알림 (추후)', SNOW, ASH),
            ('출입통제시스템\n(추후 연계)', SNOW, ASH),
        ]),
    ]

    y0 = Inches(0.82)
    LH = Inches(0.44)
    GAP = Inches(0.05)
    LBL_W = Inches(0.85)
    LBL_G = Inches(0.05)

    for idx, (lbl, items) in enumerate(layers):
        y = y0 + idx * (LH + GAP)
        _box(s, ML, y, LBL_W, LH, lbl, fill=NAVY, stk=NAVY, sz=7, b=True, fc=WHITE)
        n = len(items)
        avail = CW - LBL_W - LBL_G
        igap = Inches(0.04)
        iw = (avail - igap * (n - 1)) // n
        x = ML + LBL_W + LBL_G
        for txt, fill, stk in items:
            _box(s, x, y, iw, LH, txt, fill=fill, stk=stk, sz=7, fc=CHAR)
            x += iw + igap
        if idx < len(layers) - 1:
            _darr(s, ML + CW // 2, y + LH + Inches(0.01))

    # [C] 서술 테이블
    dy = y0 + 6 * (LH + GAP) + Inches(0.05)
    _tbl(s, ML, dy, CW, Inches(1.45), [
        ['프론트엔드 기술스택', '백엔드 기술스택', '인프라 / 보안 / 연계'],
        ['• React 18.x + TypeScript 5.x\n'
         '• Ant Design 5.24 (Form/Select/TimePicker)\n'
         '• Zustand 5 + TanStack Query 5\n'
         '• React Router 7 + Vite 5\n'
         '• React Hook Form + Zod\n'
         '• 5개 고유 페이지, 21 프로세스',
         '• Spring Boot 3.x (추후 개발)\n'
         '• MSW 2.12 Mock API (현재)\n'
         '• Faker.js 9 한국어 데이터\n'
         '• REST API 15+ endpoint\n'
         '• JWT + Redis 인증 (추후)\n'
         '• MyBatis 3.5 ORM (추후)',
         '• PostgreSQL 15 (추후 개발)\n'
         '• Linux / Tomcat 10.x (추후)\n'
         '• SSL/TLS 1.3, GPKI (추후)\n'
         '• 알림톡/SMS/이메일 (추후)\n'
         '• 출입통제시스템 (추후)\n'
         '• 사용자: 3역할, 21 프로세스'],
    ], cws=[CW // 3] * 3)


# ══════════════════════════════════════════════════
# Page 2: 목표모델 (As-Is → To-Be)
# ══════════════════════════════════════════════════

def p2_target(prs):
    s = _sl(prs)
    _title(s, '목표모델 (As-Is → To-Be)',
           '현행 문제점 진단 및 목표 시스템 기능 구성', 2)

    pairs = [
        ('수기 예약 관리\n전화/방문으로 예약, 중복예약 빈번',
         '온라인 예약 신청\nForm 10필드+부대별 회의실 필터'),
        ('예약현황 파악 불가\n누가 어떤 회의실 쓰는지 불명확',
         '회의현황 실시간 조회\nSearchForm 3필드+DataTable 전체현황'),
        ('승인 프로세스 부재\n선착순 점유, 충돌 시 조정 어려움',
         '승인/반려 워크플로우\n관리자 Popconfirm 승인/반려 처리'),
        ('회의실 정보 분산\n장비·시간대·사진 별도 관리',
         '회의실 통합 관리\nTabs 4탭(정보/시간대/장비/사진)'),
        ('내 예약 확인 불편\n별도 기록 없음, 기억에 의존',
         '내예약확인 페이지\nDataTable+상세Modal+수정/삭제'),
    ]

    y0 = Inches(1.0)
    RH = Inches(0.52)
    RG = Inches(0.03)
    AW = Inches(3.5)
    ARR_W = Inches(0.22)
    AX = ML
    ARX = ML + AW + Inches(0.15)
    TX = ARX + ARR_W + Inches(0.15)

    _box(s, AX, y0 - Inches(0.06), AW, Inches(0.18),
         '현행 시스템 (As-Is)', fill=ERR, stk=ERR, sz=8, b=True, fc=WHITE)
    _box(s, TX, y0 - Inches(0.06), AW, Inches(0.18),
         '목표 시스템 (To-Be)', fill=GRN, stk=GRN, sz=8, b=True, fc=WHITE)

    for i, (asis, tobe) in enumerate(pairs):
        y = y0 + i * (RH + RG)
        _box(s, AX, y, AW, RH, asis, fill=PEACH, stk=ERR, sz=7, fc=CHAR)
        _rarr(s, ARX, y + RH // 2, ARR_W)
        _box(s, TX, y, AW, RH, tobe, fill=MINT, stk=GRN, sz=7, fc=CHAR)

    # 우측 기능 요약
    FX = TX + AW + Inches(0.12)
    FW = ML + CW - FX
    _box(s, FX, y0, FW, 5 * RH + 4 * RG,
         '핵심 기능\n─────────\n'
         '회의예약신청  1\n내예약확인     1\n회의현황       1\n'
         '예약관리       1\n회의실관리     1\n'
         '─────────\n총 5개 고유 메뉴\n21 프로세스\n'
         '─────────\n게시판 2\n관리자 5+',
         fill=ICE, stk=TEAL, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)

    # [C] 서술 테이블
    dy = y0 + 5 * (RH + RG) + Inches(0.12)
    _tbl(s, ML, dy, CW, Inches(1.25), [
        ['개선항목', '현행 문제점', '목표 개선방향', '기대 효과', '정량 목표'],
        ['예약신청', '전화/방문 수기 예약', 'Form 10필드 온라인 신청', '예약 편의성 향상', '100% 온라인'],
        ['현황조회', '현황 파악 불가', 'SearchForm+DataTable 실시간', '회의실 현황 즉시 파악', '실시간 조회'],
        ['승인관리', '선착순 점유, 충돌', '관리자 승인/반려 워크플로우', '중복예약 방지', '충돌 0건'],
        ['회의실관리', '정보 분산 관리', 'Tabs 4탭 통합 (정보/시간/장비/사진)', '회의실 정보 일원화', '통합 100%'],
        ['내예약확인', '기록 없음', 'DataTable+상세/수정/삭제', '개인 예약 이력 관리', '이력 100%'],
    ], cws=[int(CW * r) for r in [0.10, 0.24, 0.24, 0.22, 0.10]])


# ══════════════════════════════════════════════════
# Page 3: 유스케이스
# ══════════════════════════════════════════════════

def p3_usecase(prs):
    s = _sl(prs)
    _title(s, '유스케이스 다이어그램',
           '3개 액터 x 8개 유스케이스, 21개 단위 프로세스', 3)

    y0 = Inches(0.82)

    # 좌측: 액터
    actors = [
        '일반사용자\n예약신청·확인·현황조회\nROLE_USER',
        '예약관리자\n예약 승인·반려\nROLE_MANAGER',
        '시스템관리자\n회의실·코드·권한 관리\nROLE_ADMIN',
    ]
    AX = ML
    AW = Inches(1.15)
    AH = Inches(0.55)
    AG = Inches(0.15)
    for i, atxt in enumerate(actors):
        _box(s, AX, y0 + i * (AH + AG), AW, AH, atxt,
             fill=ICE, stk=NAVY, sz=7, fc=CHAR)

    # 중앙: 시스템 바운더리
    BX = AX + AW + Inches(0.2)
    BW = Inches(6.2)
    BH = 3 * AH + 2 * AG + Inches(0.3)
    bnd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, BX, y0 - Inches(0.05), BW, BH)
    bnd.fill.solid()
    bnd.fill.fore_color.rgb = RGBColor(0xFB, 0xFB, 0xFB)
    bnd.line.color.rgb = NAVY
    bnd.line.width = Pt(2)
    _tb(s, BX + Inches(0.1), y0 - Inches(0.02), Inches(5), Inches(0.18),
        '회의실예약관리체계 (SYS16) — 바운더리', 8, True, NAVY)

    # UC 그룹
    uc_groups = [
        ('회의예약신청 (3)', 'UC-ROOM-001~003\n부대필터+회의실선택+Form(10필드)', ICE, SKY),
        ('내예약확인 (3)', 'UC-ROOM-004~006\n목록조회+상세Modal+수정/삭제', ICE, SKY),
        ('회의현황 (2)', 'UC-ROOM-007~008\nSearchForm(3필드)+DataTable 전체현황', ICE, SKY),
        ('회의예약관리 (3)', 'UC-ROOM-009~011\n승인/반려+엑셀다운로드+프린트', LBLUE, TEAL),
        ('회의실관리 (5)', 'UC-ROOM-012~016\nCRUD+시간대설정+장비관리+사진관리', LBLUE, TEAL),
        ('게시판+관리자 (5+)', 'UC-SYS 공지/질의/코드관리\n권한관리·시스템 관리', LBLUE, TEAL),
    ]
    UY0 = y0 + Inches(0.22)
    UW = Inches(2.95)
    UH = Inches(0.48)
    UG = Inches(0.03)
    for i, (grp, desc, fill, stk) in enumerate(uc_groups):
        r, c = divmod(i, 2)
        ux = BX + Inches(0.08) + c * (UW + Inches(0.1))
        uy = UY0 + r * (UH + UG)
        sh = _box(s, ux, uy, UW, UH, f'{grp}\n{desc}',
                  fill=fill, stk=stk, sz=6, b=False, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True

    # 우측: 외부 시스템
    EX = BX + BW + Inches(0.12)
    EW = ML + CW - EX
    _box(s, EX, y0 + Inches(0.2), EW, Inches(0.35),
         '출입통제시스템\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)
    _box(s, EX, y0 + Inches(0.65), EW, Inches(0.35),
         'GPKI/SSO\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)

    # [C] 서술 테이블
    dy = Inches(3.2)
    _tbl(s, ML, dy, CW, Inches(2.1), [
        ['UC-ID', '유스케이스명', '주 액터', '기능 설명 (구체적)', '우선'],
        ['UC-ROOM-001', '회의예약신청', '사용자',
         'Form 10필드 — 관리부대(Select)+회의실(Select 필터)+회의명+예약일(DatePicker)+시간(TimePicker x2)+등급+주관부서+목적+인원+참석자', '상'],
        ['UC-ROOM-004', '내예약확인', '사용자',
         'DataTable(8컬럼: 번호/회의실/예약자(군번/계급/성명)/예약일/시간/목적/상태/신청일) + 상세Modal(Descriptions) + 수정/삭제', '상'],
        ['UC-ROOM-007', '회의현황', '사용자',
         'SearchForm(3필드: 관리부대/회의실/회의일자) + DataTable(7컬럼) + StatusBadge(대기/승인/반려)', '상'],
        ['UC-ROOM-009', '회의예약관리', '관리자',
         'DataTable + Popconfirm 승인/반려 + militaryPersonColumn(군번/계급/성명) + 엑셀CSV다운로드 + 프린트', '상'],
        ['UC-ROOM-012', '회의실관리', '관리자',
         'Row(8:16) 좌:목록 우:Tabs 4탭(기본정보Descriptions/시간대Table+Switch/장비Table CRUD/사진Upload+Image)', '상'],
        ['UC-ROOM-017', '공지사항/질의응답', '사용자',
         'SimpleBoardPage 재사용: 공지사항(/sys16/board/1), 질의응답(/sys16/board/2)', '하'],
        ['UC-ROOM-018', '코드/권한관리', '관리자',
         '공통 CodeMgmtIndex, AuthGroupIndex 페이지 재사용', '하'],
    ], cws=[int(CW * r) for r in [0.10, 0.12, 0.06, 0.66, 0.05]])


# ══════════════════════════════════════════════════
# Page 4: 스토리보드 (화면흐름도)
# ══════════════════════════════════════════════════

def p4_storyboard(prs):
    s = _sl(prs)
    _title(s, '스토리보드 (화면흐름도)',
           '총 5개 고유 화면 + 게시판/관리자, 핵심 네비게이션 흐름', 4)

    y0 = Inches(0.82)

    # 로그인 → 대시보드
    _box(s, ML, y0, Inches(1.3), Inches(0.35),
         '로그인\n/login', fill=WARN, stk=WARN, sz=8, b=True, fc=WHITE)
    _rarr(s, ML + Inches(1.35), y0 + Inches(0.12), Inches(0.25))
    _box(s, ML + Inches(1.65), y0, Inches(2.0), Inches(0.35),
         '대시보드 / 홈\n/sys16', fill=NAVY, stk=NAVY, sz=8, b=True, fc=WHITE)

    _darr(s, ML + Inches(2.65), y0 + Inches(0.37))

    # 메뉴 그룹
    groups = [
        ('1. 회의예약신청', 'Form 10필드 예약', '/sys16/1/2', '1', ICE, SKY),
        ('2. 내예약확인', 'DataTable+수정/삭제', '/sys16/1/3', '1', ICE, SKY),
        ('3. 회의현황', 'SearchForm+DataTable', '/sys16/1/4', '1', ICE, SKY),
        ('4. 회의예약관리', '승인/반려+엑셀/프린트', '/sys16/1/5', '1', LBLUE, TEAL),
        ('5. 회의실관리', 'Tabs 4탭 통합관리', '/sys16/1/6', '1', LBLUE, TEAL),
        ('6. 게시판+관리자', '공지·질의·코드·권한', '/sys16/board,admin', '4+', LBLUE, TEAL),
    ]

    GY0 = y0 + Inches(0.55)
    GW = Inches(4.5)
    GH = Inches(0.58)
    GG = Inches(0.03)
    GCOL_GAP = Inches(0.15)

    for i, (name, screens, path, cnt, fill, stk) in enumerate(groups):
        r, c = divmod(i, 2)
        gx = ML + c * (GW + GCOL_GAP)
        gy = GY0 + r * (GH + GG)
        sh = _box(s, gx, gy, GW, GH,
                  f'{name} ({cnt}개)\n{screens}\n{path}',
                  fill=fill, stk=stk, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        pars = list(sh.text_frame.paragraphs)
        if len(pars) >= 3:
            for r_obj in pars[-1].runs:
                r_obj.font.color.rgb = ASH
                r_obj.font.size = Pt(8)

    # 핵심 화면 상세 (3개)
    DY = GY0 + 3 * (GH + GG) + Inches(0.08)
    det_boxes = [
        ('회의예약신청\n/sys16/1/2',
         'Card Layout (maxWidth 600)\nForm 10필드 vertical\n관리부대 Select -> 회의실 필터\n'
         'DatePicker + TimePicker x2\n회의등급(일반/비밀/대외비)\n참석인원 InputNumber',
         ICE, SKY),
        ('회의실관리\n/sys16/1/6',
         'Row(8:16) 좌:DataTable 목록\n우:Tabs 4탭 상세\n'
         '기본정보 Descriptions\n시간대 Table+Switch+TimePicker\n'
         '장비 Table CRUD\n사진 Upload+Image Gallery',
         MINT, GRN),
        ('회의예약관리\n/sys16/1/5',
         'DataTable 예약 목록\nmilitaryPersonColumn\n'
         'Popconfirm 승인/반려\nStatusBadge(대기/승인/반려)\n'
         'CSV 엑셀 다운로드\n프린트 기능',
         CREAM, WARN),
    ]
    DW = Inches(2.95)
    DH = Inches(1.1)
    DG = Inches(0.1)
    for i, (ttl, desc, fill, stk) in enumerate(det_boxes):
        dx = ML + i * (DW + DG)
        sh = _box(s, dx, DY, DW, DH,
                  f'{ttl}\n──────────\n{desc}',
                  fill=fill, stk=stk, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        sh.text_frame.paragraphs[1].runs[0].font.bold = True

    # [C] 서술 테이블 2개 좌우 배치
    dy = DY + DH + Inches(0.08)

    LW = int(CW * 0.48)
    _tbl(s, ML, dy, LW, Inches(0.95), [
        ['영역', '화면 수', 'URL 패턴', '주요 컴포넌트'],
        ['회의예약신청', '1', '/sys16/1/2', 'Form+Select+DatePicker+TimePicker'],
        ['내예약확인', '1', '/sys16/1/3', 'DataTable+Descriptions+Modal'],
        ['회의현황', '1', '/sys16/1/4', 'SearchForm+DataTable+StatusBadge'],
        ['예약관리', '1', '/sys16/1/5', 'DataTable+Popconfirm+CSV'],
        ['회의실관리+기타', '1+4', '/sys16/1/6,board,admin', 'Tabs+CrudForm+Upload+Admin'],
    ], cws=[int(LW * r) for r in [0.20, 0.10, 0.28, 0.42]])

    RX = ML + LW + Inches(0.15)
    RW = CW - LW - Inches(0.15)
    _tbl(s, RX, dy, RW, Inches(0.95), [
        ['핵심 흐름 시나리오'],
        ['시나리오 1: 회의예약 신청 흐름\n'
         '1. 로그인 -> 대시보드\n'
         '2. [회의예약신청] 메뉴 선택\n'
         '3. 관리부대 Select -> 회의실 필터링\n'
         '4. 10개 필드 입력 (회의명/일자/시간/등급 등)\n'
         '5. [예약 신청] -> 성공 메시지\n'
         '6. 내예약확인 페이지로 자동 이동'],
        ['시나리오 2: 예약 승인 흐름\n'
         '1. [회의예약관리] 메뉴 선택\n'
         '2. DataTable에서 대기 예약 확인\n'
         '3. Popconfirm [승인] or [반려] 클릭\n'
         '4. 상태 변경 -> StatusBadge 갱신'],
    ])


# ══════════════════════════════════════════════════
# Page 5: 플로우차트 (업무처리흐름)
# ══════════════════════════════════════════════════

def p5_flowchart(prs):
    s = _sl(prs)
    _title(s, '업무처리 플로우차트',
           '핵심 프로세스: 예약신청 -> 승인 -> 현황 — 4개 레인, 12단계', 5)

    y0 = Inches(0.82)

    # 4개 스윔레인
    lanes = [
        ('사용자 (신청자/관리자)', ICE, NAVY),
        ('프론트엔드 (React 18)', LBLUE, SKY),
        ('백엔드 (API/Mock)', MINT, GRN),
        ('DB / 외부', CREAM, WARN),
    ]
    LN_W = CW // 4
    LN_H = Inches(2.75)
    LN_HDR = Inches(0.22)

    for i, (name, fill, stk) in enumerate(lanes):
        lx = ML + LN_W * i
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = fill
        bg.line.color.rgb = stk
        bg.line.width = Pt(1)
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_HDR)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = stk
        hdr.line.fill.background()
        tf = hdr.text_frame
        _run(tf.paragraphs[0], name, 7, True, WHITE)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    steps = [
        (0, 0.30, '1. 예약 메뉴 접근\n회의예약신청 선택'),
        (1, 0.30, '2. 예약 폼 렌더\nForm 10필드 표시'),
        (0, 0.70, '3. 부대 선택+필드 입력\n회의실 필터 + 시간 선택'),
        (1, 0.70, '4. POST /reservations\n예약 데이터 전송'),
        (2, 0.70, '5. 예약 데이터 검증\n중복/시간 충돌 체크'),
        (3, 0.70, '6. Reservation INSERT\n예약 저장 (status=pending)'),
        (1, 1.10, '7. 내예약확인 이동\nDataTable 목록 렌더'),
        (0, 1.10, '8. 관리자: 예약관리\n대기 목록 확인'),
        (2, 1.10, '9. PATCH approve/reject\n상태 변경 처리'),
        (0, 1.50, '10. 회의현황 조회\nSearchForm 3필드 입력'),
        (1, 1.50, '11. GET /status 호출\n필터링된 목록 렌더'),
        (2, 1.90, '12. 현황 조회 API\n부대/회의실/일자 필터'),
    ]

    SBW = LN_W - Inches(0.12)
    SBH = Inches(0.40)

    for lane, y_off, txt in steps:
        sx = ML + LN_W * lane + Inches(0.06)
        sy = y0 + LN_HDR + Inches(y_off)
        _box(s, sx, sy, SBW, SBH, txt,
             fill=WHITE, stk=ASH, sw=1, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)

    # 판단 분기
    fx = ML + LN_W * 1 + Inches(0.06)
    fy = y0 + LN_HDR + Inches(1.35)
    _box(s, fx, fy, SBW, Inches(0.12),
         '시간 중복 -> 예약 불가 에러 -> 시간 재선택',
         fill=PEACH, stk=ERR, sw=1, sz=5, fc=ERR, al=PP_ALIGN.CENTER)

    # [C] 서술 테이블 2개 좌우 배치
    dy = y0 + LN_H + Inches(0.1)

    LW = int(CW * 0.52)
    _tbl(s, ML, dy, LW, Inches(1.42), [
        ['단계', '동작', '조건/상세'],
        ['1~2', '메뉴접근 -> 폼 렌더', 'Form 10필드: 관리부대/회의실/회의명/일자/시작·종료시간/등급/주관부서/목적/인원/참석자'],
        ['3~4', '입력 -> API 호출', '관리부대 Select 변경 -> 회의실 목록 필터링, POST /reservations'],
        ['5~6', '서버검증 -> DB저장', '중복예약 체크, status=pending, createdAt 자동 기록'],
        ['7~8', '내예약확인 -> 관리자', '내예약: DataTable+수정/삭제(pending만), 관리자: 전체 목록'],
        ['9', '승인/반려 처리', 'PATCH /reservations/:id/approve or reject + processedAt 기록'],
        ['10~12', '회의현황 조회', 'SearchForm(관리부대/회의실/일자) -> GET /reservations/status 필터'],
    ], cws=[int(LW * r) for r in [0.08, 0.28, 0.64]])

    RX = ML + LW + Inches(0.12)
    RW = CW - LW - Inches(0.12)
    _tbl(s, RX, dy, RW, Inches(1.42), [
        ['업무 규칙 / 예외 처리'],
        ['[업무 규칙]\n'
         'BR-ROOM-001: 관리부대 3개 (해병대사령부/1사단/2사단)\n'
         'BR-ROOM-002: 회의등급 3종 (일반/비밀/대외비)\n'
         'BR-ROOM-003: 시간선택 30분 단위 (minuteStep=30)\n'
         'BR-ROOM-004: 회의실 5개, 장비 5종 (빔프로젝터/화이트보드/마이크/스크린/TV)\n'
         'BR-ROOM-005: 요일별 운영시간 설정 (월~금 09:00~18:00 기본)\n'
         'BR-ROOM-006: pending 상태에서만 수정/삭제 가능\n'
         'BR-ROOM-007: 사진 최대 5장 Upload, 장비 CRUD'],
        ['[예외 처리]\n'
         'E-001: 세션만료(401) -> 로그인 리다이렉트\n'
         'E-002: 권한없음(403) -> 접근 제한 알림\n'
         'E-003: 시간 중복 -> 에러 메시지 표시\n'
         'E-004: 서버오류(500) -> 에러 메시지+재시도'],
    ])


# ══════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════

def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    print('[1/5] 시스템 개념도...')
    p1_concept(prs)
    print('[2/5] 목표모델...')
    p2_target(prs)
    print('[3/5] 유스케이스...')
    p3_usecase(prs)
    print('[4/5] 스토리보드...')
    p4_storyboard(prs)
    print('[5/5] 플로우차트...')
    p5_flowchart(prs)

    out_path = os.path.join(OUT, 'SYS16_회의실예약관리체계_분석설계_v5.pptx')
    prs.save(out_path)
    print(f'\n완료! 총 {len(prs.slides)}장 슬라이드')
    print(f'파일: {out_path}')


if __name__ == '__main__':
    main()
