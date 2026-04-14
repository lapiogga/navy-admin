#!/usr/bin/env python3
"""
SYS17 검열결과관리체계 기능명세서 PPTX v7 (portrait, 4페이지)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
system_level_spec_guide_v7.md 기반 핵심 5페이지:
  P1. 시스템 개념도      P2. 목표모델 (As-Is → To-Be)
  P3. 유스케이스          P4. 스토리보드 (화면흐름도)
  P5. 플로우차트 (업무처리흐름)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ═══ 색상 팔레트 ═══
NAVY  = RGBColor(0x6B, 0x8E, 0xA8)
TEAL  = RGBColor(0x8E, 0xC5, 0xC0)
SKY   = RGBColor(0xA8, 0xC8, 0xE8)
MINT  = RGBColor(0xDC, 0xF5, 0xE4)
PEACH = RGBColor(0xFC, 0xD4, 0xD0)
CREAM = RGBColor(0xFD, 0xEE, 0xC2)
ICE   = RGBColor(0xE4, 0xEE, 0xF8)
SNOW  = RGBColor(0xF7, 0xF9, 0xFC)
CHAR  = RGBColor(0x1A, 0x1A, 0x1A)
ASH   = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRN   = RGBColor(0xA4, 0xD4, 0xB4)
WARN  = RGBColor(0xF5, 0xCC, 0xA0)
ERR   = RGBColor(0xF0, 0xB2, 0xAE)
LBLUE = RGBColor(0xDC, 0xEA, 0xFE)
LSUB  = RGBColor(0xC5, 0xDC, 0xEC)
HBAR  = RGBColor(0xD4, 0xE2, 0xEC)  # 제목바 연한 배경
THBG  = RGBColor(0xD4, 0xE2, 0xEC)  # 표 헤더 연한 배경
BORD  = RGBColor(0x4A, 0x6B, 0x8C)  # 박스 테두리 (진한)
TBLB  = RGBColor(0x80, 0x80, 0x80)  # 표 테두리 (회색)

# ═══ 치수 ═══
SW = Inches(7.5)
SH = Inches(10)
ML = Inches(0.27)      # 좌우 마진
CW = Inches(6.96)      # 콘텐츠 폭 = 7.5 - 0.27*2
FONT = 'Freesentation 5 Medium'
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'ppt')


# ══════════════════════════════════════════════════
# 헬퍼 함수
# ══════════════════════════════════════════════════

def _sl(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = SNOW
    return s


def _run(p, txt, sz=9, b=False, c=CHAR):
    """단락에 서식 런 추가 (최소 9pt 보장)"""
    sz = max(sz, 10)
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = b
    r.font.color.rgb = c
    r.font.name = FONT
    return r


def _tb(s, l, t, w, h, txt, sz=9, b=False, c=CHAR, al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    _run(tb.text_frame.paragraphs[0], txt, sz, b, c)
    tb.text_frame.paragraphs[0].alignment = al
    return tb


def _title(s, title, sub, pg):
    """[A] 제목 영역 (portrait) — 연한 바 + 진한 테두리 + 검정 글자"""
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = HBAR
    bar.line.color.rgb = BORD
    bar.line.width = Pt(0.5)
    _tb(s, ML, Inches(0.06), Inches(5.5), Inches(0.36), title, 20, True, CHAR)
    if sub:
        _tb(s, ML, Inches(0.42), Inches(6.5), Inches(0.26), sub, 10, False, CHAR)


def _box(s, l, t, w, h, txt, fill=ICE, stk=BORD, sw=0.5,
         sz=9, b=False, fc=CHAR, al=PP_ALIGN.CENTER):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = stk
    sh.line.width = Pt(sw)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.margin_left = tf.margin_right = Inches(0.03)
    for i, ln in enumerate(txt.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p, ln, sz, b, fc)
        p.alignment = al
        p.space_before = p.space_after = Pt(0)
    bp = sh._element.txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', 'ctr')
    return sh


def _darr(s, cx, y, h=Inches(0.08)):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                           cx - Inches(0.06), y, Inches(0.12), h)
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _rarr(s, x, cy, w=Inches(0.18)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                           x, cy - Inches(0.05), w, Inches(0.1))
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _cell_bg(cell, hx):
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for sf in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(sf)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), val=hx)


def _cell_border(cell, color_hex='808080', width_pt=0.5):
    """셀 테두리 4면 회색 설정"""
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for border in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for existing in tcPr.findall(qn(border)):
            tcPr.remove(existing)
        ln = etree.SubElement(tcPr, qn(border))
        ln.set('w', str(int(width_pt * 12700)))
        ln.set('cap', 'flat')
        ln.set('cmpd', 'sng')
        ln.set('algn', 'ctr')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr'), val=color_hex)
        etree.SubElement(ln, qn('a:prstDash'), val='solid')
        etree.SubElement(ln, qn('a:round'))


def _tbl(s, l, t, w, h, data, cws=None):
    """테이블 (멀티라인 셀 지원, 진한 테두리 0.5pt, 연한 헤더 배경). data[0]=헤더"""
    rows, cols = len(data), len(data[0])
    sh = s.shapes.add_table(rows, cols, l, t, w, h)
    tbl = sh.table
    if cws:
        for i, cw in enumerate(cws):
            tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ''
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_top = tf.margin_bottom = Inches(0.04)
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.06)
            hdr = (r == 0)
            for i, ln in enumerate(data[r][c].split('\n')):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                _run(p, ln, 9 if not hdr else 10, hdr, CHAR)
                p.alignment = PP_ALIGN.CENTER if hdr else PP_ALIGN.LEFT
                p.space_before = p.space_after = Pt(0)
            _cell_bg(cell, 'D4E2EC' if hdr else ('FFFFFF' if r % 2 == 1 else 'F2F7FB'))
            _cell_border(cell, '000000', 0.5)
    return tbl




# ══════════════════════════════════════════════════
# Page 2: 목표모델 (As-Is → To-Be)
# ══════════════════════════════════════════════════

def p2_target(prs):
    s = _sl(prs)
    _title(s, '목표모델 (As-Is → To-Be)',
           '현행 문제점 진단 및 목표 시스템 기능 구성', 1)

    pairs = [
        ('검열결과 수기관리\n엑셀/공문 기반 과제 추적, 현황 파악 어려움',
         '통합 검열결과 관리\n계획→과제→결과→접수 전 과정 온라인 관리'),
        ('부대별 분산 관리\n각 부대 독립적 조치현황 보고, 취합 지연',
         '실시간 추진현황\nBar 차트+Progress 부대별 추진율 시각화'),
        ('결재 프로세스 부재\n조치결과 접수·반송 이력 추적 불가',
         '3단계 결재 체계\nSteps 시각화 + 접수/반송 프로세스'),
        ('검열부대 선정 비효율\n조직도 기반 부대 지정 체계 없음',
         'Tree 기반 부대 지정\n조직도 체크트리 + 선택 테이블'),
        ('과제처리 이력 미관리\n담당자·상태변경·조치내용 추적 불가',
         'Timeline 이력 관리\n과제별 처리이력 시간순 표시'),
    ]

    y0 = Inches(1.0)
    RH = Inches(0.78)
    RG = Inches(0.06)
    AW = Inches(2.6)
    ARR_W = Inches(0.22)
    AX = ML
    ARX = ML + AW + Inches(0.15)
    TX = ARX + ARR_W + Inches(0.15)

    _box(s, AX, y0 - Inches(0.06), AW, Inches(0.18),
         '현행 시스템 (As-Is)', fill=ERR, stk=ERR, sz=8, b=True, fc=WHITE)
    _box(s, TX, y0 - Inches(0.06), AW, Inches(0.18),
         '목표 시스템 (To-Be)', fill=GRN, stk=GRN, sz=8, b=True, fc=WHITE)

    for i, (asis, tobe) in enumerate(pairs):
        y = y0 + i * (RH + RG)
        _box(s, AX, y, AW, RH, asis, fill=PEACH, stk=ERR, sz=7, fc=CHAR)
        _rarr(s, ARX, y + RH // 2, ARR_W)
        _box(s, TX, y, AW, RH, tobe, fill=MINT, stk=GRN, sz=7, fc=CHAR)

    # 우측 기능 요약
    FX = TX + AW + Inches(0.12)
    FW = ML + CW - FX
    _box(s, FX, y0, FW, 5 * RH + 4 * RG,
         '핵심 기능\n─────────\n'
         '검열결과 관리  5\n게시판           2\n관리자           4\n'
         '데이터           2\n'
         '─────────\n총 13개 메뉴\n25 프로세스\n'
         '─────────\n부대 Tree 선택\n3단계 결재선',
         fill=ICE, stk=TEAL, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)

    # [C] 서술 테이블
    dy = y0 + 5 * (RH + RG) + Inches(0.12)
    _tbl(s, ML, dy, CW, Inches(1.25), [
        ['개선항목', '현행 문제점', '목표 개선방향', '기대 효과', '정량 목표'],
        ['검열계획', '수기/공문 기반 관리', 'CRUD + SearchForm(4필드)', '계획 수립 효율화', '등록 3분'],
        ['조치과제', '엑셀 취합, 추적 불가', 'Tabs 2탭 + DataTable(11컬럼)', '과제 실시간 추적', '이력 100%'],
        ['결재체계', '접수·반송 이력 없음', 'Steps 3단계 + 승인/반려', '결재 투명성 확보', '전자결재 100%'],
        ['추진현황', '수기 취합, 현황 지연', 'Bar 차트+Progress+카드 통계', '실시간 현황 파악', '자동 집계'],
        ['부대지정', '조직도 없이 수기 선정', 'Tree 체크박스 + 테이블', '부대 선정 효율화', '트리 선택'],
    ], cws=[int(CW * r) for r in [0.10, 0.24, 0.24, 0.22, 0.10]])


# ══════════════════════════════════════════════════
# Page 3: 유스케이스
# ══════════════════════════════════════════════════

def p3_usecase(prs):
    s = _sl(prs)
    _title(s, '유스케이스 다이어그램',
           '3개 액터 x 14개 유스케이스, 25개 단위 프로세스', 2)

    y0 = Inches(0.82)

    # 좌측: 액터
    actors = [
        '검열관리자\n검열계획·부대지정·결재\nROLE_INSPECTOR',
        '조치부서 담당자\n조치결과 입력·보고\nROLE_HANDLER',
        '시스템관리자\nROLE_ADMIN',
    ]
    AX = ML
    AW = Inches(1.0)
    AH = Inches(0.55)
    AG = Inches(0.15)
    for i, atxt in enumerate(actors):
        _box(s, AX, y0 + i * (AH + AG), AW, AH, atxt,
             fill=ICE, stk=NAVY, sz=7, fc=CHAR)

    # 중앙: 시스템 바운더리
    BX = AX + AW + Inches(0.2)
    BW = Inches(4.3)
    BH = 3 * AH + 2 * AG + Inches(0.3)
    bnd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, BX, y0 - Inches(0.05), BW, BH)
    bnd.fill.solid()
    bnd.fill.fore_color.rgb = RGBColor(0xFB, 0xFB, 0xFB)
    bnd.line.color.rgb = NAVY
    bnd.line.width = Pt(2)
    _tb(s, BX + Inches(0.1), y0 - Inches(0.02), Inches(5), Inches(0.18),
        '검열결과관리체계 (SYS17) — 바운더리', 8, True, NAVY)

    # UC 그룹
    uc_groups = [
        ('검열부대 지정 (2)', 'UC-INSP-001~002\n조직도 Tree 체크박스 선택+저장', ICE, SKY),
        ('검열계획 (4)', 'UC-INSP-003~006\nCRUD+SearchForm(4필드)+DataTable(8컬럼)', ICE, SKY),
        ('검열결과 (6)', 'UC-INSP-007~012\n조치과제 CRUD+결과입력+Timeline 이력', ICE, SKY),
        ('결재 (3)', 'UC-INSP-013~015\n접수대기+승인/반송+Steps 3단계', ICE, SKY),
        ('추진현황 (2)', 'UC-INSP-016~017\nBar 차트+Progress+부대별 세부', LBLUE, TEAL),
        ('관리자 (6)', 'UC-SYS 공통코드·부대관리\n사용자권한·접속로그·게시판', LBLUE, TEAL),
    ]
    UY0 = y0 + Inches(0.22)
    UW = Inches(2.0)
    UH = Inches(0.70)
    UG = Inches(0.06)
    for i, (grp, desc, fill, stk) in enumerate(uc_groups):
        r, c = divmod(i, 2)
        ux = BX + Inches(0.08) + c * (UW + Inches(0.1))
        uy = UY0 + r * (UH + UG)
        sh = _box(s, ux, uy, UW, UH, f'{grp}\n{desc}',
                  fill=fill, stk=stk, sz=6, b=False, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True

    # 우측: 외부 시스템
    EX = BX + BW + Inches(0.12)
    EW = ML + CW - EX
    _box(s, EX, y0 + Inches(0.2), EW, Inches(0.35),
         '해병대 조직도\n(부대 Tree 연동)', fill=SNOW, stk=ASH, sz=7, fc=ASH)
    _box(s, EX, y0 + Inches(0.65), EW, Inches(0.35),
         'GPKI/SSO\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)

    # [C] 서술 테이블
    dy = Inches(3.2)
    _tbl(s, ML, dy, CW, Inches(2.1), [
        ['UC-ID', '유스케이스명', '주 액터', '기능 설명 (구체적)', '우선'],
        ['UC-INSP-001', '검열부대 지정', '관리자',
         'Row(12:12) — 좌: Tree 조직도(checkable, defaultExpandAll), 우: Table(선택부대 목록)+저장', '상'],
        ['UC-INSP-003', '검열계획 CRUD', '관리자',
         'SearchForm(4필드: 연도/계획명/부대/기간)+DataTable(8컬럼)+CrudForm(7필드) Modal', '상'],
        ['UC-INSP-007', '조치과제 목록', '담당자',
         'Tabs 2탭(조치과제/조치결과)+SearchForm(5필드)+DataTable(11컬럼)+군번/계급/성명', '상'],
        ['UC-INSP-010', '조치결과 입력', '담당자',
         '상세Modal(14필드)+조치결과CrudForm(4필드: 상태/문제점/결과/첨부)+Timeline 이력', '상'],
        ['UC-INSP-013', '결재 접수대기', '관리자',
         'Tabs 2탭(접수대기/완료)+SearchForm(5필드)+DataTable(6컬럼)+Steps(3단계)', '상'],
        ['UC-INSP-016', '추진현황 종합', '관리자',
         'Tabs 2탭(종합/세부)+카드4개(총/완료/진행/미조치)+Bar 스택차트+Progress', '중'],
        ['UC-INSP-018', '게시판', '사용자',
         'SimpleBoardPage 재사용: 공지사항(/sys17/board/1), 질의응답(/sys17/board/2)', '하'],
    ], cws=[int(CW * r) for r in [0.10, 0.12, 0.06, 0.66, 0.05]])


# ══════════════════════════════════════════════════
# Page 4: 스토리보드 (화면흐름도)
# ══════════════════════════════════════════════════

def p4_storyboard(prs):
    s = _sl(prs)
    _title(s, '스토리보드 (화면흐름도)',
           '총 14개 화면, 3개 업무영역, 핵심 네비게이션 흐름', 3)

    y0 = Inches(0.82)

    # 로그인 → 대시보드
    _box(s, ML, y0, Inches(1.1), Inches(0.35),
         '로그인\n/login', fill=WARN, stk=WARN, sz=8, b=True, fc=WHITE)
    _rarr(s, ML + Inches(1.35), y0 + Inches(0.12), Inches(0.25))
    _box(s, ML + Inches(1.65), y0, Inches(2.0), Inches(0.35),
         '대시보드 / 홈\n/sys17', fill=NAVY, stk=NAVY, sz=8, b=True, fc=WHITE)

    _darr(s, ML + Inches(2.65), y0 + Inches(0.37))

    # 메뉴 그룹
    groups = [
        ('1. 검열결과 관리', '부대지정·계획·결과·결재·현황', '/sys17/1/*', '5', ICE, SKY),
        ('2. 게시판', '공지사항 · 질의응답', '/sys17/board/*', '2', ICE, SKY),
        ('3. 관리자', '공통코드·부대관리·권한·로그', '/sys17/admin/*', '4', LBLUE, TEAL),
        ('4. 데이터', '검열계획 정보 · 검열결과 정보', '/sys17/3/*', '2', LBLUE, TEAL),
    ]

    GY0 = y0 + Inches(0.55)
    GW = Inches(3.2)
    GH = Inches(0.90)
    GG = Inches(0.06)
    GCOL_GAP = Inches(0.15)

    for i, (name, screens, path, cnt, fill, stk) in enumerate(groups):
        r, c = divmod(i, 2)
        gx = ML + c * (GW + GCOL_GAP)
        gy = GY0 + r * (GH + GG)
        sh = _box(s, gx, gy, GW, GH,
                  f'{name} ({cnt}개)\n{screens}\n{path}',
                  fill=fill, stk=stk, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        pars = list(sh.text_frame.paragraphs)
        if len(pars) >= 3:
            for r_obj in pars[-1].runs:
                r_obj.font.color.rgb = ASH
                r_obj.font.size = Pt(8)

    # 핵심 화면 상세
    DY = GY0 + 2 * (GH + GG) + Inches(0.08)
    det_boxes = [
        ('검열부대 지정\n/sys17/1/2',
         'Row(12:12) 레이아웃\n좌: Tree 조직도(checkable)\n우: Table(선택부대목록)\n'
         '저장/삭제 버튼\ndefaultExpandAll\n해병대사령부 트리',
         ICE, SKY),
        ('검열결과\n/sys17/1/4',
         'Tabs 2탭 구조\n조치과제: SearchForm(5필드)\nDataTable(11컬럼)\n'
         '조치결과: SearchForm(4필드)\n군번/계급/성명 표시\n상세Modal+Timeline',
         MINT, GRN),
        ('추진현황\n/sys17/1/6',
         'Tabs 2탭 구조\n종합: 카드4개+Bar 차트\nProgress(전체추진율)\n'
         '세부: DataTable(부대별)\n확장행(과제목록)\nBar 스택 차트',
         CREAM, WARN),
    ]
    DW = Inches(2.95)
    DH = Inches(1.1)
    DG = Inches(0.1)
    for i, (ttl, desc, fill, stk) in enumerate(det_boxes):
        dx = ML + i * (DW + DG)
        sh = _box(s, dx, DY, DW, DH,
                  f'{ttl}\n──────────\n{desc}',
                  fill=fill, stk=stk, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        sh.text_frame.paragraphs[1].runs[0].font.bold = True

    # [C] 서술 테이블 2개 좌우 배치
    dy = DY + DH + Inches(0.08)

    LW = int(CW * 0.48)
    _tbl(s, ML, dy, LW, Inches(0.95), [
        ['영역', '화면 수', 'URL 패턴', '주요 컴포넌트'],
        ['검열부대 지정', '1', '/sys17/1/2', 'Tree+Table+Button'],
        ['검열계획', '1', '/sys17/1/3', 'SearchForm+DataTable+CrudForm'],
        ['검열결과', '2', '/sys17/1/4', 'Tabs+DataTable+Modal+Timeline'],
        ['결재', '1', '/sys17/1/5', 'Tabs+DataTable+Steps+Modal'],
        ['추진현황+기타', '3+6', '/sys17/1/6,board,admin', 'Bar+Progress+Expandable'],
    ], cws=[int(LW * r) for r in [0.22, 0.10, 0.28, 0.40]])

    RX = ML + LW + Inches(0.15)
    RW = CW - LW - Inches(0.15)
    _tbl(s, RX, dy, RW, Inches(0.95), [
        ['핵심 흐름 시나리오'],
        ['시나리오 1: 검열계획 수립 → 과제 추적\n'
         '① 로그인 → 대시보드\n'
         '② [검열부대 지정] 조직도 Tree 선택 → 저장\n'
         '③ [검열계획] 연도/부대 선택 → 계획 작성\n'
         '④ [검열결과] 조치과제 등록(11필드)\n'
         '⑤ 조치결과 입력 → Timeline 이력 확인\n'
         '⑥ [결재] 접수대기 → 승인/반송'],
        ['시나리오 2: 추진현황 조회\n'
         '① [추진현황] 종합 탭 → 카드 통계 확인\n'
         '② Bar 차트 부대별 현황 시각화\n'
         '③ 세부 탭 → 부대별 확장행 과제 목록'],
    ])


# ══════════════════════════════════════════════════
# Page 5: 플로우차트 (업무처리흐름)
# ══════════════════════════════════════════════════

def p5_flowchart(prs):
    s = _sl(prs)
    _title(s, '업무처리 플로우차트',
           '핵심 프로세스: 부대지정 → 계획 → 과제등록 → 결재 — 4개 레인, 12단계', 4)

    y0 = Inches(0.82)

    # 4개 스윔레인
    lanes = [
        ('검열관리자', ICE, NAVY),
        ('프론트엔드', LBLUE, SKY),
        ('백엔드', MINT, GRN),
        ('DB / 외부', CREAM, WARN),
    ]
    LGAP = Inches(0.14)  # 레인 사이 간격 (10pt)
    LN_W = (CW - 3 * LGAP) // 4
    LN_H = Inches(3.2)   # 다이어그램 높이 (겹침 방지 축소)
    LN_HDR = Inches(0.22)

    for i, (name, fill, stk) in enumerate(lanes):
        lx = ML + (LN_W + LGAP) * i
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = fill
        bg.line.color.rgb = stk
        bg.line.width = Pt(1)
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_HDR)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = stk
        hdr.line.fill.background()
        tf = hdr.text_frame
        _run(tf.paragraphs[0], name, 7, True, WHITE)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    steps = [
        (0, 0.30, '① 검열부대 지정\n조직도 Tree 체크박스'),
        (1, 0.30, '② Tree 렌더 + 선택\nRow(12:12) 레이아웃'),
        (2, 0.30, '③ PUT /sys17/units\n선택 부대 ID 저장'),
        (0, 0.70, '④ 검열계획 작성\nSearchForm 4필드 검색'),
        (1, 0.70, '⑤ CrudForm Modal\n7필드 입력 (연도/계획명/기간/부대)'),
        (2, 0.70, '⑥ POST /sys17/plans\n신규 검열계획 저장'),
        (3, 0.70, '⑦ InspectionPlan INSERT\n데이터 저장'),
        (0, 1.10, '⑧ 조치과제 등록\n11필드 + 군번/계급/성명'),
        (1, 1.10, '⑨ DataTable + Modal\n11컬럼 + CrudForm 13필드'),
        (2, 1.50, '⑩ POST /sys17/tasks\n조치과제 저장'),
        (0, 1.90, '⑪ 결재 접수/반송\nSteps 3단계 시각화'),
        (2, 1.90, '⑫ PUT approve/reject\n결재상태 변경'),
    ]

    SBW = LN_W - Inches(0.12)
    SBH = Inches(0.56)

    for lane, y_off, txt in steps:
        sx = ML + (LN_W + LGAP) * lane + Inches(0.07)
        sy = y0 + LN_HDR + Inches(y_off)
        _box(s, sx, sy, SBW, SBH, txt,
             fill=WHITE, stk=ASH, sw=1, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)

    # 판단 분기
    fx = ML + (LN_W + LGAP) * 1 + Inches(0.07)
    fy = y0 + LN_HDR + Inches(1.35)
    _box(s, fx, fy, SBW, Inches(0.12),
         '검증 실패 → 필드 에러 표시 → 재입력',
         fill=PEACH, stk=ERR, sw=1, sz=5, fc=ERR, al=PP_ALIGN.CENTER)

    # [C] 서술 테이블 2개 좌우 배치
    dy = y0 + LN_H + Inches(0.1)

    _tbl(s, ML, dy, CW, Inches(1.6), [
        ['단계', '동작', '조건/상세'],
        ['①②③', '부대지정 → Tree → 저장', 'Tree checkable, PUT /sys17/units'],
        ['④⑤⑥', '계획검색 → CrudForm → 등록', 'SearchForm 4필드, CrudForm 7필드'],
        ['⑦', 'DB 저장', 'InspectionPlan INSERT (Mock 10건)'],
        ['⑧⑨', '조치과제 등록 → DataTable', 'Tabs 2탭, 11컬럼+군번/계급/성명'],
        ['⑩', '과제 저장', 'POST /sys17/tasks (Mock 30건)'],
        ['⑪⑫', '결재 접수/반송', 'Steps 3단계 + approve/reject API'],
    ], cws=[int(CW * r) for r in [0.08, 0.28, 0.64]])

    _tbl(s, ML, dy + Inches(1.6) + Inches(0.1), CW, Inches(1.6), [
        ['업무 규칙 / 예외 처리'],
        ['[업무 규칙]\n'
         'BR-INSP-001: 검열분야 5종 (전투준비태세/교육훈련/군수/인사/정보화)\n'
         'BR-INSP-002: 처분종류 5종 (시정/주의/경고/개선/보완)\n'
         'BR-INSP-003: 진행상태 4종 (미조치/진행중/조치완료/접수완료)\n'
         'BR-INSP-004: 결재상태 4종 (접수대기/검토중/접수완료/반송)\n'
         'BR-INSP-005: 결재선 3단계 (담당자→결재자→승인권자)\n'
         'BR-INSP-006: 조직도 Tree 2뎁스 (사령부→사단/단)\n'
         'BR-INSP-007: 담당검열관 군번/계급/성명 필수 입력\n'
         'BR-INSP-008: 과제처리 이력 Timeline (시간순 역순)'],
        ['[예외 처리]\n'
         'E-001: 세션만료(401) → 로그인 리다이렉트\n'
         'E-002: 권한없음(403) → 접근 제한 알림\n'
         'E-003: 검증실패 → 필드별 에러 표시+재입력\n'
         'E-004: 서버오류(500) → 에러 메시지+재시도'],
    ])


# ══════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════

def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    print('[1/4] 목표모델...')
    p2_target(prs)
    print('[2/4] 유스케이스...')
    p3_usecase(prs)
    print('[3/4] 스토리보드...')
    p4_storyboard(prs)
    print('[4/4] 플로우차트...')
    p5_flowchart(prs)

    out_path = os.path.join(OUT, 'SYS17_검열결과관리체계_분석설계_v7.pptx')
    prs.save(out_path)
    print(f'\n완료! 총 {len(prs.slides)}장 슬라이드')
    print(f'파일: {out_path}')


if __name__ == '__main__':
    main()
