#!/usr/bin/env python3
"""
SYS18 직무기술서관리체계 기능명세서 PPTX v5
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
system_level_spec_guide_v5.md 기반 핵심 5페이지:
  P1. 시스템 개념도      P2. 목표모델 (As-Is → To-Be)
  P3. 유스케이스          P4. 스토리보드 (화면흐름도)
  P5. 플로우차트 (업무처리흐름)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ═══ 색상 팔레트 ═══
NAVY  = RGBColor(0x6B, 0x8E, 0xA8)
TEAL  = RGBColor(0x8E, 0xC5, 0xC0)
SKY   = RGBColor(0xA8, 0xC8, 0xE8)
MINT  = RGBColor(0xDC, 0xF5, 0xE4)
PEACH = RGBColor(0xFC, 0xD4, 0xD0)
CREAM = RGBColor(0xFD, 0xEE, 0xC2)
ICE   = RGBColor(0xE4, 0xEE, 0xF8)
SNOW  = RGBColor(0xF7, 0xF9, 0xFC)
CHAR  = RGBColor(0x1A, 0x1A, 0x1A)
ASH   = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRN   = RGBColor(0xA4, 0xD4, 0xB4)
WARN  = RGBColor(0xF5, 0xCC, 0xA0)
ERR   = RGBColor(0xF0, 0xB2, 0xAE)
LBLUE = RGBColor(0xDC, 0xEA, 0xFE)
LSUB  = RGBColor(0xC5, 0xDC, 0xEC)
HBAR  = RGBColor(0xD4, 0xE2, 0xEC)  # 제목바 연한 배경
THBG  = RGBColor(0xD4, 0xE2, 0xEC)  # 표 헤더 연한 배경
BORD  = RGBColor(0x4A, 0x6B, 0x8C)  # 박스 테두리 (진한)
TBLB  = RGBColor(0x80, 0x80, 0x80)  # 표 테두리 (회색)

# ═══ 치수 ═══
SW = Inches(10)
SH = Inches(5.625)
ML = Inches(0.4)
CW = Inches(9.2)
FONT = 'Freesentation 5 Medium'
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'ppt')


# ══════════════════════════════════════════════════
# 헬퍼 함수
# ══════════════════════════════════════════════════

def _sl(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = SNOW
    return s


def _run(p, txt, sz=9, b=False, c=CHAR):
    """단락에 서식 런 추가 (최소 9pt 보장)"""
    sz = max(sz, 9)
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = b
    r.font.color.rgb = c
    r.font.name = FONT
    return r


def _tb(s, l, t, w, h, txt, sz=9, b=False, c=CHAR, al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    _run(tb.text_frame.paragraphs[0], txt, sz, b, c)
    tb.text_frame.paragraphs[0].alignment = al
    return tb


def _title(s, title, sub, pg):
    """[A] 제목 영역 — 연한 바 + 진한 테두리 + 검정 글자"""
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = HBAR
    bar.line.color.rgb = BORD
    bar.line.width = Pt(0.5)
    _tb(s, ML, Inches(0.06), Inches(7.5), Inches(0.36), title, 22, True, CHAR)
    if sub:
        _tb(s, ML, Inches(0.42), Inches(8.5), Inches(0.26), sub, 11, False, CHAR)
    _tb(s, Inches(9.0), Inches(0.18), Inches(0.6), Inches(0.3),
        f'{pg}/5', 12, True, CHAR, PP_ALIGN.RIGHT)


def _box(s, l, t, w, h, txt, fill=ICE, stk=BORD, sw=0.5,
         sz=9, b=False, fc=CHAR, al=PP_ALIGN.CENTER):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = stk
    sh.line.width = Pt(sw)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.margin_left = tf.margin_right = Inches(0.03)
    for i, ln in enumerate(txt.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p, ln, sz, b, fc)
        p.alignment = al
        p.space_before = p.space_after = Pt(0)
    bp = sh._element.txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', 'ctr')
    return sh


def _darr(s, cx, y, h=Inches(0.08)):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                           cx - Inches(0.06), y, Inches(0.12), h)
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _rarr(s, x, cy, w=Inches(0.18)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                           x, cy - Inches(0.05), w, Inches(0.1))
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()


def _cell_bg(cell, hx):
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for sf in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(sf)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), val=hx)


def _cell_border(cell, color_hex='808080', width_pt=0.5):
    """셀 테두리 4면 회색 설정"""
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
    for border in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for existing in tcPr.findall(qn(border)):
            tcPr.remove(existing)
        ln = etree.SubElement(tcPr, qn(border))
        ln.set('w', str(int(width_pt * 12700)))
        ln.set('cap', 'flat')
        ln.set('cmpd', 'sng')
        ln.set('algn', 'ctr')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr'), val=color_hex)
        etree.SubElement(ln, qn('a:prstDash'), val='solid')
        etree.SubElement(ln, qn('a:round'))


def _tbl(s, l, t, w, h, data, cws=None):
    """테이블 (멀티라인 셀 지원, 진한 테두리 0.5pt, 연한 헤더 배경). data[0]=헤더"""
    rows, cols = len(data), len(data[0])
    sh = s.shapes.add_table(rows, cols, l, t, w, h)
    tbl = sh.table
    if cws:
        for i, cw in enumerate(cws):
            tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ''
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_top = tf.margin_bottom = Inches(0.04)
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.06)
            hdr = (r == 0)
            for i, ln in enumerate(data[r][c].split('\n')):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                _run(p, ln, 9 if not hdr else 10, hdr, CHAR)
                p.alignment = PP_ALIGN.CENTER if hdr else PP_ALIGN.LEFT
                p.space_before = p.space_after = Pt(0)
            _cell_bg(cell, 'D4E2EC' if hdr else ('FFFFFF' if r % 2 == 1 else 'F2F7FB'))
            _cell_border(cell, '808080', 0.5)
    return tbl


# ══════════════════════════════════════════════════
# Page 1: 시스템 개념도
# ══════════════════════════════════════════════════

def p1_concept(prs):
    s = _sl(prs)
    _title(s, '시스템 개념도',
           '직무기술서관리체계의 전체 아키텍처 및 기술 구성 — 6개 계층', 1)

    layers = [
        ('사용자\n계층', [
            ('작성자 (개인)\nJD 작성·제출·복사', ICE, NAVY),
            ('관리자\nJD 조회·검토·의견·반송', ICE, NAVY),
            ('시스템관리자\n코드·권한·표준시간 관리', ICE, NAVY),
        ]),
        ('채널\n계층', [
            ('웹 브라우저\nChrome/Edge 1920x1080', ICE, TEAL),
            ('반응형 모바일 웹\n360px~', ICE, TEAL),
            ('API Gateway\nRESTful 28 endpoints', ICE, TEAL),
        ]),
        ('프론트\n엔드', [
            ('React 18 + TS 5\nVite 5 · Router 7', LBLUE, SKY),
            ('Ant Design 5.24\nTabs·Form·Charts', LBLUE, SKY),
            ('Zustand 5 + TQ 5\n상태관리 · 서버캐싱', LBLUE, SKY),
        ]),
        ('백엔드\n계층', [
            ('Spring Boot 3.x\n(추후 개발) Java 17', MINT, GRN),
            ('MSW 2.12 Mock\nFaker.js 9 데이터', MINT, GRN),
            ('REST API 설계\nJSON/HTTPS 28개', MINT, GRN),
        ]),
        ('데이터\n계층', [
            ('PostgreSQL 15\n(추후) 정규화 설계', CREAM, WARN),
            ('Mock 데이터 38건\nJD20+진단5+결재8+시간5', CREAM, WARN),
            ('파일 스토리지\n인쇄/엑셀 출력 (추후)', CREAM, WARN),
        ]),
        ('외부\n연계', [
            ('GPKI/SSO 인증\n(추후 연계)', SNOW, ASH),
            ('인사관리시스템\n(추후 연계)', SNOW, ASH),
            ('알림톡/SMS\n(추후 연계)', SNOW, ASH),
        ]),
    ]

    y0 = Inches(0.82)
    LH = Inches(0.44)
    GAP = Inches(0.05)
    LBL_W = Inches(0.85)
    LBL_G = Inches(0.05)

    for idx, (lbl, items) in enumerate(layers):
        y = y0 + idx * (LH + GAP)
        _box(s, ML, y, LBL_W, LH, lbl, fill=NAVY, stk=NAVY, sz=7, b=True, fc=WHITE)
        n = len(items)
        avail = CW - LBL_W - LBL_G
        igap = Inches(0.04)
        iw = (avail - igap * (n - 1)) // n
        x = ML + LBL_W + LBL_G
        for txt, fill, stk in items:
            _box(s, x, y, iw, LH, txt, fill=fill, stk=stk, sz=7, fc=CHAR)
            x += iw + igap
        if idx < len(layers) - 1:
            _darr(s, ML + CW // 2, y + LH + Inches(0.01))

    # [C] 서술 테이블
    dy = y0 + 6 * (LH + GAP) + Inches(0.05)
    _tbl(s, ML, dy, CW, Inches(1.45), [
        ['프론트엔드 기술스택', '백엔드 기술스택', '인프라 / 보안 / 연계'],
        ['• React 18.x + TypeScript 5.x\n'
         '• Ant Design 5.24 (Tabs/Form/Charts)\n'
         '• Zustand 5 + TanStack Query 5\n'
         '• React Router 7 + Vite 5\n'
         '• @ant-design/charts (Bar/Column/Pie)\n'
         '• 6개 고유 페이지, 15개 화면',
         '• Spring Boot 3.x (추후 개발)\n'
         '• MSW 2.12 Mock API (현재)\n'
         '• Faker.js 9 한국어 데이터\n'
         '• REST API 28개 endpoint\n'
         '• JWT + Redis 인증 (추후)\n'
         '• MyBatis 3.5 ORM (추후)',
         '• PostgreSQL 15 (추후 개발)\n'
         '• Linux / Tomcat 10.x (추후)\n'
         '• SSL/TLS 1.3, GPKI (추후)\n'
         '• 인사관리시스템 (추후 연계)\n'
         '• 알림톡/SMS/이메일 (추후)\n'
         '• 사용자: 3역할, 47 프로세스'],
    ], cws=[CW // 3] * 3)


# ══════════════════════════════════════════════════
# Page 2: 목표모델 (As-Is → To-Be)
# ══════════════════════════════════════════════════

def p2_target(prs):
    s = _sl(prs)
    _title(s, '목표모델 (As-Is → To-Be)',
           '현행 문제점 진단 및 목표 시스템 기능 구성', 2)

    pairs = [
        ('수기 직무기술서 작성\n워드/한글 문서 개별 작성, 양식 불통일',
         '온라인 JD 작성 시스템\nTabs 3탭(개인/직책/부서) + 업무비율 자동계산'),
        ('조직진단 대상 수기관리\n대상자/제외자 목록 엑셀 관리',
         '조직진단 대상 CRUD\nSearchForm+DataTable+기간관리+대상자Select'),
        ('결재 프로세스 부재\n상급자 승인 이메일/공문 기반',
         '2단계 전자결재\n1차승인→2차승인+반려/재결재 프로세스'),
        ('관리자 통계 부재\n부대별/직급별 작성현황 파악 어려움',
         '관리자 통계 대시보드\nBar+Column+Pie 차트 3종 시각화'),
        ('표준업무시간 미적용\n직급별 기준시간 없이 무분별 배분',
         '표준업무시간 CRUD\n신분별 주간시간 + 적용기간 자동상태'),
    ]

    y0 = Inches(1.0)
    RH = Inches(0.52)
    RG = Inches(0.03)
    AW = Inches(3.5)
    ARR_W = Inches(0.22)
    AX = ML
    ARX = ML + AW + Inches(0.15)
    TX = ARX + ARR_W + Inches(0.15)

    _box(s, AX, y0 - Inches(0.06), AW, Inches(0.18),
         '현행 시스템 (As-Is)', fill=ERR, stk=ERR, sz=8, b=True, fc=WHITE)
    _box(s, TX, y0 - Inches(0.06), AW, Inches(0.18),
         '목표 시스템 (To-Be)', fill=GRN, stk=GRN, sz=8, b=True, fc=WHITE)

    for i, (asis, tobe) in enumerate(pairs):
        y = y0 + i * (RH + RG)
        _box(s, AX, y, AW, RH, asis, fill=PEACH, stk=ERR, sz=7, fc=CHAR)
        _rarr(s, ARX, y + RH // 2, ARR_W)
        _box(s, TX, y, AW, RH, tobe, fill=MINT, stk=GRN, sz=7, fc=CHAR)

    # 우측 기능 요약
    FX = TX + AW + Inches(0.12)
    FW = ML + CW - FX
    _box(s, FX, y0, FW, 5 * RH + 4 * RG,
         '핵심 기능\n─────────\n'
         'JD 관리       5\n관리자          3\n게시판          3\n'
         '─────────\n총 11개 메뉴\n47 프로세스\n'
         '─────────\n3종 JD 유형\n3종 차트 통계\n'
         '표준업무시간',
         fill=ICE, stk=TEAL, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)

    # [C] 서술 테이블
    dy = y0 + 5 * (RH + RG) + Inches(0.12)
    _tbl(s, ML, dy, CW, Inches(1.25), [
        ['개선항목', '현행 문제점', '목표 개선방향', '기대 효과', '정량 목표'],
        ['JD 작성', '수기/양식 불통일', 'Tabs 3탭 + 업무비율 자동계산', '양식 통일, 작성 효율', '온라인 100%'],
        ['조직진단', '엑셀 대상자 관리', 'CRUD+기간관리+대상/제외자', '대상자 체계적 관리', '전산화 100%'],
        ['결재체계', '이메일/공문 기반', '2단계 전자결재+반려/재결재', '결재 투명성·속도', '전자결재 100%'],
        ['관리자통계', '현황 파악 불가', 'Bar+Column+Pie 3종 차트', '실시간 현황 파악', '자동 집계'],
        ['표준시간', '기준시간 미적용', '신분별 주간시간 CRUD', '업무배분 합리화', '5개 신분'],
    ], cws=[int(CW * r) for r in [0.10, 0.24, 0.24, 0.22, 0.10]])


# ══════════════════════════════════════════════════
# Page 3: 유스케이스
# ══════════════════════════════════════════════════

def p3_usecase(prs):
    s = _sl(prs)
    _title(s, '유스케이스 다이어그램',
           '3개 액터 x 18개 유스케이스, 47개 단위 프로세스', 3)

    y0 = Inches(0.82)

    # 좌측: 액터
    actors = [
        '작성자 (개인)\nJD 작성·제출·복사\nROLE_WRITER',
        '관리자\nJD 조회·검토·의견·반송\nROLE_MANAGER',
        '시스템관리자\nROLE_ADMIN',
    ]
    AX = ML
    AW = Inches(1.15)
    AH = Inches(0.55)
    AG = Inches(0.15)
    for i, atxt in enumerate(actors):
        _box(s, AX, y0 + i * (AH + AG), AW, AH, atxt,
             fill=ICE, stk=NAVY, sz=7, fc=CHAR)

    # 중앙: 시스템 바운더리
    BX = AX + AW + Inches(0.2)
    BW = Inches(6.2)
    BH = 3 * AH + 2 * AG + Inches(0.3)
    bnd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, BX, y0 - Inches(0.05), BW, BH)
    bnd.fill.solid()
    bnd.fill.fore_color.rgb = RGBColor(0xFB, 0xFB, 0xFB)
    bnd.line.color.rgb = NAVY
    bnd.line.width = Pt(2)
    _tb(s, BX + Inches(0.1), y0 - Inches(0.02), Inches(5), Inches(0.18),
        '직무기술서관리체계 (SYS18) — 바운더리', 8, True, NAVY)

    # UC 그룹
    uc_groups = [
        ('JD 작성 (6)', 'UC-JOB-001~006\nTabs 3탭(개인/직책/부서)+CRUD+복사+인쇄', ICE, SKY),
        ('조직진단 (4)', 'UC-JOB-007~010\nCRUD+SearchForm(3필드)+기간관리', ICE, SKY),
        ('결재 (4)', 'UC-JOB-011~014\n1차/2차 승인+반려+재결재', ICE, SKY),
        ('관리자 JD조회 (5)', 'UC-JOB-015~019\n개인/부서탭+검토+의견+반송+통계', ICE, SKY),
        ('표준업무시간 (3)', 'UC-JOB-020~022\nCRUD+적용상태자동+신분별', LBLUE, TEAL),
        ('관리자+게시판 (6)', 'UC-SYS 공통코드·권한관리\n공지/질의/자료실 게시판', LBLUE, TEAL),
    ]
    UY0 = y0 + Inches(0.22)
    UW = Inches(2.95)
    UH = Inches(0.48)
    UG = Inches(0.03)
    for i, (grp, desc, fill, stk) in enumerate(uc_groups):
        r, c = divmod(i, 2)
        ux = BX + Inches(0.08) + c * (UW + Inches(0.1))
        uy = UY0 + r * (UH + UG)
        sh = _box(s, ux, uy, UW, UH, f'{grp}\n{desc}',
                  fill=fill, stk=stk, sz=6, b=False, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True

    # 우측: 외부 시스템
    EX = BX + BW + Inches(0.12)
    EW = ML + CW - EX
    _box(s, EX, y0 + Inches(0.2), EW, Inches(0.35),
         '인사관리시스템\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)
    _box(s, EX, y0 + Inches(0.65), EW, Inches(0.35),
         'GPKI/SSO\n(추후 연계)', fill=SNOW, stk=ASH, sz=7, fc=ASH)

    # [C] 서술 테이블
    dy = Inches(3.2)
    _tbl(s, ML, dy, CW, Inches(2.1), [
        ['UC-ID', '유스케이스명', '주 액터', '기능 설명 (구체적)', '우선'],
        ['UC-JOB-001', 'JD 작성(개인)', '작성자',
         'Tabs 3탭(개인/직책/부서)+SearchForm(3필드)+DataTable(8컬럼)+CRUD+복사+인쇄', '상'],
        ['UC-JOB-007', '조직진단 대상 관리', '관리자',
         'SearchForm(3필드: 진단명/부대/상태)+DataTable(10컬럼)+CrudForm(7필드: 진단명/부대/기간/대상자/제외자)', '상'],
        ['UC-JOB-011', '결재(1차/2차)', '관리자',
         'SearchForm(3필드)+DataTable(7컬럼)+승인/반려/재결재 버튼+반려사유 Modal', '상'],
        ['UC-JOB-015', 'JD 조회(관리자)', '관리자',
         'Tabs 3탭(개인/부서/통계)+SearchForm(4필드)+DataTable+검토결과/의견/반송 Modal', '상'],
        ['UC-JOB-018', '통계 차트', '관리자',
         'Bar(부대별)+Column(직급별)+Pie(업무분류별) 3종 차트+SearchForm', '중'],
        ['UC-JOB-020', '표준업무시간', '관리자',
         'DataTable(6컬럼)+CRUD Modal(신분/주간시간/적용기간)+적용상태 자동계산(적용중/예정/만료)', '중'],
        ['UC-JOB-023', '게시판', '사용자',
         'SimpleBoardPage 재사용: 공지(/sys18/board/1), 질의(/sys18/board/2), 자료실(/sys18/board/3)', '하'],
    ], cws=[int(CW * r) for r in [0.10, 0.12, 0.06, 0.66, 0.05]])


# ══════════════════════════════════════════════════
# Page 4: 스토리보드 (화면흐름도)
# ══════════════════════════════════════════════════

def p4_storyboard(prs):
    s = _sl(prs)
    _title(s, '스토리보드 (화면흐름도)',
           '총 15개 화면, 3개 업무영역, 핵심 네비게이션 흐름', 4)

    y0 = Inches(0.82)

    # 로그인 → 대시보드
    _box(s, ML, y0, Inches(1.3), Inches(0.35),
         '로그인\n/login', fill=WARN, stk=WARN, sz=8, b=True, fc=WHITE)
    _rarr(s, ML + Inches(1.35), y0 + Inches(0.12), Inches(0.25))
    _box(s, ML + Inches(1.65), y0, Inches(2.0), Inches(0.35),
         '대시보드 / 홈\n/sys18', fill=NAVY, stk=NAVY, sz=8, b=True, fc=WHITE)

    _darr(s, ML + Inches(2.65), y0 + Inches(0.37))

    # 메뉴 그룹
    groups = [
        ('1. 직무기술서 관리', '조직진단·JD작성·결재·관리자조회·표준시간', '/sys18/1/*', '5', ICE, SKY),
        ('2. 게시판', '공지사항 · 질의응답 · 자료실', '/sys18/board/*', '3', ICE, SKY),
        ('3. 관리자', '공통코드·표준업무시간·사용자권한', '/sys18/admin/*', '3', LBLUE, TEAL),
        ('메인화면', 'SubsystemHomePage + 공지/질의 링크', '/sys18', '1', LBLUE, TEAL),
    ]

    GY0 = y0 + Inches(0.55)
    GW = Inches(4.5)
    GH = Inches(0.58)
    GG = Inches(0.03)
    GCOL_GAP = Inches(0.15)

    for i, (name, screens, path, cnt, fill, stk) in enumerate(groups):
        r, c = divmod(i, 2)
        gx = ML + c * (GW + GCOL_GAP)
        gy = GY0 + r * (GH + GG)
        sh = _box(s, gx, gy, GW, GH,
                  f'{name} ({cnt}개)\n{screens}\n{path}',
                  fill=fill, stk=stk, sz=7, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        pars = list(sh.text_frame.paragraphs)
        if len(pars) >= 3:
            for r_obj in pars[-1].runs:
                r_obj.font.color.rgb = ASH
                r_obj.font.size = Pt(8)

    # 핵심 화면 상세
    DY = GY0 + 2 * (GH + GG) + Inches(0.08)
    det_boxes = [
        ('직무기술서 작성\n/sys18/1/3',
         'Tabs 3탭 구조\n개인/직책/부서 JD\nSearchForm(3필드)\n'
         'DataTable(8컬럼)\nCRUD+복사+인쇄\n군번/계급/성명 표시',
         ICE, SKY),
        ('JD 조회(관리자)\n/sys18/1/5',
         'Tabs 3탭 구조\n개인/부서/통계\n검토결과/의견/반송 Modal\n'
         'SearchForm(4필드)\nBar+Column+Pie 차트\n인쇄+엑셀 다운로드',
         MINT, GRN),
        ('조직진단 대상\n/sys18/1/2',
         'SearchForm(3필드)\nDataTable(10컬럼)\nCrudForm Modal(7필드)\n'
         '기간: RangePicker\n대상자/제외자: Select\n진단기간 후 수정 불가',
         CREAM, WARN),
    ]
    DW = Inches(2.95)
    DH = Inches(1.1)
    DG = Inches(0.1)
    for i, (ttl, desc, fill, stk) in enumerate(det_boxes):
        dx = ML + i * (DW + DG)
        sh = _box(s, dx, DY, DW, DH,
                  f'{ttl}\n──────────\n{desc}',
                  fill=fill, stk=stk, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)
        sh.text_frame.paragraphs[0].runs[0].font.bold = True
        sh.text_frame.paragraphs[1].runs[0].font.bold = True

    # [C] 서술 테이블 2개 좌우 배치
    dy = DY + DH + Inches(0.08)

    LW = int(CW * 0.48)
    _tbl(s, ML, dy, LW, Inches(0.95), [
        ['영역', '화면 수', 'URL 패턴', '주요 컴포넌트'],
        ['조직진단', '1', '/sys18/1/2', 'SearchForm+DataTable+Form Modal'],
        ['JD 작성', '1', '/sys18/1/3', 'Tabs+SearchForm+DataTable+CRUD'],
        ['결재', '1', '/sys18/1/4', 'SearchForm+DataTable+승인/반려'],
        ['관리자 JD', '1', '/sys18/1/5', 'Tabs+DataTable+Modal(검토/의견/반송)'],
        ['표준시간+기타', '2+9', '/sys18/2/2,board,admin', 'DataTable+Form+Charts'],
    ], cws=[int(LW * r) for r in [0.18, 0.10, 0.30, 0.42]])

    RX = ML + LW + Inches(0.15)
    RW = CW - LW - Inches(0.15)
    _tbl(s, RX, dy, RW, Inches(0.95), [
        ['핵심 흐름 시나리오'],
        ['시나리오 1: JD 작성 → 제출 → 결재\n'
         '① 로그인 → 대시보드\n'
         '② [조직진단 대상 관리] 진단 등록\n'
         '③ [직무기술서 작성] 개인탭 → 신규 작성\n'
         '④ 업무비율(%) 입력 → 임시저장\n'
         '⑤ [제출] → 결재 요청\n'
         '⑥ [결재] 1차승인 → 2차승인 → 완료'],
        ['시나리오 2: 관리자 조회 → 검토\n'
         '① [JD 조회(관리자)] 검색+필터\n'
         '② 검토결과 입력 (적합/수정필요/부적합)\n'
         '③ 의견보내기 또는 반송 처리'],
    ])


# ══════════════════════════════════════════════════
# Page 5: 플로우차트 (업무처리흐름)
# ══════════════════════════════════════════════════

def p5_flowchart(prs):
    s = _sl(prs)
    _title(s, '업무처리 플로우차트',
           '핵심 프로세스: 조직진단 → JD작성 → 결재 → 관리자검토 — 4개 레인, 12단계', 5)

    y0 = Inches(0.82)

    # 4개 스윔레인
    lanes = [
        ('작성자 / 관리자', ICE, NAVY),
        ('프론트엔드 (React 18)', LBLUE, SKY),
        ('백엔드 (API/Mock)', MINT, GRN),
        ('DB / 외부', CREAM, WARN),
    ]
    LN_W = CW // 4
    LN_H = Inches(2.75)
    LN_HDR = Inches(0.22)

    for i, (name, fill, stk) in enumerate(lanes):
        lx = ML + LN_W * i
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = fill
        bg.line.color.rgb = stk
        bg.line.width = Pt(1)
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y0, LN_W, LN_HDR)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = stk
        hdr.line.fill.background()
        tf = hdr.text_frame
        _run(tf.paragraphs[0], name, 7, True, WHITE)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    steps = [
        (0, 0.30, '① 조직진단 등록\n진단명/부대/기간/대상자'),
        (1, 0.30, '② Form Modal 렌더\nRangePicker+Select'),
        (2, 0.30, '③ POST /org-diagnosis\n조직진단 저장'),
        (0, 0.70, '④ JD 작성 (개인탭)\n업무목록+비율(%) 입력'),
        (1, 0.70, '⑤ Tabs + Form 렌더\n업무비율 자동합계 100%'),
        (2, 0.70, '⑥ POST /job-descs\n직무기술서 저장'),
        (3, 0.70, '⑦ JobDesc INSERT\n데이터 저장'),
        (0, 1.10, '⑧ 제출 → 결재 요청\n[결재 요청] 버튼 클릭'),
        (2, 1.10, '⑨ PUT /submit\n상태: draft→submitted'),
        (0, 1.50, '⑩ 결재 1차/2차 승인\n승인/반려/재결재 버튼'),
        (2, 1.50, '⑪ PUT approve/reject\n결재상태 변경'),
        (0, 1.90, '⑫ 관리자 검토·반송\n검토결과/의견/반송 Modal'),
    ]

    SBW = LN_W - Inches(0.12)
    SBH = Inches(0.40)

    for lane, y_off, txt in steps:
        sx = ML + LN_W * lane + Inches(0.06)
        sy = y0 + LN_HDR + Inches(y_off)
        _box(s, sx, sy, SBW, SBH, txt,
             fill=WHITE, stk=ASH, sw=1, sz=6, fc=CHAR, al=PP_ALIGN.LEFT)

    # 판단 분기
    fx = ML + LN_W * 1 + Inches(0.06)
    fy = y0 + LN_HDR + Inches(1.35)
    _box(s, fx, fy, SBW, Inches(0.12),
         '반려 → 반려사유 입력 → 재결재 요청',
         fill=PEACH, stk=ERR, sw=1, sz=5, fc=ERR, al=PP_ALIGN.CENTER)

    # [C] 서술 테이블 2개 좌우 배치
    dy = y0 + LN_H + Inches(0.1)

    LW = int(CW * 0.52)
    _tbl(s, ML, dy, LW, Inches(1.42), [
        ['단계', '동작', '조건/상세'],
        ['①②③', '조직진단 등록', 'Form 7필드, POST /org-diagnosis'],
        ['④⑤⑥', 'JD 작성 → 저장', 'Tabs 3탭, 업무비율합계=100%'],
        ['⑦', 'DB 저장', 'JobDesc INSERT (Mock 20건)'],
        ['⑧⑨', '제출 → 결재요청', 'PUT /submit, status: draft→submitted'],
        ['⑩⑪', '1차/2차 결재', 'PUT approve/reject, 반려사유 입력'],
        ['⑫', '관리자 검토·반송', '검토결과(적합/수정/부적합)+의견+반송'],
    ], cws=[int(LW * r) for r in [0.08, 0.28, 0.64]])

    RX = ML + LW + Inches(0.12)
    RW = CW - LW - Inches(0.12)
    _tbl(s, RX, dy, RW, Inches(1.42), [
        ['업무 규칙 / 예외 처리'],
        ['[업무 규칙]\n'
         'BR-JOB-001: JD 유형 3종 (개인/직책/부서)\n'
         'BR-JOB-002: JD 상태 5종 (작성중/제출/1차완료/최종완료/반려)\n'
         'BR-JOB-003: 업무비율 합계 100% 필수\n'
         'BR-JOB-004: 업무유형 4종 (정책/관리/지원/기타)\n'
         'BR-JOB-005: 결재 2단계 (1차승인→2차승인)\n'
         'BR-JOB-006: 표준업무시간 신분 5종 (장관/영관/부사관/원사/병)\n'
         'BR-JOB-007: 적용상태 자동계산 (적용중/적용예정/적용만료)\n'
         'BR-JOB-008: 진단기간 이후 조직진단 수정/삭제 불가'],
        ['[예외 처리]\n'
         'E-001: 세션만료(401) → 로그인 리다이렉트\n'
         'E-002: 권한없음(403) → 접근 제한 알림\n'
         'E-003: 검증실패 → 필드별 에러 표시+재입력\n'
         'E-004: 서버오류(500) → 에러 메시지+재시도'],
    ])


# ══════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════

def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    print('[1/5] 시스템 개념도...')
    p1_concept(prs)
    print('[2/5] 목표모델...')
    p2_target(prs)
    print('[3/5] 유스케이스...')
    p3_usecase(prs)
    print('[4/5] 스토리보드...')
    p4_storyboard(prs)
    print('[5/5] 플로우차트...')
    p5_flowchart(prs)

    out_path = os.path.join(OUT, 'SYS18_직무기술서관리체계_분석설계_v5.pptx')
    prs.save(out_path)
    print(f'\n완료! 총 {len(prs.slides)}장 슬라이드')
    print(f'파일: {out_path}')


if __name__ == '__main__':
    main()
